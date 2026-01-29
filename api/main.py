"""
AgentForge - API Gateway v3.2
Complete API with Wizard-based Agent Creation, RAG Support, and Multi-Provider Configuration
"""

import os
import json

# Production mode - set to False to reduce logging
DEBUG_MODE = os.environ.get('DEBUG_MODE', 'false').lower() == 'true'

def debug_log(*args, **kwargs):
    """Only print if DEBUG_MODE is enabled"""
    if DEBUG_MODE:
        print(*args, **kwargs)
import uuid
import hashlib
import asyncio
import httpx
import re
import yaml
from typing import Dict, Any, List, Optional, Set
from datetime import datetime, timedelta
from contextlib import asynccontextmanager
from pathlib import Path
from urllib.parse import urljoin, urlparse
from abc import ABC, abstractmethod
from enum import Enum

from fastapi import FastAPI, HTTPException, UploadFile, File, Form, Request, Depends
from fastapi.responses import FileResponse, StreamingResponse
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import HTMLResponse
from pydantic import BaseModel, Field
import uvicorn

# Security Module
# Initialize User and get_current_user at module level
User = None
get_current_user = None

try:
    from core.security import security_state, User
    from api.security import router as security_router, get_current_user, get_current_user_optional
    SECURITY_AVAILABLE = True
    print("✅ Security module available")
except ImportError:
    SECURITY_AVAILABLE = False
    print("⚠️ Security module not installed - running without authentication")
    # Fallback for get_current_user if security module not available
    class User:
        pass
    def get_current_user():
        return None
    def get_current_user_optional():
        return None

# Health Check Module
try:
    from api.health import router as health_router
    HEALTH_CHECK_AVAILABLE = True
    print("✅ Health check module available")
except ImportError:
    HEALTH_CHECK_AVAILABLE = False
    print("⚠️ Health check module not available")

# Access Control Service
ACCESS_CONTROL_AVAILABLE = False
AccessControlService = None
try:
    from api.modules.access_control.service import AccessControlService
    ACCESS_CONTROL_AVAILABLE = True
    print("✅ Access Control Service available")
except ImportError as e:
    print(f"⚠️ Access Control Service not available: {e}")

# Instruction Enforcer Service
INSTRUCTION_ENFORCER_AVAILABLE = False
InstructionEnforcer = None
try:
    from core.llm.instruction_enforcer import InstructionEnforcer
    INSTRUCTION_ENFORCER_AVAILABLE = True
    print("✅ Instruction Enforcer Service available")
except ImportError as e:
    print(f"⚠️ Instruction Enforcer Service not available: {e}")

# Check for Playwright availability
PLAYWRIGHT_AVAILABLE = False
try:
    from playwright.async_api import async_playwright
    PLAYWRIGHT_AVAILABLE = True
    print("✅ Playwright available - dynamic site rendering enabled")
except ImportError:
    print("⚠️ Playwright not installed - dynamic sites will use basic scraping")

# Check for ChromaDB availability
CHROMADB_AVAILABLE = False
try:
    import chromadb
    from chromadb.config import Settings as ChromaSettings
    CHROMADB_AVAILABLE = True
    print("✅ ChromaDB available - vector search enabled")
except ImportError:
    print("⚠️ ChromaDB not installed - using keyword search")
    print("   Install with: pip install chromadb")

# Check for Sentence Transformers availability
SENTENCE_TRANSFORMERS_AVAILABLE = False
try:
    from sentence_transformers import SentenceTransformer
    SENTENCE_TRANSFORMERS_AVAILABLE = True
    print("✅ Sentence Transformers available - local embeddings enabled")
except ImportError:
    print("⚠️ Sentence Transformers not installed - using API embeddings only")
    print("   Install with: pip install sentence-transformers")


# ============================================================================
# Provider Enums
# ============================================================================

class LLMProvider(str, Enum):
    OPENAI = "openai"
    AZURE_OPENAI = "azure_openai"
    ANTHROPIC = "anthropic"
    GOOGLE = "google"
    XAI = "xai"
    GROQ = "groq"
    MISTRAL = "mistral"
    DEEPSEEK = "deepseek"
    TOGETHER = "together"
    PERPLEXITY = "perplexity"
    COHERE = "cohere"
    AWS_BEDROCK = "aws_bedrock"
    OLLAMA = "ollama"
    LMSTUDIO = "lmstudio"
    CUSTOM = "custom"

class EmbeddingProvider(str, Enum):
    OPENAI = "openai"
    AZURE_OPENAI = "azure_openai"
    SENTENCE_TRANSFORMERS = "sentence_transformers"
    OLLAMA = "ollama"
    COHERE = "cohere"
    CUSTOM = "custom"

class VectorDBProvider(str, Enum):
    CHROMADB = "chromadb"
    PINECONE = "pinecone"
    QDRANT = "qdrant"
    MILVUS = "milvus"
    WEAVIATE = "weaviate"
    MEMORY = "memory"  # In-memory keyword search (current)


# ============================================================================
# Settings Models
# ============================================================================

class LLMConfig(BaseModel):
    provider: LLMProvider = LLMProvider.OPENAI
    model: str = "gpt-4o"
    api_key: str = ""
    api_base: str = ""  # For Azure, Ollama, Custom
    api_version: str = ""  # For Azure
    temperature: float = 0.7
    max_tokens: int = 4096
    # Ollama specific
    ollama_host: str = "http://localhost:11434"
    # Azure specific
    azure_deployment: str = ""
    # AWS Bedrock specific
    aws_region: str = "us-east-1"
    aws_access_key: str = ""
    aws_secret_key: str = ""

class EmbeddingConfig(BaseModel):
    provider: EmbeddingProvider = EmbeddingProvider.OPENAI
    model: str = "text-embedding-3-small"
    api_key: str = ""
    api_base: str = ""
    # Sentence Transformers specific
    local_model: str = "all-MiniLM-L6-v2"
    # Ollama specific
    ollama_host: str = "http://localhost:11434"
    ollama_model: str = "nomic-embed-text"
    # Dimensions
    dimensions: int = 1536

class VectorDBConfig(BaseModel):
    provider: VectorDBProvider = VectorDBProvider.MEMORY
    # ChromaDB specific
    chroma_path: str = "./data/chromadb"
    chroma_collection: str = "agentforge"
    # Pinecone specific
    pinecone_api_key: str = ""
    pinecone_environment: str = ""
    pinecone_index: str = ""
    # Qdrant specific
    qdrant_host: str = "localhost"
    qdrant_port: int = 6333
    qdrant_api_key: str = ""
    qdrant_collection: str = "agentforge"
    # Milvus specific
    milvus_host: str = "localhost"
    milvus_port: int = 19530
    milvus_collection: str = "agentforge"

class GoogleConfig(BaseModel):
    """Google API Keys configuration"""
    gemini_key: str = ""  # From AI Studio - for Gemini models
    cloud_key: str = ""   # From Cloud Console - for Vision, Maps, etc.

class LLMProviderConfig(BaseModel):
    """Configuration for a single LLM provider"""
    provider: str
    name: str
    api_key: str = ""
    api_base: str = ""
    resource: str = ""
    models: List[str] = []

class SystemSettings(BaseModel):
    # LLM Settings (default provider)
    llm: LLMConfig = Field(default_factory=LLMConfig)
    # All configured LLM providers
    llm_providers: List[LLMProviderConfig] = Field(default_factory=list)
    # Google Services
    google: GoogleConfig = Field(default_factory=GoogleConfig)
    # Embedding Settings
    embedding: EmbeddingConfig = Field(default_factory=EmbeddingConfig)
    # Vector DB Settings
    vector_db: VectorDBConfig = Field(default_factory=VectorDBConfig)
    # General Settings
    chunk_size: int = 1000
    chunk_overlap: int = 200
    search_top_k: int = 5
    # Feature Flags
    enable_rag: bool = True
    enable_web_scraping: bool = True
    enable_api_tools: bool = True
    # UI Settings
    app_name: str = "AgentForge"
    app_logo: str = "🔥"
    theme: str = "dark"


# ============================================================================
# Provider Abstractions
# ============================================================================

class BaseLLMProvider(ABC):
    @abstractmethod
    async def generate(self, messages: List[Dict], **kwargs) -> str:
        pass
    
    @abstractmethod
    def get_available_models(self) -> List[str]:
        pass

class BaseEmbeddingProvider(ABC):
    @abstractmethod
    async def embed(self, texts: List[str]) -> List[List[float]]:
        pass
    
    @abstractmethod
    def get_dimensions(self) -> int:
        pass

class BaseVectorDB(ABC):
    @abstractmethod
    async def add_documents(self, documents: List[Dict], embeddings: List[List[float]]):
        pass
    
    @abstractmethod
    async def search(self, query_embedding: List[float], top_k: int = 5, filter: Dict = None) -> List[Dict]:
        pass
    
    @abstractmethod
    async def delete(self, ids: List[str]):
        pass


# ============================================================================
# LLM Provider Implementations
# ============================================================================

class OpenAILLM(BaseLLMProvider):
    def __init__(self, config: LLMConfig):
        self.config = config
        self.api_key = config.api_key or os.environ.get("OPENAI_API_KEY", "")
    
    async def generate(self, messages: List[Dict], **kwargs) -> str:
        try:
            from openai import AsyncOpenAI
            client = AsyncOpenAI(api_key=self.api_key)
            response = await client.chat.completions.create(
                model=kwargs.get("model", self.config.model),
                messages=messages,
                temperature=kwargs.get("temperature", self.config.temperature),
                max_tokens=kwargs.get("max_tokens", self.config.max_tokens)
            )
            return response.choices[0].message.content or ""
        except Exception as e:
            raise Exception(f"OpenAI error: {str(e)}")
    
    def get_available_models(self) -> List[str]:
        return ["gpt-4o", "gpt-4o-mini", "gpt-4-turbo", "gpt-4", "gpt-3.5-turbo"]


class AzureOpenAILLM(BaseLLMProvider):
    def __init__(self, config: LLMConfig):
        self.config = config
        self.api_key = config.api_key or os.environ.get("AZURE_OPENAI_API_KEY", "")
        self.api_base = config.api_base or os.environ.get("AZURE_OPENAI_ENDPOINT", "")
        self.api_version = config.api_version or "2024-02-15-preview"
    
    async def generate(self, messages: List[Dict], **kwargs) -> str:
        try:
            from openai import AsyncAzureOpenAI
            client = AsyncAzureOpenAI(
                api_key=self.api_key,
                api_version=self.api_version,
                azure_endpoint=self.api_base
            )
            response = await client.chat.completions.create(
                model=kwargs.get("model", self.config.azure_deployment),
                messages=messages,
                temperature=kwargs.get("temperature", self.config.temperature),
                max_tokens=kwargs.get("max_tokens", self.config.max_tokens)
            )
            return response.choices[0].message.content or ""
        except Exception as e:
            raise Exception(f"Azure OpenAI error: {str(e)}")
    
    def get_available_models(self) -> List[str]:
        return ["gpt-4o", "gpt-4", "gpt-35-turbo"]  # Deployment names


class AnthropicLLM(BaseLLMProvider):
    def __init__(self, config: LLMConfig):
        self.config = config
        self.api_key = config.api_key or os.environ.get("ANTHROPIC_API_KEY", "")
    
    async def generate(self, messages: List[Dict], **kwargs) -> str:
        try:
            print(f"[AnthropicLLM] Calling with {len(messages)} messages, model: {kwargs.get('model', self.config.model)}")
            import anthropic
            client = anthropic.AsyncAnthropic(api_key=self.api_key)
            # Convert messages format
            system = ""
            msgs = []
            for m in messages:
                if m["role"] == "system":
                    system = m["content"]
                elif m["role"] == "tool":
                    # Convert OpenAI tool response to Anthropic tool_result format
                    # Anthropic uses "user" role with tool_result content blocks
                    tool_call_id = m.get("tool_call_id", "")
                    tool_content = m.get("content", "")
                    msgs.append({
                        "role": "user",
                        "content": [{
                            "type": "tool_result",
                            "tool_use_id": tool_call_id,
                            "content": tool_content
                        }]
                    })
                elif m["role"] in ["user", "assistant"]:
                    # Handle both string and list content (for Anthropic block format)
                    content = m.get("content", "")
                    if isinstance(content, list):
                        # Already in Anthropic block format
                        msgs.append({"role": m["role"], "content": content})
                else:
                        # String content
                        msgs.append({"role": m["role"], "content": content})
            
            # Build request kwargs - system must be a list for new Anthropic SDK
            request_kwargs = {
                "model": kwargs.get("model", self.config.model),
                "max_tokens": kwargs.get("max_tokens", self.config.max_tokens),
                "messages": msgs
            }
            # Only add system if we have one (as a list of content blocks)
            if system:
                request_kwargs["system"] = [{"type": "text", "text": system}]
            
            response = await client.messages.create(**request_kwargs)
            text = response.content[0].text
            print(f"[AnthropicLLM] ✅ Got response: {len(text)} chars")
            return text
        except Exception as e:
            print(f"[AnthropicLLM] ❌ Error: {e}")
            import traceback
            traceback.print_exc()
            raise Exception(f"Anthropic error: {str(e)}")
    
    def get_available_models(self) -> List[str]:
        return ["claude-sonnet-4-20250514", "claude-opus-4-20250514", "claude-3-5-sonnet-20241022", "claude-3-5-haiku-20241022"]


class OllamaLLM(BaseLLMProvider):
    def __init__(self, config: LLMConfig):
        self.config = config
        self.host = config.ollama_host or "http://localhost:11434"
    
    async def generate(self, messages: List[Dict], **kwargs) -> str:
        try:
            async with httpx.AsyncClient(timeout=120.0) as client:
                response = await client.post(
                    f"{self.host}/api/chat",
                    json={
                        "model": kwargs.get("model", self.config.model),
                        "messages": messages,
                        "stream": False,
                        "options": {
                            "temperature": kwargs.get("temperature", self.config.temperature)
                        }
                    }
                )
                response.raise_for_status()
                return response.json()["message"]["content"]
        except Exception as e:
            raise Exception(f"Ollama error: {str(e)}")
    
    def get_available_models(self) -> List[str]:
        try:
            import requests
            response = requests.get(f"{self.host}/api/tags", timeout=5)
            if response.status_code == 200:
                models = response.json().get("models", [])
                return [m["name"] for m in models]
        except:
            pass
        return ["llama3.2", "llama3.1", "mistral", "codellama", "phi3", "gemma2"]


class CustomLLM(BaseLLMProvider):
    """Custom LLM provider for OpenAI-compatible APIs"""
    def __init__(self, config: LLMConfig):
        self.config = config
        self.api_key = config.api_key
        self.api_base = config.api_base
    
    async def generate(self, messages: List[Dict], **kwargs) -> str:
        try:
            from openai import AsyncOpenAI
            client = AsyncOpenAI(api_key=self.api_key, base_url=self.api_base)
            response = await client.chat.completions.create(
                model=kwargs.get("model", self.config.model),
                messages=messages,
                temperature=kwargs.get("temperature", self.config.temperature),
                max_tokens=kwargs.get("max_tokens", self.config.max_tokens)
            )
            return response.choices[0].message.content or ""
        except Exception as e:
            raise Exception(f"Custom LLM error: {str(e)}")
    
    def get_available_models(self) -> List[str]:
        return [self.config.model] if self.config.model else []


class OpenAICompatibleLLM(BaseLLMProvider):
    """LLM provider for OpenAI-compatible APIs (Groq, xAI, Mistral, DeepSeek, Together, Perplexity, etc)"""
    
    # API base URLs for each provider
    PROVIDER_CONFIGS = {
        "xai": {
            "base_url": "https://api.x.ai/v1",
            "models": ["grok-2", "grok-2-mini", "grok-beta"]
        },
        "groq": {
            "base_url": "https://api.groq.com/openai/v1",
            "models": ["llama-3.3-70b-versatile", "llama-3.3-70b-specdec", "llama-3.1-8b-instant", "gemma2-9b-it", "mixtral-8x7b-32768"]
        },
        "mistral": {
            "base_url": "https://api.mistral.ai/v1",
            "models": ["mistral-large-2411", "mistral-small-2503", "open-mistral-nemo", "codestral-2501", "ministral-8b-2410"]
        },
        "deepseek": {
            "base_url": "https://api.deepseek.com/v1",
            "models": ["deepseek-chat", "deepseek-coder", "deepseek-reasoner"]
        },
        "together": {
            "base_url": "https://api.together.xyz/v1",
            "models": ["meta-llama/Llama-3.3-70B-Instruct-Turbo", "meta-llama/Llama-3.1-405B-Instruct-Turbo", "mistralai/Mixtral-8x22B-Instruct-v0.1"]
        },
        "perplexity": {
            "base_url": "https://api.perplexity.ai",
            "models": ["sonar", "sonar-pro", "sonar-reasoning"]
        },
        "lmstudio": {
            "base_url": "http://localhost:1234/v1",
            "models": ["local-model"]
        }
    }
    
    def __init__(self, config: LLMConfig, provider_name: str = None):
        self.config = config
        self.provider_name = provider_name or config.provider.value
        self.api_key = config.api_key
        
        # Get provider-specific base URL or use custom
        provider_config = self.PROVIDER_CONFIGS.get(self.provider_name, {})
        self.api_base = config.api_base or provider_config.get("base_url", "")
    
    async def generate(self, messages: List[Dict], **kwargs) -> str:
        try:
            from openai import AsyncOpenAI
            client = AsyncOpenAI(api_key=self.api_key, base_url=self.api_base)
            response = await client.chat.completions.create(
                model=kwargs.get("model", self.config.model),
                messages=messages,
                temperature=kwargs.get("temperature", self.config.temperature),
                max_tokens=kwargs.get("max_tokens", self.config.max_tokens)
            )
            return response.choices[0].message.content or ""
        except Exception as e:
            raise Exception(f"{self.provider_name} LLM error: {str(e)}")
    
    def get_available_models(self) -> List[str]:
        provider_config = self.PROVIDER_CONFIGS.get(self.provider_name, {})
        return provider_config.get("models", [self.config.model])


class GoogleLLM(BaseLLMProvider):
    """Google Gemini LLM provider - uses REST API for async compatibility"""
    def __init__(self, config: LLMConfig):
        self.config = config
        self.api_key = config.api_key
    
    async def generate(self, messages: List[Dict], **kwargs) -> str:
        """Generate response using Google Gemini REST API"""
        import httpx
        
        # Build the prompt from messages
        contents = []
        system_instruction = ""
        
        for msg in messages:
            role = msg.get("role", "user")
            content = msg.get("content", "")
            
            if role == "system":
                system_instruction = content
            elif role == "user":
                contents.append({"role": "user", "parts": [{"text": content}]})
            elif role == "assistant":
                contents.append({"role": "model", "parts": [{"text": content}]})
        
        # If we have a system instruction, prepend it to the first user message
        if system_instruction and contents:
            first_user_text = contents[0]["parts"][0]["text"]
            contents[0]["parts"][0]["text"] = f"{system_instruction}\n\n{first_user_text}"
        
        # Ensure we have at least one message
        if not contents:
            contents = [{"role": "user", "parts": [{"text": "Hello"}]}]
        
        # Fix model name format - Google API uses specific names
        model_name = self.config.model or "gemini-1.5-flash"
        # Map common names to actual API model names that support function calling
        model_mapping = {
            "gemini-1.5-pro": "gemini-1.5-flash",  # 1.5-pro requires specific version
            "gemini-1.5-flash": "gemini-1.5-flash", 
            "gemini-pro": "gemini-1.5-flash",
            "gemini-2.0-flash": "gemini-1.5-flash",  # 2.0 not stable yet
        }
        api_model = model_mapping.get(model_name, "gemini-1.5-flash")
        print(f"[GoogleLLM] Using API model: {api_model}")
        url = f"https://generativelanguage.googleapis.com/v1beta/models/{api_model}:generateContent"
        
        payload = {
            "contents": contents,
            "generationConfig": {
                "temperature": kwargs.get("temperature", getattr(self.config, 'temperature', 0.7)),
                "maxOutputTokens": kwargs.get("max_tokens", getattr(self.config, 'max_tokens', 4096))
            }
        }
        
        try:
            async with httpx.AsyncClient() as client:
                print(f"[GoogleLLM] Calling {model_name} with {len(contents)} messages...")
                response = await client.post(
                    url,
                    params={"key": self.api_key},
                    json=payload,
                    timeout=60.0
                )
                
                if response.status_code != 200:
                    error_text = response.text
                    print(f"[GoogleLLM] ❌ API error {response.status_code}: {error_text[:200]}")
                    return f"Error: Google API returned {response.status_code}"
                
                result = response.json()
                
                # Extract text from response
                candidates = result.get("candidates", [])
                if not candidates:
                    print(f"[GoogleLLM] ⚠️ No candidates in response: {result}")
                    # Check for blocked content
                    if result.get("promptFeedback", {}).get("blockReason"):
                        return f"Response blocked: {result['promptFeedback']['blockReason']}"
                    return ""
                
                content = candidates[0].get("content", {})
                parts = content.get("parts", [])
                if not parts:
                    print(f"[GoogleLLM] ⚠️ No parts in content: {content}")
                    return ""
                
                text = parts[0].get("text", "")
                print(f"[GoogleLLM] ✅ Got response: {len(text)} chars")
                return text
                
        except httpx.TimeoutException:
            print(f"[GoogleLLM] ❌ Request timeout")
            return "Error: Request timed out"
        except Exception as e:
            print(f"[GoogleLLM] ❌ Error: {e}")
            import traceback
            traceback.print_exc()
            return f"Error: {str(e)}"
    
    def get_available_models(self) -> List[str]:
        return ["gemini-2.0-flash", "gemini-1.5-pro", "gemini-1.5-flash"]


class CohereLLM(BaseLLMProvider):
    """Cohere LLM provider - uses native Cohere API"""
    def __init__(self, config: LLMConfig):
        self.config = config
        self.api_key = config.api_key
    
    async def generate(self, messages: List[Dict], **kwargs) -> str:
        try:
            print(f"[CohereLLM] Calling with {len(messages)} messages, model: {kwargs.get('model', self.config.model)}")
            import httpx
            
            # Convert messages to Cohere format
            chat_history = []
            user_message = ""
            preamble = ""
            
            for msg in messages:
                if msg["role"] == "system":
                    preamble = msg["content"]
                elif msg["role"] == "user":
                    user_message = msg["content"]
                elif msg["role"] == "assistant":
                    if user_message:
                        chat_history.append({"role": "USER", "message": user_message})
                        user_message = ""
                    chat_history.append({"role": "CHATBOT", "message": msg["content"]})
            
            # The last user message
            if not user_message and messages:
                for msg in reversed(messages):
                    if msg["role"] == "user":
                        user_message = msg["content"]
                        break
            
            async with httpx.AsyncClient(timeout=60.0) as client:
                response = await client.post(
                    "https://api.cohere.com/v2/chat",
                    headers={
                        "Authorization": f"Bearer {self.api_key}",
                        "Content-Type": "application/json",
                        "Accept": "application/json"
                    },
                    json={
                        "model": kwargs.get("model", self.config.model),
                        "messages": [{"role": "user", "content": user_message}],
                    }
                )
                
                if response.status_code != 200:
                    error_detail = response.text
                    raise Exception(f"Error code: {response.status_code} - {error_detail}")
                
                result = response.json()
                # Cohere v2 response format
                if "message" in result and "content" in result["message"]:
                    content = result["message"]["content"]
                    if isinstance(content, list) and len(content) > 0:
                        text = content[0].get("text", "")
                        print(f"[CohereLLM] ✅ Got response: {len(text)} chars")
                        return text
                    print(f"[CohereLLM] ✅ Got response: {len(str(content))} chars")
                    return str(content)
                text = result.get("text", "")
                print(f"[CohereLLM] ✅ Got response: {len(text)} chars")
                return text
                
        except Exception as e:
            print(f"[CohereLLM] ❌ Error: {e}")
            raise Exception(f"Cohere LLM error: {str(e)}")
    
    def get_available_models(self) -> List[str]:
        return ["command-a-03-2025", "command-r-plus-08-2024", "command-r-08-2024", "command-r7b-12-2024"]


# ============================================================================
# Embedding Provider Implementations
# ============================================================================

class OpenAIEmbedding(BaseEmbeddingProvider):
    def __init__(self, config: EmbeddingConfig):
        self.config = config
        self.api_key = config.api_key or os.environ.get("OPENAI_API_KEY", "")
        self._dimensions = {"text-embedding-3-small": 1536, "text-embedding-3-large": 3072, "text-embedding-ada-002": 1536}
    
    async def embed(self, texts: List[str]) -> List[List[float]]:
        try:
            from openai import AsyncOpenAI
            client = AsyncOpenAI(api_key=self.api_key)
            response = await client.embeddings.create(
                model=self.config.model,
                input=texts
            )
            return [item.embedding for item in response.data]
        except Exception as e:
            raise Exception(f"OpenAI Embedding error: {str(e)}")
    
    def get_dimensions(self) -> int:
        return self._dimensions.get(self.config.model, 1536)


class SentenceTransformersEmbedding(BaseEmbeddingProvider):
    def __init__(self, config: EmbeddingConfig):
        self.config = config
        self.model = None
        self._dimensions = {
            "all-MiniLM-L6-v2": 384,
            "all-mpnet-base-v2": 768,
            "paraphrase-multilingual-MiniLM-L12-v2": 384,
            "multi-qa-MiniLM-L6-cos-v1": 384
        }
    
    def _load_model(self):
        if self.model is None and SENTENCE_TRANSFORMERS_AVAILABLE:
            from sentence_transformers import SentenceTransformer
            self.model = SentenceTransformer(self.config.local_model)
    
    async def embed(self, texts: List[str]) -> List[List[float]]:
        try:
            self._load_model()
            if self.model is None:
                raise Exception("Sentence Transformers not available")
            embeddings = self.model.encode(texts, convert_to_numpy=True)
            return embeddings.tolist()
        except Exception as e:
            raise Exception(f"Sentence Transformers error: {str(e)}")
    
    def get_dimensions(self) -> int:
        return self._dimensions.get(self.config.local_model, 384)


class OllamaEmbedding(BaseEmbeddingProvider):
    def __init__(self, config: EmbeddingConfig):
        self.config = config
        self.host = config.ollama_host or "http://localhost:11434"
    
    async def embed(self, texts: List[str]) -> List[List[float]]:
        try:
            embeddings = []
            async with httpx.AsyncClient(timeout=60.0) as client:
                for text in texts:
                    response = await client.post(
                        f"{self.host}/api/embeddings",
                        json={"model": self.config.ollama_model, "prompt": text}
                    )
                    response.raise_for_status()
                    embeddings.append(response.json()["embedding"])
            return embeddings
        except Exception as e:
            raise Exception(f"Ollama Embedding error: {str(e)}")
    
    def get_dimensions(self) -> int:
        return self.config.dimensions or 768


# ============================================================================
# Vector DB Implementations
# ============================================================================

class MemoryVectorDB(BaseVectorDB):
    """In-memory keyword search (no embeddings required)"""
    def __init__(self, config: VectorDBConfig):
        self.config = config
        self.documents: List[Dict] = []
    
    async def add_documents(self, documents: List[Dict], embeddings: List[List[float]] = None):
        for doc in documents:
            self.documents.append(doc)
    
    async def search(self, query_embedding: List[float] = None, top_k: int = 5, filter: Dict = None, query_text: str = "") -> List[Dict]:
        # Keyword search fallback
        if not query_text:
            return []
        query_words = query_text.lower().split()
        results = []
        for doc in self.documents:
            if filter:
                match = all(doc.get(k) == v for k, v in filter.items())
                if not match:
                    continue
            text_lower = doc.get('text', '').lower()
            score = sum(1 for word in query_words if word in text_lower)
            if score > 0:
                results.append({**doc, "score": score / len(query_words)})
        results.sort(key=lambda x: x['score'], reverse=True)
        return results[:top_k]
    
    async def delete(self, ids: List[str]):
        self.documents = [d for d in self.documents if d.get('id') not in ids]


class ChromaVectorDB(BaseVectorDB):
    """ChromaDB vector database - runs locally, no external dependencies"""
    def __init__(self, config: VectorDBConfig, embedding_dim: int = 384):
        self.config = config
        self.client = None
        self.collection = None
        self.embedding_dim = embedding_dim
    
    def _init_client(self):
        if self.client is None and CHROMADB_AVAILABLE:
            import chromadb
            from chromadb.config import Settings
            
            persist_path = self.config.chroma_path or "./data/chromadb"
            os.makedirs(persist_path, exist_ok=True)
            
            self.client = chromadb.PersistentClient(
                path=persist_path,
                settings=Settings(anonymized_telemetry=False)
            )
            self.collection = self.client.get_or_create_collection(
                name=self.config.chroma_collection or "agentforge",
                metadata={"hnsw:space": "cosine"}
            )
    
    async def add_documents(self, documents: List[Dict], embeddings: List[List[float]]):
        self._init_client()
        if self.collection is None:
            raise Exception("ChromaDB not initialized")
        
        ids = [doc.get('id', str(uuid.uuid4())) for doc in documents]
        texts = [doc.get('text', '') for doc in documents]
        metadatas = [{k: v for k, v in doc.items() if k not in ['id', 'text', 'embedding']} for doc in documents]
        
        self.collection.add(
            ids=ids,
            embeddings=embeddings,
            documents=texts,
            metadatas=metadatas
        )
    
    async def search(self, query_embedding: List[float], top_k: int = 5, filter: Dict = None, query_text: str = "") -> List[Dict]:
        self._init_client()
        if self.collection is None:
            return []
        
        where_filter = filter if filter else None
        
        results = self.collection.query(
            query_embeddings=[query_embedding],
            n_results=top_k,
            where=where_filter
        )
        
        documents = []
        if results and results['documents']:
            for i, doc in enumerate(results['documents'][0]):
                documents.append({
                    "id": results['ids'][0][i] if results['ids'] else "",
                    "text": doc,
                    "score": 1 - (results['distances'][0][i] if results['distances'] else 0),
                    **(results['metadatas'][0][i] if results['metadatas'] else {})
                })
        return documents
    
    async def delete(self, ids: List[str]):
        self._init_client()
        if self.collection:
            self.collection.delete(ids=ids)
    
    async def delete_by_filter(self, filter: Dict):
        """Delete documents matching filter"""
        self._init_client()
        if self.collection:
            self.collection.delete(where=filter)


# ============================================================================
# Provider Factory
# ============================================================================

class ProviderFactory:
    @staticmethod
    def get_llm_provider(config: LLMConfig) -> BaseLLMProvider:
        # Direct mapping for main providers
        direct_providers = {
            LLMProvider.OPENAI: OpenAILLM,
            LLMProvider.AZURE_OPENAI: AzureOpenAILLM,
            LLMProvider.ANTHROPIC: AnthropicLLM,
            LLMProvider.OLLAMA: OllamaLLM,
            LLMProvider.GOOGLE: GoogleLLM,
            LLMProvider.COHERE: CohereLLM,
            LLMProvider.CUSTOM: CustomLLM,
        }
        
        # OpenAI-compatible providers
        openai_compatible = {
            LLMProvider.XAI,
            LLMProvider.GROQ,
            LLMProvider.MISTRAL,
            LLMProvider.DEEPSEEK,
            LLMProvider.TOGETHER,
            LLMProvider.PERPLEXITY,
            LLMProvider.LMSTUDIO,
        }
        
        if config.provider in direct_providers:
            return direct_providers[config.provider](config)
        elif config.provider in openai_compatible:
            return OpenAICompatibleLLM(config, config.provider.value)
        else:
            # Default to OpenAI
            return OpenAILLM(config)
        
        if config.provider in direct_providers:
            return direct_providers[config.provider](config)
        elif config.provider in openai_compatible:
            return OpenAICompatibleLLM(config, config.provider.value)
        else:
            # Default to OpenAI
            return OpenAILLM(config)
    
    @staticmethod
    def get_embedding_provider(config: EmbeddingConfig) -> BaseEmbeddingProvider:
        providers = {
            EmbeddingProvider.OPENAI: OpenAIEmbedding,
            EmbeddingProvider.SENTENCE_TRANSFORMERS: SentenceTransformersEmbedding,
            EmbeddingProvider.OLLAMA: OllamaEmbedding,
        }
        provider_class = providers.get(config.provider, OpenAIEmbedding)
        return provider_class(config)
    
    @staticmethod
    def get_vector_db(config: VectorDBConfig, embedding_dim: int = 384) -> BaseVectorDB:
        if config.provider == VectorDBProvider.CHROMADB and CHROMADB_AVAILABLE:
            return ChromaVectorDB(config, embedding_dim)
        # Default to memory
        return MemoryVectorDB(config)
    
    @staticmethod
    def get_kb_embedding_provider(kb_config, global_config: EmbeddingConfig) -> BaseEmbeddingProvider:
        """Get embedding provider for a specific Knowledge Base"""
        if kb_config.get('use_global', True):
            return ProviderFactory.get_embedding_provider(global_config)
        
        # Create custom config from KB settings
        emb_config = EmbeddingConfig(
            provider=EmbeddingProvider(kb_config.get('provider', 'openai')),
            model=kb_config.get('model', 'text-embedding-3-small'),
            local_model=kb_config.get('local_model', 'all-MiniLM-L6-v2'),
            api_key=kb_config.get('api_key', '')
        )
        return ProviderFactory.get_embedding_provider(emb_config)
    
    @staticmethod
    def get_kb_vector_db(kb_config, global_config: VectorDBConfig, collection_id: str, embedding_dim: int = 384) -> BaseVectorDB:
        """Get vector DB for a specific Knowledge Base with its own collection"""
        if kb_config.get('use_global', True):
            # Use global provider but with KB-specific collection
            config = VectorDBConfig(
                provider=global_config.provider,
                chroma_path=global_config.chroma_path,
                chroma_collection=collection_id,  # KB-specific collection
                pinecone_api_key=global_config.pinecone_api_key,
                pinecone_index=global_config.pinecone_index,
            )
        else:
            # Use KB-specific provider
            provider = kb_config.get('provider', 'chromadb')
            config = VectorDBConfig(
                provider=VectorDBProvider(provider),
                chroma_path=global_config.chroma_path,
                chroma_collection=collection_id,
                pinecone_api_key=kb_config.get('pinecone_api_key', ''),
            )
        
        return ProviderFactory.get_vector_db(config, embedding_dim)


# ============================================================================
# Data Models
# ============================================================================

class APIInputParameter(BaseModel):
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    name: str
    description: str = ""
    data_type: str = "string"
    required: bool = False
    location: str = "query"


class APIEndpointConfig(BaseModel):
    base_url: str = ""
    http_method: str = "GET"
    endpoint_path: str = ""
    auth_type: str = "none"
    auth_value: str = ""
    api_key_name: str = "X-API-Key"
    api_key_location: str = "header"
    headers: Dict[str, str] = {}
    input_parameters: List[APIInputParameter] = []


class KnowledgeBaseEmbeddingConfig(BaseModel):
    use_global: bool = True
    provider: str = "openai"
    model: str = "text-embedding-3-small"
    local_model: str = "all-MiniLM-L6-v2"
    api_key: str = ""

class KnowledgeBaseVectorDBConfig(BaseModel):
    use_global: bool = True
    provider: str = "chromadb"
    pinecone_api_key: str = ""

class KnowledgeBaseConfig(BaseModel):
    """Configuration for a Knowledge Base tool with its own RAG settings"""
    collection_id: str = ""
    # Embedding settings
    embedding: KnowledgeBaseEmbeddingConfig = Field(default_factory=KnowledgeBaseEmbeddingConfig)
    # Vector DB settings
    vector_db: KnowledgeBaseVectorDBConfig = Field(default_factory=KnowledgeBaseVectorDBConfig)
    # RAG Settings
    chunk_size: int = 1000
    chunk_overlap: int = 200
    top_k: int = 5
    search_type: str = "hybrid"  # hybrid, semantic, keyword
    similarity_threshold: float = 0.7
    # Advanced
    reranking: str = "none"  # none, cohere, cross-encoder
    context_window: int = 4000
    include_metadata: bool = True
    auto_reindex: bool = False


class ToolConfiguration(BaseModel):
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    type: str  # document, website, api, knowledge, etc.
    name: str
    description: str = ""
    config: Dict[str, Any] = {}
    api_config: Optional[APIEndpointConfig] = None
    kb_config: Optional[KnowledgeBaseConfig] = None  # For knowledge base tools
    documents: List[Dict[str, Any]] = []
    scraped_pages: List[Dict[str, Any]] = []
    created_at: str = Field(default_factory=lambda: datetime.utcnow().isoformat())
    is_active: bool = True
    # Access Control
    owner_id: str = ""  # User who created/owns the tool
    access_type: str = "owner_only"  # 'owner_only', 'authenticated', 'specific_users', 'public'
    allowed_user_ids: List[str] = []  # Users who can view/use the tool
    allowed_group_ids: List[str] = []  # Groups who can view/use the tool
    can_edit_user_ids: List[str] = []  # Users who can edit the tool
    can_delete_user_ids: List[str] = []  # Users who can delete the tool
    can_execute_user_ids: List[str] = []  # Users who can use tool in agents


class TaskInstruction(BaseModel):
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    text: str


class TaskDefinition(BaseModel):
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    name: str
    description: str = ""
    instructions: List[TaskInstruction] = []


class AgentPersonality(BaseModel):
    tone: str = "professional"
    voice: str = "balanced"
    languages: List[str] = ["English"]
    traits: List[str] = []
    # New personality sliders (1-10 scale)
    creativity: int = 5      # 1=Factual only, 10=Highly creative
    length: int = 5          # 1=Ultra brief, 10=Exhaustive
    formality: int = 7       # 1=Very casual, 10=Executive formal
    empathy: int = 6         # 1=Direct/matter-of-fact, 10=Deeply caring
    proactivity: int = 5     # 1=Reactive only, 10=Proactive advisor
    confidence: int = 7      # 1=Very cautious, 10=Authoritative


class AgentGuardrails(BaseModel):
    """Safety and accuracy guardrails for the agent"""
    # Anti-hallucination settings
    anti_hallucination: bool = True
    cite_sources: bool = True
    admit_uncertainty: bool = True
    verify_facts: bool = True
    no_speculation: bool = False
    
    # Content boundaries
    avoid_topics: List[str] = []
    focus_topics: List[str] = []
    
    # Response limits
    max_length: str = "medium"  # short, medium, long, unlimited
    language: str = "user"  # user, en, ar, multi
    
    # Escalation rules
    escalate_angry: bool = True
    escalate_complex: bool = True
    escalate_request: bool = True
    escalate_sensitive: bool = False
    
    # PII protection
    pii_protection: bool = True
    mask_pii: bool = True
    no_store_pii: bool = True


class AgentData(BaseModel):
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    name: str
    icon: str = "🤖"
    goal: str = ""
    personality: AgentPersonality = Field(default_factory=AgentPersonality)
    guardrails: AgentGuardrails = Field(default_factory=AgentGuardrails)
    tasks: List[TaskDefinition] = []
    tool_ids: List[str] = []
    model_id: str = "gpt-4o"
    created_at: str = Field(default_factory=lambda: datetime.utcnow().isoformat())
    updated_at: str = Field(default_factory=lambda: datetime.utcnow().isoformat())
    status: str = "draft"
    is_active: bool = True
    # Memory - stores conversation summaries for continuity
    memory: List[Dict[str, Any]] = []  # [{timestamp, summary, key_facts}]
    memory_enabled: bool = True  # Whether to use memory across conversations
    # Ownership - who owns this agent
    owner_id: Optional[str] = None
    created_by: Optional[str] = None

    class Config:
        protected_namespaces = ()


class ConversationMessage(BaseModel):
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    role: str
    content: str
    timestamp: str = Field(default_factory=lambda: datetime.utcnow().isoformat())
    tool_calls: List[Dict[str, Any]] = []
    sources: List[Dict[str, Any]] = []


class ConversationAccessCache(BaseModel):
    """Cached access control for a conversation session"""
    user_id: str
    user_name: str = ""
    user_role: str = ""
    user_groups: List[str] = []
    accessible_task_names: List[str] = []
    denied_task_names: List[str] = []
    accessible_tool_ids: List[str] = []
    denied_tool_ids: List[str] = []
    loaded_at: str = Field(default_factory=lambda: datetime.utcnow().isoformat())


class Conversation(BaseModel):
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    agent_id: str
    user_id: Optional[str] = None  # Owner of this conversation (for privacy)
    title: str = "New Conversation"
    messages: List[ConversationMessage] = []
    created_at: str = Field(default_factory=lambda: datetime.utcnow().isoformat())
    updated_at: str = Field(default_factory=lambda: datetime.utcnow().isoformat())
    # Cached access control - loaded once when conversation starts
    access_cache: Optional[ConversationAccessCache] = None


class Document(BaseModel):
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    tool_id: str
    filename: str
    original_name: str
    file_type: str
    file_size: int
    content: str = ""
    chunks: List[Dict[str, Any]] = []
    status: str = "pending"
    error_message: str = ""
    created_at: str = Field(default_factory=lambda: datetime.utcnow().isoformat())


class ScrapedPage(BaseModel):
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    tool_id: str
    url: str
    title: str = ""
    content: str = ""
    chunks: List[Dict[str, Any]] = []
    status: str = "pending"
    error_message: str = ""
    created_at: str = Field(default_factory=lambda: datetime.utcnow().isoformat())


# Request Models
class ChatRequest(BaseModel):
    message: str
    conversation_id: Optional[str] = None
    timezone: Optional[str] = None  # e.g., "Africa/Cairo", "America/New_York"


class ChatResponse(BaseModel):
    response: str
    conversation_id: str
    tool_calls: List[Dict[str, Any]] = []
    sources: List[Dict[str, Any]] = []
    formatted: bool = True


class CreateAgentRequest(BaseModel):
    name: str
    goal: str
    personality: Optional[Dict[str, Any]] = None
    guardrails: Optional[Dict[str, Any]] = None
    tasks: List[Dict[str, Any]] = []
    tool_ids: List[str] = []
    model_id: str = "gpt-4o"
    status: str = "draft"

    class Config:
        protected_namespaces = ()


class UpdateAgentRequest(BaseModel):
    name: Optional[str] = None
    goal: Optional[str] = None
    personality: Optional[Dict[str, Any]] = None
    guardrails: Optional[Dict[str, Any]] = None
    tasks: Optional[List[Dict[str, Any]]] = None
    tool_ids: Optional[List[str]] = None
    model_id: Optional[str] = None
    icon: Optional[str] = None
    status: Optional[str] = None
    is_active: Optional[bool] = None  # For publishing/unpublishing agents
    is_published: Optional[bool] = None  # Alternative field for publishing

    class Config:
        protected_namespaces = ()


class CreateToolRequest(BaseModel):
    type: str
    name: str
    description: str = ""
    config: Dict[str, Any] = {}
    api_config: Optional[Dict[str, Any]] = None
    # Access Control
    access_type: str = "owner_only"  # 'owner_only', 'authenticated', 'specific_users', 'public'
    allowed_user_ids: List[str] = []
    allowed_group_ids: List[str] = []
    can_edit_user_ids: List[str] = []
    can_delete_user_ids: List[str] = []
    can_execute_user_ids: List[str] = []


class ScrapeRequest(BaseModel):
    url: str
    recursive: bool = False
    max_pages: int = 10


class APITestRequest(BaseModel):
    base_url: str
    http_method: str = "GET"
    endpoint_path: str = ""
    auth_type: str = "none"
    auth_value: str = ""
    api_key_name: str = "X-API-Key"
    api_key_location: str = "header"
    headers: Dict[str, str] = {}
    parameters: Dict[str, Any] = {}


# ============================================================================
# Document Processing
# ============================================================================

class DocumentProcessor:
    @staticmethod
    async def extract_text(file_path: str, file_type: str) -> str:
        try:
            if file_type in ['txt', 'md', 'text']:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read()
            elif file_type == 'pdf':
                try:
                    import fitz
                    doc = fitz.open(file_path)
                    text = "".join([page.get_text() for page in doc])
                    doc.close()
                    return text
                except:
                    return "[PDF extraction failed]"
            elif file_type in ['docx', 'doc']:
                try:
                    from docx import Document as DocxDocument
                    doc = DocxDocument(file_path)
                    return "\n".join([p.text for p in doc.paragraphs])
                except:
                    return "[DOCX extraction failed]"
            elif file_type in ['xlsx', 'xls']:
                try:
                    import openpyxl
                    wb = openpyxl.load_workbook(file_path, read_only=True)
                    text = ""
                    for sheet in wb.sheetnames:
                        ws = wb[sheet]
                        text += f"\n=== Sheet: {sheet} ===\n"
                        for row in ws.iter_rows(values_only=True):
                            text += " | ".join([str(c) if c else "" for c in row]) + "\n"
                    return text
                except:
                    return "[Excel extraction failed]"
            elif file_type in ['pptx', 'ppt']:
                try:
                    from pptx import Presentation
                    prs = Presentation(file_path)
                    text = ""
                    for i, slide in enumerate(prs.slides):
                        text += f"\n=== Slide {i+1} ===\n"
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text += shape.text + "\n"
                    return text
                except:
                    return "[PowerPoint extraction failed]"
            elif file_type == 'csv':
                import csv
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    reader = csv.reader(f)
                    return "\n".join([" | ".join(row) for row in reader])
            else:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read()
        except Exception as e:
            return f"[Error: {str(e)}]"
    
    @staticmethod
    def chunk_text(text: str, chunk_size: int = 1000, overlap: int = 200) -> List[Dict]:
        chunks = []
        start = 0
        chunk_id = 0
        while start < len(text):
            end = start + chunk_size
            chunk_text = text[start:end]
            if end < len(text):
                for sep in ['\n\n', '\n', '. ', ' ']:
                    idx = chunk_text.rfind(sep)
                    if idx > chunk_size * 0.5:
                        chunk_text = chunk_text[:idx + len(sep)]
                        end = start + idx + len(sep)
                        break
            if chunk_text.strip():
                chunks.append({"id": chunk_id, "text": chunk_text.strip(), "start": start, "end": end})
                chunk_id += 1
            start = end - overlap
            if start < 0: start = 0
            if start >= len(text): break
        return chunks


# ============================================================================
# Website Scraper with Playwright Support for Dynamic Sites
# ============================================================================

# List of known dynamic sites that require JavaScript rendering
DYNAMIC_SITES = [
    'oracle.com',
    'azure.microsoft.com',
    'cloud.google.com',
    'aws.amazon.com',
    'vercel.com',
    'netlify.com',
    'render.com',
    'railway.app',
    'supabase.com',
    'planetscale.com',
    'firebase.google.com',
    'digitalocean.com',
    'linode.com',
    'vultr.com',
    'heroku.com',
    'cloudflare.com',
    'fastly.com',
    'akamai.com',
]

def should_use_playwright(url: str) -> bool:
    """Determine if a URL needs JavaScript rendering based on known dynamic sites."""
    domain = urlparse(url).netloc.lower()
    return any(site in domain for site in DYNAMIC_SITES)


async def fetch_with_playwright(url: str, wait_time: int = 10000, scroll: bool = True) -> str:
    """
    Fetch page content using Playwright headless browser.
    This renders JavaScript and returns the fully loaded HTML.
    
    Args:
        url: The URL to fetch
        wait_time: Time to wait for JS to render (ms)
        scroll: Whether to scroll the page to trigger lazy loading
    
    Returns:
        The fully rendered HTML content
    """
    if not PLAYWRIGHT_AVAILABLE:
        raise Exception("Playwright not available")
    
    async with async_playwright() as p:
        print(f"🎭 [1/7] Launching Chromium browser...")
        browser = await p.chromium.launch(
            headless=True,
            args=['--no-sandbox', '--disable-dev-shm-usage', '--disable-gpu']
        )
        
        print(f"🎭 [2/7] Creating browser context...")
        context = await browser.new_context(
            user_agent='Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36',
            viewport={'width': 1920, 'height': 1080},
            java_script_enabled=True
        )
        page = await context.new_page()
        
        # Block unnecessary resources to speed up loading
        await page.route("**/*.{png,jpg,jpeg,gif,svg,ico,woff,woff2,ttf}", lambda route: route.abort())
        await page.route("**/analytics**", lambda route: route.abort())
        await page.route("**/tracking**", lambda route: route.abort())
        await page.route("**/ads**", lambda route: route.abort())
        
        print(f"🎭 [3/7] Navigating to {url[:60]}...")
        # Use 'domcontentloaded' instead of 'networkidle' - much faster for heavy sites
        try:
            await page.goto(url, wait_until='domcontentloaded', timeout=120000)
        except Exception as e:
            print(f"⚠️ Initial load warning: {str(e)[:100]}, continuing anyway...")
        
        print(f"🎭 [4/7] Waiting for JavaScript to render ({wait_time/1000}s)...")
        await page.wait_for_timeout(wait_time)
        
        # Try to click "Expand All" or similar buttons if they exist
        try:
            expand_buttons = await page.query_selector_all('button:has-text("Expand"), button:has-text("Show All"), button:has-text("View All"), [aria-expanded="false"]')
            if expand_buttons:
                print(f"🎭 [4.5/7] Found {len(expand_buttons)} expandable sections, clicking...")
                for btn in expand_buttons[:20]:
                    try:
                        await btn.click()
                        await page.wait_for_timeout(300)
                    except:
                        pass
        except:
            pass
        
        if scroll:
            print(f"🎭 [5/7] Scrolling through entire page to load all content...")
            # Get total page height
            total_height = await page.evaluate('document.body.scrollHeight')
            print(f"📏 Initial page height: {total_height:,}px")
            
            # Scroll in smaller increments for long pages (every 800px)
            scroll_step = 800
            current_position = 0
            scroll_count = 0
            max_scrolls = 100  # Increased safety limit for very long pages
            
            while current_position < total_height and scroll_count < max_scrolls:
                current_position += scroll_step
                try:
                    await page.evaluate(f'window.scrollTo(0, {current_position})')
                    await page.wait_for_timeout(400)  # Slightly faster
                except:
                    pass
                scroll_count += 1
                
                # Check if page height increased (more content loaded)
                try:
                    new_height = await page.evaluate('document.body.scrollHeight')
                    if new_height > total_height:
                        total_height = new_height
                except:
                    pass
                
                if scroll_count % 15 == 0:
                    print(f"🎭 [5/7] Scrolled {scroll_count} times, position: {current_position:,}px / {total_height:,}px")
            
            print(f"🎭 [5/7] Completed {scroll_count} scroll operations, final height: {total_height:,}px")
            
            # Scroll back to top
            await page.evaluate('window.scrollTo(0, 0)')
            await page.wait_for_timeout(500)
        
        # Final wait for any remaining lazy-loaded content
        print(f"🎭 [6/7] Final wait for lazy-loaded content...")
        await page.wait_for_timeout(3000)
        
        print(f"🎭 [7/7] Extracting rendered HTML...")
        html_content = await page.content()
        
        await browser.close()
        
        print(f"✅ Playwright fetched {len(html_content):,} bytes")
        return html_content


class WebsiteScraper:
    def __init__(self, base_url: str, max_pages: int = 10):
        self.base_url = base_url
        self.max_pages = max_pages
        self.visited: Set[str] = set()
        self.pages: List[Dict] = []
        self.use_playwright = PLAYWRIGHT_AVAILABLE and should_use_playwright(base_url)
        
        if self.use_playwright:
            print(f"🎭 Detected dynamic site - will use Playwright for JavaScript rendering")
    
    async def scrape(self, recursive: bool = False) -> List[Dict]:
        await self._scrape_page(self.base_url, recursive)
        return self.pages
    
    async def _scrape_page(self, url: str, recursive: bool):
        if url in self.visited or len(self.pages) >= self.max_pages:
            return
        self.visited.add(url)
        
        try:
            # Use Playwright for dynamic sites, httpx for static
            if self.use_playwright:
                print(f"🎭 Scraping with Playwright: {url[:60]}...")
                html = await fetch_with_playwright(url)
            else:
                print(f"📄 Scraping with httpx: {url[:60]}...")
                async with httpx.AsyncClient(timeout=30.0, follow_redirects=True) as client:
                    response = await client.get(url, headers={"User-Agent": "Mozilla/5.0 (compatible; AgentForge/1.0)"})
                    if response.status_code != 200:
                        return
                    html = response.text
            
            title_match = re.search(r'<title[^>]*>([^<]+)</title>', html, re.IGNORECASE)
            title = title_match.group(1).strip() if title_match else url
            text = self._extract_text(html)
            
            if text.strip():
                # Extract tables specifically for pricing data
                tables = self._extract_tables(html)
                full_content = text
                if tables:
                    full_content += "\n\n" + tables
                    print(f"📊 Extracted {tables.count('[TABLE]')} tables with pricing data")
                
                self.pages.append({"url": url, "title": title, "content": full_content})
                print(f"✅ Scraped: {title[:50]}... ({len(full_content):,} chars)")
            
            if recursive and len(self.pages) < self.max_pages:
                links = self._extract_links(html, url)
                for link in links[:10]:
                    if len(self.pages) >= self.max_pages:
                        break
                    await self._scrape_page(link, recursive)
                    
        except Exception as e:
            print(f"❌ Scrape error for {url}: {e}")
    
    def _extract_text(self, html: str) -> str:
        """Extract clean text from HTML."""
        html = re.sub(r'<script[^>]*>.*?</script>', '', html, flags=re.DOTALL | re.IGNORECASE)
        html = re.sub(r'<style[^>]*>.*?</style>', '', html, flags=re.DOTALL | re.IGNORECASE)
        html = re.sub(r'<nav[^>]*>.*?</nav>', '', html, flags=re.DOTALL | re.IGNORECASE)
        html = re.sub(r'<footer[^>]*>.*?</footer>', '', html, flags=re.DOTALL | re.IGNORECASE)
        html = re.sub(r'<header[^>]*>.*?</header>', '', html, flags=re.DOTALL | re.IGNORECASE)
        text = re.sub(r'<[^>]+>', ' ', html)
        text = re.sub(r'\s+', ' ', text)
        return text.strip()
    
    def _extract_tables(self, html: str) -> str:
        """Extract tables from HTML preserving structure (important for pricing data)."""
        tables_text = []
        table_pattern = r'<table[^>]*>(.*?)</table>'
        tables = re.findall(table_pattern, html, re.DOTALL | re.IGNORECASE)
        
        for table in tables:
            rows = []
            # Extract rows
            row_pattern = r'<tr[^>]*>(.*?)</tr>'
            for row_match in re.finditer(row_pattern, table, re.DOTALL | re.IGNORECASE):
                row_content = row_match.group(1)
                # Extract cells (th or td)
                cells = re.findall(r'<t[hd][^>]*>(.*?)</t[hd]>', row_content, re.DOTALL | re.IGNORECASE)
                # Clean cell content
                clean_cells = []
                for cell in cells:
                    cell_text = re.sub(r'<[^>]+>', ' ', cell)
                    cell_text = re.sub(r'\s+', ' ', cell_text).strip()
                    if cell_text:
                        clean_cells.append(cell_text)
                if clean_cells:
                    rows.append(' | '.join(clean_cells))
            
            if rows:
                tables_text.append("[TABLE]\n" + '\n'.join(rows) + "\n[/TABLE]")
        
        return '\n\n'.join(tables_text)
    
    def _extract_links(self, html: str, base_url: str) -> List[str]:
        links = []
        base_domain = urlparse(base_url).netloc
        for match in re.finditer(r'href=["\']([^"\']+)["\']', html, re.IGNORECASE):
            href = match.group(1)
            if href.startswith('#') or href.startswith('javascript:') or href.startswith('mailto:'):
                continue
            full_url = urljoin(base_url, href)
            parsed = urlparse(full_url)
            if parsed.netloc == base_domain and full_url not in self.visited:
                links.append(full_url)
        return list(set(links))


# ============================================================================
# OpenAPI Parser
# ============================================================================

class OpenAPIParser:
    @staticmethod
    def parse(content: str, file_type: str) -> Dict[str, Any]:
        try:
            if file_type in ['yaml', 'yml']:
                spec = yaml.safe_load(content)
            else:
                spec = json.loads(content)
            result = {"base_url": "", "endpoints": []}
            if 'servers' in spec and spec['servers']:
                result["base_url"] = spec['servers'][0].get('url', '')
            elif 'host' in spec:
                scheme = spec.get('schemes', ['https'])[0]
                base_path = spec.get('basePath', '')
                result["base_url"] = f"{scheme}://{spec['host']}{base_path}"
            if 'paths' in spec:
                for path, methods in spec['paths'].items():
                    for method, details in methods.items():
                        if method.upper() in ['GET', 'POST', 'PUT', 'DELETE', 'PATCH']:
                            endpoint = {"path": path, "method": method.upper(), "summary": details.get('summary', ''), "description": details.get('description', ''), "parameters": []}
                            params = details.get('parameters', [])
                            for param in params:
                                endpoint["parameters"].append({
                                    "name": param.get('name', ''),
                                    "description": param.get('description', ''),
                                    "location": param.get('in', 'query'),
                                    "required": param.get('required', False),
                                    "data_type": param.get('schema', {}).get('type', 'string') if 'schema' in param else param.get('type', 'string')
                                })
                            result["endpoints"].append(endpoint)
            return result
        except Exception as e:
            return {"error": str(e)}


# ============================================================================
# Demo Kit System - Tools, Knowledge Bases, and Demo Assets
# ============================================================================

class DemoAPITool(BaseModel):
    """A mock API tool for demo purposes"""
    id: str = ""
    name: str
    description: str
    method: str = "GET"
    endpoint: str
    parameters: List[Dict[str, Any]] = []
    sample_request: Optional[Dict[str, Any]] = None
    sample_response: Dict[str, Any] = {}
    created_at: str = ""

class DemoKnowledgeBase(BaseModel):
    """A knowledge base document for demo purposes"""
    id: str = ""
    name: str
    description: str
    content: str  # The actual content
    sections: List[Dict[str, str]] = []  # {title, content}
    created_at: str = ""

class DemoAsset(BaseModel):
    """A sample document/image for demo testing (invoices, receipts, etc.)"""
    id: str = ""
    name: str
    description: str
    asset_type: str  # "invoice", "receipt", "document", "image"
    file_path: Optional[str] = None  # Path to generated file
    file_url: Optional[str] = None  # URL to download
    preview_data: Optional[str] = None  # Base64 or text preview
    metadata: Dict[str, Any] = {}
    created_at: str = ""

class DemoKit(BaseModel):
    """A complete demo kit with APIs, Knowledge Bases, and Assets"""
    id: str = ""
    name: str
    description: str
    original_prompt: str  # What the user asked for
    apis: List[DemoAPITool] = []
    knowledge_bases: List[DemoKnowledgeBase] = []
    assets: List[DemoAsset] = []
    created_at: str = ""
    updated_at: str = ""

class GenerateDemoKitRequest(BaseModel):
    """Request to generate a demo kit from natural language"""
    description: str
    kit_name: Optional[str] = None


# ============================================================================
# Application State
# ============================================================================

class AppState:
    def __init__(self):
        self.agents: Dict[str, AgentData] = {}
        self.tools: Dict[str, ToolConfiguration] = {}
        self.documents: Dict[str, Document] = {}
        self.scraped_pages: Dict[str, ScrapedPage] = {}
        self.conversations: Dict[str, Conversation] = {}
        self.document_chunks: List[Dict] = []
        self.settings: SystemSettings = SystemSettings()
        # Demo Kit System
        self.demo_kits: Dict[str, DemoKit] = {}
        # OAuth Tokens Storage (tool_id -> tokens)
        self.oauth_tokens: Dict[str, Dict[str, Any]] = {}
        # Integration Settings (Google, Microsoft OAuth credentials)
        self.integrations: Dict[str, Dict[str, str]] = {}
        # Providers (initialized lazily)
        self._llm_provider: Optional[BaseLLMProvider] = None
        self._embedding_provider: Optional[BaseEmbeddingProvider] = None
        self._vector_db: Optional[BaseVectorDB] = None
    
    def get_llm_provider(self) -> BaseLLMProvider:
        if self._llm_provider is None:
            llm_config = self.settings.llm
            
            # If provider is Google and no api_key, use gemini_key from google settings
            if llm_config.provider == LLMProvider.GOOGLE and not llm_config.api_key:
                if self.settings.google and self.settings.google.gemini_key:
                    llm_config = LLMConfig(
                        provider=llm_config.provider,
                        model=llm_config.model,
                        api_key=self.settings.google.gemini_key,
                        api_base=llm_config.api_base,
                        temperature=llm_config.temperature,
                        max_tokens=llm_config.max_tokens
                    )
            
            # If no api_key in llm config, try to find it in llm_providers
            if not llm_config.api_key and self.settings.llm_providers:
                provider_data = next(
                    (p for p in self.settings.llm_providers if p.provider == llm_config.provider.value),
                    None
                )
                if provider_data:
                    llm_config = LLMConfig(
                        provider=llm_config.provider,
                        model=llm_config.model,
                        api_key=provider_data.api_key,
                        api_base=provider_data.api_base or llm_config.api_base,
                        temperature=llm_config.temperature,
                        max_tokens=llm_config.max_tokens
                    )
            
            self._llm_provider = ProviderFactory.get_llm_provider(llm_config)
        return self._llm_provider
    
    def get_llm_provider_for(self, provider_name: str, model: str = None) -> BaseLLMProvider:
        """Get LLM provider for a specific provider name (useful for multi-provider support)"""
        # Find provider in configured providers
        provider_data = next(
            (p for p in self.settings.llm_providers if p.provider == provider_name),
            None
        )
        
        if not provider_data:
            raise ValueError(f"Provider {provider_name} is not configured")
        
        # Create config
        try:
            provider_enum = LLMProvider(provider_name)
        except ValueError:
            provider_enum = LLMProvider.CUSTOM
        
        config = LLMConfig(
            provider=provider_enum,
            model=model or (provider_data.models[0] if provider_data.models else ""),
            api_key=provider_data.api_key,
            api_base=provider_data.api_base,
        )
        
        return ProviderFactory.get_llm_provider(config)
    
    def get_embedding_provider(self) -> BaseEmbeddingProvider:
        if self._embedding_provider is None:
            self._embedding_provider = ProviderFactory.get_embedding_provider(self.settings.embedding)
        return self._embedding_provider
    
    def get_vector_db(self) -> BaseVectorDB:
        if self._vector_db is None:
            embedding_dim = self.get_embedding_provider().get_dimensions()
            self._vector_db = ProviderFactory.get_vector_db(self.settings.vector_db, embedding_dim)
        return self._vector_db
    
    def reset_providers(self):
        """Reset providers when settings change"""
        self._llm_provider = None
        self._embedding_provider = None
        self._vector_db = None
        
    def save_to_disk(self):
        data_dir = os.environ.get("DATA_PATH", "data")
        os.makedirs(data_dir, exist_ok=True)
        
        # Save agents to database first
        try:
            from database.services import AgentService
            # Get default org_id and owner_id
            org_id = "org_default"
            owner_id = None
            created_by = None
            
            # Try to get owner_id from security_state if available
            try:
                if SECURITY_AVAILABLE and security_state.users:
                    # Use first user as default owner (or super admin if available)
                    super_admin = next((u for u in security_state.users.values() if u.role_ids and any(r == "super_admin" or "super" in r.lower() for r in u.role_ids)), None)
                    if super_admin:
                        owner_id = super_admin.id
                        created_by = super_admin.id
                    else:
                        first_user = next(iter(security_state.users.values()), None)
                        if first_user:
                            owner_id = first_user.id
                            created_by = first_user.id
            except Exception:
                pass
            
            # Save each agent to database
            for agent_id, agent in self.agents.items():
                try:
                    agent_dict = agent.dict()
                    # Get org_id and owner_id from agent_dict if available (from database)
                    agent_org_id = agent_dict.get('org_id', org_id)
                    agent_owner_id = agent_dict.get('owner_id', owner_id)
                    agent_created_by = agent_dict.get('created_by', created_by or agent_owner_id)
                    
                    # Remove extra fields that AgentData doesn't have
                    agent_dict.pop('org_id', None)
                    agent_dict.pop('owner_id', None)
                    agent_dict.pop('created_by', None)
                    agent_dict.pop('shared_with_user_ids', None)
                    agent_dict.pop('shared_with_role_ids', None)
                    agent_dict.pop('usage_count', None)
                    agent_dict.pop('last_used_at', None)
                    agent_dict.pop('context_window', None)
                    agent_dict.pop('version', None)
                    agent_dict.pop('parent_version_id', None)
                    agent_dict.pop('published_at', None)
                    agent_dict.pop('extra_metadata', None)
                    
                    # Save to database
                    AgentService.save_agent(
                        agent_dict,
                        org_id=agent_org_id,
                        owner_id=agent_owner_id or "00000000-0000-0000-0000-000000000000",  # Fallback UUID
                        created_by=agent_created_by or agent_owner_id or "00000000-0000-0000-0000-000000000000",
                        updated_by=agent_created_by or agent_owner_id
                    )
                except Exception as e:
                    print(f"⚠️  [DATABASE ERROR] Failed to save agent '{agent.name}' (ID: {agent_id[:8]}...): {e}")
                    import traceback
                    traceback.print_exc()
                    continue
            
            # Also save to JSON as backup
            with open(os.path.join(data_dir, "agents.json"), "w") as f:
                json.dump({k: v.dict() for k, v in self.agents.items()}, f, indent=2)
        except Exception as db_error:
            print(f"❌ [DATABASE ERROR] Failed to save agents: {type(db_error).__name__}: {str(db_error)}")
            import traceback
            traceback.print_exc()
            # Fallback to JSON only
            with open(os.path.join(data_dir, "agents.json"), "w") as f:
                json.dump({k: v.dict() for k, v in self.agents.items()}, f, indent=2)
        
        # Save tools to database
        try:
            from database.services import ToolService
            for tool_id, tool in self.tools.items():
                tool_dict = tool.dict()
                # Convert api_config to dict if it's an object
                if hasattr(tool, 'api_config') and tool.api_config:
                    tool_dict['api_config'] = tool.api_config.dict() if hasattr(tool.api_config, 'dict') else tool.api_config
                ToolService.create_tool(tool_dict, "org_default", tool_dict.get('owner_id', 'system'))
        except Exception as e:
            print(f"⚠️  [DATABASE] Failed to save tools: {e}")
            # Fallback to JSON
            with open(os.path.join(data_dir, "tools.json"), "w") as f:
                json.dump({k: v.dict() for k, v in self.tools.items()}, f, indent=2)
        
        # Save other data to JSON (documents, scraped_pages, conversations)
        for name, data in [("documents.json", self.documents), ("scraped_pages.json", self.scraped_pages), ("conversations.json", self.conversations)]:
            with open(os.path.join(data_dir, name), "w") as f:
                json.dump({k: v.dict() for k, v in data.items()}, f, indent=2)
        with open(os.path.join(data_dir, "chunks_index.json"), "w") as f:
            json.dump(self.document_chunks, f)
        # Save settings
        with open(os.path.join(data_dir, "settings.json"), "w") as f:
            json.dump(self.settings.dict(), f, indent=2)
        # Save integrations (OAuth credentials)
        with open(os.path.join(data_dir, "integrations.json"), "w") as f:
            json.dump(self.integrations, f, indent=2)
    
    def load_from_disk(self):
        data_dir = os.environ.get("DATA_PATH", "data")
        
        # Load integrations first (OAuth credentials)
        integrations_path = os.path.join(data_dir, "integrations.json")
        if os.path.exists(integrations_path):
            try:
                with open(integrations_path) as f:
                    self.integrations = json.load(f)
                    # Set environment variables from saved integrations
                    for provider, creds in self.integrations.items():
                        if provider == 'google':
                            if creds.get('client_id'):
                                os.environ['GOOGLE_CLIENT_ID'] = creds['client_id']
                            if creds.get('client_secret'):
                                os.environ['GOOGLE_CLIENT_SECRET'] = creds['client_secret']
                        elif provider == 'microsoft':
                            if creds.get('client_id'):
                                os.environ['MICROSOFT_CLIENT_ID'] = creds['client_id']
                            if creds.get('client_secret'):
                                os.environ['MICROSOFT_CLIENT_SECRET'] = creds['client_secret']
                    print(f"✅ Loaded integrations: {list(self.integrations.keys())}")
            except Exception as e:
                print(f"⚠️ Error loading integrations: {e}")
        
        # Load settings from database first
        db_settings_loaded = False
        try:
            from database.services import SystemSettingsService
            db_settings = SystemSettingsService.get_system_setting("system_settings")
            if db_settings:
                self.settings = SystemSettings(**db_settings)
                llm_provider = self.settings.llm.provider.value if hasattr(self.settings.llm.provider, 'value') else str(self.settings.llm.provider)
                vector_db_provider = self.settings.vector_db.provider.value if hasattr(self.settings.vector_db.provider, 'value') else str(self.settings.vector_db.provider)
                db_settings_loaded = True
        except Exception as db_error:
            print(f"❌ [DATABASE ERROR] Failed to load platform settings: {type(db_error).__name__}: {str(db_error)}")
            import traceback
            traceback.print_exc()
            print("📂 Loading settings from files (database unavailable)")
        
        # Load settings from JSON only if database loading failed
        if not db_settings_loaded:
            settings_path = os.path.join(data_dir, "settings.json")
            if os.path.exists(settings_path):
                try:
                    with open(settings_path) as f:
                        settings_data = json.load(f)
                        self.settings = SystemSettings(**settings_data)
                        llm_provider = self.settings.llm.provider.value if hasattr(self.settings.llm.provider, 'value') else str(self.settings.llm.provider)
                        vector_db_provider = self.settings.vector_db.provider.value if hasattr(self.settings.vector_db.provider, 'value') else str(self.settings.vector_db.provider)
                        print(f"✅ Loaded settings from file: LLM={llm_provider}, VectorDB={vector_db_provider}")
                        if self.settings.llm_providers:
                            print(f"✅ Loaded {len(self.settings.llm_providers)} LLM providers from file: {[p.name for p in self.settings.llm_providers]}")
                except Exception as e:
                    print(f"⚠️ Error loading settings: {e}, using defaults")
                    import traceback
                    traceback.print_exc()
        
        # Load agents from database first
        db_agents_loaded = False
        try:
            from database.services import AgentService
            db_agents = AgentService.get_all_agents("org_default")
            if db_agents:
                for agent_dict in db_agents:
                    try:
                        # Convert dictionary to AgentData - include owner_id and created_by
                        # These fields are now in AgentData model
                        agent_data = AgentData(**{
                            k: v for k, v in agent_dict.items() 
                            if k in AgentData.__fields__ or k in ('owner_id', 'created_by')
                        })
                        self.agents[agent_data.id] = agent_data
                    except Exception as e:
                        print(f"⚠️  Error converting agent from database: {e}")
                        import traceback
                        traceback.print_exc()
                        continue
                db_agents_loaded = True
                print(f"✅ Loaded {len(self.agents)} agents from database")
        except Exception as db_error:
            print(f"❌ [DATABASE ERROR] Failed to load agents: {type(db_error).__name__}: {str(db_error)}")
            import traceback
            traceback.print_exc()
            print("📂 Loading agents from files (database unavailable)")
        
        # Load agents from JSON only if database loading failed
        if not db_agents_loaded:
            agents_path = os.path.join(data_dir, "agents.json")
            if os.path.exists(agents_path):
                try:
                    with open(agents_path) as f:
                        data = json.load(f)
                        for k, v in data.items():
                            try:
                                if 'tasks' in v:
                                    tasks = []
                                    for t in v.get('tasks', []):
                                        instructions = []
                                        for i in t.get('instructions', []):
                                            if isinstance(i, dict): instructions.append(TaskInstruction(**i))
                                            elif isinstance(i, str): instructions.append(TaskInstruction(text=i))
                                        tasks.append(TaskDefinition(id=t.get('id', str(uuid.uuid4())), name=t.get('name', ''), description=t.get('description', ''), instructions=instructions))
                                    v['tasks'] = tasks
                                    if 'personality' in v and isinstance(v['personality'], dict):
                                        v['personality'] = AgentPersonality(**v['personality'])
                                self.agents[k] = AgentData(**v)
                            except Exception as e:
                                print(f"Error loading agents.json item {k}: {e}")
                except Exception as e:
                    print(f"Error loading agents.json: {e}")
        
        # Load tools from database first
        db_tools_loaded = False
        try:
            from database.services import ToolService
            db_tools = ToolService.get_all_tools()
            if db_tools:
                for tool_dict in db_tools:
                    try:
                        # Convert api_config if present
                        if 'api_config' in tool_dict and tool_dict['api_config']:
                            api_cfg = tool_dict['api_config']
                            if 'input_parameters' in api_cfg:
                                params = [APIInputParameter(**p) for p in api_cfg['input_parameters']]
                                api_cfg['input_parameters'] = params
                            tool_dict['api_config'] = APIEndpointConfig(**api_cfg)
                        
                        # Migration: Set defaults for access control on existing tools
                        if not tool_dict.get('owner_id'):
                            tool_dict['owner_id'] = tool_dict.get('created_by', 'system')
                        if not tool_dict.get('access_type'):
                            # Default to authenticated so existing tools work for logged-in users
                            tool_dict['access_type'] = 'authenticated'
                        
                        # Debug: Log access control fields from database
                        print(f"   📦 Loading tool '{tool_dict.get('name')}': access_type={tool_dict.get('access_type')}, allowed_users={tool_dict.get('allowed_user_ids')}, allowed_groups={tool_dict.get('allowed_group_ids')}")
                        
                        tool = ToolConfiguration(**{k: v for k, v in tool_dict.items() if k in ToolConfiguration.__fields__})
                        self.tools[tool.id] = tool
                    except Exception as e:
                        print(f"⚠️  Error converting tool from database: {e}")
                db_tools_loaded = True
                print(f"✅ Loaded {len(self.tools)} tools from database")
        except Exception as db_error:
            print(f"⚠️  [DATABASE] Failed to load tools: {db_error}")
        
        # Load other data from JSON (tools only if DB failed, documents, etc.)
        for name, cls, container in [("documents.json", Document, self.documents), ("scraped_pages.json", ScrapedPage, self.scraped_pages), ("conversations.json", Conversation, self.conversations)] + ([] if db_tools_loaded else [("tools.json", ToolConfiguration, self.tools)]):
            path = os.path.join(data_dir, name)
            if os.path.exists(path):
                try:
                    with open(path) as f:
                        data = json.load(f)
                        for k, v in data.items():
                            try:
                                if cls == ToolConfiguration and 'api_config' in v and v['api_config']:
                                    api_cfg = v['api_config']
                                    if 'input_parameters' in api_cfg:
                                        params = [APIInputParameter(**p) for p in api_cfg['input_parameters']]
                                        api_cfg['input_parameters'] = params
                                    v['api_config'] = APIEndpointConfig(**api_cfg)
                                
                                # Migration: Set defaults for access control on existing tools (JSON loaded)
                                if cls == ToolConfiguration:
                                    if not v.get('owner_id'):
                                        v['owner_id'] = v.get('created_by', 'system')
                                    if not v.get('access_type'):
                                        v['access_type'] = 'authenticated'
                                
                                container[k] = cls(**v)
                            except Exception as e:
                                print(f"Error loading {name} item {k}: {e}")
                except Exception as e:
                    print(f"Error loading {name}: {e}")
        chunks_path = os.path.join(data_dir, "chunks_index.json")
        if os.path.exists(chunks_path):
            with open(chunks_path) as f:
                self.document_chunks = json.load(f)


app_state = AppState()


def resolve_tool_id(tid: str) -> Optional['ToolConfiguration']:
    """Resolve a tool ID to its configuration, handling prefixed IDs"""
    # Try direct lookup first
    if tid in app_state.tools:
        return app_state.tools[tid]
    
    # Try without prefix (api:xxx, kb:xxx)
    clean_tid = tid.replace('api:', '').replace('kb:', '')
    if clean_tid in app_state.tools:
        return app_state.tools[clean_tid]
    
    return None


def get_agent_tools(agent: 'AgentData') -> List['ToolConfiguration']:
    """Get all tools for an agent, resolving prefixed IDs"""
    tools = []
    for tid in agent.tool_ids:
        tool = resolve_tool_id(tid)
        if tool:
            tools.append(tool)
    return tools


async def search_knowledge_base(query: str, kb_tool: ToolConfiguration, top_k: int = None) -> List[Dict]:
    """Search a specific Knowledge Base tool with its own RAG configuration"""
    try:
        kb_config = kb_tool.config
        collection_id = kb_config.get('collection_id', kb_tool.id)
        search_type = kb_config.get('search_type', 'hybrid')
        top_k = top_k or kb_config.get('top_k', 5)
        similarity_threshold = kb_config.get('similarity_threshold', 0.7)
        
        # If keyword only search
        if search_type == 'keyword':
            return search_documents_keyword(query, [kb_tool.id], top_k)
        
        # Get KB-specific or global embedding provider
        emb_config = kb_config.get('embedding', {})
        embedding_provider = ProviderFactory.get_kb_embedding_provider(
            emb_config, 
            app_state.settings.embedding
        )
        
        # Get query embedding
        query_embedding = (await embedding_provider.embed([query]))[0]
        
        # Get KB-specific vector DB with its collection
        vdb_config = kb_config.get('vector_db', {})
        embedding_dim = embedding_provider.get_dimensions()
        vector_db = ProviderFactory.get_kb_vector_db(
            vdb_config,
            app_state.settings.vector_db,
            collection_id,
            embedding_dim
        )
        
        # Search
        results = await vector_db.search(
            query_embedding=query_embedding,
            top_k=top_k * 2 if search_type == 'hybrid' else top_k,  # Get more for hybrid
            filter={"tool_id": kb_tool.id},
            query_text=query
        )
        
        # Filter by similarity threshold
        results = [r for r in results if r.get('score', 0) >= similarity_threshold]
        
        # If hybrid, combine with keyword search
        if search_type == 'hybrid':
            keyword_results = search_documents_keyword(query, [kb_tool.id], top_k)
            # Merge and dedupe
            seen_texts = set(r.get('text', '')[:100] for r in results)
            for kr in keyword_results:
                if kr.get('text', '')[:100] not in seen_texts:
                    results.append(kr)
            results = results[:top_k]
        
        return results[:top_k]
        
    except Exception as e:
        print(f"⚠️ KB search failed for {kb_tool.name}: {e}, falling back to keyword")
        return search_documents_keyword(query, [kb_tool.id], top_k or 5)


async def search_documents_rag(query: str, tool_ids: List[str] = None, top_k: int = 5) -> List[Dict]:
    """RAG-based search - routes to KB-specific search for knowledge tools"""
    try:
        all_results = []
        kb_tool_ids = set()
        
        # Check if any tools are Knowledge Bases
        if tool_ids:
            for tool_id in tool_ids:
                tool = app_state.tools.get(tool_id)
                if tool and tool.type == 'knowledge':
                    kb_tool_ids.add(tool_id)
                    # Search this KB with its own config
                    kb_results = await search_knowledge_base(query, tool)
                    all_results.extend(kb_results)
        
        # Search non-KB tools with global RAG
        non_kb_tool_ids = [tid for tid in (tool_ids or []) if tid not in kb_tool_ids]
        
        if non_kb_tool_ids or not tool_ids:
            # Check if RAG is enabled and providers are available
            if not app_state.settings.enable_rag:
                all_results.extend(search_documents_keyword(query, non_kb_tool_ids or None, top_k))
            elif app_state.settings.vector_db.provider == VectorDBProvider.MEMORY:
                all_results.extend(search_documents_keyword(query, non_kb_tool_ids or None, top_k))
            else:
                # Get embedding for query
                embedding_provider = app_state.get_embedding_provider()
                query_embedding = (await embedding_provider.embed([query]))[0]
                
                # Search vector DB
                vector_db = app_state.get_vector_db()
                filter_dict = {"tool_id": {"$in": non_kb_tool_ids}} if non_kb_tool_ids else None
                
                results = await vector_db.search(
                    query_embedding=query_embedding,
                    top_k=top_k,
                    filter=filter_dict,
                    query_text=query
                )
                all_results.extend(results)
        
        # Sort by score and return top_k
        all_results.sort(key=lambda x: x.get('score', 0), reverse=True)
        return all_results[:top_k]
        
    except Exception as e:
        print(f"⚠️ RAG search failed: {e}, falling back to keyword search")
        return search_documents_keyword(query, tool_ids, top_k)


def search_documents_keyword(query: str, tool_ids: List[str] = None, top_k: int = 5) -> List[Dict]:
    """Improved keyword-based search with TF-IDF style scoring"""
    import re
    import math
    
    # Stop words to filter out
    stop_words = {'the', 'a', 'an', 'is', 'are', 'was', 'were', 'be', 'been', 'being', 
                  'have', 'has', 'had', 'do', 'does', 'did', 'will', 'would', 'could', 
                  'should', 'may', 'might', 'must', 'shall', 'can', 'need', 'dare',
                  'to', 'of', 'in', 'for', 'on', 'with', 'at', 'by', 'from', 'as',
                  'and', 'but', 'if', 'or', 'because', 'what', 'which', 'who', 'this',
                  'that', 'i', 'me', 'my', 'we', 'our', 'you', 'your', 'it', 'its'}
    
    query_lower = query.lower().strip()
    query_words = [w for w in re.findall(r'\b\w+\b', query_lower) if w not in stop_words and len(w) > 1]
    
    if not query_words:
        query_words = re.findall(r'\b\w+\b', query_lower)
    
    if not query_words:
        return []
    
    # Get relevant chunks
    relevant_chunks = [c for c in app_state.document_chunks 
                       if not tool_ids or c.get('tool_id') in tool_ids]
    
    if not relevant_chunks:
        return []
    
    # Calculate IDF
    doc_count = len(relevant_chunks)
    idf = {}
    for word in query_words:
        docs_with_word = sum(1 for c in relevant_chunks if word in c.get('text', '').lower())
        idf[word] = math.log((doc_count + 1) / (docs_with_word + 1)) + 1
    
    results = []
    for chunk in relevant_chunks:
        text = chunk.get('text', '').lower()
        text_words = re.findall(r'\b\w+\b', text)
        text_word_count = len(text_words) if text_words else 1
        
        score = 0
        matched_words = []
        
        for word in query_words:
            tf = text.count(word) / text_word_count
            word_score = tf * idf.get(word, 1)
            if word in text:
                matched_words.append(word)
                score += word_score
        
        # Coverage bonus
        coverage = len(matched_words) / len(query_words) if query_words else 0
        score *= (1 + coverage)
        
        # Phrase match bonus
        if len(query_words) > 1 and ' '.join(query_words) in text:
            score *= 2
        
        if score > 0 and (coverage >= 0.3 or len(query_words) == 1):
            results.append({
                "text": chunk.get('text', ''),
                "score": min(score, 1.0),
                "source": chunk.get('source', 'Unknown'),
                "tool_id": chunk.get('tool_id'),
                "type": chunk.get('type', 'document')
            })
    
    results.sort(key=lambda x: x['score'], reverse=True)
    return results[:top_k]


def search_documents(query: str, tool_ids: List[str] = None, top_k: int = 5) -> List[Dict]:
    """Synchronous wrapper for backward compatibility"""
    import asyncio
    try:
        loop = asyncio.get_event_loop()
        if loop.is_running():
            # Already in async context, use keyword search
            return search_documents_keyword(query, tool_ids, top_k)
        return loop.run_until_complete(search_documents_rag(query, tool_ids, top_k))
    except:
        return search_documents_keyword(query, tool_ids, top_k)


async def call_llm(messages: List[Dict], model_id: str = None) -> Dict:
    """Call LLM using configured provider - auto-detects provider from model name"""
    try:
        model = model_id or app_state.settings.llm.model
        
        # Auto-detect provider based on model name
        model_lower = model.lower()
        print(f"[LLM] Model requested: {model}")
        
        if model_lower.startswith('claude'):
            # Use Anthropic for Claude models
            print(f"[LLM] Detected Claude model, using Anthropic provider")
            print(f"[LLM] Available providers: {[p.provider for p in app_state.settings.llm_providers]}")
            provider_data = next(
                (p for p in app_state.settings.llm_providers if p.provider == 'anthropic'),
                None
            )
            print(f"[LLM] Anthropic provider found: {provider_data is not None}")
            if provider_data:
                print(f"[LLM] Anthropic API key set: {bool(provider_data.api_key)}")
            if provider_data and provider_data.api_key:
                config = LLMConfig(
                    provider=LLMProvider.ANTHROPIC,
                    model=model,
                    api_key=provider_data.api_key,
                )
                llm_provider = ProviderFactory.get_llm_provider(config)
            else:
                return {"content": "Error: Anthropic API key not configured. Please add it in Settings."}
        elif model_lower.startswith('gpt') or model_lower.startswith('o1') or model_lower.startswith('o3'):
            # Use OpenAI for GPT/O1/O3 models
            print(f"[LLM] Detected OpenAI model, using OpenAI provider")
            provider_data = next(
                (p for p in app_state.settings.llm_providers if p.provider == 'openai'),
                None
            )
            if provider_data and provider_data.api_key:
                config = LLMConfig(
                    provider=LLMProvider.OPENAI,
                    model=model,
                    api_key=provider_data.api_key,
                )
                llm_provider = ProviderFactory.get_llm_provider(config)
            else:
                return {"content": "Error: OpenAI API key not configured. Please add it in Settings."}
        elif model_lower.startswith('mistral') or model_lower.startswith('open-mistral') or model_lower.startswith('codestral') or model_lower.startswith('ministral'):
            # Use Mistral for Mistral models
            print(f"[LLM] Detected Mistral model, using Mistral provider")
            provider_data = next(
                (p for p in app_state.settings.llm_providers if p.provider == 'mistral'),
                None
            )
            if provider_data and provider_data.api_key:
                config = LLMConfig(
                    provider=LLMProvider.MISTRAL,
                    model=model,
                    api_key=provider_data.api_key,
                    api_base="https://api.mistral.ai/v1"
                )
                llm_provider = OpenAICompatibleLLM(config, "mistral")
            else:
                return {"content": "Error: Mistral API key not configured. Please add it in Settings."}
        elif model_lower.startswith('llama') or model_lower.startswith('gemma'):
            # Use Groq for Llama/Gemma models (NOT mixtral - that goes to Mistral)
            print(f"[LLM] Detected Groq model, using Groq provider")
            provider_data = next(
                (p for p in app_state.settings.llm_providers if p.provider == 'groq'),
                None
            )
            if provider_data and provider_data.api_key:
                config = LLMConfig(
                    provider=LLMProvider.GROQ,
                    model=model,
                    api_key=provider_data.api_key,
                    api_base="https://api.groq.com/openai/v1"
                )
                llm_provider = OpenAICompatibleLLM(config, "groq")
            else:
                return {"content": "Error: Groq API key not configured. Please add it in Settings."}
        elif model_lower.startswith('gemini'):
            # Use Google for Gemini models
            print(f"[LLM] Detected Gemini model, using Google provider")
            provider_data = next(
                (p for p in app_state.settings.llm_providers if p.provider == 'google'),
                None
            )
            api_key = provider_data.api_key if provider_data else None
            print(f"[LLM] Google provider from settings: {provider_data is not None}, API key set: {bool(api_key)}")
            if not api_key and app_state.settings.google:
                api_key = app_state.settings.google.gemini_key
                print(f"[LLM] Fallback to settings.google.gemini_key: {bool(api_key)}")
            if api_key:
                config = LLMConfig(
                    provider=LLMProvider.GOOGLE,
                    model=model,
                    api_key=api_key,
                )
                llm_provider = ProviderFactory.get_llm_provider(config)
            else:
                print(f"[LLM] ❌ No Google API key found!")
                return {"content": "Error: Google API key not configured. Please add it in Settings."}
        else:
            # Use default provider from settings
            print(f"[LLM] Using default provider from settings")
            llm_provider = app_state.get_llm_provider()
        
        print(f"[LLM] Calling provider.generate()...")
        content = await llm_provider.generate(messages, model=model)
        print(f"[LLM] Response received: {len(content) if content else 0} chars")
        return {"content": content}
    except Exception as e:
        print(f"[LLM] ❌ Error: {e}")
        import traceback
        traceback.print_exc()
        return {"content": f"Error: {str(e)}"}


# ============================================================================
# TOOL CALLING SYSTEM - Real tool execution for agents
# ============================================================================

def build_tool_definitions(tools: List['ToolConfiguration']) -> List[Dict]:
    """Build OpenAI-compatible tool definitions from agent tools"""
    tool_defs = []
    used_names = set()  # Track used names to ensure uniqueness
    
    def make_unique_name(base_name: str, tool_id: str) -> str:
        """Generate a unique function name"""
        # Clean the full tool ID to be a valid function name (alphanumeric and underscore only)
        import re
        clean_id = re.sub(r'[^a-zA-Z0-9]', '_', tool_id)
        # Remove consecutive underscores and trim
        clean_id = re.sub(r'_+', '_', clean_id).strip('_')
        name = f"{base_name}_{clean_id}"
        
        # Ensure uniqueness (should already be unique with full ID, but just in case)
        if name in used_names:
            counter = 1
            while f"{name}_{counter}" in used_names:
                counter += 1
            name = f"{name}_{counter}"
        
        used_names.add(name)
        return name
    
    # Additional safeguard: Track function names to ensure final uniqueness
    final_function_names = set()
    
    for tool in tools:
        if tool.type == 'email':
            func_name = make_unique_name("send_email", tool.id)
            tool_defs.append({
                "type": "function",
                "function": {
                    "name": func_name,
                    "description": f"Send an email using {tool.name}. {tool.description or ''}",
                    "parameters": {
                        "type": "object",
                        "properties": {
                            "to": {"type": "string", "description": "Recipient email address"},
                            "subject": {"type": "string", "description": "Email subject line"},
                            "body": {"type": "string", "description": "Email message body"}
                        },
                        "required": ["to", "subject", "body"]
                    }
                },
                "_tool_id": tool.id,
                "_tool_type": "email"
            })
        
        elif tool.type == 'api':
            # Build parameters from API config
            params_props = {}
            required_params = []
            
            if tool.api_config and tool.api_config.input_parameters:
                for param in tool.api_config.input_parameters:
                    param_type = "string"
                    if param.data_type in ['integer', 'number']:
                        param_type = param.data_type
                    elif param.data_type == 'boolean':
                        param_type = 'boolean'
                    
                    params_props[param.name] = {
                        "type": param_type,
                        "description": param.description or f"Parameter: {param.name}"
                    }
                    if param.required:
                        required_params.append(param.name)
            
            func_name = make_unique_name("call_api", tool.id)
            tool_defs.append({
                "type": "function",
                "function": {
                    "name": func_name,
                    "description": f"Call {tool.name} API. {tool.description or ''}",
                    "parameters": {
                        "type": "object",
                        "properties": params_props if params_props else {
                            "query": {"type": "string", "description": "Query or request parameters"}
                        },
                        "required": required_params if required_params else []
                    }
                },
                "_tool_id": tool.id,
                "_tool_type": "api"
            })
        
        elif tool.type == 'websearch':
            func_name = make_unique_name("web_search", tool.id)
            tool_defs.append({
                "type": "function",
                "function": {
                    "name": func_name,
                    "description": f"Search the web using {tool.name}. {tool.description or ''}",
                    "parameters": {
                        "type": "object",
                        "properties": {
                            "query": {"type": "string", "description": "Search query"}
                        },
                        "required": ["query"]
                    }
                },
                "_tool_id": tool.id,
                "_tool_type": "websearch"
            })
        
        elif tool.type == 'webhook':
            func_name = make_unique_name("trigger_webhook", tool.id)
            tool_defs.append({
                "type": "function",
                "function": {
                    "name": func_name,
                    "description": f"Trigger webhook: {tool.name}. {tool.description or ''}",
                    "parameters": {
                        "type": "object",
                        "properties": {
                            "data": {"type": "object", "description": "Data to send to webhook"}
                        },
                        "required": []
                    }
                },
                "_tool_id": tool.id,
                "_tool_type": "webhook"
            })
        
        elif tool.type == 'slack':
            func_name = make_unique_name("send_slack", tool.id)
            tool_defs.append({
                "type": "function",
                "function": {
                    "name": func_name,
                    "description": f"Send Slack message using {tool.name}. {tool.description or ''}",
                    "parameters": {
                        "type": "object",
                        "properties": {
                            "message": {"type": "string", "description": "Message to send"},
                            "channel": {"type": "string", "description": "Channel name (optional)"}
                        },
                        "required": ["message"]
                    }
                },
                "_tool_id": tool.id,
                "_tool_type": "slack"
            })
        
        elif tool.type in ['knowledge', 'document', 'website']:
            func_name = make_unique_name("search_kb", tool.id)
            tool_defs.append({
                "type": "function",
                "function": {
                    "name": func_name,
                    "description": f"Search knowledge base: {tool.name}. {tool.description or ''}",
                    "parameters": {
                        "type": "object",
                        "properties": {
                            "query": {"type": "string", "description": "Search query"}
                        },
                        "required": ["query"]
                    }
                },
                "_tool_id": tool.id,
                "_tool_type": "knowledge"
            })
        
        elif tool.type == 'database':
            func_name = make_unique_name("query_db", tool.id)
            tool_defs.append({
                "type": "function",
                "function": {
                    "name": func_name,
                    "description": f"Query database: {tool.name}. {tool.description or ''}",
                    "parameters": {
                        "type": "object",
                        "properties": {
                            "query": {"type": "string", "description": "SQL query or natural language query"}
                        },
                        "required": ["query"]
                    }
                },
                "_tool_id": tool.id,
                "_tool_type": "database"
            })
    
    # Final validation: Ensure all function names are unique
    seen_func_names = {}
    for i, tool_def in enumerate(tool_defs):
        func_name = tool_def['function']['name']
        if func_name in seen_func_names:
            # Duplicate found - make it unique
            counter = 1
            new_name = f"{func_name}_{counter}"
            while new_name in seen_func_names:
                counter += 1
                new_name = f"{func_name}_{counter}"
            tool_def['function']['name'] = new_name
            func_name = new_name
            print(f"⚠️  Duplicate tool name detected, renamed to: {new_name}")
        seen_func_names[func_name] = i
    
    print(f"📦 Built {len(tool_defs)} tool definitions: {[t['function']['name'] for t in tool_defs]}")
    return tool_defs


async def execute_tool(tool_id: str, tool_type: str, arguments: Dict) -> Dict:
    """Execute a tool and return the result"""
    print(f"\n🔧 EXECUTING TOOL")
    print(f"   Tool ID: {tool_id}")
    print(f"   Tool Type: {tool_type}")
    print(f"   Arguments: {arguments}")
    
    tool = resolve_tool_id(tool_id)
    if not tool:
        print(f"   ❌ Tool not found: {tool_id}")
        return {"success": False, "error": f"Tool {tool_id} not found"}
    
    print(f"   ✅ Tool resolved: {tool.name} (actual ID: {tool.id})")
    
    try:
        if tool_type == 'email':
            # Send email
            config = tool.config or {}
            provider = config.get('provider', '')
            
            if provider == 'google':
                access_token = config.get('access_token')
                refresh_token = config.get('refresh_token')
                
                # Refresh token if needed
                if refresh_token:
                    new_token = await EmailService.refresh_google_token(refresh_token)
                    if new_token:
                        access_token = new_token
                        tool.config['access_token'] = new_token
                        app_state.save_to_disk()
                
                if access_token:
                    result = await EmailService.send_gmail(
                        access_token=access_token,
                        to=arguments.get('to', ''),
                        subject=arguments.get('subject', ''),
                        body=arguments.get('body', '')
                    )
                    if result.get('success'):
                        return {
                            "success": True,
                            "message": f"Email sent successfully to {arguments.get('to')}",
                            "message_id": result.get('message_id')
                        }
                    else:
                        return {"success": False, "error": result.get('error', 'Failed to send email')}
                else:
                    return {"success": False, "error": "No valid access token"}
            
            elif provider == 'sendgrid':
                api_key = config.get('apiKey', config.get('api_key'))
                from_email = config.get('fromEmail', config.get('from_email', config.get('email')))
                
                result = await EmailService.send_sendgrid(
                    api_key=api_key,
                    to=arguments.get('to', ''),
                    subject=arguments.get('subject', ''),
                    body=arguments.get('body', ''),
                    from_email=from_email
                )
                return result
            
            else:
                return {"success": False, "error": f"Unsupported email provider: {provider}"}
        
        elif tool_type == 'api':
            # Call external API
            if tool.api_config:
                base_url = tool.api_config.base_url
                endpoint = tool.api_config.endpoint_path
                method = tool.api_config.http_method
                
                # Replace parameters in endpoint
                url = f"{base_url.rstrip('/')}/{endpoint.lstrip('/')}"
                for key, value in arguments.items():
                    url = url.replace(f"{{{key}}}", str(value))
                
                # Prepare headers
                headers = dict(tool.api_config.headers) if tool.api_config.headers else {}
                
                # Add auth
                if tool.api_config.auth_type == 'bearer':
                    headers['Authorization'] = f"Bearer {tool.api_config.auth_value}"
                elif tool.api_config.auth_type == 'api_key':
                    if tool.api_config.api_key_location == 'header':
                        headers[tool.api_config.api_key_name] = tool.api_config.auth_value
                
                async with httpx.AsyncClient(timeout=30.0) as client:
                    if method == 'GET':
                        response = await client.get(url, headers=headers, params=arguments)
                    elif method == 'POST':
                        response = await client.post(url, headers=headers, json=arguments)
                    elif method == 'PUT':
                        response = await client.put(url, headers=headers, json=arguments)
                    elif method == 'DELETE':
                        response = await client.delete(url, headers=headers)
                    else:
                        response = await client.get(url, headers=headers)
                    
                    if response.status_code < 400:
                        try:
                            data = response.json()
                        except:
                            data = response.text
                        return {"success": True, "data": data}
                    else:
                        return {"success": False, "error": f"API returned {response.status_code}"}
            
            # Check for mock response (demo mode)
            if tool.config and tool.config.get('mock_response'):
                return {"success": True, "data": tool.config['mock_response']}
            
            return {"success": False, "error": "API not configured"}
        
        elif tool_type == 'knowledge':
            # Search knowledge base
            query = arguments.get('query', '')
            results = search_documents(query, [tool_id], top_k=5)
            
            if results:
                return {
                    "success": True,
                    "results": [{"text": r['text'][:500], "source": r['source']} for r in results]
                }
            return {"success": True, "results": [], "message": "No relevant documents found"}
        
        elif tool_type == 'webhook':
            # Trigger webhook
            webhook_url = tool.config.get('url', '')
            if webhook_url:
                async with httpx.AsyncClient(timeout=30.0) as client:
                    response = await client.post(webhook_url, json=arguments.get('data', {}))
                    return {"success": response.status_code < 400, "status_code": response.status_code}
            return {"success": False, "error": "Webhook URL not configured"}
        
        elif tool_type == 'slack':
            # Send Slack message
            webhook_url = tool.config.get('webhook_url', '')
            if webhook_url:
                async with httpx.AsyncClient(timeout=30.0) as client:
                    response = await client.post(webhook_url, json={
                        "text": arguments.get('message', ''),
                        "channel": arguments.get('channel', '')
                    })
                    return {"success": response.status_code < 400}
            return {"success": False, "error": "Slack webhook not configured"}
        
        elif tool_type == 'websearch':
            # Web search (would need actual implementation with search API)
            query = arguments.get('query', '')
            return {
                "success": True,
                "message": f"Web search for '{query}' - requires search API integration",
                "results": []
            }
        
        elif tool_type == 'database':
            # Database query (would need actual DB connection)
            query = arguments.get('query', '')
            return {
                "success": True,
                "message": f"Database query - requires actual DB connection",
                "query": query
            }
        
        else:
            return {"success": False, "error": f"Unknown tool type: {tool_type}"}
    
    except Exception as e:
        print(f"   ❌ Tool execution error: {e}")
        import traceback
        traceback.print_exc()
        return {"success": False, "error": str(e)}


async def call_llm_with_tools(messages: List[Dict], tools: List[Dict], model_id: str = None) -> Dict:
    """Call LLM with tool calling support"""
    model = model_id or app_state.settings.llm.model
    model_lower = model.lower()
    
    print(f"\n🤖 LLM CALL WITH TOOLS")
    print(f"   Model: {model}")
    print(f"   Model lower: {model_lower}")
    print(f"   Tools count: {len(tools)}")
    print(f"   Checking conditions: gpt={model_lower.startswith('gpt')}, claude={model_lower.startswith('claude')}, gemini={model_lower.startswith('gemini')}")
    
    # Extract just the function definitions (without our metadata)
    openai_tools = []
    tool_mapping = {}  # Map function name to tool info
    
    for t in tools:
        func_name = t['function']['name']
        openai_tools.append({
            "type": "function",
            "function": t['function']
        })
        tool_mapping[func_name] = {
            "tool_id": t.get('_tool_id'),
            "tool_type": t.get('_tool_type')
        }
    
    try:
        if model_lower.startswith('gpt') or model_lower.startswith('o1') or model_lower.startswith('o3'):
            # OpenAI with function calling
            provider_data = next(
                (p for p in app_state.settings.llm_providers if p.provider == 'openai'),
                None
            )
            
            if not provider_data or not provider_data.api_key:
                return {"content": "Error: OpenAI API key not configured", "tool_calls": []}
            
            async with httpx.AsyncClient(timeout=60.0) as client:
                response = await client.post(
                    "https://api.openai.com/v1/chat/completions",
                    headers={
                        "Authorization": f"Bearer {provider_data.api_key}",
                        "Content-Type": "application/json"
                    },
                    json={
                        "model": model,
                        "messages": messages,
                        "tools": openai_tools if openai_tools else None,
                        "tool_choice": "auto" if openai_tools else None
                    }
                )
                
                if response.status_code != 200:
                    error = response.json()
                    return {"content": f"Error: {error}", "tool_calls": []}
                
                data = response.json()
                choice = data['choices'][0]
                message = choice['message']
                
                # Check for tool calls
                if message.get('tool_calls'):
                    tool_calls = []
                    for tc in message['tool_calls']:
                        func_name = tc['function']['name']
                        try:
                            args = json.loads(tc['function']['arguments'])
                        except:
                            args = {}
                        
                        tool_info = tool_mapping.get(func_name, {})
                        tool_calls.append({
                            "id": tc['id'],
                            "name": func_name,
                            "arguments": args,
                            "tool_id": tool_info.get('tool_id'),
                            "tool_type": tool_info.get('tool_type')
                        })
                    
                    return {
                        "content": message.get('content', ''),
                        "tool_calls": tool_calls,
                        "finish_reason": choice.get('finish_reason')
                    }
                
                return {
                    "content": message.get('content', ''),
                    "tool_calls": [],
                    "finish_reason": choice.get('finish_reason')
                }
        
        elif model_lower.startswith('claude'):
            # Anthropic with tool use
            provider_data = next(
                (p for p in app_state.settings.llm_providers if p.provider == 'anthropic'),
                None
            )
            
            if not provider_data or not provider_data.api_key:
                return {"content": "Error: Anthropic API key not configured", "tool_calls": []}
            
            # Convert tools to Anthropic format
            anthropic_tools = []
            for t in tools:
                anthropic_tools.append({
                    "name": t['function']['name'],
                    "description": t['function']['description'],
                    "input_schema": t['function']['parameters']
                })
            
            # Convert messages to Anthropic format
            system_content = ""
            anthropic_messages = []
            for msg in messages:
                if msg['role'] == 'system':
                    system_content = msg['content'] if isinstance(msg['content'], str) else str(msg['content'])
                elif msg['role'] == 'tool':
                    # Convert OpenAI tool response to Anthropic tool_result format
                    # This shouldn't happen if using correct format, but handle it just in case
                    anthropic_messages.append({
                        "role": "user",
                        "content": [
                            {
                                "type": "tool_result",
                                "tool_use_id": msg.get('tool_call_id', ''),
                                "content": msg.get('content', '')
                            }
                        ]
                    })
                elif msg['role'] == 'assistant' and msg.get('tool_calls'):
                    # Convert OpenAI assistant tool_calls to Anthropic format
                    content_blocks = []
                    if msg.get('content'):
                        content_blocks.append({"type": "text", "text": msg['content']})
                    for tc in msg['tool_calls']:
                        content_blocks.append({
                            "type": "tool_use",
                            "id": tc.get('id', ''),
                            "name": tc.get('function', {}).get('name', ''),
                            "input": json.loads(tc.get('function', {}).get('arguments', '{}'))
                        })
                    anthropic_messages.append({
                        "role": "assistant",
                        "content": content_blocks
                    })
                elif msg['role'] in ['user', 'assistant']:
                    # Regular message - check if content is already in Anthropic format
                    content = msg.get('content', '')
                    if isinstance(content, list):
                        # Already in Anthropic block format
                        anthropic_messages.append({"role": msg['role'], "content": content})
                    else:
                        # String content
                        anthropic_messages.append({"role": msg['role'], "content": content})
            
            async with httpx.AsyncClient(timeout=60.0) as client:
                request_body = {
                    "model": model,
                    "max_tokens": 4096,
                    "messages": anthropic_messages
                }
                if system_content:
                    request_body["system"] = system_content
                if anthropic_tools:
                    request_body["tools"] = anthropic_tools
                
                print(f"   📤 Sending to Anthropic: {len(anthropic_messages)} messages, {len(anthropic_tools)} tools")
                
                response = await client.post(
                    "https://api.anthropic.com/v1/messages",
                    headers={
                        "x-api-key": provider_data.api_key,
                        "Content-Type": "application/json",
                        "anthropic-version": "2023-06-01"
                    },
                    json=request_body
                )
                
                if response.status_code != 200:
                    error = response.json()
                    print(f"   ❌ Anthropic API Error: {error}")
                    return {"content": f"Error: {error}", "tool_calls": []}
                
                data = response.json()
                
                # Parse Anthropic response
                content_text = ""
                tool_calls = []
                
                for block in data.get('content', []):
                    if block['type'] == 'text':
                        content_text += block['text']
                    elif block['type'] == 'tool_use':
                        func_name = block['name']
                        tool_info = tool_mapping.get(func_name, {})
                        tool_calls.append({
                            "id": block['id'],
                            "name": func_name,
                            "arguments": block['input'],
                            "tool_id": tool_info.get('tool_id'),
                            "tool_type": tool_info.get('tool_type')
                        })
                
                return {
                    "content": content_text,
                    "tool_calls": tool_calls,
                    "finish_reason": data.get('stop_reason')
                }
        
        elif model_lower.startswith('gemini'):
            # Google Gemini with function calling
            print(f"   🔧 ENTERING Gemini function calling branch!")
            provider_data = next(
                (p for p in app_state.settings.llm_providers if p.provider == 'google'),
                None
            )
            
            if not provider_data or not provider_data.api_key:
                return {"content": "Error: Google API key not configured", "tool_calls": []}
            
            # Map model name to actual API model (gemini-1.5-flash is most stable for function calling)
            # Google API requires specific model names
            model_mapping = {
                "gemini-pro": "gemini-1.5-flash",
                "gemini-1.5-flash": "gemini-1.5-flash", 
                "gemini-1.5-pro": "gemini-1.5-flash",  # Use flash for reliability
                "gemini-2.0-flash": "gemini-1.5-flash",
                "gemini-flash": "gemini-1.5-flash"
            }
            api_model = model_mapping.get(model_lower, "gemini-1.5-flash")
            print(f"   📍 Mapped model {model_lower} -> {api_model} (for function calling)")
            
            print(f"   🔧 Gemini function calling with model: {api_model}")
            
            # Convert messages to Gemini format
            gemini_messages = []
            system_instruction = ""
            for msg in messages:
                if msg['role'] == 'system':
                    system_instruction = msg['content']
                else:
                    gemini_messages.append({
                        "role": "user" if msg['role'] == 'user' else "model",
                        "parts": [{"text": msg['content']}]
                    })
            
            # Convert tools to Gemini format
            gemini_tools = []
            if openai_tools:
                function_declarations = []
                for t in openai_tools:
                    func = t['function']
                    function_declarations.append({
                        "name": func['name'],
                        "description": func.get('description', ''),
                        "parameters": func.get('parameters', {"type": "object", "properties": {}})
                    })
                gemini_tools = [{"function_declarations": function_declarations}]
            
            async with httpx.AsyncClient(timeout=60.0) as client:
                request_body = {
                    "contents": gemini_messages,
                    "generationConfig": {
                        "temperature": 0.7,
                        "maxOutputTokens": 4096
                    }
                }
                
                if system_instruction:
                    request_body["systemInstruction"] = {"parts": [{"text": system_instruction}]}
                
                if gemini_tools:
                    request_body["tools"] = gemini_tools
                
                response = await client.post(
                    f"https://generativelanguage.googleapis.com/v1beta/models/{api_model}:generateContent?key={provider_data.api_key}",
                    headers={"Content-Type": "application/json"},
                    json=request_body
                )
                
                if response.status_code != 200:
                    error = response.text
                    print(f"   ❌ Gemini error: {error}")
                    return {"content": f"Error: {error}", "tool_calls": []}
                
                data = response.json()
                
                if 'candidates' not in data or not data['candidates']:
                    return {"content": "No response from Gemini", "tool_calls": []}
                
                candidate = data['candidates'][0]
                content_parts = candidate.get('content', {}).get('parts', [])
                
                content_text = ""
                tool_calls = []
                
                for part in content_parts:
                    if 'text' in part:
                        content_text += part['text']
                    elif 'functionCall' in part:
                        func_call = part['functionCall']
                        func_name = func_call['name']
                        tool_info = tool_mapping.get(func_name, {})
                        tool_calls.append({
                            "id": f"gemini_{func_name}_{id(part)}",
                            "name": func_name,
                            "arguments": func_call.get('args', {}),
                            "tool_id": tool_info.get('tool_id'),
                            "tool_type": tool_info.get('tool_type')
                        })
                
                print(f"   ✅ Gemini response: {len(content_text)} chars, {len(tool_calls)} tool calls")
                
                return {
                    "content": content_text,
                    "tool_calls": tool_calls,
                    "finish_reason": candidate.get('finishReason')
                }
        
        else:
            # Truly unsupported models
            print(f"   ⚠️  Model {model} doesn't support tools, falling back to basic LLM call")
            result = await call_llm(messages, model_id)
            print(f"   📤 Fallback LLM result: {len(result.get('content', '') or '')} chars")
            if not result.get('content'):
                print(f"   ⚠️  Empty content from fallback! Full result: {result}")
            return {"content": result.get('content', ''), "tool_calls": []}
    
    except Exception as e:
        print(f"   ❌ LLM call error: {e}")
        import traceback
        traceback.print_exc()
        return {"content": f"Error: {str(e)}", "tool_calls": []}


# ============================================================================
# AGENT CHAT SYSTEM - Handles all agent conversations
# ============================================================================

def get_current_datetime_for_user(timezone_str: str = None) -> str:
    """Get current date/time formatted for the user's timezone"""
    from datetime import datetime, timezone
    
    try:
        if timezone_str:
            # Use zoneinfo (Python 3.9+) instead of pytz
            from zoneinfo import ZoneInfo
            tz = ZoneInfo(timezone_str)
            current_date = datetime.now(tz)
        else:
            # Default to UTC if no timezone provided
            current_date = datetime.now(timezone.utc)
        
        return f"""=== CURRENT DATE & TIME ===
Today: {current_date.strftime('%A, %B %d, %Y')}
Time: {current_date.strftime('%I:%M %p')} ({timezone_str or 'UTC'})
"""
    except Exception:
        # Fallback to server time
        current_date = datetime.now()
        return f"""=== CURRENT DATE & TIME ===
Today: {current_date.strftime('%A, %B %d, %Y')}
Time: {current_date.strftime('%I:%M %p')} (Server Time)
"""


async def update_agent_memory(agent: AgentData, conversation: 'Conversation') -> None:
    """
    Update agent memory with a summary of the conversation.
    Called after a conversation ends or has significant content.
    """
    if not agent.memory_enabled:
        return
    
    # Only update if conversation has enough messages
    if len(conversation.messages) < 4:  # At least 2 exchanges
        return
    
    try:
        # Build conversation text for summarization
        conv_text = "\n".join([
            f"{msg.role}: {msg.content[:200]}" 
            for msg in conversation.messages[-10:]  # Last 10 messages
        ])
        
        # Use LLM to summarize
        messages = [
            {
                "role": "system",
                "content": "Summarize this conversation in 1-2 sentences, focusing on key facts, decisions, or user preferences learned. Be concise."
            },
            {
                "role": "user",
                "content": conv_text
            }
        ]
        
        result = await call_llm(messages, agent.model_id)
        summary = result.get("content", "")
        
        if summary:
            # Add to agent memory
            agent.memory.append({
                "timestamp": datetime.utcnow().isoformat(),
                "summary": summary[:500],  # Limit length
                "conversation_id": conversation.id
            })
            
            # Keep only last 20 memories
            if len(agent.memory) > 20:
                agent.memory = agent.memory[-20:]
            
            agent.updated_at = datetime.utcnow().isoformat()
            app_state.save_to_disk()
            
    except Exception as e:
        print(f"[Memory] Error updating memory: {e}")


def get_language_instruction(language_setting: str = 'user') -> str:
    """Get language/dialect matching instruction for system prompt"""
    lang_map = {
        'user': 'RESPOND IN THE USER\'S LANGUAGE AND DIALECT. Detect the language from the user\'s message, not from attached documents/images. Mirror their exact dialect, tone, and style.',
        'en': 'Respond in English only',
        'ar': 'Respond in Modern Standard Arabic (فصحى) only', 
        'multi': 'Match the language and dialect the user writes in'
    }
    return lang_map.get(language_setting, lang_map['user'])


async def process_test_agent_chat(agent: AgentData, message: str, conversation: Conversation, attachments: List[Dict] = None, timezone: str = None) -> Dict:
    """
    Agent Chat with File Processing Support
    
    Processes conversations using the agent's full configuration:
    - Tasks & Instructions
    - Personality settings
    - Guardrails
    - Knowledge bases
    - Available APIs/Tools
    """
    
    # 1. Process attachments with REAL extraction
    extracted_content = []
    if attachments:
        for attachment in attachments:
            file_path = attachment.get('path')
            file_name = attachment.get('name', '')
            file_type = attachment.get('type', '')
            
            if file_path and os.path.exists(file_path):
                content = None
                
                if file_type.startswith('image/') or file_name.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.webp')):
                    content = await extract_content_from_image(file_path)
                elif file_type == 'application/pdf' or file_name.lower().endswith('.pdf'):
                    content = extract_text_from_pdf(file_path)
                elif file_name.lower().endswith(('.docx', '.doc')):
                    content = extract_text_from_docx(file_path)
                elif file_name.lower().endswith(('.txt', '.csv', '.md')):
                    try:
                        with open(file_path, 'r', encoding='utf-8') as f:
                            content = f.read()
                    except:
                        pass
                
                if content:
                    extracted_content.append({"filename": file_name, "content": content})
    
    # 2. Get agent tools (resolves prefixed IDs)
    agent_tools = get_agent_tools(agent)
    
    # 3. Search Knowledge Bases
    context = ""
    sources = []
    
    # Search regular knowledge bases
    tool_ids = [t.id for t in agent_tools]
    search_results = search_documents(message, tool_ids, top_k=5)
    if search_results:
        context = "\n\n=== KNOWLEDGE BASE ===\n"
        for i, result in enumerate(search_results):
            context += f"\n[Source {i+1}: {result['source']}]\n{result['text'][:800]}\n"
            sources.append({"source": result['source'], "type": result['type'], "relevance": round(result['score'] * 100)})
    
    # Add Knowledge Base content with inline sections
    for tool in agent_tools:
        if tool.type == 'knowledge' and tool.config.get('sections'):
            kb_sections = tool.config.get('sections', [])
            if kb_sections:
                context += f"\n\n=== {tool.name.upper()} ===\n"
                for section in kb_sections:
                    context += f"\n## {section.get('title', '')}\n{section.get('content', '')}\n"
                sources.append({"source": tool.name, "type": "knowledge_base", "relevance": 95})
    
    # 4. Build APIs context
    apis_context = ""
    api_tools = [t for t in agent_tools if t.type == 'api']
    if api_tools:
        apis_context = "\n\n=== AVAILABLE APIs ===\n"
        for api in api_tools:
            method = api.config.get('method') or (api.api_config.http_method if api.api_config else 'GET')
            endpoint = api.config.get('endpoint') or (api.api_config.endpoint_path if api.api_config else '')
            apis_context += f"• **{api.name}**: {api.description}\n"
            apis_context += f"  [{method}] {endpoint}\n"
    
    # 5. Add uploaded files
    files_context = ""
    if extracted_content:
        files_context = "\n\n=== UPLOADED FILES ===\n"
        for ec in extracted_content:
            files_context += f"\n--- {ec['filename']} ---\n{ec['content']}\n"
            sources.append({"source": ec['filename'], "type": "uploaded_file", "relevance": 100})
    
    # 6. Build FULL system prompt with ALL configuration
    g = agent.guardrails
    p = agent.personality
    
    # Get current date/time for user's timezone
    date_info = get_current_datetime_for_user(timezone)
    
    # Personality descriptions
    creativity_desc = "(Factual only)" if p.creativity <= 3 else "(Creative)" if p.creativity >= 7 else "(Balanced)"
    length_desc = "(Brief responses)" if p.length <= 3 else "(Detailed)" if p.length >= 7 else "(Moderate)"
    formality_desc = "(Casual)" if p.formality <= 3 else "(Formal)" if p.formality >= 7 else "(Professional)"
    empathy_desc = "(Direct)" if p.empathy <= 3 else "(Empathetic)" if p.empathy >= 7 else "(Supportive)"
    
    system_prompt = f"""You are {agent.name}.

=== LANGUAGE ===
{get_language_instruction(g.language)}

{date_info}
=== YOUR GOAL ===
{agent.goal}
"""

    # Add memory context if enabled and has memories
    if agent.memory_enabled and agent.memory:
        system_prompt += "\n=== MEMORY (Previous Conversations) ===\n"
        recent_memories = agent.memory[-5:]
        for mem in recent_memories:
            system_prompt += f"• {mem.get('summary', '')}\n"
        system_prompt += "\nUse this context for continuity but don't reference it explicitly unless relevant.\n"

    system_prompt += f"""
=== YOUR PERSONALITY ===
• Tone: {p.tone}
• Voice: {p.voice}
• Creativity: {p.creativity}/10 {creativity_desc}
• Response Length: {p.length}/10 {length_desc}
• Formality: {p.formality}/10 {formality_desc}
• Empathy: {p.empathy}/10 {empathy_desc}

=== YOUR TASKS ==="""

    # Add tasks with instructions
    for task in agent.tasks:
        system_prompt += f"\n\n### {task.name}\n{task.description}"
        if task.instructions:
            system_prompt += "\n**Steps:**"
            for i, inst in enumerate(task.instructions, 1):
                inst_text = inst.text if hasattr(inst, 'text') else str(inst)
                system_prompt += f"\n{i}. {inst_text}"

    # Add guardrails
    system_prompt += "\n\n=== GUARDRAILS (MUST FOLLOW) ==="
    
    if g.anti_hallucination:
        system_prompt += "\n\n**ACCURACY:**"
        if g.cite_sources:
            system_prompt += "\n• Cite sources using [Source X] format when using knowledge base"
        if g.admit_uncertainty:
            system_prompt += "\n• If information is not in knowledge base and you cannot deduce it, say so"
        if g.verify_facts:
            system_prompt += "\n• Use facts from the knowledge base - don't make up company-specific data"
        if g.no_speculation:
            system_prompt += "\n• Don't speculate about company policies or data you don't have"
    
    if g.avoid_topics:
        system_prompt += f"\n\n**AVOID TOPICS:** {', '.join(g.avoid_topics)}"
    
    if g.focus_topics:
        system_prompt += f"\n\n**FOCUS ON:** {', '.join(g.focus_topics)}"
    
    # Response length from guardrails
    length_map = {
        'short': '1-2 short paragraphs MAX',
        'medium': '2-3 paragraphs',
        'long': 'Detailed as needed',
        'unlimited': 'Thorough coverage'
    }
    system_prompt += f"\n\n**RESPONSE LENGTH:** {length_map.get(g.max_length, '2-3 paragraphs')}"
    
    # Language/dialect matching
    system_prompt += f"\n\n**LANGUAGE:** {get_language_instruction(g.language)}"
    
    if g.escalate_angry or g.escalate_complex or g.escalate_request:
        system_prompt += "\n\n**ESCALATION:** Offer human help if user is frustrated or issue is complex"
    
    if g.pii_protection:
        system_prompt += "\n\n**PRIVACY:** Don't ask for or repeat sensitive personal info"

    # Add context
    system_prompt += context
    system_prompt += apis_context
    system_prompt += files_context
    
    # Build tool definitions for action tools
    action_tools = [t for t in agent_tools if t.type in ['email', 'api', 'webhook', 'slack', 'websearch', 'database']]
    tool_definitions = build_tool_definitions(action_tools) if action_tools else []
    
    # Add available tools to system prompt
    if action_tools:
        system_prompt += "\n\n=== AVAILABLE TOOLS ===\nYou have access to the following tools. Use them when appropriate:\n"
        for t in action_tools:
            system_prompt += f"• **{t.name}** ({t.type}): {t.description or 'No description'}\n"
        system_prompt += "\nWhen you need to use a tool, the system will automatically execute it for you.\n"
    
    # Reminder to follow task instructions
    system_prompt += """

=== IMPORTANT ===
• Follow the task instructions above carefully
• Be smart: calculate dates, times, durations when user provides relative terms (tomorrow, next week, etc.)
• Only ask for information that is truly missing - don't re-ask what user already provided
• Use the knowledge base for facts, but apply common sense for calculations"""

    # 7. Build conversation
    llm_messages = [{"role": "system", "content": system_prompt}]
    for msg in conversation.messages[-10:]:
        llm_messages.append({"role": msg.role, "content": msg.content})
    llm_messages.append({"role": "user", "content": message})
    
    # Detect if using Claude/Anthropic
    model_id = agent.model_id or ''
    is_anthropic = model_id.lower().startswith('claude')
    
    # 8. Call LLM with tools if available
    tool_calls_made = []
    if tool_definitions:
        print(f"\n🤖 TEST AGENT CHAT WITH TOOLS")
        print(f"   Agent: {agent.name}")
        print(f"   Model: {model_id} (Anthropic: {is_anthropic})")
        print(f"   Tools: {[t.name for t in action_tools]}")
        
        # Tool calling loop
        max_iterations = 5
        for iteration in range(max_iterations):
            print(f"\n   📍 Iteration {iteration + 1}/{max_iterations}")
            
            result = await call_llm_with_tools(llm_messages, tool_definitions, agent.model_id)
            
            if result.get('tool_calls'):
                for tc in result['tool_calls']:
                    print(f"\n   🔧 Tool call: {tc['name']}")
                    
                    # Check if it's the embedded security tool
                    if tc['name'] == 'check_user_permissions':
                        print(f"   🔐 Executing SECURITY TOOL")
                        task_name = tc.get('arguments', {}).get('task_name')
                        tool_result = execute_security_tool(current_user, access_control, agent, task_name)
                        tool_result = json.dumps(tool_result, indent=2)
                    else:
                        tool_result = await execute_tool(
                            tc.get('tool_id', ''),
                            tc.get('tool_type', ''),
                            tc.get('arguments', {})
                        )
                    
                    tool_calls_made.append({
                        "tool": tc['name'],
                        "arguments": tc['arguments'],
                        "result": tool_result
                    })
                    
                    # Format messages based on provider
                    if is_anthropic:
                        # Anthropic format: assistant with tool_use, then user with tool_result
                        llm_messages.append({
                            "role": "assistant",
                            "content": [
                                {"type": "text", "text": result.get('content', '') or "I'll use a tool to help with this."},
                                {
                                    "type": "tool_use",
                                    "id": tc['id'],
                                    "name": tc['name'],
                                    "input": tc['arguments']
                                }
                            ] if result.get('content') else [
                                {
                                    "type": "tool_use",
                                    "id": tc['id'],
                                    "name": tc['name'],
                                    "input": tc['arguments']
                                }
                            ]
                        })
                        llm_messages.append({
                            "role": "user",
                            "content": [
                                {
                                    "type": "tool_result",
                                    "tool_use_id": tc['id'],
                                    "content": json.dumps(tool_result)
                                }
                            ]
                        })
                    else:
                        # OpenAI format
                        llm_messages.append({
                            "role": "assistant",
                            "content": result.get('content', ''),
                            "tool_calls": [{
                                "id": tc['id'],
                                "type": "function",
                                "function": {
                                    "name": tc['name'],
                                    "arguments": json.dumps(tc['arguments'])
                                }
                            }]
                        })
                        llm_messages.append({
                            "role": "tool",
                            "tool_call_id": tc['id'],
                            "content": json.dumps(tool_result)
                        })
            else:
                return {
                    "content": result.get('content', ''),
                    "sources": sources,
                    "files_processed": len(extracted_content),
                    "tool_calls": tool_calls_made
                }
        
        return {
            "content": result.get('content', ''),
            "sources": sources,
            "files_processed": len(extracted_content),
            "tool_calls": tool_calls_made
        }
    else:
        result = await call_llm(llm_messages, agent.model_id)
        return {
            "content": result["content"],
            "sources": sources,
            "files_processed": len(extracted_content),
            "tool_calls": []
        }


async def extract_content_from_image(file_path: str) -> str:
    """Extract content from image using vision model - generic extraction"""
    try:
        import base64
        
        with open(file_path, "rb") as f:
            image_data = base64.b64encode(f.read()).decode('utf-8')
        
        ext = os.path.splitext(file_path)[1].lower()
        mime_types = {'.png': 'image/png', '.jpg': 'image/jpeg', '.jpeg': 'image/jpeg', '.gif': 'image/gif', '.webp': 'image/webp'}
        mime_type = mime_types.get(ext, 'image/png')
        
        messages = [
            {
                "role": "user",
                "content": [
                    {
                        "type": "image_url",
                        "image_url": {"url": f"data:{mime_type};base64,{image_data}"}
                    },
                    {
                        "type": "text",
                        "text": "Extract all text and information from this image. Preserve the structure and include all details, numbers, dates, and any other visible content."
                    }
                ]
            }
        ]
        
        result = await call_llm(messages, "gpt-4o")
        return result.get("content", "")
        
    except Exception as e:
        print(f"[Image Extract] Error: {e}")
        return ""


def extract_text_from_pdf(file_path: str) -> str:
    """Extract text from PDF"""
    try:
        import fitz
        doc = fitz.open(file_path)
        text = ""
        for page in doc:
            text += page.get_text()
        doc.close()
        return text
    except Exception as e:
        print(f"[PDF Extract] Error: {e}")
        return ""


def extract_text_from_docx(file_path: str) -> str:
    """Extract text from DOCX"""
    try:
        from docx import Document
        doc = Document(file_path)
        return "\n".join([para.text for para in doc.paragraphs])
    except Exception as e:
        print(f"[DOCX Extract] Error: {e}")
        return ""


# ========================================================================
# EMBEDDED SECURITY TOOL - Used by LLM to check user permissions
# ========================================================================
def build_security_tool_definition():
    """
    Build the embedded security tool that LLM uses to check permissions.
    This is an INTERNAL tool - not visible to users but used by the agent.
    """
    return {
        "type": "function",
        "function": {
            "name": "check_user_permissions",
            "description": "IMPORTANT: Call this tool FIRST before executing any task to verify the current user's identity and permissions. Returns the user's name, role, groups, and which tasks they can access.",
            "parameters": {
                "type": "object",
                "properties": {
                    "task_name": {
                        "type": "string",
                        "description": "The name of the task you want to execute (optional - if provided, checks if user can access this specific task)"
                    }
                },
                "required": []
            }
        }
    }


def execute_security_tool(current_user, access_control, agent, task_name: str = None) -> Dict:
    """
    Execute the embedded security tool.
    Returns user context and permission status.
    """
    # Build user info
    user_name = "Anonymous"
    user_role = "User"
    user_groups = []
    
    if current_user:
        if hasattr(current_user, 'first_name') and current_user.first_name:
            user_name = current_user.first_name
            if hasattr(current_user, 'last_name') and current_user.last_name:
                user_name += f" {current_user.last_name}"
        elif hasattr(current_user, 'email'):
            user_name = current_user.email.split('@')[0].replace('.', ' ').title()
        
        # Get role name
        if hasattr(current_user, 'role_ids') and current_user.role_ids and SECURITY_AVAILABLE and security_state:
            for rid in current_user.role_ids:
                role = security_state.roles.get(rid)
                if role:
                    user_role = role.name
                    break
        
        # Get group names
        if hasattr(current_user, 'group_ids') and current_user.group_ids and SECURITY_AVAILABLE and security_state:
            for gid in current_user.group_ids:
                group = security_state.groups.get(gid)
                if group:
                    user_groups.append(group.name)
    
    # Get permissions
    denied_tasks = []
    accessible_tasks = [t.name for t in agent.tasks]
    
    if access_control and hasattr(access_control, 'denied_tasks'):
        denied_tasks = access_control.denied_tasks or []
        accessible_tasks = [t for t in accessible_tasks if t not in denied_tasks]
    
    # Check specific task if requested
    task_allowed = True
    if task_name:
        task_allowed = task_name not in denied_tasks
    
    result = {
        "user": {
            "name": user_name,
            "role": user_role,
            "groups": user_groups
        },
        "permissions": {
            "accessible_tasks": accessible_tasks,
            "restricted_tasks": denied_tasks,
            "has_full_access": len(denied_tasks) == 0
        }
    }
    
    if task_name:
        result["task_check"] = {
            "task_name": task_name,
            "allowed": task_allowed,
            "reason": None if task_allowed else f"Task '{task_name}' is not available for {user_name}'s account"
        }
    
    return result


async def process_agent_chat(agent: AgentData, message: str, conversation: Conversation, timezone: str = None, access_control = None, current_user = None) -> Dict:
    """
    Process agent chat with access control enforcement.
    
    Args:
        agent: The agent to use
        message: User's message
        conversation: Current conversation
        timezone: User's timezone
        access_control: AccessCheckResult with denied_tasks and denied_tools
        current_user: The authenticated user making the request
    """
    agent_tools = get_agent_tools(agent)  # Resolves prefixed IDs
    
    # ========================================================================
    # ACCESS CONTROL: Filter tools based on permissions (Level 3)
    # ========================================================================
    denied_tool_ids = []
    if access_control and hasattr(access_control, 'denied_tools') and access_control.denied_tools:
        denied_tool_ids = access_control.denied_tools
        print(f"🔐 Filtering out {len(denied_tool_ids)} denied tools")
        agent_tools = [t for t in agent_tools if t.id not in denied_tool_ids]
    
    tool_ids = [t.id for t in agent_tools]
    search_results = search_documents(message, tool_ids, top_k=5)
    context = ""
    sources = []
    if search_results:
        context = "\n\nRELEVANT INFORMATION FROM KNOWLEDGE BASE:\n"
        for i, result in enumerate(search_results):
            context += f"\n[Source {i+1}: {result['source']}]\n{result['text'][:800]}\n"
            sources.append({"source": result['source'], "type": result['type'], "relevance": round(result['score'] * 100)})
    
    # Build system prompt with guardrails and FULL personality
    g = agent.guardrails
    p = agent.personality
    
    # Get current date/time for user's timezone
    date_info = get_current_datetime_for_user(timezone)
    
    # Personality descriptions based on sliders
    creativity_desc = "(Factual only)" if p.creativity <= 3 else "(Creative)" if p.creativity >= 7 else "(Balanced)"
    length_desc = "(Brief)" if p.length <= 3 else "(Detailed)" if p.length >= 7 else "(Moderate)"
    formality_desc = "(Casual)" if p.formality <= 3 else "(Formal)" if p.formality >= 7 else "(Professional)"
    empathy_desc = "(Direct)" if p.empathy <= 3 else "(Empathetic)" if p.empathy >= 7 else "(Supportive)"
    
    # Build tool definitions for action tools (non-knowledge)
    # These are already filtered by access control above
    action_tools = [t for t in agent_tools if t.type in ['email', 'api', 'webhook', 'slack', 'websearch', 'database']]
    tool_definitions = build_tool_definitions(action_tools) if action_tools else []
    
    # ========================================================================
    # ADD EMBEDDED SECURITY TOOL (LLM uses this to check permissions)
    # ========================================================================
    security_tool = build_security_tool_definition()
    tool_definitions = [security_tool] + tool_definitions  # Security tool FIRST
    print(f"🔐 Added embedded security tool. Total tools: {len(tool_definitions)}")
    
    # Add available tools to system prompt
    tools_description = ""
    tools_description = "\n\n=== AVAILABLE TOOLS ===\n"
    tools_description += "**🔐 check_user_permissions** (SECURITY): ALWAYS call this FIRST before executing any task. Returns the current user's identity and what they can access.\n"
    if action_tools:
        tools_description += "\nOther available tools:\n"
        for t in action_tools:
            tools_description += f"• **{t.name}** ({t.type}): {t.description or 'No description'}\n"
    tools_description += "\nWhen you need to use a tool, the system will automatically execute it for you.\n"
    
    # ========================================================================
    # ACCESS CONTROL: Filter tasks based on permissions (Level 2)
    # ========================================================================
    denied_task_names = []
    print(f"🔐 [PROCESS_CHAT] access_control param: {access_control}")
    print(f"🔐 [PROCESS_CHAT] current_user param: {current_user.email if current_user else None}")
    if access_control and hasattr(access_control, 'denied_tasks') and access_control.denied_tasks:
        denied_task_names = access_control.denied_tasks  # Now contains task NAMES, not IDs
        print(f"🔐 Denied tasks (by name): {denied_task_names}")
        
        # ========================================================================
        # CRITICAL: Also filter out tools associated with denied tasks
        # When a task is denied, the LLM should NOT have access to its tools
        # ========================================================================
        if denied_task_names:
            # Build a list of task-related keywords for tool filtering
            denied_keywords = []
            for task_name in denied_task_names:
                # Extract keywords from task name (lowercase, split by spaces)
                words = task_name.lower().replace('-', ' ').replace('_', ' ').split()
                denied_keywords.extend([w for w in words if len(w) > 3])
            
            # Filter out tools whose names contain denied task keywords
            original_tool_count = len(action_tools)
            filtered_action_tools = []
            for tool in action_tools:
                tool_name_lower = (tool.name or '').lower()
                tool_desc_lower = (tool.description or '').lower()
                
                # Check if tool matches any denied task keyword
                is_denied = False
                for keyword in denied_keywords:
                    if keyword in tool_name_lower or keyword in tool_desc_lower:
                        is_denied = True
                        print(f"   🔐 Tool '{tool.name}' FILTERED (matches denied keyword '{keyword}')")
                        break
                
                if not is_denied:
                    filtered_action_tools.append(tool)
            
            action_tools = filtered_action_tools
            print(f"🔐 Tools filtered: {original_tool_count} → {len(action_tools)}")
            
            # Rebuild tool definitions with filtered tools
            tool_definitions = build_tool_definitions(action_tools) if action_tools else []
            
            # Rebuild tools description
            tools_description = ""
            if action_tools:
                tools_description = "\n\n=== AVAILABLE TOOLS ===\nYou have access to the following tools. Use them when appropriate:\n"
                for t in action_tools:
                    tools_description += f"• **{t.name}** ({t.type}): {t.description or 'No description'}\n"
                tools_description += "\nWhen you need to use a tool, the system will automatically execute it for you.\n"
    else:
        print(f"🔐 [PROCESS_CHAT] No denied tasks - access_control={access_control is not None}, has denied_tasks={hasattr(access_control, 'denied_tasks') if access_control else False}")
    
    # Get tasks that user has access to - compare by NAME since IDs can change
    accessible_tasks = [task for task in agent.tasks if task.name not in denied_task_names]
    print(f"🔐 Accessible tasks: {[t.name for t in accessible_tasks]}")
    
    # ========================================================================
    # BUILD USER CONTEXT (Best Practice: Personalization with Security)
    # ========================================================================
    user_context = ""
    if current_user:
        # Get user's display name
        user_name = ""
        if hasattr(current_user, 'first_name') and current_user.first_name:
            user_name = current_user.first_name
            if hasattr(current_user, 'last_name') and current_user.last_name:
                user_name += f" {current_user.last_name}"
        elif hasattr(current_user, 'display_name') and current_user.display_name:
            user_name = current_user.display_name
        elif hasattr(current_user, 'email') and current_user.email:
            user_name = current_user.email.split('@')[0].replace('.', ' ').title()
        
        # Get user's groups (for context, not for permission decisions)
        user_groups = []
        if hasattr(current_user, 'group_ids') and current_user.group_ids:
            # Try to get group names from security_state
            try:
                if SECURITY_AVAILABLE and security_state:
                    for gid in current_user.group_ids:
                        group = security_state.groups.get(gid)
                        if group:
                            user_groups.append(group.name)
            except:
                pass
        
        # Get user's role (for context)
        user_role = ""
        if hasattr(current_user, 'role_ids') and current_user.role_ids:
            try:
                if SECURITY_AVAILABLE and security_state:
                    for rid in current_user.role_ids:
                        role = security_state.roles.get(rid)
                        if role:
                            user_role = role.name
                            break
            except:
                pass
        
        user_context = f"\n=== CURRENT USER ==="
        user_context += f"\n• Name: {user_name}" if user_name else ""
        user_context += f"\n• Role: {user_role}" if user_role else ""
        user_context += f"\n• Groups: {', '.join(user_groups)}" if user_groups else ""
        
        # Security instruction
        user_context += """

**USER CONTEXT RULES:**
- You may greet the user by name for a personalized experience
- NEVER reveal the user's permissions, role, or group membership
- NEVER discuss what other users can or cannot do
- NEVER explain the permission system or how access is determined
- Focus on what YOU can help with, not on what the user is allowed to do
"""
        
        print(f"👤 User context: name={user_name}, role={user_role}, groups={user_groups}")
    
    system_prompt = f"""You are {agent.name}.
{user_context}
=== LANGUAGE ===
{get_language_instruction(g.language)}

{date_info}
=== GOAL ===
{agent.goal}
"""

    # Add memory context if enabled and has memories
    if agent.memory_enabled and agent.memory:
        system_prompt += "\n=== MEMORY (Previous Conversations) ===\n"
        # Get last 5 memories
        recent_memories = agent.memory[-5:]
        for mem in recent_memories:
            system_prompt += f"• {mem.get('summary', '')}\n"
        system_prompt += "\nUse this context to provide continuity but don't reference it explicitly unless relevant.\n"

    system_prompt += f"""
=== PERSONALITY ===
• Tone: {p.tone}
• Voice: {p.voice}
• Creativity: {p.creativity}/10 {creativity_desc}
• Response Length: {p.length}/10 {length_desc}
• Formality: {p.formality}/10 {formality_desc}
• Empathy: {p.empathy}/10 {empathy_desc}

=== TASKS ==="""
    
    # Add instruction enforcement if available
    if INSTRUCTION_ENFORCER_AVAILABLE and InstructionEnforcer:
        # Detect language for enforcement
        enforce_lang = 'ar' if any('\u0600' <= c <= '\u06FF' for c in (request.message if hasattr(request, 'message') else '')) else 'en'
        phrases = InstructionEnforcer.ENFORCEMENT_PHRASES.get(enforce_lang, InstructionEnforcer.ENFORCEMENT_PHRASES['en'])
        
        system_prompt += f"""

=== ⚠️ INSTRUCTION ENFORCEMENT ===
{phrases['critical']}
{phrases['mandatory']}
{phrases['no_skip']}
"""
    
    # Only include tasks the user has access to - with enforced instructions
    for task in accessible_tasks:
        system_prompt += f"\n\n### TASK: {task.name}\n{task.description}"
        if task.instructions:
            if INSTRUCTION_ENFORCER_AVAILABLE and InstructionEnforcer:
                enforce_lang = 'ar' if g.language == 'ar' else 'en'
                system_prompt += InstructionEnforcer.format_instructions_for_prompt(
                    [inst.text for inst in task.instructions],
                    task.name,
                    enforce_lang
                )
            else:
                system_prompt += "\n**⚠️ MANDATORY STEPS (Execute ALL in order):**"
                for i, inst in enumerate(task.instructions, 1):
                    system_prompt += f"\n  [{i}] ✓ {inst.text}"
    
    # Add information about restricted tasks if any
    if denied_task_names and len(denied_task_names) > 0:
        system_prompt += "\n\n=== RESTRICTED TASKS (NOT AVAILABLE FOR THIS USER) ==="
        system_prompt += "\nThe following tasks are NOT available for this user due to permission settings:"
        for task_name in denied_task_names:
            system_prompt += f"\n• {task_name}"
        
        system_prompt += """

**🚫 CRITICAL SECURITY RULE - RESTRICTED TASKS:**
When the user asks for ANYTHING related to these restricted tasks, you MUST:
1. Politely decline the request
2. Explain this information is not available for their account
3. Suggest they contact their administrator if they need access
4. Do NOT provide any information, data, or details related to these tasks

⚠️ IMPORTANT: You do NOT have access to data for restricted tasks.
- DO NOT make up, guess, or hallucinate any data
- DO NOT provide partial information or workarounds
- DO NOT say "I can look that up" or similar - you CANNOT access this data
- ANY request related to restricted tasks (even specific items within them) must be declined

Example response:
"I'm sorry, but I don't have access to employee information with your current permissions. This data is restricted for your account. Please contact your administrator if you need access to this feature."

This applies to ALL variations of the request - whether asking for "all data", "specific items", or "just one piece of information"."""
    
    # Add security tool instructions
    system_prompt += """

=== 🔐 MANDATORY SECURITY PROTOCOL ===
You have access to a special security tool called 'check_user_permissions'.

**CRITICAL RULES:**
1. You ALREADY know the current user from the USER CONTEXT section above
2. If you're unsure about a user's permissions, call 'check_user_permissions' 
3. NEVER ask the user about their role, department, or permissions
4. If a task is in RESTRICTED TASKS, decline immediately without checking
5. For tasks that require data access, verify permissions internally using the tool

**REMEMBER:** The user's identity and permissions are KNOWN to you. Never ask them to verify their identity or role."""
    
    # Add tools description
    system_prompt += tools_description
    
    # Add guardrails to system prompt
    system_prompt += "\n\n=== GUARDRAILS (MUST FOLLOW) ==="
    
    # Anti-hallucination guardrails
    if g.anti_hallucination:
        system_prompt += "\n\n**ACCURACY:**"
        if g.cite_sources:
            system_prompt += "\n• Cite sources using [Source X] format when using knowledge base"
        if g.admit_uncertainty:
            system_prompt += "\n• If information is not in knowledge base and you cannot deduce it, say so"
        if g.verify_facts:
            system_prompt += "\n• Use facts from the knowledge base - don't make up company-specific data"
        if g.no_speculation:
            system_prompt += "\n• Don't speculate about company policies or data you don't have"
    
    # Content boundaries
    if g.avoid_topics:
        system_prompt += f"\n\n**AVOID:** {', '.join(g.avoid_topics)}"
    
    if g.focus_topics:
        system_prompt += f"\n\n**FOCUS ON:** {', '.join(g.focus_topics)}"
    
    # Response limits
    length_guide = {
        'short': '1-2 paragraphs MAX',
        'medium': '2-3 paragraphs',
        'long': 'Detailed as needed',
        'unlimited': 'Thorough coverage'
    }
    system_prompt += f"\n\n**RESPONSE LENGTH:** {length_guide.get(g.max_length, '2-3 paragraphs')}"
    
    # Language/dialect matching
    system_prompt += f"\n\n**LANGUAGE:** {get_language_instruction(g.language)}"
    
    # Escalation rules
    if g.escalate_angry or g.escalate_complex or g.escalate_request or g.escalate_sensitive:
        system_prompt += "\n\n**ESCALATION:** Offer human assistance if user is frustrated or issue is complex"
    
    # PII protection
    if g.pii_protection:
        system_prompt += "\n\n**PRIVACY:** Don't ask for or repeat sensitive personal info"
    
    # Add context
    system_prompt += context + """

=== IMPORTANT ===
• Follow the task instructions above carefully
• Be smart: calculate dates, times, durations when user provides relative terms
• Only ask for information that is truly missing
• Use markdown for formatting when helpful"""
    
    llm_messages = [{"role": "system", "content": system_prompt}]
    for msg in conversation.messages[-10:]:
        llm_messages.append({"role": msg.role, "content": msg.content})
    llm_messages.append({"role": "user", "content": message})
    
    # Detect if using Claude/Anthropic
    model_id = agent.model_id or ''
    is_anthropic = model_id.lower().startswith('claude')
    
    # If we have action tools, use tool-calling LLM
    tool_calls_made = []
    if tool_definitions:
        print(f"\n🤖 AGENT CHAT WITH TOOLS")
        print(f"   Agent: {agent.name}")
        print(f"   Model: {model_id} (Anthropic: {is_anthropic})")
        print(f"   Tools: {[t.name for t in action_tools]}")
        
        # Tool calling loop (max 5 iterations to prevent infinite loops)
        max_iterations = 5
        for iteration in range(max_iterations):
            print(f"\n   📍 Iteration {iteration + 1}/{max_iterations}")
            
            result = await call_llm_with_tools(llm_messages, tool_definitions, agent.model_id)
            
            if result.get('tool_calls'):
                # Execute each tool call
                for tc in result['tool_calls']:
                    print(f"\n   🔧 Tool call: {tc['name']}")
                    
                    # Check if it's the embedded security tool
                    if tc['name'] == 'check_user_permissions':
                        print(f"   🔐 Executing SECURITY TOOL")
                        task_name = tc.get('arguments', {}).get('task_name')
                        tool_result = execute_security_tool(current_user, access_control, agent, task_name)
                        tool_result = json.dumps(tool_result, indent=2)
                    else:
                        tool_result = await execute_tool(
                            tc.get('tool_id', ''),
                            tc.get('tool_type', ''),
                            tc.get('arguments', {})
                        )
                    
                    tool_calls_made.append({
                        "tool": tc['name'],
                        "arguments": tc['arguments'],
                        "result": tool_result
                    })
                    
                    # Format messages based on provider
                    if is_anthropic:
                        # Anthropic format: assistant with tool_use, then user with tool_result
                        llm_messages.append({
                            "role": "assistant",
                            "content": [
                                {"type": "text", "text": result.get('content', '') or "I'll use a tool to help with this."},
                                {
                                    "type": "tool_use",
                                    "id": tc['id'],
                                    "name": tc['name'],
                                    "input": tc['arguments']
                                }
                            ] if result.get('content') else [
                                {
                                    "type": "tool_use",
                                    "id": tc['id'],
                                    "name": tc['name'],
                                    "input": tc['arguments']
                                }
                            ]
                        })
                        llm_messages.append({
                            "role": "user",
                            "content": [
                                {
                                    "type": "tool_result",
                                    "tool_use_id": tc['id'],
                                    "content": json.dumps(tool_result)
                                }
                            ]
                        })
                    else:
                        # OpenAI format
                        llm_messages.append({
                            "role": "assistant",
                            "content": result.get('content', ''),
                            "tool_calls": [{
                                "id": tc['id'],
                                "type": "function",
                                "function": {
                                    "name": tc['name'],
                                    "arguments": json.dumps(tc['arguments'])
                                }
                            }]
                        })
                        llm_messages.append({
                            "role": "tool",
                            "tool_call_id": tc['id'],
                            "content": json.dumps(tool_result)
                        })
            else:
                # No more tool calls, we have final response
                final_content = result.get('content', '')
                if tool_calls_made:
                    # Add tool execution summary
                    tool_summary = "\n\n---\n**Actions performed:**\n"
                    for tc in tool_calls_made:
                        # Handle both dict and string results
                        tool_result = tc['result']
                        if isinstance(tool_result, str):
                            try:
                                tool_result = json.loads(tool_result)
                            except:
                                tool_result = {}
                        status = "✅" if isinstance(tool_result, dict) and tool_result.get('success') else "❌"
                        tool_summary += f"- {status} {tc['tool']}\n"
                    # Don't append summary to content, just log it
                    print(f"   📋 Tool Summary: {len(tool_calls_made)} tools executed")
                
                return {
                    "content": final_content,
                    "sources": sources,
                    "tool_calls": tool_calls_made
                }
        
        # Max iterations reached
        return {
            "content": result.get('content', 'I apologize, but I encountered an issue processing your request.'),
            "sources": sources,
            "tool_calls": tool_calls_made
        }
    else:
        # No tools, use regular LLM call
        result = await call_llm(llm_messages, agent.model_id)
        return {"content": result["content"], "sources": sources, "tool_calls": []}


# ============================================================================
# FastAPI Application
# ============================================================================

@asynccontextmanager
async def lifespan(app: FastAPI):
    try:
        print("🔥 Starting AgentForge v3.1...")
        app_state.load_from_disk()
        upload_dir = os.environ.get("UPLOAD_PATH", "data/uploads")
        os.makedirs(upload_dir, exist_ok=True)
        print(f"✅ Loaded {len(app_state.agents)} agents, {len(app_state.tools)} tools")
        
        # Load Security State
        if SECURITY_AVAILABLE:
            security_state.load_from_disk()
            print(f"✅ Security module loaded - {len(security_state.users)} users, {len(security_state.roles)} roles")
        
        # Test endpoints to catch any import/runtime errors
        print("🧪 Testing endpoints...")
        try:
            test_response = await root()
            print(f"✅ Root endpoint test: {test_response}")
        except Exception as test_error:
            print(f"❌❌❌ ROOT ENDPOINT TEST FAILED: {test_error}")
            import traceback
            traceback.print_exc()
            raise
        
        yield
        
        print("💾 Saving...")
        app_state.save_to_disk()
        
        # Save Security State
        if SECURITY_AVAILABLE:
            security_state.save_to_disk()
            print("✅ Security state saved")
    except Exception as e:
        print(f"❌❌❌ STARTUP ERROR: {e}")
        import traceback
        traceback.print_exc()
        raise


app = FastAPI(title="AgentForge", version="3.1.0", lifespan=lifespan)
app.add_middleware(CORSMiddleware, allow_origins=["*"], allow_credentials=True, allow_methods=["*"], allow_headers=["*"])

# Include Security Router
if SECURITY_AVAILABLE:
    app.include_router(security_router)
    print("✅ Security API endpoints registered")

# Include Health Check Router
if HEALTH_CHECK_AVAILABLE:
    app.include_router(health_router)
    print("✅ Health check endpoints registered")

# Include Access Control Module Router
try:
    from api.modules.access_control import router as access_control_router
    app.include_router(access_control_router)
    print("✅ Access Control module registered")
except ImportError as e:
    print(f"⚠️ Access Control module not available: {e}")

# Include Lab Module Router (Test Data Generator)
try:
    from api.modules.lab import lab_router
    app.include_router(lab_router)
    print("✅ Lab module registered")
except ImportError as e:
    print(f"⚠️ Lab module not available: {e}")

# Conversation Management Module
try:
    from api.modules.conversations import conversations_router
    app.include_router(conversations_router)
    print("✅ Conversation Management module registered")
except ImportError as e:
    print(f"⚠️ Conversation Management module not available: {e}")


@app.get("/")
async def root():
    try:
        return {"name": "AgentForge", "version": "3.2.0"}
    except Exception as e:
        print(f"❌ ROOT ENDPOINT ERROR: {e}")
        import traceback
        traceback.print_exc()
        raise


@app.get("/demo")
async def demo_redirect():
    """Redirect to Demo Lab page"""
    from starlette.responses import RedirectResponse
    return RedirectResponse(url="/#demo")


@app.get("/agent-hub.png")
async def get_agent_hub_icon():
    """Serve Agent Hub icon"""
    icon_paths = ["agent-hub.png", "/app/agent-hub.png", "ui/agent-hub.png"]
    for path in icon_paths:
        if os.path.exists(path):
            return FileResponse(path, media_type="image/png")
    raise HTTPException(404, "Icon not found")


@app.get("/agent-studio.png")
async def get_agent_studio_icon():
    """Serve Agent Studio icon"""
    icon_paths = ["agent-studio.png", "/app/agent-studio.png", "ui/agent-studio.png"]
    for path in icon_paths:
        if os.path.exists(path):
            return FileResponse(path, media_type="image/png")
    raise HTTPException(404, "Icon not found")


@app.get("/tools-icon.png")
async def get_tools_icon():
    """Serve Tools icon"""
    icon_paths = ["tools-icon.png", "/app/tools-icon.png", "ui/tools-icon.png"]
    for path in icon_paths:
        if os.path.exists(path):
            return FileResponse(path, media_type="image/png")
    raise HTTPException(404, "Icon not found")


@app.get("/health")
async def health():
    try:
        return {"status": "healthy", "agents": len(app_state.agents), "tools": len(app_state.tools)}
    except Exception as e:
        print(f"❌ HEALTH ENDPOINT ERROR: {e}")
        import traceback
        traceback.print_exc()
        raise


@app.get("/api/debug/agents-ownership")
async def debug_agents_ownership():
    """
    TEMPORARY DEBUG ENDPOINT: Show all agents and their owners.
    Remove this in production!
    """
    try:
        from database.services import AgentService
        db_agents = AgentService.get_all_agents("org_default")
        
        # Also get users for reference
        users_info = {}
        try:
            if SECURITY_AVAILABLE:
                for user_id, user in security_state.users.items():
                    users_info[user_id] = {
                        "email": user.email,
                        "name": user.get_display_name() if hasattr(user, 'get_display_name') else user.email
                    }
        except:
            pass
        
        return {
            "agents": [
                {
                    "id": str(a.get('id')),
                    "name": a.get('name'),
                    "status": a.get('status'),
                    "owner_id": str(a.get('owner_id')) if a.get('owner_id') else None,
                    "owner_email": users_info.get(str(a.get('owner_id')), {}).get('email', 'Unknown'),
                    "created_by": str(a.get('created_by')) if a.get('created_by') else None,
                    "created_by_email": users_info.get(str(a.get('created_by')), {}).get('email', 'Unknown'),
                }
                for a in (db_agents or [])
            ],
            "users": [
                {"id": uid, "email": u.get("email")}
                for uid, u in users_info.items()
            ],
            "total_agents": len(db_agents or []),
            "total_users": len(users_info)
        }
    except Exception as e:
        return {"error": str(e)}


# Agent Endpoints

@app.get("/api/agents/accessible")
async def list_accessible_agents(current_user: User = Depends(get_current_user)):
    """
    Get PUBLISHED agents that the current user has access to.
    Used by End User Portal to show only accessible agents.
    
    OWNERSHIP-BASED ACCESS:
    - Agents are PRIVATE by default
    - Owner always sees their agents
    - Others only see if explicitly granted access via Access Control
    """
    if not current_user:
        raise HTTPException(401, "Authentication required")
    
    org_id = current_user.org_id or "org_default"
    user_id = str(current_user.id)
    user_role_ids = getattr(current_user, 'role_ids', []) or []
    user_group_ids = get_user_group_ids(user_id) if user_id else []
    
    # Get all published agents with ownership info
    all_agents = []
    agent_ownership = {}  # Map agent_id -> owner_id
    
    try:
        from database.services import AgentService
        db_agents = AgentService.get_all_agents(org_id)
        if db_agents:
            for agent_dict in db_agents:
                try:
                    # Store ownership info
                    agent_id = agent_dict.get('id')
                    agent_ownership[agent_id] = {
                        'owner_id': agent_dict.get('owner_id'),
                        'created_by': agent_dict.get('created_by')
                    }
                    
                    agent_dict_clean = {k: v for k, v in agent_dict.items() if k in AgentData.__fields__}
                    agent_data = AgentData(**agent_dict_clean)
                    if agent_data.status == 'published':
                        all_agents.append(agent_data)
                except Exception as e:
                    continue
    except Exception as e:
        print(f"⚠️  [DATABASE ERROR] Failed to load agents: {e}")
        all_agents = [a for a in app_state.agents.values() if a.status == 'published']
    
    if not all_agents:
        all_agents = [a for a in app_state.agents.values() if a.status == 'published']
    
    # OWNERSHIP-BASED FILTERING: Private by default
    accessible_agents = []
    
    for agent in all_agents:
        # Check if user is the OWNER
        ownership_info = agent_ownership.get(agent.id, {})
        owner_id = ownership_info.get('owner_id')
        created_by = ownership_info.get('created_by')
        
        is_owner = (str(owner_id) == user_id) if owner_id else False
        is_creator = (str(created_by) == user_id) if created_by else False
        
        if is_owner or is_creator:
            # Owner/Creator always sees their published agents
            accessible_agents.append(agent)
            continue
        
        # Check if user is a DELEGATED ADMIN
        is_delegated_admin = False
        if ACCESS_CONTROL_AVAILABLE and AccessControlService:
            try:
                perm_result = AccessControlService.check_agent_permission(
                    user_id=user_id,
                    user_role_ids=user_role_ids,
                    user_group_ids=user_group_ids,
                    agent_id=agent.id,
                    org_id=org_id,
                    permission='any_admin'
                )
                is_delegated_admin = perm_result.get('has_permission', False)
                if is_delegated_admin:
                    print(f"   🔑 [ACCESSIBLE] User is DELEGATED ADMIN for agent {agent.id[:8]}...")
                    accessible_agents.append(agent)
                    continue
            except Exception as e:
                print(f"⚠️  Delegated admin check failed for agent {agent.id}: {e}")
        
        # For non-owners and non-admins, check Access Control permissions
        if ACCESS_CONTROL_AVAILABLE and AccessControlService:
            try:
                print(f"   🔐 [ACCESSIBLE] Checking access for agent {agent.id[:8]}... user={user_id[:8]}...")
                access_result = AccessControlService.check_user_access(
                    user_id=user_id,
                    user_role_ids=user_role_ids,
                    user_group_ids=user_group_ids,
                    agent_id=agent.id,
                    org_id=org_id
                )
                print(f"   📋 [ACCESSIBLE] Result: has_access={access_result.has_access}, reason={access_result.reason}")
                if access_result.has_access:
                    accessible_agents.append(agent)
                    print(f"   ✅ [ACCESSIBLE] Agent added to list")
            except Exception as e:
                print(f"⚠️  Access check failed for agent {agent.id}: {e}")
                import traceback
                traceback.print_exc()
                # FAIL SECURE: Don't show agent on error
                continue
        else:
            print(f"   ⚠️ [ACCESSIBLE] Access Control not available")
        # If no access control configured, agent remains private (not shown)
    
    return {
        "agents": [
            {
                "id": a.id,
                "name": a.name,
                "icon": a.icon,
                "goal": a.goal[:100] + "..." if len(a.goal) > 100 else a.goal,
                "description": a.goal[:150] if a.goal else "",
                "tasks_count": len(a.tasks),
                "tools_count": len(a.tool_ids),
                "status": a.status,
                "created_at": a.created_at
            }
            for a in accessible_agents
        ]
    }


@app.get("/api/agents")
async def list_agents(status: Optional[str] = None, current_user: User = Depends(get_current_user)):
    """
    List agents - OWNERSHIP BASED ACCESS CONTROL
    
    - Users only see agents they OWN or have been GRANTED access to
    - Agents are private by default until owner grants access
    """
    # Get user info for filtering
    user_id = str(current_user.id) if current_user else None
    org_id = current_user.org_id if current_user else "org_default"
    
    print(f"📋 [LIST_AGENTS] Request from user: {current_user.email if current_user else 'None'}, user_id: {user_id}, org_id: {org_id}, status: {status}")
    user_role_ids = getattr(current_user, 'role_ids', []) or []
    # IMPORTANT: Use get_user_group_ids to get groups from member_ids (source of truth)
    user_group_ids = get_user_group_ids(user_id) if user_id else []
    print(f"📋 [LIST_AGENTS] User groups (from member_ids): {user_group_ids}")
    
    # Try to load from database first, then fallback to in-memory
    all_agents = []
    agent_ownership = {}  # Map agent_id -> owner_id
    
    try:
        from database.services import AgentService
        db_agents = AgentService.get_all_agents(org_id)
        if db_agents:
            for agent_dict in db_agents:
                try:
                    # Store ownership info before cleaning
                    agent_id = agent_dict.get('id')
                    agent_ownership[agent_id] = {
                        'owner_id': agent_dict.get('owner_id'),
                        'created_by': agent_dict.get('created_by')
                    }
                    
                    # Remove extra fields that AgentData doesn't have
                    agent_dict_clean = {k: v for k, v in agent_dict.items() if k in AgentData.__fields__}
                    agent_data = AgentData(**agent_dict_clean)
                    all_agents.append(agent_data)
                    # Update in-memory cache
                    app_state.agents[agent_data.id] = agent_data
                except Exception as e:
                    print(f"⚠️  Error converting agent from database: {e}")
                    continue
    except Exception as e:
        print(f"⚠️  [DATABASE ERROR] Failed to load agents from database: {e}, using in-memory")
        all_agents = list(app_state.agents.values())
    
    # Fallback to in-memory if database loading failed or returned no agents
    if not all_agents:
        all_agents = list(app_state.agents.values())
    
    # Filter by status if specified
    if status:
        all_agents = [a for a in all_agents if a.status == status]
    
    # OWNERSHIP-BASED FILTERING: Only show agents user owns OR has been granted access
    visible_agents = []
    
    print(f"📋 [LIST_AGENTS] Filtering {len(all_agents)} agents for user {user_id[:8] if user_id else 'None'}...")
    
    for agent in all_agents:
        # Check if user is the OWNER
        ownership_info = agent_ownership.get(agent.id, {})
        owner_id = ownership_info.get('owner_id')
        created_by = ownership_info.get('created_by')
        
        is_owner = (str(owner_id) == user_id) if owner_id else False
        is_creator = (str(created_by) == user_id) if created_by else False
        
        print(f"   🔍 Agent {agent.id[:8]}... '{agent.name}': owner={owner_id}, created_by={created_by}, user_id={user_id}, is_owner={is_owner}, is_creator={is_creator}")
        
        if is_owner or is_creator:
            # Owner/Creator always sees their agents
            visible_agents.append((agent, str(owner_id) if owner_id else str(created_by)))
            continue
        
        # Check if user is a DELEGATED ADMIN for this agent (not just chat access)
        # Admin Portal should only show agents user can MANAGE, not just chat with
        if ACCESS_CONTROL_AVAILABLE and AccessControlService:
            try:
                print(f"   🔐 Checking MANAGEMENT permission for agent {agent.id[:8]}...")
                # Check if user has ANY management permission (not just full_admin)
                mgmt_result = AccessControlService.check_agent_permission(
                    user_id=user_id,
                    user_role_ids=user_role_ids,
                    user_group_ids=user_group_ids,
                    agent_id=agent.id,
                    org_id=org_id,
                    permission='any_admin'  # Any delegated admin permission
                )
                if mgmt_result.get('has_permission', False):
                    visible_agents.append((agent, str(owner_id) if owner_id else None))
                    print(f"   ✅ Agent added to visible list - user is delegated admin")
                else:
                    print(f"   ⛔ User is not owner or delegated admin - NOT shown in admin portal")
            except Exception as e:
                # On error, don't show the agent (fail secure)
                print(f"⚠️  Management check error for agent {agent.id}: {e}")
                continue
        else:
            print(f"   ⚠️ Access Control not available")
    
    print(f"✅ [LIST_AGENTS] Returning {len(visible_agents)} visible agents out of {len(all_agents)} total")
    
    return {
        "agents": [
            {
                "id": a.id,
                "name": a.name,
                "icon": a.icon,
                "goal": a.goal[:100] + "..." if len(a.goal) > 100 else a.goal,
                "tasks_count": len(a.tasks),
                "tools_count": len(a.tool_ids),
                "status": a.status,
                "created_at": a.created_at,
                "owner_id": owner_id,
                "is_owner": str(agent_ownership.get(a.id, {}).get('owner_id')) == user_id or str(agent_ownership.get(a.id, {}).get('created_by')) == user_id
            }
            for a, owner_id in visible_agents
        ]
    }


@app.get("/api/agents/{agent_id}")
async def get_agent(agent_id: str, current_user: User = Depends(get_current_user)):
    # Try to get agent from database first, then fallback to in-memory
    agent = None
    owner_id = None
    created_by = None
    admin_ids = []
    
    try:
        from database.services import AgentService
        org_id = current_user.org_id if current_user else "org_default"
        agent_dict = AgentService.get_agent_by_id(agent_id, org_id)
        if agent_dict:
            # Preserve ownership info before cleaning
            owner_id = agent_dict.get('owner_id')
            created_by = agent_dict.get('created_by')
            admin_ids = agent_dict.get('admin_ids') or []
            
            # Remove extra fields that AgentData doesn't have
            agent_dict_clean = {k: v for k, v in agent_dict.items() if k in AgentData.__fields__}
            agent = AgentData(**agent_dict_clean)
            app_state.agents[agent.id] = agent  # Update in-memory cache
    except Exception as e:
        print(f"⚠️  [DATABASE ERROR] Failed to load agent from database: {e}, using in-memory")
    
    # Fallback to in-memory if not found in database
    if not agent:
        if agent_id not in app_state.agents:
            raise HTTPException(404, "Agent not found")
        agent = app_state.agents[agent_id]
    
    # Resolve tools - handle prefixed IDs (api:xxx, kb:xxx)
    tools = []
    for tid in agent.tool_ids:
        # Try direct lookup first
        if tid in app_state.tools:
            tools.append(app_state.tools[tid].dict())
        else:
            # Try without prefix
            clean_tid = tid.replace('api:', '').replace('kb:', '')
            if clean_tid in app_state.tools:
                tools.append(app_state.tools[clean_tid].dict())
    
    # Include ownership info in response
    response = {**agent.dict(), "tools": tools}
    if owner_id:
        response["owner_id"] = owner_id
    if created_by:
        response["created_by"] = created_by
    if admin_ids:
        response["admin_ids"] = admin_ids
    
    return response


@app.get("/api/agents/{agent_id}/my-permissions")
async def get_my_agent_permissions(agent_id: str, current_user: User = Depends(get_current_user)):
    """
    Get the current user's permissions for a specific agent.
    Returns:
    - is_owner: True if user owns this agent
    - is_admin: True if user is a delegated admin
    - permissions: List of specific permissions (or ['full_admin'] if owner)
    - denied_task_names: List of tasks the user cannot access (for chat)
    """
    if not current_user:
        raise HTTPException(401, "Authentication required")
    
    user_id = str(current_user.id)
    org_id = current_user.org_id if current_user else "org_default"
    user_role_ids = getattr(current_user, 'role_ids', []) or []
    # IMPORTANT: Use get_user_group_ids to get groups from member_ids (source of truth)
    user_group_ids = get_user_group_ids(user_id) if user_id else []
    
    print(f"🔐 [MY-PERMS] user_id={user_id[:8]}..., user_group_ids={user_group_ids}")
    
    # Check if user is owner
    is_owner = False
    owner_id = None
    
    try:
        from database.services import AgentService
        agent_dict = AgentService.get_agent_by_id(agent_id, org_id)
        if agent_dict:
            owner_id = str(agent_dict.get('owner_id')) if agent_dict.get('owner_id') else None
            created_by = str(agent_dict.get('created_by')) if agent_dict.get('created_by') else None
            is_owner = (user_id == owner_id) or (user_id == created_by)
    except Exception as e:
        print(f"⚠️  Error checking ownership: {e}")
    
    if is_owner:
        # Owner has all permissions
        return {
            "is_owner": True,
            "is_admin": True,
            "permissions": ["full_admin"],
            "denied_task_names": [],
            "can_edit_basic_info": True,
            "can_edit_personality": True,
            "can_edit_model": True,
            "can_edit_guardrails": True,
            "can_manage_tasks": True,
            "can_manage_tools": True,
            "can_manage_knowledge": True,
            "can_manage_access": True,
            "can_manage_task_permissions": True,
            "can_publish": True,
            "can_delete": True
        }
    
    # Check if user is delegated admin
    permissions = []
    denied_task_names = []
    is_admin = False
    
    if ACCESS_CONTROL_AVAILABLE and AccessControlService:
        try:
            # Get admin permissions from policy
            from api.modules.access_control.service import get_session, normalize_org_id
            from database.models.agent_access import AgentAccessPolicy
            import json
            
            org_id = normalize_org_id(org_id)
            
            with get_session() as session:
                admin_policy = session.query(AgentAccessPolicy).filter(
                    AgentAccessPolicy.agent_id == agent_id,
                    AgentAccessPolicy.org_id == org_id,
                    AgentAccessPolicy.access_type == 'agent_admin',
                    AgentAccessPolicy.is_active == True
                ).first()
                
                if admin_policy:
                    # Check if user is in admin list
                    user_is_admin = user_id in (admin_policy.user_ids or [])
                    group_is_admin = any(g in (admin_policy.group_ids or []) for g in user_group_ids)
                    
                    if user_is_admin or group_is_admin:
                        is_admin = True
                        permissions = []  # Start with empty - must get from config
                        
                        # Get specific permissions from description
                        if admin_policy.description:
                            try:
                                admin_config = json.loads(admin_policy.description)
                                print(f"🔐 [MY-PERMS] Parsed admin_config: {admin_config}")
                                print(f"🔐 [MY-PERMS] Looking for user_id: {user_id}")
                                
                                if user_id in admin_config:
                                    entity_config = admin_config[user_id]
                                    print(f"🔐 [MY-PERMS] Found user config: {entity_config}")
                                    if isinstance(entity_config, dict):
                                        permissions = entity_config.get('permissions', [])
                                        denied_task_names = entity_config.get('denied_task_names', [])
                                    elif isinstance(entity_config, list):
                                        permissions = entity_config
                                else:
                                    # Check groups
                                    for group_id in user_group_ids:
                                        if group_id in admin_config:
                                            entity_config = admin_config[group_id]
                                            print(f"🔐 [MY-PERMS] Found group config for {group_id}: {entity_config}")
                                            if isinstance(entity_config, dict):
                                                permissions = entity_config.get('permissions', [])
                                                denied_task_names = entity_config.get('denied_task_names', [])
                                            elif isinstance(entity_config, list):
                                                permissions = entity_config
                                            break
                            except Exception as e:
                                print(f"⚠️  Error parsing admin permissions: {e}")
                        else:
                            # No description stored yet - means legacy data
                            print(f"⚠️  [MY-PERMS] No description in policy, defaulting to full_admin (legacy)")
                            permissions = ['full_admin']
        except Exception as e:
            print(f"⚠️  Error checking admin permissions: {e}")
    
    print(f"🔐 [MY-PERMS] Final permissions for user {user_id[:8]}...: {permissions}")
    
    if not is_admin:
        raise HTTPException(403, "You don't have management access to this agent")
    
    # Build permission flags
    has_full_admin = 'full_admin' in permissions
    
    return {
        "is_owner": False,
        "is_admin": True,
        "permissions": permissions,
        "denied_task_names": denied_task_names,
        "can_edit_basic_info": has_full_admin or 'edit_basic_info' in permissions,
        "can_edit_personality": has_full_admin or 'edit_personality' in permissions,
        "can_edit_model": has_full_admin or 'edit_model' in permissions,
        "can_edit_guardrails": has_full_admin or 'edit_guardrails' in permissions,
        "can_manage_tasks": has_full_admin or 'manage_tasks' in permissions,
        "can_manage_tools": has_full_admin or 'manage_tools' in permissions,
        "can_manage_knowledge": has_full_admin or 'manage_knowledge' in permissions,
        "can_manage_access": has_full_admin or 'manage_access' in permissions,
        "can_manage_task_permissions": has_full_admin or 'manage_task_permissions' in permissions,
        "can_publish": has_full_admin or 'publish_agent' in permissions,
        "can_delete": 'delete_agent' in permissions  # Never included in full_admin
    }


# ============================================================================
# EMAIL SERVICE
# ============================================================================

class EmailService:
    """Service for sending emails via Gmail API, Microsoft Graph, or SMTP"""
    
    @staticmethod
    async def send_gmail(
        access_token: str,
        to: str,
        subject: str,
        body: str,
        html: bool = False
    ) -> Dict[str, Any]:
        """Send email via Gmail API"""
        import base64
        from email.mime.text import MIMEText
        from email.mime.multipart import MIMEMultipart
        
        # Create message
        message = MIMEMultipart('alternative')
        message['to'] = to
        message['subject'] = subject
        
        if html:
            part = MIMEText(body, 'html')
        else:
            part = MIMEText(body, 'plain')
        message.attach(part)
        
        # Encode message
        raw_message = base64.urlsafe_b64encode(message.as_bytes()).decode('utf-8')
        
        # Send via Gmail API
        async with httpx.AsyncClient() as client:
            print(f"📧 Sending email via Gmail API to: {to}")
            print(f"   Subject: {subject}")
            print(f"   Token (first 20 chars): {access_token[:20]}...")
            
            response = await client.post(
                "https://gmail.googleapis.com/gmail/v1/users/me/messages/send",
                headers={
                    "Authorization": f"Bearer {access_token}",
                    "Content-Type": "application/json"
                },
                json={"raw": raw_message}
            )
            
            print(f"   Gmail API Response Status: {response.status_code}")
            
            if response.status_code == 200:
                result = response.json()
                print(f"   ✅ Email sent! Message ID: {result.get('id')}")
                return {
                    "success": True,
                    "message_id": result.get("id"),
                    "thread_id": result.get("threadId")
                }
            else:
                error_data = response.json()
                error_msg = error_data.get("error", {}).get("message", "Failed to send email")
                print(f"   ❌ Gmail API Error: {error_msg}")
                print(f"   Full error: {error_data}")
                return {
                    "success": False,
                    "error": error_msg,
                    "details": error_data
                }
    
    @staticmethod
    async def refresh_google_token(refresh_token: str) -> Optional[str]:
        """Refresh Google access token"""
        client_id = os.environ.get("GOOGLE_CLIENT_ID")
        client_secret = os.environ.get("GOOGLE_CLIENT_SECRET")
        
        if not client_id or not client_secret:
            print("⚠️ Cannot refresh token: GOOGLE_CLIENT_ID or GOOGLE_CLIENT_SECRET not set")
            return None
        
        print(f"🔄 Refreshing Google access token...")
        
        async with httpx.AsyncClient() as client:
            response = await client.post(
                "https://oauth2.googleapis.com/token",
                data={
                    "client_id": client_id,
                    "client_secret": client_secret,
                    "refresh_token": refresh_token,
                    "grant_type": "refresh_token"
                }
            )
            
            print(f"   Token refresh response: {response.status_code}")
            
            if response.status_code == 200:
                tokens = response.json()
                print(f"   ✅ Token refreshed successfully")
                return tokens.get("access_token")
            else:
                error_data = response.json()
                print(f"   ❌ Token refresh failed: {error_data}")
                return None
    
    @staticmethod
    async def send_smtp(
        host: str,
        port: int,
        username: str,
        password: str,
        to: str,
        subject: str,
        body: str,
        from_email: str = None,
        html: bool = False
    ) -> Dict[str, Any]:
        """Send email via SMTP"""
        import smtplib
        from email.mime.text import MIMEText
        from email.mime.multipart import MIMEMultipart
        
        try:
            # Create message
            message = MIMEMultipart('alternative')
            message['From'] = from_email or username
            message['To'] = to
            message['Subject'] = subject
            
            if html:
                part = MIMEText(body, 'html')
            else:
                part = MIMEText(body, 'plain')
            message.attach(part)
            
            # Send via SMTP
            with smtplib.SMTP(host, port) as server:
                server.starttls()
                server.login(username, password)
                server.send_message(message)
            
            return {
                "success": True,
                "message": f"Email sent to {to}"
            }
        except Exception as e:
            return {
                "success": False,
                "error": str(e)
            }
    
    @staticmethod
    async def send_sendgrid(
        api_key: str,
        to: str,
        subject: str,
        body: str,
        from_email: str,
        html: bool = False
    ) -> Dict[str, Any]:
        """Send email via SendGrid API"""
        async with httpx.AsyncClient() as client:
            response = await client.post(
                "https://api.sendgrid.com/v3/mail/send",
                headers={
                    "Authorization": f"Bearer {api_key}",
                    "Content-Type": "application/json"
                },
                json={
                    "personalizations": [{"to": [{"email": to}]}],
                    "from": {"email": from_email},
                    "subject": subject,
                    "content": [{
                        "type": "text/html" if html else "text/plain",
                        "value": body
                    }]
                }
            )
            
            if response.status_code in [200, 201, 202]:
                return {
                    "success": True,
                    "message": f"Email sent to {to} via SendGrid"
                }
            else:
                return {
                    "success": False,
                    "error": f"SendGrid error: {response.status_code}"
                }


# ============================================================================
# OAUTH ENDPOINTS
# ============================================================================

@app.get("/api/oauth/google/email/url")
async def google_oauth_url():
    """Get Google OAuth URL for email integration"""
    client_id = os.environ.get("GOOGLE_CLIENT_ID")
    if not client_id:
        raise HTTPException(400, "Google OAuth not configured. Set GOOGLE_CLIENT_ID and GOOGLE_CLIENT_SECRET environment variables.")
    
    from urllib.parse import urlencode
    
    # Get the redirect URI based on the request
    redirect_uri = f"{os.environ.get('APP_URL', 'http://localhost:8000')}/api/oauth/google/callback"
    
    params = {
        "client_id": client_id,
        "redirect_uri": redirect_uri,
        "response_type": "code",
        "scope": "https://www.googleapis.com/auth/gmail.send https://www.googleapis.com/auth/userinfo.email",
        "access_type": "offline",
        "prompt": "consent"
    }
    
    url = f"https://accounts.google.com/o/oauth2/v2/auth?{urlencode(params)}"
    return {"url": url}


@app.get("/api/oauth/google/callback")
async def google_oauth_callback(code: str = None, error: str = None):
    """Handle Google OAuth callback"""
    if error:
        return HTMLResponse(f"""
            <html><body>
                <script>
                    window.opener.postMessage({{ type: 'google-oauth-error', error: '{error}' }}, '*');
                    window.close();
                </script>
            </body></html>
        """)
    
    if not code:
        raise HTTPException(400, "No authorization code received")
    
    client_id = os.environ.get("GOOGLE_CLIENT_ID")
    client_secret = os.environ.get("GOOGLE_CLIENT_SECRET")
    redirect_uri = f"{os.environ.get('APP_URL', 'http://localhost:8000')}/api/oauth/google/callback"
    
    # Exchange code for tokens
    import httpx
    async with httpx.AsyncClient() as client:
        token_response = await client.post(
            "https://oauth2.googleapis.com/token",
            data={
                "client_id": client_id,
                "client_secret": client_secret,
                "code": code,
                "grant_type": "authorization_code",
                "redirect_uri": redirect_uri
            }
        )
        
        if token_response.status_code != 200:
            return HTMLResponse(f"""
                <html><body>
                    <script>
                        window.opener.postMessage({{ type: 'google-oauth-error', error: 'Token exchange failed' }}, '*');
                        window.close();
                    </script>
                </body></html>
            """)
        
        tokens = token_response.json()
        
        # Get user info
        user_response = await client.get(
            "https://www.googleapis.com/oauth2/v2/userinfo",
            headers={"Authorization": f"Bearer {tokens['access_token']}"}
        )
        
        user_info = user_response.json()
        email = user_info.get("email", "unknown")
        
        # Store tokens temporarily with email as key
        # These will be saved to tool config when tool is created
        access_token = tokens.get("access_token", "")
        refresh_token = tokens.get("refresh_token", "")
        
        # Return success to popup with both tokens
        return HTMLResponse(f"""
            <html><body>
                <script>
                    window.opener.postMessage({{
                        type: 'google-oauth-success',
                        email: '{email}',
                        access_token: '{access_token}',
                        refresh_token: '{refresh_token}'
                    }}, '*');
                    window.close();
                </script>
            </body></html>
        """)


@app.get("/api/oauth/microsoft/email/url")
async def microsoft_oauth_url():
    """Get Microsoft OAuth URL for email integration"""
    client_id = os.environ.get("MICROSOFT_CLIENT_ID")
    if not client_id:
        raise HTTPException(400, "Microsoft OAuth not configured")
    
    from urllib.parse import urlencode
    redirect_uri = f"{os.environ.get('APP_URL', 'http://localhost:8000')}/api/oauth/microsoft/callback"
    
    params = {
        "client_id": client_id,
        "redirect_uri": redirect_uri,
        "response_type": "code",
        "scope": "https://graph.microsoft.com/Mail.Send User.Read offline_access",
        "response_mode": "query"
    }
    
    url = f"https://login.microsoftonline.com/common/oauth2/v2.0/authorize?{urlencode(params)}"
    return {"url": url}


@app.get("/AgentForge_Logo.png")
async def get_logo():
    """Serve the AgentForge logo"""
    logo_paths = [
        "AgentForge_Logo.png",
        "data/AgentForge_Logo.png",
        "/mnt/user-data/uploads/AgentForge_Logo.png",
        os.path.join(os.environ.get("UPLOAD_PATH", "data/uploads"), "AgentForge_Logo.png")
    ]
    
    for path in logo_paths:
        if os.path.exists(path):
            return FileResponse(path, media_type="image/png")
    
    raise HTTPException(404, "Logo not found")


@app.post("/api/agents")
async def create_agent(request: CreateAgentRequest, current_user: User = Depends(get_current_user)):
    tasks = []
    for t in request.tasks:
        instructions = []
        for i in t.get('instructions', []):
            if isinstance(i, dict): instructions.append(TaskInstruction(**i))
            elif isinstance(i, str): instructions.append(TaskInstruction(text=i))
        tasks.append(TaskDefinition(name=t.get('name', ''), description=t.get('description', ''), instructions=instructions))
    
    # Parse guardrails
    guardrails = AgentGuardrails()
    if request.guardrails:
        g = request.guardrails
        guardrails = AgentGuardrails(
            anti_hallucination=g.get('antiHallucination', True),
            cite_sources=g.get('citeSources', True),
            admit_uncertainty=g.get('admitUncertainty', True),
            verify_facts=g.get('verifyFacts', True),
            no_speculation=g.get('noSpeculation', False),
            avoid_topics=[t.strip() for t in g.get('avoidTopics', '').split('\n') if t.strip()] if isinstance(g.get('avoidTopics'), str) else g.get('avoidTopics', []),
            focus_topics=[t.strip() for t in g.get('focusTopics', '').split('\n') if t.strip()] if isinstance(g.get('focusTopics'), str) else g.get('focusTopics', []),
            max_length=g.get('maxLength', 'medium'),
            language=g.get('language', 'user'),
            escalate_angry=g.get('escalateAngry', True),
            escalate_complex=g.get('escalateComplex', True),
            escalate_request=g.get('escalateRequest', True),
            escalate_sensitive=g.get('escalateSensitive', False),
            pii_protection=g.get('piiProtection', True),
            mask_pii=g.get('maskPii', True),
            no_store_pii=g.get('noStorePii', True)
        )
    
    agent = AgentData(
        name=request.name, 
        goal=request.goal, 
        personality=AgentPersonality(**request.personality) if request.personality else AgentPersonality(),
        guardrails=guardrails,
        tasks=tasks, 
        tool_ids=request.tool_ids, 
        model_id=request.model_id, 
        status=request.status
    )
    
    # Get org_id and owner_id from the AUTHENTICATED USER
    org_id = current_user.org_id if current_user else "org_default"
    owner_id = current_user.id if current_user else None
    created_by = current_user.id if current_user else None
    
    # Log ownership info
    if current_user:
        print(f"👤 Creating agent with owner_id={current_user.id[:8]}... (user: {current_user.email})")
    
    # Save to database using AgentService
    try:
        from database.services import AgentService
        agent_dict = agent.dict()
        saved_agent = AgentService.save_agent(
            agent_dict,
            org_id=org_id,
            owner_id=owner_id or "00000000-0000-0000-0000-000000000000",  # Fallback UUID
            created_by=created_by or owner_id or "00000000-0000-0000-0000-000000000000",
            updated_by=created_by or owner_id
        )
        # Update agent with saved data (including ID from database)
        agent.id = saved_agent['id']
        app_state.agents[agent.id] = agent
    except Exception as e:
        print(f"⚠️  [DATABASE ERROR] Failed to save agent to database: {e}, saving to in-memory only")
        import traceback
        traceback.print_exc()
        app_state.agents[agent.id] = agent
        app_state.save_to_disk()  # Will try to save to database in save_to_disk()
    
    return {
        "status": "success", 
        "agent_id": agent.id, 
        "agent": agent.dict(),
        "owner_id": str(owner_id) if owner_id else None,
        "created_by": str(created_by) if created_by else None
    }


class GenerateAgentConfigRequest(BaseModel):
    goal: str
    update_mode: Optional[bool] = False  # True when updating existing config

@app.post("/api/agents/generate-config")
async def generate_agent_config(request: GenerateAgentConfigRequest):
    """Use AI to generate a complete agent configuration based on a goal description - 100% LLM powered, no hardcoding"""
    try:
        goal = request.goal
        update_mode = request.update_mode or False
        
        if not goal or len(goal.strip()) < 10:
            raise HTTPException(400, "Please provide a more detailed goal description (at least 10 characters)")
        
        # Get configured LLM providers to include in prompt
        available_models = []
        tool_capable_models = []  # Models that support function calling
        
        for provider in app_state.settings.llm_providers:
            for model in provider.models:
                model_lower = model.lower()
                supports_tools = (
                    model_lower.startswith('gpt') or 
                    model_lower.startswith('o1') or 
                    model_lower.startswith('o3') or
                    model_lower.startswith('claude') or
                    model_lower.startswith('gemini')  # Gemini now supports tools!
                )
                available_models.append({
                    "id": model,
                    "provider": provider.name,
                    "supports_tools": supports_tools
                })
                if supports_tools:
                    tool_capable_models.append(model)
        
        # Build available models string for the prompt with tool support info
        if available_models:
            models_list = "\n".join([
                f"- {m['id']} ({m['provider']}) {'✅ SUPPORTS TOOLS' if m['supports_tools'] else '⚠️ NO TOOL SUPPORT'}" 
                for m in available_models
            ])
            if tool_capable_models:
                tool_models_note = f"\n\n⚠️ CRITICAL: For agents using tools (APIs, databases, etc.), you MUST choose one of these models: {', '.join(tool_capable_models)}"
            else:
                tool_models_note = "\n\n⚠️ WARNING: No tool-capable models are configured. Agent may not work properly with tools."
        else:
            models_list = "- gpt-4o (OpenAI) - default if no models configured"
            tool_models_note = ""
        
        # Find the best available LLM to use for generation
        generation_llm = await _get_generation_llm()
        if not generation_llm:
            raise HTTPException(503, "No LLM configured. Please add at least one LLM provider in Settings.")
        
        # Build the dynamic prompt
        system_prompt = f"""You are an expert AI agent architect. Your job is to design the perfect agent configuration based on the user's goal.

IMPORTANT: You must analyze the goal deeply and generate ALL configurations dynamically. Do NOT use generic/default values - every field must be specifically tailored to this agent's purpose.

AVAILABLE AI MODELS (only recommend from this list):
{models_list}
{tool_models_note}

Return a JSON object with these fields:

1. "name": A unique, descriptive name (2-4 words) that captures the agent's purpose

2. "icon": A single emoji that best represents this agent's role

3. "goal": Enhanced version of the goal (2-3 clear sentences). Make it specific and actionable.

4. "model": Choose the BEST model from the available list above. 
   
   🚨 CRITICAL RULE: If the agent needs ANY tools (API calls, database queries, web search, email, etc.), 
   you MUST select a model marked with "✅ SUPPORTS TOOLS". Models without tool support CANNOT use tools!
   
   Tool-capable models: OpenAI (gpt-*), Anthropic (claude-*), Google (gemini-*)
   
   Consider:
   - Agent uses tools → Choose gpt-4o, gpt-4o-mini, claude-sonnet-4, gemini-1.5-pro, etc.
   - Complex reasoning/analysis → gpt-4o, claude-sonnet-4, gemini-1.5-pro
   - High volume/simple tasks → gpt-4o-mini, gemini-1.5-flash
   - Creative/writing tasks → claude models
   - Code tasks → gpt-4o or claude

5. "modelReason": Explain WHY this specific model is best for THIS agent (be specific, not generic)

6. "personality": Object with 6 traits (each 1-10), specifically calibrated for this agent:
   - "creativity": 1=Only verified facts, 10=Freely brainstorm ideas
   - "length": 1=One-line answers, 10=Comprehensive explanations  
   - "formality": 1=Casual chat style, 10=Executive/legal formal
   - "empathy": 1=Pure facts, no emotion, 10=Deeply supportive
   - "proactivity": 1=Only answer what's asked, 10=Anticipate needs & suggest
   - "confidence": 1=Many disclaimers, 10=Authoritative expert
   
   Think about: Who uses this agent? What situations? What tone fits?

7. "personalityReason": 2-3 sentences explaining WHY these specific values fit this agent

8. "personalityDescriptions": Object with description for each trait's current value:
   - "creativityDesc": What this creativity level means for THIS agent
   - "lengthDesc": What this length level means for THIS agent
   - "formalityDesc": What this formality level means for THIS agent
   - "empathyDesc": What this empathy level means for THIS agent
   - "proactivityDesc": What this proactivity level means for THIS agent
   - "confidenceDesc": What this confidence level means for THIS agent

9. "tasks": Array of 3-5 tasks this agent should handle. Each task:
   - "name": Task name (2-4 words)
   - "description": What this accomplishes (1 sentence)
   - "instructions": Array of 4-6 SMART, SPECIFIC instructions (not generic!)
   
   BAD instruction: "Help the user"
   GOOD instruction: "When user uploads a document, extract key data automatically instead of asking them to type it"

10. "suggestedTools": Array of tools this agent needs:
    - "knowledge" (documents/FAQs)
    - "database" (structured data)
    - "email" (send emails)
    - "calendar" (scheduling)
    - "crm" (customer data)
    - "websearch" (internet search)
    - "code" (run code)
    - "slack" (messaging)
    - "ocr" (read images/PDFs)
    - "api" (external APIs)

11. "guardrails": Safety settings object:
    - "antiHallucination": Should agent strictly avoid making things up?
    - "citeSources": Should agent cite information sources?
    - "escalateComplex": Should complex issues go to humans?
    - "escalateSensitive": Should sensitive topics go to humans?
    - "piiProtection": Should agent protect personal information?
    
    Think about: What could go wrong? What's risky for this use case?

12. "guardrailsReason": Why these guardrail settings fit this agent

Remember: EVERY value must be specifically chosen for THIS agent's goal. No defaults!"""

        user_message = f"Design the perfect agent configuration for this goal:\n\n{goal}"
        
        # Call the LLM
        try:
            response_text = await generation_llm.generate(
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": user_message}
                ],
                temperature=0.7,
                max_tokens=4000
            )
            
            # Parse JSON from response
            # Handle potential markdown code blocks
            response_text = response_text.strip()
            if response_text.startswith("```json"):
                response_text = response_text[7:]
            if response_text.startswith("```"):
                response_text = response_text[3:]
            if response_text.endswith("```"):
                response_text = response_text[:-3]
            
            config = json.loads(response_text.strip())
            
            # Validate required fields
            required_fields = ["name", "icon", "goal", "model", "personality", "tasks"]
            for field in required_fields:
                if field not in config:
                    raise ValueError(f"Missing required field: {field}")
            
            # Ensure model is from available models (or use first available)
            if available_models:
                available_ids = [m["id"] for m in available_models]
                if config["model"] not in available_ids:
                    config["model"] = available_ids[0]
                    config["modelReason"] = f"Selected {available_ids[0]} as it's available in your configured providers."
            
            print(f"[Agent Config] ✅ Generated 100% dynamic config for: {goal[:50]}...")
            return config
            
        except json.JSONDecodeError as e:
            print(f"[Agent Config] JSON parse error: {e}")
            print(f"[Agent Config] Response was: {response_text[:500]}...")
            raise HTTPException(500, "AI returned invalid configuration format. Please try again.")
            
    except HTTPException:
        raise
    except Exception as e:
        print(f"[Agent Config] Error: {e}")
        import traceback
        traceback.print_exc()
        raise HTTPException(500, f"Failed to generate agent configuration: {str(e)}")


async def _get_generation_llm():
    """Get the best available LLM for agent generation from configured providers"""
    
    print(f"[Generation LLM] Checking configured providers...")
    print(f"[Generation LLM] Number of providers: {len(app_state.settings.llm_providers)}")
    
    # Check configured providers first
    priority_providers = ["openai", "anthropic", "google", "groq", "mistral", "deepseek"]
    
    for provider_name in priority_providers:
        for provider in app_state.settings.llm_providers:
            provider_type = provider.provider.lower() if provider.provider else ""
            print(f"[Generation LLM] Checking provider: {provider.name} ({provider_type})")
            
            if provider_type == provider_name and provider.api_key:
                print(f"[Generation LLM] Found matching provider: {provider.name}")
                try:
                    # Use direct API call for reliability
                    if provider_type == "openai":
                        model = provider.models[0] if provider.models else "gpt-4o"
                        return OpenAIDirectLLM(api_key=provider.api_key, model=model)
                    elif provider_type == "anthropic":
                        model = provider.models[0] if provider.models else "claude-sonnet-4-20250514"
                        return AnthropicDirectLLM(api_key=provider.api_key, model=model)
                    else:
                        # Try ProviderFactory for other providers
                        llm = await ProviderFactory.create(
                            LLMProvider(provider.provider),
                            api_key=provider.api_key,
                            api_base=provider.api_base,
                            model=provider.models[0] if provider.models else None
                        )
                        return llm
                except Exception as e:
                    print(f"[Generation LLM] Failed to create {provider_name}: {e}")
                    continue
    
    # Fallback to environment variables only if no configured providers
    openai_key = os.getenv("OPENAI_API_KEY")
    if openai_key:
        print("[Generation LLM] Using OpenAI from environment variable as fallback")
        return OpenAIDirectLLM(api_key=openai_key, model="gpt-4o")
    
    print("[Generation LLM] No LLM provider found!")
    return None


class OpenAIDirectLLM:
    """Simple OpenAI wrapper for generation"""
    def __init__(self, api_key: str, model: str = "gpt-4o"):
        self.api_key = api_key
        self.model = model
    
    async def generate(self, messages: list, temperature: float = 0.7, max_tokens: int = 4000) -> str:
        async with httpx.AsyncClient(timeout=120.0) as client:
            response = await client.post(
                "https://api.openai.com/v1/chat/completions",
                headers={
                    "Authorization": f"Bearer {self.api_key}",
                    "Content-Type": "application/json"
                },
                json={
                    "model": self.model,
                    "messages": messages,
                    "max_tokens": max_tokens,
                    "temperature": temperature,
                    "response_format": {"type": "json_object"}
                }
            )
            
            if response.status_code != 200:
                raise Exception(f"OpenAI API error: {response.status_code} - {response.text}")
            
            result = response.json()
            return result["choices"][0]["message"]["content"]


class AnthropicDirectLLM:
    """Simple Anthropic wrapper for generation"""
    def __init__(self, api_key: str, model: str = "claude-sonnet-4-20250514"):
        self.api_key = api_key
        self.model = model
    
    async def generate(self, messages: list, temperature: float = 0.7, max_tokens: int = 4000) -> str:
        # Convert messages format
        system_msg = ""
        user_messages = []
        for msg in messages:
            if msg["role"] == "system":
                system_msg = msg["content"]
            else:
                user_messages.append(msg)
        
        async with httpx.AsyncClient(timeout=120.0) as client:
            response = await client.post(
                "https://api.anthropic.com/v1/messages",
                headers={
                    "x-api-key": self.api_key,
                    "Content-Type": "application/json",
                    "anthropic-version": "2023-06-01"
                },
                json={
                    "model": self.model,
                    "max_tokens": max_tokens,
                    "system": system_msg,
                    "messages": user_messages
                }
            )
            
            if response.status_code != 200:
                raise Exception(f"Anthropic API error: {response.status_code} - {response.text}")
            
            result = response.json()
            return result["content"][0]["text"]


class ConfigureToolsRequest(BaseModel):
    goal: str
    agent_name: Optional[str] = None
    tasks: Optional[List[Dict[str, Any]]] = None


@app.post("/api/agents/configure-tools")
async def configure_tools_for_agent(request: ConfigureToolsRequest):
    """AI-powered tool recommendations based on agent goal and tasks"""
    try:
        goal_lower = request.goal.lower()
        tools = []
        demo_suggestions = {
            "apis": [],
            "documents": [],
            "images": []
        }
        
        # Customer Support
        if any(kw in goal_lower for kw in ["customer", "support", "help", "service", "ticket"]):
            tools = ["knowledge", "database", "email"]
            demo_suggestions = {
                "apis": [
                    {"name": "Customer Orders API", "endpoint": "/api/orders", "description": "Track order status"},
                    {"name": "Tickets API", "endpoint": "/api/tickets", "description": "Support ticket management"}
                ],
                "documents": [
                    {"name": "FAQ Document", "type": "faq"},
                    {"name": "Return Policy", "type": "policy"},
                    {"name": "Product Manual", "type": "manual"}
                ],
                "images": [
                    {"name": "Product Images", "count": 5}
                ]
            }
        
        # Sales
        elif any(kw in goal_lower for kw in ["sales", "sell", "lead", "crm", "deal"]):
            tools = ["crm", "database", "email", "knowledge"]
            demo_suggestions = {
                "apis": [
                    {"name": "Products Catalog API", "endpoint": "/api/products", "description": "Product listings"},
                    {"name": "Pricing API", "endpoint": "/api/pricing", "description": "Price lookups"},
                    {"name": "Leads API", "endpoint": "/api/leads", "description": "Lead management"}
                ],
                "documents": [
                    {"name": "Product Catalog", "type": "catalog"},
                    {"name": "Pricing Guide", "type": "pricing"},
                    {"name": "Sales Playbook", "type": "playbook"}
                ],
                "images": [
                    {"name": "Product Photos", "count": 10}
                ]
            }
        
        # HR
        elif any(kw in goal_lower for kw in ["hr", "employee", "staff", "hiring", "human resource"]):
            tools = ["knowledge", "calendar", "database"]
            demo_suggestions = {
                "apis": [
                    {"name": "Employees API", "endpoint": "/api/employees", "description": "Employee directory"},
                    {"name": "Leave API", "endpoint": "/api/leave", "description": "Leave management"}
                ],
                "documents": [
                    {"name": "Employee Handbook", "type": "handbook"},
                    {"name": "Benefits Guide", "type": "benefits"},
                    {"name": "Leave Policy", "type": "policy"}
                ],
                "images": []
            }
        
        # Research
        elif any(kw in goal_lower for kw in ["research", "analysis", "report", "data"]):
            tools = ["websearch", "knowledge", "code"]
            demo_suggestions = {
                "apis": [
                    {"name": "Data API", "endpoint": "/api/data", "description": "Data access"}
                ],
                "documents": [
                    {"name": "Research Template", "type": "template"},
                    {"name": "Analysis Framework", "type": "framework"}
                ],
                "images": []
            }
        
        # Coding
        elif any(kw in goal_lower for kw in ["code", "developer", "programming", "software"]):
            tools = ["code", "websearch", "knowledge"]
            demo_suggestions = {
                "apis": [],
                "documents": [
                    {"name": "API Documentation", "type": "api_doc"},
                    {"name": "Code Standards", "type": "standards"}
                ],
                "images": []
            }
        
        # Default
        else:
            tools = ["knowledge", "websearch"]
            demo_suggestions = {
                "apis": [
                    {"name": "Sample API", "endpoint": "/api/sample", "description": "Sample data"}
                ],
                "documents": [
                    {"name": "Knowledge Base", "type": "kb"}
                ],
                "images": []
            }
        
        return {
            "tools": tools,
            "demo_suggestions": demo_suggestions
        }
        
    except Exception as e:
        print(f"[Configure Tools] Error: {e}")
        raise HTTPException(500, str(e))


class DemoDataRequest(BaseModel):
    type: str  # apis, documents, images
    agent_goal: str
    agent_name: Optional[str] = None
    tasks: Optional[List[Dict[str, Any]]] = None


@app.post("/api/demo-lab/generate-for-agent")
async def generate_demo_data_for_agent(request: DemoDataRequest):
    """Generate demo data (mock APIs, documents, images) based on agent configuration"""
    try:
        goal_lower = request.agent_goal.lower()
        items = []
        
        if request.type == "apis":
            # Create mock APIs based on agent goal
            if "customer" in goal_lower or "order" in goal_lower:
                items.append({
                    "name": "Orders API",
                    "endpoint": "/api/mock/orders",
                    "method": "GET",
                    "sample_response": {"orders": [{"id": "ORD-001", "status": "shipped", "total": 99.99}]}
                })
            if "product" in goal_lower:
                items.append({
                    "name": "Products API",
                    "endpoint": "/api/mock/products",
                    "method": "GET",
                    "sample_response": {"products": [{"id": "P001", "name": "Sample Product", "price": 49.99}]}
                })
            if "employee" in goal_lower or "hr" in goal_lower:
                items.append({
                    "name": "Employees API",
                    "endpoint": "/api/mock/employees",
                    "method": "GET",
                    "sample_response": {"employees": [{"id": "E001", "name": "John Doe", "department": "Engineering"}]}
                })
            
            if not items:
                items.append({
                    "name": "Sample Data API",
                    "endpoint": "/api/mock/data",
                    "method": "GET",
                    "sample_response": {"data": [{"id": 1, "value": "sample"}]}
                })
        
        elif request.type == "documents":
            # Generate document recommendations
            if "customer" in goal_lower or "support" in goal_lower:
                items = [
                    {"name": "Customer FAQ", "type": "faq", "pages": 5},
                    {"name": "Return Policy", "type": "policy", "pages": 2},
                    {"name": "Product Guide", "type": "guide", "pages": 10}
                ]
            elif "sales" in goal_lower or "product" in goal_lower:
                items = [
                    {"name": "Product Catalog", "type": "catalog", "pages": 20},
                    {"name": "Pricing Sheet", "type": "pricing", "pages": 3}
                ]
            elif "hr" in goal_lower or "employee" in goal_lower:
                items = [
                    {"name": "Employee Handbook", "type": "handbook", "pages": 30},
                    {"name": "Benefits Guide", "type": "benefits", "pages": 10}
                ]
            else:
                items = [
                    {"name": "Knowledge Base", "type": "kb", "pages": 15}
                ]
        
        elif request.type == "images":
            # Generate image placeholders
            if "product" in goal_lower:
                items = [
                    {"name": f"Product Image {i+1}", "url": f"/demo/images/product_{i+1}.png"} 
                    for i in range(5)
                ]
            else:
                items = [
                    {"name": "Placeholder Image", "url": "/demo/images/placeholder.png"}
                ]
        
        return {
            "status": "success",
            "type": request.type,
            "count": len(items),
            "items": items
        }
        
    except Exception as e:
        print(f"[Demo Data] Error: {e}")
        raise HTTPException(500, str(e))


# ============================================================
# DEMO KIT SYSTEM - AI-Powered Tool Generation
# ============================================================

@app.post("/api/demo-kits/generate")
async def generate_demo_kit(request: GenerateDemoKitRequest, current_user: User = Depends(get_current_user_optional)):
    """Generate a complete demo kit from natural language description"""
    try:
        description = request.description
        
        # Get owner_id from current user (for access control on created tools)
        owner_id = str(current_user.id) if current_user else "system"
        
        # Get LLM for generation
        generation_llm = await _get_generation_llm()
        if not generation_llm:
            raise HTTPException(503, "No LLM configured. Please add at least one LLM provider in Settings.")
        
        # Build the prompt
        system_prompt = """You are an expert at designing demo environments for AI agents. 
Based on the user's description, generate a complete demo kit with:

1. **APIs**: Mock API endpoints the agent will need to call
2. **Knowledge Bases**: Documents with information the agent needs to answer questions
3. **Demo Assets**: Sample files (invoices, receipts, documents) for testing

Return a JSON object with this structure:
{
    "kit_name": "Short descriptive name",
    "kit_description": "Brief description of what this kit provides",
    "apis": [
        {
            "name": "API name",
            "description": "What this API does",
            "method": "GET or POST",
            "endpoint": "/api/mock/...",
            "parameters": [
                {"name": "param_name", "type": "string", "required": true, "description": "..."}
            ],
            "sample_request": {"key": "value"} or null for GET,
            "sample_response": {"realistic": "response data"}
        }
    ],
    "knowledge_bases": [
        {
            "name": "KB name",
            "description": "What info this contains",
            "sections": [
                {"title": "Section Title", "content": "Detailed content..."}
            ]
        }
    ],
    "assets": [
        {
            "name": "Asset name",
            "description": "What this is",
            "asset_type": "invoice" | "receipt" | "document" | "image",
            "metadata": {
                "fields": {"field_name": "sample_value"}
            }
        }
    ]
}

IMPORTANT:
- Generate REALISTIC sample data (real-looking names, amounts, dates)
- APIs should have proper RESTful design
- Knowledge base content should be comprehensive (not placeholder text)
- Assets should match what the agent would actually process"""

        user_message = f"Generate a demo kit for this use case:\n\n{description}"
        
        # Call LLM
        response_text = await generation_llm.generate(
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_message}
            ],
            temperature=0.7,
            max_tokens=8000
        )
        
        # Parse response
        response_text = response_text.strip()
        if response_text.startswith("```json"):
            response_text = response_text[7:]
        if response_text.startswith("```"):
            response_text = response_text[3:]
        if response_text.endswith("```"):
            response_text = response_text[:-3]
        
        kit_data = json.loads(response_text.strip())
        
        # Create Demo Kit
        kit_id = f"kit_{uuid.uuid4().hex[:8]}"
        now = datetime.now().isoformat()
        
        # Process APIs
        apis = []
        for i, api_data in enumerate(kit_data.get("apis", [])):
            api = DemoAPITool(
                id=f"{kit_id}_api_{i}",
                name=api_data.get("name", f"API {i+1}"),
                description=api_data.get("description", ""),
                method=api_data.get("method", "GET"),
                endpoint=api_data.get("endpoint", f"/api/mock/endpoint_{i}"),
                parameters=api_data.get("parameters", []),
                sample_request=api_data.get("sample_request"),
                sample_response=api_data.get("sample_response", {}),
                created_at=now
            )
            apis.append(api)
        
        # Process Knowledge Bases
        knowledge_bases = []
        for i, kb_data in enumerate(kit_data.get("knowledge_bases", [])):
            # Combine sections into full content
            sections = kb_data.get("sections", [])
            full_content = "\n\n".join([
                f"## {s.get('title', '')}\n{s.get('content', '')}" 
                for s in sections
            ])
            
            kb = DemoKnowledgeBase(
                id=f"{kit_id}_kb_{i}",
                name=kb_data.get("name", f"Knowledge Base {i+1}"),
                description=kb_data.get("description", ""),
                content=full_content,
                sections=sections,
                created_at=now
            )
            knowledge_bases.append(kb)
        
        # Process Assets
        assets = []
        for i, asset_data in enumerate(kit_data.get("assets", [])):
            asset = DemoAsset(
                id=f"{kit_id}_asset_{i}",
                name=asset_data.get("name", f"Asset {i+1}"),
                description=asset_data.get("description", ""),
                asset_type=asset_data.get("asset_type", "document"),
                metadata=asset_data.get("metadata", {}),
                created_at=now
            )
            assets.append(asset)
        
        # Create the kit
        demo_kit = DemoKit(
            id=kit_id,
            name=request.kit_name or kit_data.get("kit_name", "Demo Kit"),
            description=kit_data.get("kit_description", description[:100]),
            original_prompt=description,
            apis=apis,
            knowledge_bases=knowledge_bases,
            assets=assets,
            created_at=now,
            updated_at=now
        )
        
        # Save to state
        app_state.demo_kits[kit_id] = demo_kit
        
        # Also create ToolConfiguration entries for APIs and KBs
        # Get the actual server URL for the base_url
        server_base_url = os.environ.get("PUBLIC_URL", "http://localhost:8000")
        
        for api in apis:
            # Create proper api_config - same structure as user-created APIs
            api_config = APIEndpointConfig(
                base_url=server_base_url,
                endpoint_path=api.endpoint,
                http_method=api.method,
                auth_type="none",
                auth_value="",
                api_key_name="",
                api_key_location="header",
                headers={},
                input_parameters=[
                    APIInputParameter(
                        name=p.get("name", ""),
                        data_type=p.get("type", "string"),
                        required=p.get("required", False),
                        description=p.get("description", ""),
                        location=p.get("location", "query")
                    ) for p in (api.parameters or [])
                ]
            )
            
            # Create tool with same structure as user-created tools
            # Demo Kit tools are PUBLIC by default so everyone can use them
            tool = ToolConfiguration(
                id=api.id,
                type="api",
                name=api.name,
                description=api.description,
                api_config=api_config,
                config={
                    "mock_response": api.sample_response,
                    "sample_request": api.sample_request
                },
                # Access Control - Demo tools are PUBLIC by default
                owner_id=owner_id,
                access_type="public"
            )
            app_state.tools[api.id] = tool
        
        for kb in knowledge_bases:
            # Create knowledge tool - same structure as user-created KBs
            # Demo Kit tools are PUBLIC by default so everyone can use them
            tool = ToolConfiguration(
                id=kb.id,
                type="knowledge",
                name=kb.name,
                description=kb.description,
                config={
                    "content": kb.content,
                    "sections": kb.sections
                },
                # Access Control - Demo tools are PUBLIC by default
                owner_id=owner_id,
                access_type="public"
            )
            app_state.tools[kb.id] = tool
        
        print(f"[Demo Kit] ✅ Generated kit '{demo_kit.name}' with {len(apis)} APIs, {len(knowledge_bases)} KBs, {len(assets)} assets")
        
        return demo_kit.dict()
        
    except json.JSONDecodeError as e:
        print(f"[Demo Kit] JSON parse error: {e}")
        raise HTTPException(500, "Failed to parse AI response. Please try again.")
    except HTTPException:
        raise
    except Exception as e:
        print(f"[Demo Kit] Error: {e}")
        import traceback
        traceback.print_exc()
        raise HTTPException(500, f"Failed to generate demo kit: {str(e)}")


@app.get("/api/demo-kits")
async def list_demo_kits():
    """List all demo kits"""
    kits = []
    for kit in app_state.demo_kits.values():
        kits.append({
            "id": kit.id,
            "name": kit.name,
            "description": kit.description,
            "api_count": len(kit.apis),
            "kb_count": len(kit.knowledge_bases),
            "asset_count": len(kit.assets),
            "created_at": kit.created_at
        })
    return {"kits": kits}


@app.get("/api/demo-kits/{kit_id}")
async def get_demo_kit(kit_id: str):
    """Get a specific demo kit"""
    if kit_id not in app_state.demo_kits:
        raise HTTPException(404, "Demo kit not found")
    return app_state.demo_kits[kit_id].dict()


@app.delete("/api/demo-kits/{kit_id}")
async def delete_demo_kit(kit_id: str):
    """Delete a demo kit"""
    if kit_id not in app_state.demo_kits:
        raise HTTPException(404, "Demo kit not found")
    
    # Also delete associated tools
    kit = app_state.demo_kits[kit_id]
    for api in kit.apis:
        if api.id in app_state.tools:
            del app_state.tools[api.id]
    for kb in kit.knowledge_bases:
        if kb.id in app_state.tools:
            del app_state.tools[kb.id]
    
    del app_state.demo_kits[kit_id]
    return {"status": "deleted", "id": kit_id}


class UpdateDemoAPIRequest(BaseModel):
    name: Optional[str] = None
    description: Optional[str] = None
    method: Optional[str] = None
    endpoint: Optional[str] = None
    parameters: Optional[List[Dict[str, Any]]] = None
    sample_request: Optional[Dict[str, Any]] = None
    sample_response: Optional[Dict[str, Any]] = None


@app.put("/api/demo-kits/{kit_id}/api/{api_id}")
async def update_demo_api(kit_id: str, api_id: str, request: UpdateDemoAPIRequest):
    """Update an API in a demo kit"""
    if kit_id not in app_state.demo_kits:
        raise HTTPException(404, "Demo kit not found")
    
    kit = app_state.demo_kits[kit_id]
    api_idx = next((i for i, a in enumerate(kit.apis) if a.id == api_id), None)
    
    if api_idx is None:
        raise HTTPException(404, "API not found in kit")
    
    api = kit.apis[api_idx]
    
    # Update fields
    if request.name is not None:
        api.name = request.name
    if request.description is not None:
        api.description = request.description
    if request.method is not None:
        api.method = request.method
    if request.endpoint is not None:
        api.endpoint = request.endpoint
    if request.parameters is not None:
        api.parameters = request.parameters
    if request.sample_request is not None:
        api.sample_request = request.sample_request
    if request.sample_response is not None:
        api.sample_response = request.sample_response
    
    kit.apis[api_idx] = api
    kit.updated_at = datetime.now().isoformat()
    
    # Also update in tools
    if api.id in app_state.tools:
        tool = app_state.tools[api.id]
        tool.name = api.name
        tool.description = api.description
        tool.config.update({
            "method": api.method,
            "endpoint": api.endpoint,
            "parameters": api.parameters,
            "sample_request": api.sample_request,
            "sample_response": api.sample_response
        })
    
    return api.dict()


class UpdateDemoKBRequest(BaseModel):
    name: Optional[str] = None
    description: Optional[str] = None
    content: Optional[str] = None
    sections: Optional[List[Dict[str, str]]] = None


@app.put("/api/demo-kits/{kit_id}/kb/{kb_id}")
async def update_demo_kb(kit_id: str, kb_id: str, request: UpdateDemoKBRequest):
    """Update a knowledge base in a demo kit"""
    if kit_id not in app_state.demo_kits:
        raise HTTPException(404, "Demo kit not found")
    
    kit = app_state.demo_kits[kit_id]
    kb_idx = next((i for i, k in enumerate(kit.knowledge_bases) if k.id == kb_id), None)
    
    if kb_idx is None:
        raise HTTPException(404, "Knowledge base not found in kit")
    
    kb = kit.knowledge_bases[kb_idx]
    
    # Update fields
    if request.name is not None:
        kb.name = request.name
    if request.description is not None:
        kb.description = request.description
    if request.content is not None:
        kb.content = request.content
    if request.sections is not None:
        kb.sections = request.sections
        # Regenerate content from sections
        kb.content = "\n\n".join([
            f"## {s.get('title', '')}\n{s.get('content', '')}" 
            for s in request.sections
        ])
    
    kit.knowledge_bases[kb_idx] = kb
    kit.updated_at = datetime.now().isoformat()
    
    # Also update in tools
    if kb.id in app_state.tools:
        tool = app_state.tools[kb.id]
        tool.name = kb.name
        tool.description = kb.description
        tool.config.update({
            "content": kb.content,
            "sections": kb.sections
        })
    
    return kb.dict()


class UpdateDemoAssetRequest(BaseModel):
    name: Optional[str] = None
    description: Optional[str] = None
    asset_type: Optional[str] = None
    metadata: Optional[Dict[str, Any]] = None


@app.put("/api/demo-kits/{kit_id}/asset/{asset_id}")
async def update_demo_asset(kit_id: str, asset_id: str, request: UpdateDemoAssetRequest):
    """Update an asset in a demo kit"""
    if kit_id not in app_state.demo_kits:
        raise HTTPException(404, "Demo kit not found")
    
    kit = app_state.demo_kits[kit_id]
    asset_idx = next((i for i, a in enumerate(kit.assets) if a.id == asset_id), None)
    
    if asset_idx is None:
        raise HTTPException(404, "Asset not found in kit")
    
    asset = kit.assets[asset_idx]
    
    # Update fields
    if request.name is not None:
        asset.name = request.name
    if request.description is not None:
        asset.description = request.description
    if request.asset_type is not None:
        asset.asset_type = request.asset_type
    if request.metadata is not None:
        asset.metadata = request.metadata
    
    kit.assets[asset_idx] = asset
    kit.updated_at = datetime.now().isoformat()
    
    return asset.dict()


@app.get("/api/demo-kits/tools/all")
async def get_all_demo_tools():
    """Get all tools (APIs + Knowledge Bases) from all kits for agent creation"""
    tools = {
        "apis": [],
        "knowledge_bases": []
    }
    
    for kit in app_state.demo_kits.values():
        for api in kit.apis:
            tools["apis"].append({
                **api.dict(),
                "kit_id": kit.id,
                "kit_name": kit.name
            })
        for kb in kit.knowledge_bases:
            tools["knowledge_bases"].append({
                **kb.dict(),
                "kit_id": kit.id,
                "kit_name": kit.name
            })
    
    return tools


@app.post("/api/demo-kits/{kit_id}/api/{api_id}/test")
async def test_demo_api(kit_id: str, api_id: str, request_body: Dict[str, Any] = {}):
    """Test a mock API endpoint"""
    if kit_id not in app_state.demo_kits:
        raise HTTPException(404, "Demo kit not found")
    
    kit = app_state.demo_kits[kit_id]
    api = next((a for a in kit.apis if a.id == api_id), None)
    
    if not api:
        raise HTTPException(404, "API not found in kit")
    
    # Return the sample response (in a real implementation, this could be dynamic)
    return {
        "status": "success",
        "api": api.name,
        "method": api.method,
        "endpoint": api.endpoint,
        "request": request_body,
        "response": api.sample_response
    }


@app.post("/api/demo-kits/{kit_id}/assets/{asset_id}/generate")
async def generate_demo_asset_file(kit_id: str, asset_id: str):
    """Generate an actual file for a demo asset (invoice image, receipt, etc.)"""
    if kit_id not in app_state.demo_kits:
        raise HTTPException(404, "Demo kit not found")
    
    kit = app_state.demo_kits[kit_id]
    asset = next((a for a in kit.assets if a.id == asset_id), None)
    
    if not asset:
        raise HTTPException(404, "Asset not found in kit")
    
    try:
        # Generate the file based on asset type
        from PIL import Image, ImageDraw, ImageFont
        import io
        import base64
        
        # Create output directory
        output_dir = os.path.join(os.environ.get("DATA_PATH", "data"), "demo_assets")
        os.makedirs(output_dir, exist_ok=True)
        
        filename = f"{asset.id}.png"
        filepath = os.path.join(output_dir, filename)
        
        if asset.asset_type in ["invoice", "receipt"]:
            # Generate invoice/receipt image
            img = Image.new('RGB', (600, 800), color='white')
            draw = ImageDraw.Draw(img)
            
            # Try to use a font, fallback to default
            try:
                font_large = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 24)
                font_medium = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 16)
                font_small = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 12)
            except:
                font_large = font_medium = font_small = ImageFont.load_default()
            
            # Header
            title = "INVOICE" if asset.asset_type == "invoice" else "RECEIPT"
            draw.text((250, 30), title, fill='black', font=font_large)
            
            # Company info
            draw.text((50, 80), "Demo Company Inc.", fill='black', font=font_medium)
            draw.text((50, 100), "123 Business Street", fill='gray', font=font_small)
            draw.text((50, 115), "City, State 12345", fill='gray', font=font_small)
            
            # Invoice details
            metadata = asset.metadata.get("fields", {})
            y = 160
            draw.line((50, y, 550, y), fill='lightgray')
            y += 20
            
            draw.text((50, y), f"Date: {metadata.get('date', '2024-12-28')}", fill='black', font=font_small)
            draw.text((300, y), f"Invoice #: {metadata.get('invoice_number', 'INV-001')}", fill='black', font=font_small)
            y += 40
            
            # Items header
            draw.text((50, y), "Description", fill='black', font=font_medium)
            draw.text((350, y), "Qty", fill='black', font=font_medium)
            draw.text((420, y), "Price", fill='black', font=font_medium)
            draw.text((500, y), "Total", fill='black', font=font_medium)
            y += 25
            draw.line((50, y, 550, y), fill='lightgray')
            y += 10
            
            # Sample items
            items = metadata.get("items", [
                {"description": "Product/Service 1", "qty": 2, "price": 50.00},
                {"description": "Product/Service 2", "qty": 1, "price": 75.00},
                {"description": "Product/Service 3", "qty": 3, "price": 25.00}
            ])
            
            subtotal = 0
            for item in items:
                desc = item.get("description", "Item")
                qty = item.get("qty", 1)
                price = item.get("price", 0)
                total = qty * price
                subtotal += total
                
                draw.text((50, y), desc[:30], fill='black', font=font_small)
                draw.text((350, y), str(qty), fill='black', font=font_small)
                draw.text((420, y), f"${price:.2f}", fill='black', font=font_small)
                draw.text((500, y), f"${total:.2f}", fill='black', font=font_small)
                y += 25
            
            # Totals
            y += 20
            draw.line((350, y, 550, y), fill='lightgray')
            y += 10
            tax = subtotal * 0.1
            total = subtotal + tax
            
            draw.text((350, y), "Subtotal:", fill='black', font=font_small)
            draw.text((500, y), f"${subtotal:.2f}", fill='black', font=font_small)
            y += 20
            draw.text((350, y), "Tax (10%):", fill='black', font=font_small)
            draw.text((500, y), f"${tax:.2f}", fill='black', font=font_small)
            y += 20
            draw.text((350, y), "TOTAL:", fill='black', font=font_medium)
            draw.text((500, y), f"${total:.2f}", fill='black', font=font_medium)
            
            # Footer
            draw.text((50, 750), "Thank you for your business!", fill='gray', font=font_small)
            
            img.save(filepath)
            
        else:
            # Generic document image
            img = Image.new('RGB', (600, 400), color='white')
            draw = ImageDraw.Draw(img)
            draw.text((200, 180), asset.name, fill='black')
            img.save(filepath)
        
        # Update asset with file info
        asset.file_path = filepath
        asset.file_url = f"/api/demo-assets/{asset.id}/download"
        
        # Generate preview
        with open(filepath, "rb") as f:
            asset.preview_data = base64.b64encode(f.read()).decode()
        
        return {
            "status": "success",
            "asset_id": asset.id,
            "file_url": asset.file_url,
            "preview": f"data:image/png;base64,{asset.preview_data[:100]}..."
        }
        
    except Exception as e:
        print(f"[Demo Asset] Error generating file: {e}")
        raise HTTPException(500, f"Failed to generate asset: {str(e)}")


@app.get("/api/demo-assets/{asset_id}/download")
async def download_demo_asset(asset_id: str):
    """Download a generated demo asset file"""
    # Find the asset
    for kit in app_state.demo_kits.values():
        for asset in kit.assets:
            if asset.id == asset_id and asset.file_path:
                if os.path.exists(asset.file_path):
                    return FileResponse(
                        asset.file_path,
                        filename=f"{asset.name.replace(' ', '_')}.png",
                        media_type="image/png"
                    )
    
    raise HTTPException(404, "Asset file not found")

class UnifiedDemoRequest(BaseModel):
    """Request for unified demo setup"""
    goal: str
    agent_name: Optional[str] = None
    tasks: Optional[List[Dict[str, Any]]] = None
    personality: Optional[Dict[str, str]] = None


def generate_hr_demo_data():
    """Generate HR domain demo data"""
    employees = [
        {
            "id": "EMP001", "name": "Ahmed Hassan", "email": "ahmed.hassan@company.com",
            "department": "Engineering", "position": "Senior Developer", "hire_date": "2021-03-15",
            "manager": "Sara Ahmed",
            "vacation_balance": {"total": 21, "used": 8, "remaining": 13},
            "sick_leave": {"total": 10, "used": 2, "remaining": 8}
        },
        {
            "id": "EMP002", "name": "Sara Ahmed", "email": "sara.ahmed@company.com",
            "department": "Engineering", "position": "Engineering Manager", "hire_date": "2019-06-01",
            "manager": "Mohamed Ali",
            "vacation_balance": {"total": 25, "used": 10, "remaining": 15},
            "sick_leave": {"total": 10, "used": 0, "remaining": 10}
        }
    ]
    
    policies = {
        "vacation_policy": """## Vacation Policy
- Full-time employees receive 21 days of paid vacation per year
- Managers receive 25-30 days based on seniority
- Submit requests at least 2 weeks in advance
- Maximum 10 consecutive days without VP approval
- Up to 5 unused days can carry over to next year""",
        
        "sick_leave_policy": """## Sick Leave Policy
- All employees receive 10 days of paid sick leave per year
- For 1-2 days: Self-certification is sufficient
- For 3+ days: Medical certificate required
- 2 days can be used as mental health days""",
        
        "benefits_overview": """## Employee Benefits
- Full medical coverage for employee and dependents (dental & vision included)
- Company matches 5% of salary to retirement fund
- Life insurance: 2x annual salary
- Education allowance: $2,000/year
- Transportation allowance: $200/month
- Mobile phone allowance: $50/month"""
    }
    
    return {"employees": employees, "policies": policies, "domain": "hr"}


def generate_ecommerce_demo_data():
    """Generate E-commerce domain demo data"""
    products = [
        {"id": "PROD001", "name": "Wireless Bluetooth Headphones", "category": "Electronics",
         "price": 79.99, "stock": 150, "description": "Premium noise-canceling headphones", "rating": 4.5},
        {"id": "PROD002", "name": "Smart Watch Pro", "category": "Electronics",
         "price": 299.99, "stock": 75, "description": "Advanced fitness tracking with GPS", "rating": 4.8}
    ]
    
    orders = [
        {"id": "ORD-2024-001", "customer": "John Smith", "status": "delivered",
         "total": 79.99, "tracking": "1Z999AA10123456784", "delivery_date": "2024-01-18"},
        {"id": "ORD-2024-002", "customer": "Sarah Johnson", "status": "shipped",
         "total": 349.97, "tracking": "1Z999BB20234567895", "estimated_delivery": "2024-01-25"},
        {"id": "ORD-2024-003", "customer": "Mike Wilson", "status": "processing",
         "total": 159.98, "estimated_delivery": "2024-01-28"}
    ]
    
    policies = {
        "return_policy": """## Return Policy
- 30-day returns on all items
- Items must be unused and in original packaging
- Free return shipping for defective items
- Refunds processed within 3-5 business days""",
        
        "shipping_info": """## Shipping Information
- Standard (5-7 days): Free over $50, otherwise $5.99
- Express (2-3 days): $12.99
- Next Day: $24.99
- All orders include tracking"""
    }
    
    return {"products": products, "orders": orders, "policies": policies, "domain": "ecommerce"}


def generate_support_demo_data():
    """Generate Customer Support domain demo data"""
    faqs = [
        {"q": "How do I reset my password?", "a": "Click 'Forgot Password' on login page, enter your email, check inbox for reset link, create new password. Link expires in 24 hours."},
        {"q": "What are your support hours?", "a": "Phone: Mon-Fri 9AM-6PM EST. Email/Chat: 24/7. Average response under 2 hours."},
        {"q": "Do you offer refunds?", "a": "Yes, 30-day money-back guarantee on all plans. Contact support within 30 days for full refund."},
        {"q": "How do I cancel my subscription?", "a": "Go to Account Settings > Subscription > Cancel. Access continues until end of billing period."}
    ]
    
    troubleshooting = [
        {"issue": "App not loading", "solution": "1) Clear cache/cookies, 2) Try incognito mode, 3) Disable extensions, 4) Try different browser"},
        {"issue": "Payment failed", "solution": "1) Verify card details, 2) Check sufficient funds, 3) Try different payment method, 4) Contact bank"}
    ]
    
    return {"faqs": faqs, "troubleshooting": troubleshooting, "domain": "support"}


def generate_test_prompts(domain: str, data: dict) -> list:
    """Generate test prompts that CAN BE ANSWERED from the demo data"""
    
    if domain == "hr":
        emp = data["employees"][0]
        return [
            {"prompt": "What is my vacation balance?", 
             "expected": f"You have {emp['vacation_balance']['remaining']} vacation days remaining out of {emp['vacation_balance']['total']} total days.",
             "category": "Leave Balance"},
            {"prompt": "How many sick days do I have?",
             "expected": f"You have {emp['sick_leave']['remaining']} sick days remaining out of {emp['sick_leave']['total']} total.",
             "category": "Leave Balance"},
            {"prompt": "What is the vacation policy?",
             "expected": "Full-time employees receive 21 days of paid vacation per year. Requests should be submitted 2 weeks in advance.",
             "category": "Policy Information"},
            {"prompt": "What benefits do we have?",
             "expected": "Benefits include medical coverage, 5% retirement match, life insurance, and allowances for education, transportation, and mobile.",
             "category": "Benefits Inquiry"},
            {"prompt": "Who is my manager?",
             "expected": f"Your manager is {emp['manager']}.",
             "category": "Employee Info"}
        ]
    
    elif domain == "ecommerce":
        order = data["orders"][1]
        prod = data["products"][0]
        return [
            {"prompt": "Where is my order?",
             "expected": f"Your order is currently {order['status']}. Tracking: {order['tracking']}. Estimated delivery: {order['estimated_delivery']}.",
             "category": "Order Status"},
            {"prompt": f"What's the price of {prod['name']}?",
             "expected": f"The {prod['name']} is ${prod['price']} with a {prod['rating']}/5 rating.",
             "category": "Product Info"},
            {"prompt": "What is your return policy?",
             "expected": "We offer 30-day returns. Items must be unused and in original packaging. Refunds within 3-5 business days.",
             "category": "Policy"},
            {"prompt": "How long does shipping take?",
             "expected": "Standard: 5-7 days (free over $50). Express: 2-3 days ($12.99). Next day: $24.99.",
             "category": "Shipping"}
        ]
    
    else:  # support
        return [
            {"prompt": "How do I reset my password?",
             "expected": data["faqs"][0]["a"],
             "category": "Account Help"},
            {"prompt": "What are your support hours?",
             "expected": data["faqs"][1]["a"],
             "category": "Support Info"},
            {"prompt": "Can I get a refund?",
             "expected": data["faqs"][2]["a"],
             "category": "Billing"},
            {"prompt": "The app is not loading",
             "expected": data["troubleshooting"][0]["solution"],
             "category": "Technical"}
        ]


@app.post("/api/demo/unified-setup")
async def unified_demo_setup(request: UnifiedDemoRequest):
    """
    FULLY DYNAMIC Demo Setup - LLM generates EVERYTHING:
    - Analyzes objective & tasks
    - Decides what tools to create
    - Generates realistic mock data
    - Creates test prompts
    - Generates demo assets descriptions
    
    NO HARDCODING - Everything is AI-generated based on context!
    """
    try:
        goal = request.goal
        tasks = request.tasks or []
        agent_name = request.agent_name or "Demo Agent"
        
        # Build context for LLM
        def get_instructions_text(task):
            instructions = task.get('instructions', [])
            if not instructions:
                return ""
            texts = []
            for i in instructions:
                if isinstance(i, str):
                    texts.append(i)
                elif isinstance(i, dict):
                    texts.append(i.get('text', ''))
            return ', '.join(texts)
        
        tasks_text = "\n".join([
            f"- Task: {t.get('name', 'Unnamed')}\n  Description: {t.get('description', '')}\n  Instructions: {get_instructions_text(t)}"
            for t in tasks if isinstance(t, dict)
        ]) if tasks else "No specific tasks defined"
        
        # Ask LLM to design the complete demo setup
        design_prompt = f"""You are an expert at designing demo environments for AI agents.

AGENT NAME: {agent_name}
AGENT OBJECTIVE: {goal}

TASKS:
{tasks_text}

Based on this agent's objective and tasks, design a COMPLETE demo environment.
You must create realistic, contextual mock data that makes sense for this specific agent.

Respond with a JSON object containing:

{{
    "domain": "detected domain (e.g., hr, ecommerce, healthcare, finance, support, etc.)",
    "analysis": "brief analysis of what this agent needs",
    
    "api_tools": [
        {{
            "name": "API name (e.g., Employee Leave Balance API)",
            "description": "What this API does",
            "endpoint": "/api/v1/path/to/endpoint",
            "method": "GET or POST",
            "parameters": [
                {{"name": "param_name", "type": "string/number/boolean", "required": true/false, "description": "param description"}}
            ],
            "mock_response": {{
                // Realistic mock data that this API would return
                // Use realistic names, numbers, dates relevant to the domain
            }}
        }}
    ],
    
    "knowledge_base": {{
        "name": "Knowledge Base name",
        "description": "What information it contains",
        "documents": [
            {{
                "title": "Document title",
                "type": "policy/faq/guide/record/etc",
                "content": "Full realistic content of the document with specific details, numbers, procedures"
            }}
        ]
    }},
    
    "demo_assets": [
        {{
            "type": "image",
            "name": "filename.png or filename.jpg",
            "description": "What this image shows",
            "purpose": "For OCR testing / upload testing / etc",
            "content_description": "Detailed description of what should be in the image (e.g., receipt with store name, date, items, total)",
            "sample_text": "The actual text that would appear in the image for OCR"
        }}
    ],
    
    "test_prompts": [
        {{
            "prompt": "A realistic question a user would ask",
            "expected_tool": "Which tool should handle this",
            "expected_response_summary": "Brief summary of expected response",
            "category": "Category name"
        }}
    ],
    
    "demo_script": [
        "Step-by-step description of what was created and why"
    ]
}}

IMPORTANT RULES:
1. Generate REALISTIC data with actual names, numbers, dates
2. Make the data internally consistent (e.g., vacation balance matches policy limits)
3. Create data that allows the agent to actually answer the test prompts
4. For tasks involving document/receipt processing: create IMAGE assets (png/jpg) NOT spreadsheets
   - Receipts should be images for OCR testing
   - Invoices should be images for OCR testing
   - Forms should be images that users can upload
5. Include "sample_text" for image assets - this is the text that would be extracted via OCR
6. Think about what a real user would ask and upload to this agent

Respond ONLY with valid JSON, no markdown formatting."""

        # Call LLM to design the demo
        design_messages = [{"role": "user", "content": design_prompt}]
        design_result = await call_llm(design_messages, "gpt-4o")
        
        # Parse LLM response
        try:
            response_text = design_result["content"]
            # Clean up response
            if "```json" in response_text:
                response_text = response_text.split("```json")[1].split("```")[0]
            elif "```" in response_text:
                response_text = response_text.split("```")[1].split("```")[0]
            
            demo_design = json.loads(response_text.strip())
        except Exception as parse_error:
            print(f"[Demo] Failed to parse LLM response: {parse_error}")
            print(f"[Demo] Raw response: {design_result['content'][:500]}")
            raise HTTPException(500, f"Failed to parse AI response: {str(parse_error)}")
        
        # Now create the actual tools based on LLM's design
        created_tools = []
        tool_ids = []
        
        # Helper function to get unique tool name
        def get_unique_tool_name(base_name: str) -> str:
            existing_names = [t.name.lower() for t in app_state.tools.values()]
            if base_name.lower() not in existing_names:
                return base_name
            # Add suffix to make unique
            counter = 2
            while f"{base_name} ({counter})".lower() in existing_names:
                counter += 1
            return f"{base_name} ({counter})"
        
        # Create API tools - same structure as user-created APIs
        server_base_url = os.environ.get("PUBLIC_URL", "http://localhost:8000")
        
        for api_spec in demo_design.get("api_tools", []):
            tool_name = get_unique_tool_name(api_spec.get("name", "Demo API"))
            api_tool = ToolConfiguration(
                type="api",
                name=tool_name,
                description=api_spec.get("description", ""),
                config={
                    "mock_response": api_spec.get("mock_response", {})
                },
                api_config=APIEndpointConfig(
                    base_url=server_base_url,
                    http_method=api_spec.get("method", "GET"),
                    endpoint_path=api_spec.get("endpoint", "/api/demo"),
                    auth_type="none",
                    auth_value="",
                    api_key_name="",
                    api_key_location="header",
                    headers={},
                    input_parameters=[
                        APIInputParameter(
                            name=p.get("name", "param"),
                            data_type=p.get("type", "string"),
                            required=p.get("required", False),
                            description=p.get("description", ""),
                            location=p.get("location", "query")
                        ) for p in api_spec.get("parameters", [])
                    ]
                )
            )
            app_state.tools[api_tool.id] = api_tool
            tool_ids.append(api_tool.id)
            created_tools.append({
                "id": api_tool.id,
                "type": "api",
                "name": api_tool.name,
                "description": api_tool.description,
                "endpoint": api_spec.get("endpoint"),
                "method": api_spec.get("method"),
                "mock_data": api_spec.get("mock_response")
            })
        
        # Create Knowledge Base
        kb_spec = demo_design.get("knowledge_base", {})
        if kb_spec:
            kb_name = get_unique_tool_name(kb_spec.get("name", "Demo Knowledge Base"))
            kb_id = f"demo-kb-{uuid.uuid4().hex[:8]}"
            kb_tool = ToolConfiguration(
                id=kb_id,
                type="knowledge",
                name=kb_name,
                description=kb_spec.get("description", ""),
                config={"collection_id": kb_id}
            )
            app_state.tools[kb_id] = kb_tool
            tool_ids.append(kb_id)
            
            # Index documents
            documents = kb_spec.get("documents", [])
            for doc in documents:
                chunk = {
                    "chunk_id": f"chunk-{uuid.uuid4().hex[:8]}",
                    "text": f"## {doc.get('title', 'Document')}\n\n{doc.get('content', '')}",
                    "source": doc.get("title", "Document"),
                    "type": doc.get("type", "document"),
                    "tool_id": kb_id
                }
                app_state.document_chunks.append(chunk)
            
            created_tools.append({
                "id": kb_id,
                "type": "knowledge",
                "name": kb_tool.name,
                "description": kb_tool.description,
                "documents_count": len(documents),
                "documents": [{"title": d.get("title"), "type": d.get("type")} for d in documents]
            })
        
        # Generate actual demo asset files
        generated_assets = []
        demo_assets_dir = os.path.join(os.environ.get("UPLOAD_PATH", "data/uploads"), "demo_assets")
        os.makedirs(demo_assets_dir, exist_ok=True)
        
        for asset in demo_design.get("demo_assets", []):
            asset_id = f"demo-asset-{uuid.uuid4().hex[:8]}"
            asset_name = asset.get("name", f"{asset_id}.png")
            asset_path = os.path.join(demo_assets_dir, asset_name)
            
            # Generate image with text content for OCR testing
            try:
                from PIL import Image, ImageDraw, ImageFont
                
                # Create image
                img_width, img_height = 600, 800
                img = Image.new('RGB', (img_width, img_height), color='white')
                draw = ImageDraw.Draw(img)
                
                # Try to use a font, fallback to default
                try:
                    font_large = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 24)
                    font_medium = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 16)
                    font_small = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 12)
                except:
                    font_large = ImageFont.load_default()
                    font_medium = font_large
                    font_small = font_large
                
                # Draw content
                y_pos = 30
                
                # Title/Header
                title = asset.get("description", "Demo Document")[:50]
                draw.text((30, y_pos), title, fill='black', font=font_large)
                y_pos += 50
                
                # Draw a line
                draw.line([(30, y_pos), (570, y_pos)], fill='gray', width=2)
                y_pos += 20
                
                # Sample text content for OCR
                sample_text = asset.get("sample_text", asset.get("content_description", "Sample content for testing"))
                
                # Word wrap and draw text
                lines = []
                words = sample_text.split()
                current_line = ""
                for word in words:
                    test_line = current_line + " " + word if current_line else word
                    if len(test_line) < 60:
                        current_line = test_line
                    else:
                        lines.append(current_line)
                        current_line = word
                if current_line:
                    lines.append(current_line)
                
                for line in lines[:30]:  # Limit lines
                    draw.text((30, y_pos), line, fill='black', font=font_medium)
                    y_pos += 25
                
                # Add footer
                y_pos = img_height - 50
                draw.text((30, y_pos), f"Demo Asset - Generated for testing", fill='gray', font=font_small)
                
                # Save image
                img.save(asset_path, 'PNG')
                
                generated_assets.append({
                    "id": asset_id,
                    "name": asset_name,
                    "type": asset.get("type", "image"),
                    "description": asset.get("description", ""),
                    "purpose": asset.get("purpose", "Demo testing"),
                    "sample_text": asset.get("sample_text", ""),
                    "download_url": f"/api/demo/assets/{asset_name}",
                    "generated": True
                })
                
            except ImportError:
                # PIL not available, create a text-based placeholder
                generated_assets.append({
                    "id": asset_id,
                    "name": asset_name,
                    "type": asset.get("type", "image"),
                    "description": asset.get("description", ""),
                    "purpose": asset.get("purpose", "Demo testing"),
                    "sample_text": asset.get("sample_text", ""),
                    "download_url": None,
                    "generated": False,
                    "note": "Image generation requires PIL library"
                })
        
        app_state.save_to_disk()
        
        return {
            "status": "success",
            "domain": demo_design.get("domain", "general"),
            "analysis": demo_design.get("analysis", ""),
            "tools": created_tools,
            "tool_ids": tool_ids,
            "demo_assets": generated_assets if generated_assets else demo_design.get("demo_assets", []),
            "test_prompts": demo_design.get("test_prompts", []),
            "demo_script": demo_design.get("demo_script", []),
            "summary": {
                "apis_created": len([t for t in created_tools if t["type"] == "api"]),
                "knowledge_bases": len([t for t in created_tools if t["type"] == "knowledge"]),
                "documents_indexed": sum(t.get("documents_count", 0) for t in created_tools if t["type"] == "knowledge"),
                "demo_assets": len(generated_assets) if generated_assets else len(demo_design.get("demo_assets", [])),
                "test_prompts": len(demo_design.get("test_prompts", []))
            }
        }
        
    except HTTPException:
        raise
    except Exception as e:
        print(f"[Unified Demo] Error: {e}")
        import traceback
        traceback.print_exc()
        raise HTTPException(500, str(e))


@app.get("/api/demo/assets/{filename}")
async def download_demo_asset(filename: str):
    """Download a generated demo asset file"""
    demo_assets_dir = os.path.join(os.environ.get("UPLOAD_PATH", "data/uploads"), "demo_assets")
    file_path = os.path.join(demo_assets_dir, filename)
    
    if not os.path.exists(file_path):
        raise HTTPException(404, "Asset not found")
    
    # Determine content type
    if filename.lower().endswith('.png'):
        media_type = "image/png"
    elif filename.lower().endswith('.jpg') or filename.lower().endswith('.jpeg'):
        media_type = "image/jpeg"
    elif filename.lower().endswith('.pdf'):
        media_type = "application/pdf"
    else:
        media_type = "application/octet-stream"
    
    return FileResponse(file_path, media_type=media_type, filename=filename)


@app.post("/api/demo/refine-unified")
async def refine_unified_demo(request: Dict[str, Any]):
    """Refine unified demo based on user feedback"""
    try:
        kb_id = request.get("knowledge_base_id")
        refinement = request.get("refinement_prompt", "")
        goal = request.get("goal", "")
        
        if not kb_id:
            raise HTTPException(400, "Missing knowledge_base_id")
        
        # Re-run unified setup with refined goal
        new_goal = f"{goal}. {refinement}"
        
        # Create a new request and call unified setup
        new_request = UnifiedDemoRequest(goal=new_goal)
        return await unified_demo_setup(new_request)
        
    except Exception as e:
        print(f"[Refine Unified] Error: {e}")
        raise HTTPException(500, str(e))


@app.put("/api/agents/{agent_id}")
async def update_agent(agent_id: str, request: UpdateAgentRequest, current_user: User = Depends(get_current_user)):
    if not current_user:
        raise HTTPException(401, "Authentication required")
    
    user_id = str(current_user.id)
    org_id = current_user.org_id if current_user else "org_default"
    user_role_ids = getattr(current_user, 'role_ids', []) or []
    user_group_ids = get_user_group_ids(user_id) if user_id else []
    
    # Try to get agent from database first, then fallback to in-memory
    agent = None
    owner_id = None
    is_owner = False
    
    try:
        from database.services import AgentService
        agent_dict = AgentService.get_agent_by_id(agent_id, org_id)
        if agent_dict:
            owner_id = str(agent_dict.get('owner_id')) if agent_dict.get('owner_id') else None
            created_by = str(agent_dict.get('created_by')) if agent_dict.get('created_by') else None
            is_owner = (user_id == owner_id) or (user_id == created_by)
            
            # Remove extra fields that AgentData doesn't have
            agent_dict_clean = {k: v for k, v in agent_dict.items() if k in AgentData.__fields__}
            agent = AgentData(**agent_dict_clean)
            app_state.agents[agent.id] = agent  # Update in-memory cache
    except Exception as e:
        print(f"⚠️  [DATABASE ERROR] Failed to load agent from database: {e}, using in-memory")
    
    # Fallback to in-memory if not found in database
    if not agent:
        if agent_id not in app_state.agents:
            raise HTTPException(404, "Agent not found")
        agent = app_state.agents[agent_id]
    
    # Get user's permissions for this agent
    permissions = ['full_admin'] if is_owner else []
    
    if not is_owner and ACCESS_CONTROL_AVAILABLE and AccessControlService:
        try:
            from api.modules.access_control.service import get_session, normalize_org_id
            from database.models.agent_access import AgentAccessPolicy
            import json
            
            norm_org_id = normalize_org_id(org_id)
            
            with get_session() as session:
                admin_policy = session.query(AgentAccessPolicy).filter(
                    AgentAccessPolicy.agent_id == agent_id,
                    AgentAccessPolicy.org_id == norm_org_id,
                    AgentAccessPolicy.access_type == 'agent_admin',
                    AgentAccessPolicy.is_active == True
                ).first()
                
                if admin_policy:
                    user_is_admin = user_id in (admin_policy.user_ids or [])
                    group_is_admin = any(g in (admin_policy.group_ids or []) for g in user_group_ids)
                    
                    if user_is_admin or group_is_admin:
                        permissions = []  # Start with empty - must get from config
                        if admin_policy.description:
                            try:
                                admin_config = json.loads(admin_policy.description)
                                print(f"🔐 [UPDATE] Parsed admin_config: {admin_config}")
                                
                                if user_id in admin_config:
                                    entity_config = admin_config[user_id]
                                    print(f"🔐 [UPDATE] Found user config: {entity_config}")
                                    if isinstance(entity_config, dict):
                                        permissions = entity_config.get('permissions', [])
                                    elif isinstance(entity_config, list):
                                        permissions = entity_config
                                else:
                                    # Check groups
                                    for group_id in user_group_ids:
                                        if group_id in admin_config:
                                            entity_config = admin_config[group_id]
                                            if isinstance(entity_config, dict):
                                                permissions = entity_config.get('permissions', [])
                                            elif isinstance(entity_config, list):
                                                permissions = entity_config
                                            break
                            except Exception as parse_err:
                                print(f"⚠️  [UPDATE] Error parsing permissions: {parse_err}")
                        else:
                            # No description = legacy data, grant full_admin
                            print(f"⚠️  [UPDATE] No description in policy, defaulting to full_admin (legacy)")
                            permissions = ['full_admin']
                        print(f"🔐 [UPDATE] Final permissions: {permissions}")
        except Exception as e:
            print(f"⚠️  Error checking update permissions: {e}")
    
    if not permissions:
        raise HTTPException(403, "You don't have permission to edit this agent")
    
    has_full_admin = 'full_admin' in permissions
    
    # Permission checks for each field update
    def check_perm(perm_name, field_name):
        if not has_full_admin and perm_name not in permissions:
            raise HTTPException(403, f"You don't have permission to edit {field_name} (requires: {perm_name})")
    
    # Update agent fields with permission checks
    if request.name is not None:
        check_perm('edit_basic_info', 'name')
        agent.name = request.name
    if request.goal is not None:
        check_perm('edit_basic_info', 'goal')
        agent.goal = request.goal
    if request.icon is not None:
        check_perm('edit_basic_info', 'icon')
        agent.icon = request.icon
    if request.personality is not None:
        check_perm('edit_personality', 'personality')
        agent.personality = AgentPersonality(**request.personality)
    if request.guardrails is not None:
        check_perm('edit_guardrails', 'guardrails')
        g = request.guardrails
        agent.guardrails = AgentGuardrails(
            anti_hallucination=g.get('antiHallucination', True),
            cite_sources=g.get('citeSources', True),
            admit_uncertainty=g.get('admitUncertainty', True),
            verify_facts=g.get('verifyFacts', True),
            no_speculation=g.get('noSpeculation', False),
            avoid_topics=[t.strip() for t in g.get('avoidTopics', '').split('\n') if t.strip()] if isinstance(g.get('avoidTopics'), str) else g.get('avoidTopics', []),
            focus_topics=[t.strip() for t in g.get('focusTopics', '').split('\n') if t.strip()] if isinstance(g.get('focusTopics'), str) else g.get('focusTopics', []),
            max_length=g.get('maxLength', 'medium'),
            language=g.get('language', 'user'),
            escalate_angry=g.get('escalateAngry', True),
            escalate_complex=g.get('escalateComplex', True),
            escalate_request=g.get('escalateRequest', True),
            escalate_sensitive=g.get('escalateSensitive', False),
            pii_protection=g.get('piiProtection', True),
            mask_pii=g.get('maskPii', True),
            no_store_pii=g.get('noStorePii', True)
        )
    if request.tasks is not None:
        check_perm('manage_tasks', 'tasks')
        tasks = []
        for t in request.tasks:
            instructions = []
            for i in t.get('instructions', []):
                if isinstance(i, dict): instructions.append(TaskInstruction(**i))
                elif isinstance(i, str): instructions.append(TaskInstruction(text=i))
            tasks.append(TaskDefinition(id=t.get('id', str(uuid.uuid4())), name=t.get('name', ''), description=t.get('description', ''), instructions=instructions))
        agent.tasks = tasks
    if request.tool_ids is not None:
        check_perm('manage_tools', 'tools')
        agent.tool_ids = request.tool_ids
    if request.model_id is not None:
        check_perm('edit_model', 'model')
        agent.model_id = request.model_id
    if request.status is not None:
        # Publishing requires publish_agent permission
        if request.status == 'published':
            check_perm('publish_agent', 'publish status')
        agent.status = request.status
    # Handle publishing/unpublishing
    if request.is_active is not None:
        if request.is_active:
            check_perm('publish_agent', 'publish')
        agent.is_active = request.is_active
        if request.is_active:
            agent.status = "published"
    if request.is_published is not None:
        if request.is_published:
            check_perm('publish_agent', 'publish')
        agent.is_active = request.is_published
        if request.is_published:
            agent.status = "published"
    agent.updated_at = datetime.utcnow().isoformat()
    
    # Save to database using AgentService
    try:
        from database.services import AgentService
        agent_dict = agent.dict()
        # Get org_id and updated_by from current_user
        org_id = current_user.org_id if current_user else "org_default"
        updated_by = current_user.id if current_user else None
        
        AgentService.update_agent(
            agent_id=agent.id,
            agent_data=agent_dict,
            org_id=org_id,
            updated_by=updated_by
        )
        app_state.agents[agent.id] = agent  # Update in-memory cache
    except Exception as e:
        print(f"⚠️  [DATABASE ERROR] Failed to update agent in database: {e}, saving to in-memory only")
        import traceback
        traceback.print_exc()
        app_state.agents[agent.id] = agent
        app_state.save_to_disk()  # Will try to save to database in save_to_disk()
    
    return {"status": "success", "agent": agent.dict()}


@app.delete("/api/agents/{agent_id}")
async def delete_agent(agent_id: str, current_user: User = Depends(get_current_user)):
    """
    Delete an agent. Only the owner can delete.
    NOTE: Admins are treated as normal users - they cannot delete unless they own the agent.
    """
    if not current_user:
        raise HTTPException(401, "Authentication required")
    
    user_id = str(current_user.id)
    
    # Check if agent exists and get ownership info
    agent_dict = None
    owner_id = None
    created_by = None
    
    try:
        from database.services import AgentService
        agent_dict = AgentService.get_agent_by_id(agent_id, "org_default")
        if agent_dict:
            owner_id = str(agent_dict.get('owner_id')) if agent_dict.get('owner_id') else None
            created_by = str(agent_dict.get('created_by')) if agent_dict.get('created_by') else None
    except Exception as e:
        print(f"⚠️  Error loading agent from database: {e}")
    
    # Also check in-memory
    if not agent_dict and agent_id in app_state.agents:
        agent = app_state.agents[agent_id]
        owner_id = getattr(agent, 'owner_id', None)
        created_by = getattr(agent, 'created_by', None)
        agent_dict = agent.dict()
    
    if not agent_dict:
        raise HTTPException(404, "Agent not found")
    
    # PERMISSION CHECK: Only owner can delete
    # NOTE: Admins/SuperAdmins are treated as normal users - no special privileges
    is_owner = (user_id == owner_id) or (user_id == created_by)
    
    print(f"🗑️ [DELETE AGENT] User '{user_id[:8]}...' attempting to delete agent '{agent_id[:8]}...': owner='{owner_id}', created_by='{created_by}', is_owner={is_owner}")
    
    if not is_owner:
        # Check if user has been granted delete permission via agent access policy
        has_delete_permission = False
        try:
            from database.models.agent_access import AgentAccessPolicy
            from database.config import get_db_session
            
            with get_db_session() as session:
                # Check for delete permission in access policies
                policy = session.query(AgentAccessPolicy).filter(
                    AgentAccessPolicy.agent_id == agent_id,
                    AgentAccessPolicy.is_active == True
                ).first()
                
                if policy:
                    # Check if user is in user_ids with delete permission
                    if user_id in (policy.user_ids or []):
                        # Check if delete is in their permissions
                        if policy.description:
                            try:
                                config = json.loads(policy.description)
                                permissions = config.get('permissions', [])
                                if 'delete_agent' in permissions or 'full_admin' in permissions:
                                    has_delete_permission = True
                            except:
                                pass
        except Exception as e:
            print(f"⚠️  Error checking delete permission: {e}")
        
        if not has_delete_permission:
            print(f"🚫 [DELETE AGENT] DENIED - User '{user_id[:8]}...' is not owner and has no delete permission")
            raise HTTPException(403, "You don't have permission to delete this agent. Only the owner can delete.")
    
    # Delete from database
    try:
        from database.services import AgentService
        AgentService.delete_agent(agent_id, org_id="org_default", deleted_by=user_id)
        print(f"✅ [DELETE AGENT] Deleted from database by '{user_id[:8]}...'")
    except Exception as e:
        print(f"⚠️  [DATABASE ERROR] Failed to delete agent from database: {e}, deleting from in-memory only")
        import traceback
        traceback.print_exc()
    
    # Delete from in-memory
    if agent_id in app_state.agents:
        del app_state.agents[agent_id]
    app_state.conversations = {k: v for k, v in app_state.conversations.items() if v.agent_id != agent_id}
    app_state.save_to_disk()
    
    return {"status": "success"}


# ============================================================================
# Settings Endpoints
# ============================================================================

@app.get("/api/settings")
async def get_settings():
    """Get all system settings (loads from database if available)"""
    # Try to load from database first
    try:
        from database.services import SystemSettingsService
        db_settings = SystemSettingsService.get_system_setting("system_settings")
        if db_settings:
            # Update app_state with database settings
            app_state.settings = SystemSettings(**db_settings)
            llm_provider = app_state.settings.llm.provider.value if hasattr(app_state.settings.llm.provider, 'value') else str(app_state.settings.llm.provider)
            vector_db_provider = app_state.settings.vector_db.provider.value if hasattr(app_state.settings.vector_db.provider, 'value') else str(app_state.settings.vector_db.provider)
    except Exception as e:
        print(f"⚠️  [API ERROR] Database load failed: {e}, using in-memory settings")
        import traceback
        traceback.print_exc()
    
    settings = app_state.settings.dict()
    
    # Add default models for each configured provider if not set
    provider_default_models = {
        "openai": ["gpt-4o", "gpt-4o-mini", "gpt-4-turbo", "gpt-4", "gpt-3.5-turbo"],
        "anthropic": ["claude-sonnet-4-20250514", "claude-3-5-sonnet-20241022", "claude-3-5-haiku-20241022"],
        "google": ["gemini-2.0-flash", "gemini-1.5-pro", "gemini-1.5-flash"],
        "groq": ["llama-3.3-70b-versatile", "llama-3.1-8b-instant", "gemma2-9b-it"],
        "mistral": ["mistral-large-2411", "mistral-small-2503", "open-mistral-nemo"],
        "cohere": ["command-a-03-2025", "command-r-plus-08-2024", "command-r-08-2024"],
        "xai": ["grok-2", "grok-2-mini"],
        "deepseek": ["deepseek-chat", "deepseek-coder"],
        "together": ["meta-llama/Llama-3.3-70B-Instruct-Turbo"],
        "perplexity": ["sonar", "sonar-pro", "sonar-reasoning"],
    }
    
    for provider in settings.get('llm_providers', []):
        if not provider.get('models') or len(provider.get('models', [])) == 0:
            provider['models'] = provider_default_models.get(provider.get('provider'), [])
    
    # Mask sensitive values
    if settings['llm']['api_key']:
        settings['llm']['api_key'] = '***' + settings['llm']['api_key'][-4:] if len(settings['llm']['api_key']) > 4 else '****'
    if settings['embedding']['api_key']:
        settings['embedding']['api_key'] = '***' + settings['embedding']['api_key'][-4:] if len(settings['embedding']['api_key']) > 4 else '****'
    if settings['vector_db']['pinecone_api_key']:
        settings['vector_db']['pinecone_api_key'] = '***' + settings['vector_db']['pinecone_api_key'][-4:] if len(settings['vector_db']['pinecone_api_key']) > 4 else '****'
    return {"settings": settings}


@app.put("/api/settings")
async def update_settings(request: Dict[str, Any]):
    """Update system settings"""
    try:
        # Update LLM settings
        if 'llm' in request:
            llm_data = request['llm']
            for key, value in llm_data.items():
                if key == 'provider' and value:
                    try:
                        app_state.settings.llm.provider = LLMProvider(value)
                    except ValueError:
                        app_state.settings.llm.provider = LLMProvider.CUSTOM
                elif value and not (key == 'api_key' and value.startswith('***')):
                    setattr(app_state.settings.llm, key, value)
        
        # Update LLM Providers array (for multi-provider support)
        if 'llm_providers' in request:
            providers_data = request['llm_providers']
            app_state.settings.llm_providers = [
                LLMProviderConfig(**p) if isinstance(p, dict) else p
                for p in providers_data
            ]
            print(f"✅ Saved {len(app_state.settings.llm_providers)} LLM providers")
        
        # Update Google settings
        if 'google' in request:
            google_data = request['google']
            if 'gemini_key' in google_data:
                app_state.settings.google.gemini_key = google_data['gemini_key']
            if 'cloud_key' in google_data:
                app_state.settings.google.cloud_key = google_data['cloud_key']
        
        # Update Embedding settings
        if 'embedding' in request:
            emb_data = request['embedding']
            for key, value in emb_data.items():
                if value and not (key == 'api_key' and value.startswith('***')):
                    setattr(app_state.settings.embedding, key, value)
        
        # Update Vector DB settings
        if 'vector_db' in request:
            vdb_data = request['vector_db']
            for key, value in vdb_data.items():
                if value is not None and not (key.endswith('api_key') and str(value).startswith('***')):
                    setattr(app_state.settings.vector_db, key, value)
        
        # Update general settings
        for key in ['chunk_size', 'chunk_overlap', 'search_top_k', 'enable_rag', 'enable_web_scraping', 'enable_api_tools', 'app_name', 'app_logo', 'theme']:
            if key in request:
                setattr(app_state.settings, key, request[key])
        
        # Reset providers to pick up new settings
        app_state.reset_providers()
        
        # Save to database
        print(f"💾 [API] Updating platform settings in database...")
        try:
            from database.services import SystemSettingsService
            # Save all settings as JSON in system_settings table
            settings_dict = app_state.settings.dict()
            SystemSettingsService.set_system_setting(
                key="system_settings",
                value=settings_dict,
                value_type="json",
                category="platform",
                description="Complete platform settings including LLM, Embedding, Vector DB, RAG, Features, Theme"
            )
            # Handle provider as enum or string
            llm_provider = app_state.settings.llm.provider.value if hasattr(app_state.settings.llm.provider, 'value') else str(app_state.settings.llm.provider)
            vector_db_provider = app_state.settings.vector_db.provider.value if hasattr(app_state.settings.vector_db.provider, 'value') else str(app_state.settings.vector_db.provider)
            print(f"✅ [API] Platform settings updated successfully in database (LLM: {llm_provider}, VectorDB: {vector_db_provider}, Theme: {app_state.settings.theme})")
        except Exception as e:
            print(f"⚠️  [API ERROR] Database save failed: {e}, saving to disk only")
            import traceback
            traceback.print_exc()
        
        # Save to disk (backup)
        app_state.save_to_disk()
        
        return {"status": "success", "message": "Settings updated successfully"}
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(400, f"Failed to update settings: {str(e)}")


@app.get("/api/settings/providers")
async def get_available_providers():
    """Get list of available providers - returns ONLY configured providers from database/settings"""
    
    # Get configured LLM providers from settings (loaded from database)
    configured_llm_providers = []
    
    # Provider metadata for display
    provider_meta = {
        "openai": {"name": "OpenAI", "default_models": ["gpt-4o", "gpt-4o-mini", "gpt-4-turbo", "gpt-4", "gpt-3.5-turbo"]},
        "azure_openai": {"name": "Azure OpenAI", "default_models": ["gpt-4o", "gpt-4", "gpt-35-turbo"]},
        "anthropic": {"name": "Anthropic Claude", "default_models": ["claude-sonnet-4-20250514", "claude-3-5-sonnet-20241022", "claude-3-5-haiku-20241022"]},
        "ollama": {"name": "Ollama (Local)", "default_models": ["llama3.2", "llama3.1", "mistral", "codellama", "phi3", "gemma2"]},
        "google": {"name": "Google Gemini", "default_models": ["gemini-2.0-flash", "gemini-1.5-pro", "gemini-1.5-flash"]},
        "xai": {"name": "xAI (Grok)", "default_models": ["grok-beta", "grok-2"]},
        "groq": {"name": "Groq", "default_models": ["llama-3.3-70b-versatile", "llama-3.1-8b-instant", "gemma2-9b-it"]},
        "mistral": {"name": "Mistral AI", "default_models": ["mistral-large-2411", "mistral-small-2503", "open-mistral-nemo"]},
        "deepseek": {"name": "DeepSeek", "default_models": ["deepseek-chat", "deepseek-coder"]},
        "together": {"name": "Together AI", "default_models": ["meta-llama/Llama-3-70b-chat-hf"]},
        "perplexity": {"name": "Perplexity", "default_models": ["sonar", "sonar-pro"]},
        "lmstudio": {"name": "LM Studio (Local)", "default_models": []},
        "cohere": {"name": "Cohere", "default_models": ["command-r-plus", "command-r"]},
        "custom": {"name": "Custom (OpenAI Compatible)", "default_models": []},
    }
    
    # Check configured providers from app_state.settings.llm_providers
    for provider in app_state.settings.llm_providers:
        provider_id = provider.provider
        has_api_key = bool(provider.api_key)
        meta = provider_meta.get(provider_id, {"name": provider_id.title(), "default_models": []})
        
        configured_llm_providers.append({
            "id": provider_id,
            "name": meta["name"],
            "available": has_api_key,  # Only available if API key is set
            "configured": True,
            "models": meta["default_models"]
        })
    
    # Add unconfigured providers as "available: False" so users know they exist
    configured_ids = {p["id"] for p in configured_llm_providers}
    for provider_id, meta in provider_meta.items():
        if provider_id not in configured_ids:
            configured_llm_providers.append({
                "id": provider_id,
                "name": meta["name"],
                "available": False,
                "configured": False,
                "models": meta["default_models"]
            })
    
    return {
        "llm_providers": configured_llm_providers,
        "embedding_providers": [
            {"id": "openai", "name": "OpenAI Embeddings", "available": True, "models": ["text-embedding-3-small", "text-embedding-3-large", "text-embedding-ada-002"]},
            {"id": "sentence_transformers", "name": "Sentence Transformers (Local)", "available": SENTENCE_TRANSFORMERS_AVAILABLE, "models": ["all-MiniLM-L6-v2", "all-mpnet-base-v2", "paraphrase-multilingual-MiniLM-L12-v2"]},
            {"id": "ollama", "name": "Ollama Embeddings", "available": True, "models": ["nomic-embed-text", "mxbai-embed-large"]},
            {"id": "azure_openai", "name": "Azure OpenAI Embeddings", "available": True, "models": ["text-embedding-ada-002"]},
            {"id": "cohere", "name": "Cohere", "available": False, "models": ["embed-english-v3.0", "embed-multilingual-v3.0"]},
        ],
        "vector_db_providers": [
            {"id": "memory", "name": "In-Memory (Keyword Search)", "available": True, "description": "Simple keyword search, no setup required"},
            {"id": "chromadb", "name": "ChromaDB (Local)", "available": CHROMADB_AVAILABLE, "description": "Local vector database, runs anywhere"},
            {"id": "pinecone", "name": "Pinecone (Cloud)", "available": True, "description": "Managed cloud vector database"},
            {"id": "qdrant", "name": "Qdrant (Self-hosted/Cloud)", "available": False, "description": "High-performance vector database"},
            {"id": "milvus", "name": "Milvus (Self-hosted)", "available": False, "description": "Open-source vector database"},
            {"id": "weaviate", "name": "Weaviate (Self-hosted/Cloud)", "available": False, "description": "AI-native vector database"},
        ],
        "features": {
            "playwright_available": PLAYWRIGHT_AVAILABLE,
            "chromadb_available": CHROMADB_AVAILABLE,
            "sentence_transformers_available": SENTENCE_TRANSFORMERS_AVAILABLE,
        }
    }


@app.post("/api/settings/test-llm")
async def test_llm_connection(request: Dict[str, Any]):
    """Test LLM connection - tests ALL models for the provider IN PARALLEL"""
    import asyncio
    
    provider_name = request.get('provider', 'openai')
    api_key = request.get('api_key', '')
    api_base = request.get('api_base', '')
    
    print(f"\n{'='*50}")
    print(f"[Test LLM] Testing ALL models for: {provider_name} (parallel)")
    print(f"{'='*50}")
    
    # Models to test for each provider (verified working Jan 2026)
    provider_models = {
        "openai": ["gpt-4o-mini", "gpt-4o", "gpt-3.5-turbo"],
        "anthropic": ["claude-opus-4-5-20251101", "claude-sonnet-4-20250514", "claude-3-5-haiku-20241022"],
        "google": ["gemini-2.0-flash", "gemini-1.5-flash", "gemini-1.5-pro"],
        "groq": ["llama-3.3-70b-versatile", "llama-3.1-8b-instant"],  # mixtral deprecated
        "mistral": ["mistral-small-2503", "mistral-large-2411"],  # open-mistral-nemo slow/deprecated
        "cohere": ["command-r-08-2024", "command-r-plus-08-2024"],
        "xai": ["grok-2", "grok-2-mini"],
        "deepseek": ["deepseek-chat", "deepseek-coder"],
        "together": ["meta-llama/Llama-3.3-70B-Instruct-Turbo"],
        "perplexity": ["sonar", "sonar-pro"],
    }
    
    models_to_test = provider_models.get(provider_name, ["gpt-4o-mini"])
    
    async def test_single_model(model: str) -> dict:
        """Test a single model with timeout"""
        try:
            llm_config = LLMConfig(
                provider=LLMProvider(provider_name) if provider_name in [p.value for p in LLMProvider] else LLMProvider.CUSTOM,
                model=model,
                api_key=api_key,
                api_base=api_base
            )
            
            if provider_name in ['groq', 'xai', 'mistral', 'deepseek', 'together', 'perplexity', 'lmstudio']:
                provider = OpenAICompatibleLLM(llm_config, provider_name)
            else:
                provider = ProviderFactory.get_llm_provider(llm_config)
            
            # 30 second timeout per model
            response = await asyncio.wait_for(
                provider.generate([{"role": "user", "content": "Hi"}], max_tokens=5, model=model),
                timeout=30.0
            )
            
            print(f"[Test LLM] ✅ {model}: OK")
            return {"model": model, "status": "success", "response": response[:30] if response else ""}
            
        except asyncio.TimeoutError:
            print(f"[Test LLM] ⏱️ {model}: Timeout")
            return {"model": model, "status": "error", "message": "Timeout (30s)"}
        except Exception as e:
            print(f"[Test LLM] ❌ {model}: {str(e)[:60]}")
            return {"model": model, "status": "error", "message": str(e)}
    
    # Run all tests in parallel
    results = await asyncio.gather(*[test_single_model(m) for m in models_to_test])
    
    passed = [r for r in results if r['status'] == 'success']
    failed = [r for r in results if r['status'] == 'error']
    
    print(f"[Test LLM] Result: {len(passed)}/{len(results)} passed\n")
    
    if len(failed) == 0:
        return {"status": "success", "message": f"All {len(results)} models OK!", "results": results}
    elif len(passed) > 0:
        return {"status": "partial", "message": f"{len(passed)}/{len(results)} passed. Failed: {', '.join([r['model'] for r in failed])}", "results": results}
    else:
        return {"status": "error", "message": f"All models failed: {failed[0]['message']}", "results": results}


@app.post("/api/settings/test-embedding")
async def test_embedding_connection(request: Dict[str, Any]):
    """Test embedding provider connection"""
    try:
        # Create temporary config
        emb_config = EmbeddingConfig(**request)
        provider = ProviderFactory.get_embedding_provider(emb_config)
        
        # Test with a simple text
        embeddings = await provider.embed(["Hello, AgentForge!"])
        
        return {
            "status": "success", 
            "dimensions": len(embeddings[0]),
            "sample": embeddings[0][:5]  # First 5 values
        }
    except Exception as e:
        return {"status": "error", "message": str(e)}


# ============================================================================
# ORGANIZATION BRANDING - Customization for End User Portal
# ============================================================================

@app.get("/api/organization/branding")
async def get_organization_branding(current_user: User = Depends(get_current_user)):
    """Get organization branding settings for End User Portal"""
    try:
        from database.base import get_db_session
        from database.models import Organization
        
        org_id = current_user.org_id if current_user else None
        
        with get_db_session() as db:
            if org_id:
                org = db.query(Organization).filter(Organization.id == org_id).first()
                if org:
                    branding = org.branding or {}
                    return {
                        "logo_url": org.logo_url,
                        "favicon_url": getattr(org, 'favicon_url', None),
                        "name": org.name,
                        **branding
                    }
        
        return {"logo_url": None, "name": "AgentForge"}
    except Exception as e:
        print(f"⚠️  Failed to get branding: {e}")
        return {"logo_url": None, "name": "AgentForge"}


@app.put("/api/organization/branding")
async def update_organization_branding(
    request: Dict[str, Any],
    current_user: User = Depends(get_current_user)
):
    """Update organization branding settings (Admin only)"""
    # Check if user has admin permissions
    if not current_user or not any(r in ['super_admin', 'admin', 'org_admin'] for r in (current_user.roles or [])):
        raise HTTPException(403, "Admin access required")
    
    try:
        from database.base import get_db_session
        from database.models import Organization
        
        org_id = current_user.org_id
        
        with get_db_session() as db:
            org = db.query(Organization).filter(Organization.id == org_id).first()
            if not org:
                raise HTTPException(404, "Organization not found")
            
            # Update direct fields
            if 'logo_url' in request:
                org.logo_url = request['logo_url']
            if 'favicon_url' in request:
                org.favicon_url = request['favicon_url']
            if 'name' in request:
                org.name = request['name']
            
            # Update branding JSON
            branding = org.branding or {}
            branding_fields = [
                'primary_color', 'secondary_color',
                'banner_enabled', 'banner_text', 'banner_bg_color', 'banner_text_color',
                'chat_welcome_title', 'chat_welcome_message',
                'theme', 'custom_css'
            ]
            for field in branding_fields:
                if field in request:
                    branding[field] = request[field]
            
            org.branding = branding
            db.commit()
            
            return {"status": "success", "message": "Branding updated"}
    except HTTPException:
        raise
    except Exception as e:
        print(f"❌ Failed to update branding: {e}")
        raise HTTPException(500, f"Failed to update branding: {str(e)}")


# ============================================================================
# INTEGRATIONS SETTINGS
# ============================================================================

@app.post("/api/settings/integrations")
async def save_integration_settings(request: Dict[str, Any]):
    """Save OAuth integration credentials (Google, Microsoft)"""
    provider = request.get('provider')
    client_id = request.get('clientId')
    client_secret = request.get('clientSecret')
    
    if not provider or not client_id or not client_secret:
        raise HTTPException(400, "Missing required fields")
    
    # Store in environment for immediate use
    if provider == 'google':
        os.environ['GOOGLE_CLIENT_ID'] = client_id
        os.environ['GOOGLE_CLIENT_SECRET'] = client_secret
    elif provider == 'microsoft':
        os.environ['MICROSOFT_CLIENT_ID'] = client_id
        os.environ['MICROSOFT_CLIENT_SECRET'] = client_secret
    else:
        raise HTTPException(400, f"Unknown provider: {provider}")
    
    # Save to AppState for persistence
    app_state.integrations[provider] = {
        'client_id': client_id,
        'client_secret': client_secret
    }
    
    # Save to database
    print(f"💾 [API] Saving {provider} integration to database...")
    try:
        from database.services import SystemSettingsService
        integrations_key = f"integrations_{provider}"
        SystemSettingsService.set_system_setting(
            key=integrations_key,
            value={"client_id": client_id, "client_secret": client_secret},
            value_type="json",
            category="integrations",
            description=f"{provider.title()} OAuth integration credentials"
        )
        print(f"✅ [API] {provider.title()} integration saved successfully to database")
    except Exception as e:
        print(f"⚠️  [API ERROR] Database save failed: {e}, saving to disk only")
        import traceback
        traceback.print_exc()
    
    # Save to disk immediately (backup)
    app_state.save_to_disk()
    
    print(f"✅ Saved {provider} integration credentials")
    
    return {"status": "success", "message": f"{provider} integration saved"}


@app.get("/api/settings/integrations")
async def get_integration_settings():
    """Get current integration status"""
    return {
        "google": {
            "configured": bool(os.environ.get('GOOGLE_CLIENT_ID') or app_state.integrations.get('google', {}).get('client_id'))
        },
        "microsoft": {
            "configured": bool(os.environ.get('MICROSOFT_CLIENT_ID') or app_state.integrations.get('microsoft', {}).get('client_id'))
        }
    }


@app.post("/api/settings/reindex")
async def reindex_knowledge_base():
    """Reindex all documents in the knowledge base with current embedding settings"""
    try:
        if app_state.settings.vector_db.provider == VectorDBProvider.MEMORY:
            return {"status": "success", "message": "Memory provider doesn't require reindexing"}
        
        # Get providers
        embedding_provider = app_state.get_embedding_provider()
        vector_db = app_state.get_vector_db()
        
        # Clear existing index
        if hasattr(vector_db, 'delete_by_filter'):
            await vector_db.delete_by_filter({})
        
        # Reindex all chunks
        total_chunks = len(app_state.document_chunks)
        indexed = 0
        batch_size = 100
        
        for i in range(0, total_chunks, batch_size):
            batch = app_state.document_chunks[i:i+batch_size]
            texts = [c.get('text', '') for c in batch]
            
            # Get embeddings
            embeddings = await embedding_provider.embed(texts)
            
            # Add to vector DB
            documents = [{
                'id': c.get('chunk_id', str(uuid.uuid4())),
                'text': c.get('text', ''),
                'source': c.get('source', ''),
                'tool_id': c.get('tool_id', ''),
                'type': c.get('type', 'document')
            } for c in batch]
            
            await vector_db.add_documents(documents, embeddings)
            indexed += len(batch)
        
        return {"status": "success", "message": f"Reindexed {indexed} chunks"}
    except Exception as e:
        raise HTTPException(500, f"Reindexing failed: {str(e)}")


# Tool Endpoints

def get_user_group_ids(user_id: str) -> List[str]:
    """
    Get all group IDs that a user belongs to by checking group.member_ids.
    This is the source of truth since user.group_ids may not be in sync.
    """
    if not user_id or not SECURITY_AVAILABLE:
        print(f"   👥 [GET_GROUPS] No user_id or security not available")
        return []
    
    try:
        group_ids = []
        total_groups = len(security_state.groups)
        print(f"   👥 [GET_GROUPS] Checking {total_groups} groups for user {user_id[:8]}...")
        
        for group in security_state.groups.values():
            member_ids = getattr(group, 'member_ids', []) or []
            user_ids = getattr(group, 'user_ids', []) or []
            all_members = list(set(member_ids + user_ids))
            
            # Convert all to strings for comparison
            all_members_str = [str(m) for m in all_members]
            user_id_str = str(user_id)
            
            print(f"      📁 Group '{group.name}' (id={group.id[:8]}...): members={all_members_str}")
            
            if user_id_str in all_members_str:
                print(f"      ✅ User IS in group '{group.name}'!")
                group_ids.append(group.id)
        
        print(f"   👥 [GET_GROUPS] Found {len(group_ids)} groups for user")
        return group_ids
    except Exception as e:
        print(f"⚠️ Error getting user groups: {e}")
        import traceback
        traceback.print_exc()
        return []

def check_tool_access(tool: ToolConfiguration, user_id: str, user_group_ids: List[str] = None, permission: str = 'view', user_role_ids: List[str] = None, user_obj = None) -> bool:
    """
    Check if a user has access to a tool.
    
    permission: 'view', 'edit', 'delete', 'execute'
    
    NOTE: Admins are treated as normal users for tools/agents.
    They only have access based on what the owner grants them.
    """
    tool_owner = getattr(tool, 'owner_id', None) or ''
    tool_access = getattr(tool, 'access_type', 'owner_only') or 'owner_only'
    
    if not user_id:
        result = tool_access == 'public'
        print(f"   🔐 [ACCESS] No user, tool '{tool.name}' access_type='{tool_access}' -> {result}")
        return result
    
    # Owner always has full access (compare as strings)
    is_owner = str(tool_owner) == str(user_id) if tool_owner else False
    if is_owner:
        print(f"   🔐 [ACCESS] User is OWNER of '{tool.name}' -> True")
        return True
    
    # Check by access_type
    if tool_access == 'public':
        if permission == 'view' or permission == 'execute':
            print(f"   🔐 [ACCESS] Tool '{tool.name}' is PUBLIC, permission='{permission}' -> True")
            return True
    
    if tool_access == 'authenticated':
        if permission == 'view' or permission == 'execute':
            print(f"   🔐 [ACCESS] Tool '{tool.name}' is AUTHENTICATED, permission='{permission}' -> True")
            return True
    
    if tool_access == 'specific_users':
        # Get permission lists
        allowed_users = getattr(tool, 'allowed_user_ids', []) or []
        allowed_groups = getattr(tool, 'allowed_group_ids', []) or []
        can_edit_ids = getattr(tool, 'can_edit_user_ids', []) or []
        can_delete_ids = getattr(tool, 'can_delete_user_ids', []) or []
        can_execute_ids = getattr(tool, 'can_execute_user_ids', []) or []
        
        print(f"   📋 [PERMS] Tool '{tool.name}': allowed_users={allowed_users}, allowed_groups={allowed_groups}")
        print(f"   📋 [PERMS] can_edit={can_edit_ids}, can_delete={can_delete_ids}, can_execute={can_execute_ids}")
        
        # Check if user is directly in allowed list (for view permission)
        user_in_allowed = user_id in allowed_users
        
        # Check if any of user's groups are in allowed list
        group_in_allowed = False
        matched_group_id = None
        if user_group_ids:
            for group_id in user_group_ids:
                if group_id in allowed_groups:
                    group_in_allowed = True
                    matched_group_id = group_id
                    break
        
        # User has base access if in allowed_users or in an allowed_group
        has_base_access = user_in_allowed or group_in_allowed
        
        if not has_base_access:
            print(f"   🔐 [ACCESS] User '{user_id[:8]}...' NOT in allowed lists -> False")
            return False
        
        # VIEW: Anyone in allowed list can view
        if permission == 'view':
            print(f"   🔐 [ACCESS] User in allowed list -> view=True")
            return True
        
        # EXECUTE: Requires explicit can_execute permission
        if permission == 'execute':
            # Check if user is directly in can_execute_user_ids
            if user_id in can_execute_ids:
                print(f"   🔐 [ACCESS] User '{user_id[:8]}...' in can_execute_user_ids -> execute=True")
                return True
            # Check if any of user's groups are in can_execute_user_ids (format: "group:id")
            if user_group_ids:
                for group_id in user_group_ids:
                    if f"group:{group_id}" in can_execute_ids:
                        print(f"   🔐 [ACCESS] Group '{group_id}' in can_execute -> execute=True")
                        return True
            print(f"   🔐 [ACCESS] User NOT in can_execute_user_ids -> execute=False")
            return False
        
        # EDIT: Requires explicit can_edit permission
        if permission == 'edit':
            # Check if user is directly in can_edit_user_ids
            if user_id in can_edit_ids:
                print(f"   🔐 [ACCESS] User '{user_id[:8]}...' in can_edit_user_ids -> edit=True")
                return True
            # Check if any of user's groups are in can_edit_user_ids (format: "group:id")
            if user_group_ids:
                for group_id in user_group_ids:
                    if f"group:{group_id}" in can_edit_ids:
                        print(f"   🔐 [ACCESS] Group '{group_id}' in can_edit -> edit=True")
                        return True
            print(f"   🔐 [ACCESS] User NOT in can_edit_user_ids -> edit=False")
            return False
        
        # DELETE: Requires explicit can_delete permission
        if permission == 'delete':
            # Check if user is directly in can_delete_user_ids
            if user_id in can_delete_ids:
                print(f"   🔐 [ACCESS] User '{user_id[:8]}...' in can_delete_user_ids -> delete=True")
                return True
            # Check if any of user's groups are in can_delete_user_ids (format: "group:id")
            if user_group_ids:
                for group_id in user_group_ids:
                    if f"group:{group_id}" in can_delete_ids:
                        print(f"   🔐 [ACCESS] Group '{group_id}' in can_delete -> delete=True")
                        return True
            print(f"   🔐 [ACCESS] User NOT in can_delete_user_ids -> delete=False")
            return False
    
    print(f"   🔐 [ACCESS] Tool '{tool.name}' (owner='{tool_owner}', access='{tool_access}'), user='{user_id}', permission='{permission}' -> False")
    return False


@app.get("/api/tools/accessible")
async def list_accessible_tools(current_user: User = Depends(get_current_user)):
    """
    Get tools that the current user can USE in agents (execute permission).
    Used when selecting tools for an agent.
    """
    if not current_user:
        raise HTTPException(401, "Authentication required")
    
    user_id = str(current_user.id)
    # Get user's groups dynamically
    user_group_ids = get_user_group_ids(user_id)
    
    accessible_tools = []
    
    for tool in app_state.tools.values():
        if check_tool_access(tool, user_id, user_group_ids, 'execute'):
            accessible_tools.append({
                **tool.dict(),
                "documents_count": len([d for d in app_state.documents.values() if d.tool_id == tool.id]),
                "pages_count": len([p for p in app_state.scraped_pages.values() if p.tool_id == tool.id]),
                "can_edit": check_tool_access(tool, user_id, user_group_ids, 'edit'),
                "can_delete": check_tool_access(tool, user_id, user_group_ids, 'delete'),
                "is_owner": tool.owner_id == user_id
            })
    
    return {"tools": accessible_tools}


@app.get("/api/tools")
async def list_tools(current_user: User = Depends(get_current_user_optional)):
    """
    List tools - with ownership-based filtering.
    Users only see tools they own or have been granted view access to.
    """
    user_id = str(current_user.id) if current_user else None
    # Get user's groups dynamically (not from user.group_ids which may be stale)
    user_group_ids = get_user_group_ids(user_id) if user_id else []
    
    print(f"📋 [LIST_TOOLS] User: {user_id}, groups: {user_group_ids}, total tools in memory: {len(app_state.tools)}")
    
    viewable_tools = []
    
    for tool in app_state.tools.values():
        # Debug: log tool ownership info
        tool_owner_id = getattr(tool, 'owner_id', None) or ''
        is_owner = (str(tool_owner_id) == str(user_id)) if user_id and tool_owner_id else False
        print(f"   🔧 Tool '{tool.name}' (id={tool.id[:8]}...): owner_id='{tool_owner_id}', user_id='{user_id}', is_owner={is_owner}, access_type={getattr(tool, 'access_type', 'N/A')}")
        
        # Check if user can view this tool
        if check_tool_access(tool, user_id, user_group_ids, 'view'):
            can_edit = check_tool_access(tool, user_id, user_group_ids, 'edit')
            can_delete = check_tool_access(tool, user_id, user_group_ids, 'delete')
            print(f"      ✅ Viewable: can_edit={can_edit}, can_delete={can_delete}, is_owner={is_owner}")
            viewable_tools.append({
                **tool.dict(),
                "documents_count": len([d for d in app_state.documents.values() if d.tool_id == tool.id]),
                "pages_count": len([p for p in app_state.scraped_pages.values() if p.tool_id == tool.id]),
                "can_edit": can_edit,
                "can_delete": can_delete,
                "can_execute": check_tool_access(tool, user_id, user_group_ids, 'execute'),
                "is_owner": is_owner
            })
    
    print(f"📋 [LIST_TOOLS] Returning {len(viewable_tools)} viewable tools")
    return {"tools": viewable_tools}


@app.post("/api/tools")
async def create_tool(request: CreateToolRequest, current_user: User = Depends(get_current_user_optional)):
    # Check for duplicate name
    existing_names = [t.name.lower() for t in app_state.tools.values()]
    if request.name.lower() in existing_names:
        raise HTTPException(400, f"A tool with name '{request.name}' already exists. Please use a unique name.")
    
    # Get owner_id from current user
    owner_id = str(current_user.id) if current_user else "system"
    
    api_config = None
    kb_config = None
    
    if request.type == 'api' and request.api_config:
        params = [APIInputParameter(**p) for p in request.api_config.get('input_parameters', [])]
        api_config = APIEndpointConfig(base_url=request.api_config.get('base_url', ''), http_method=request.api_config.get('http_method', 'GET'), endpoint_path=request.api_config.get('endpoint_path', ''), auth_type=request.api_config.get('auth_type', 'none'), auth_value=request.api_config.get('auth_value', ''), api_key_name=request.api_config.get('api_key_name', 'X-API-Key'), api_key_location=request.api_config.get('api_key_location', 'header'), headers=request.api_config.get('headers', {}), input_parameters=params)
    
    # Handle Knowledge Base tool
    if request.type == 'knowledge' and request.config:
        cfg = request.config
        kb_config = KnowledgeBaseConfig(
            collection_id=cfg.get('collection_id', ''),
            embedding=KnowledgeBaseEmbeddingConfig(
                use_global=cfg.get('embedding', {}).get('use_global', True),
                provider=cfg.get('embedding', {}).get('provider', 'openai'),
                model=cfg.get('embedding', {}).get('model', 'text-embedding-3-small'),
                local_model=cfg.get('embedding', {}).get('local_model', 'all-MiniLM-L6-v2'),
                api_key=cfg.get('embedding', {}).get('api_key', '')
            ),
            vector_db=KnowledgeBaseVectorDBConfig(
                use_global=cfg.get('vector_db', {}).get('use_global', True),
                provider=cfg.get('vector_db', {}).get('provider', 'chromadb'),
                pinecone_api_key=cfg.get('vector_db', {}).get('pinecone_api_key', '')
            ),
            chunk_size=cfg.get('chunk_size', 1000),
            chunk_overlap=cfg.get('chunk_overlap', 200),
            top_k=cfg.get('top_k', 5),
            search_type=cfg.get('search_type', 'hybrid'),
            similarity_threshold=cfg.get('similarity_threshold', 0.7),
            reranking=cfg.get('reranking', 'none'),
            context_window=cfg.get('context_window', 4000),
            include_metadata=cfg.get('include_metadata', True),
            auto_reindex=cfg.get('auto_reindex', False)
        )
    
    tool = ToolConfiguration(
        type=request.type, 
        name=request.name, 
        description=request.description, 
        config=request.config, 
        api_config=api_config,
        kb_config=kb_config,
        # Access Control
        owner_id=owner_id,
        access_type=request.access_type,
        allowed_user_ids=request.allowed_user_ids,
        allowed_group_ids=request.allowed_group_ids,
        can_edit_user_ids=request.can_edit_user_ids,
        can_delete_user_ids=request.can_delete_user_ids,
        can_execute_user_ids=request.can_execute_user_ids
    )
    
    # Set collection_id if not provided
    if request.type == 'knowledge':
        if not tool.config.get('collection_id'):
            tool.config['collection_id'] = f"kb_{tool.id.replace('-', '_')}"
        
        # Initialize vector collection for this KB
        try:
            emb_config = tool.config.get('embedding', {})
            vdb_config = tool.config.get('vector_db', {})
            collection_id = tool.config['collection_id']
            
            embedding_provider = ProviderFactory.get_kb_embedding_provider(
                emb_config, 
                app_state.settings.embedding
            )
            embedding_dim = embedding_provider.get_dimensions()
            
            vector_db = ProviderFactory.get_kb_vector_db(
                vdb_config,
                app_state.settings.vector_db,
                collection_id,
                embedding_dim
            )
            
            print(f"✅ Initialized Knowledge Base '{tool.name}' with collection '{collection_id}'")
        except Exception as e:
            print(f"⚠️ Warning: Could not initialize KB collection: {e}")
    
    app_state.tools[tool.id] = tool
    
    # Save to database
    try:
        from database.services import ToolService
        tool_dict = tool.dict()
        if tool.api_config:
            tool_dict['api_config'] = tool.api_config.dict() if hasattr(tool.api_config, 'dict') else tool.api_config
        ToolService.create_tool(tool_dict, "org_default", "system")
    except Exception as e:
        print(f"⚠️  [DATABASE] Failed to save tool: {e}")
    
    app_state.save_to_disk()
    return {"status": "success", "tool_id": tool.id, "tool": tool.dict()}


@app.get("/api/tools/{tool_id}")
async def get_tool(tool_id: str):
    if tool_id not in app_state.tools:
        raise HTTPException(404, "Tool not found")
    tool = app_state.tools[tool_id]
    
    # Debug: Log access control data
    print(f"📋 [GET TOOL] Tool '{tool.name}' access_type={tool.access_type}, allowed_users={tool.allowed_user_ids}, allowed_groups={tool.allowed_group_ids}")
    
    # Get regular documents
    documents = [d.dict() for d in app_state.documents.values() if d.tool_id == tool_id]
    
    # Get scraped pages
    pages = [p.dict() for p in app_state.scraped_pages.values() if p.tool_id == tool_id]
    
    # For knowledge/document tools, also get chunks as virtual documents (for demo/text/table entries)
    demo_documents = []
    if tool.type in ['knowledge', 'document']:
        chunks = [c for c in app_state.document_chunks if c.get('tool_id') == tool_id]
        
        # Group chunks by source
        sources = {}
        for chunk in chunks:
            source = chunk.get('source', 'Unknown')
            is_table = chunk.get('is_table', False) or chunk.get('type') == 'table'
            
            if source not in sources:
                sources[source] = {
                    'id': f"demo-doc-{chunk.get('chunk_id', '')}",
                    'original_name': source,
                    'file_type': 'table' if is_table else chunk.get('type', 'document'),
                    'file_size': len(chunk.get('text', '')),
                    'chunks_count': 0,
                    'content': '',
                    'status': 'ready',
                    'is_demo': True,
                    'is_table': is_table,
                    'rows_count': chunk.get('rows_count', 0) if is_table else 0
                }
            sources[source]['chunks_count'] += 1
            sources[source]['content'] += chunk.get('text', '') + '\n\n'
            sources[source]['file_size'] = len(sources[source]['content'])
            # Update is_table if any chunk is a table
            if is_table:
                sources[source]['is_table'] = True
                sources[source]['rows_count'] = chunk.get('rows_count', 0)
        
        demo_documents = list(sources.values())
    
    # Combine regular documents with demo documents
    all_documents = documents + [d for d in demo_documents if not any(doc['original_name'] == d['original_name'] for doc in documents)]
    
    return {
        **tool.dict(), 
        "documents": all_documents, 
        "scraped_pages": pages,
        "has_mock_response": bool(tool.config.get('mock_response')) if tool.config else False
    }


@app.get("/api/tools/{tool_id}/data")
async def get_tool_data(tool_id: str):
    """Get all extracted data and chunks for a tool with search support."""
    if tool_id not in app_state.tools:
        raise HTTPException(404, "Tool not found")
    
    tool = app_state.tools[tool_id]
    
    # Get all chunks for this tool
    chunks = [c for c in app_state.document_chunks if c.get('tool_id') == tool_id]
    
    # Get unique sources
    sources = list(set(c.get('source', 'Unknown') for c in chunks))
    
    # Calculate total characters
    total_chars = sum(len(c.get('text', '')) for c in chunks)
    
    # Format chunks for response
    formatted_chunks = [{
        'id': c.get('chunk_id', ''),
        'text': c.get('text', ''),
        'source': c.get('source', 'Unknown'),
        'type': c.get('type', 'unknown')
    } for c in chunks]
    
    return {
        "id": tool.id,
        "name": tool.name,
        "type": tool.type,
        "description": tool.description,
        "chunks": formatted_chunks,
        "total_chunks": len(chunks),
        "total_chars": total_chars,
        "sources": sources
    }


class UpdateDemoDocRequest(BaseModel):
    source: str
    content: str


@app.put("/api/tools/{tool_id}/demo-document")
async def update_demo_document(tool_id: str, request: UpdateDemoDocRequest):
    """Update content of a demo document (chunks)"""
    if tool_id not in app_state.tools:
        raise HTTPException(404, "Tool not found")
    
    # Find and update chunks with this source
    updated_count = 0
    for i, chunk in enumerate(app_state.document_chunks):
        if chunk.get('tool_id') == tool_id and chunk.get('source') == request.source:
            # Remove old chunk
            app_state.document_chunks[i] = None
            updated_count += 1
    
    # Remove None entries
    app_state.document_chunks = [c for c in app_state.document_chunks if c is not None]
    
    # Add new chunk with updated content
    new_chunk = {
        "chunk_id": f"chunk-{uuid.uuid4().hex[:8]}",
        "text": request.content,
        "source": request.source,
        "type": "document",
        "tool_id": tool_id
    }
    app_state.document_chunks.append(new_chunk)
    
    app_state.save_to_disk()
    return {"status": "success", "updated_chunks": updated_count}


@app.post("/api/tools/{tool_id}/demo-document")
async def add_demo_document(tool_id: str, request: UpdateDemoDocRequest):
    """Add a new demo document to knowledge base"""
    if tool_id not in app_state.tools:
        raise HTTPException(404, "Tool not found")
    
    # Add new chunk
    new_chunk = {
        "chunk_id": f"chunk-{uuid.uuid4().hex[:8]}",
        "text": request.content,
        "source": request.source,
        "type": "document",
        "tool_id": tool_id
    }
    app_state.document_chunks.append(new_chunk)
    
    app_state.save_to_disk()
    return {"status": "success", "chunk_id": new_chunk["chunk_id"]}


@app.delete("/api/tools/{tool_id}/demo-document/{source}")
async def delete_demo_document(tool_id: str, source: str):
    """Delete a demo document from knowledge base"""
    if tool_id not in app_state.tools:
        raise HTTPException(404, "Tool not found")
    
    # URL decode the source
    from urllib.parse import unquote
    source = unquote(source)
    
    # Remove chunks with this source
    original_count = len(app_state.document_chunks)
    app_state.document_chunks = [
        c for c in app_state.document_chunks 
        if not (c.get('tool_id') == tool_id and c.get('source') == source)
    ]
    deleted_count = original_count - len(app_state.document_chunks)
    
    app_state.save_to_disk()
    return {"status": "success", "deleted_chunks": deleted_count}


# ========== TABLE ENTRY ENDPOINTS ==========

class TableEntryRequest(BaseModel):
    source: str = None
    content: str
    table_data: Dict[str, Any]  # {headers: [], rows: [[]]}


@app.post("/api/tools/{tool_id}/table-entry")
async def add_table_entry(tool_id: str, request: TableEntryRequest):
    """Add a new table entry to knowledge base"""
    if tool_id not in app_state.tools:
        raise HTTPException(404, "Tool not found")
    
    # Add new chunk with table data
    new_chunk = {
        "chunk_id": f"table-{uuid.uuid4().hex[:8]}",
        "text": request.content,
        "source": request.source,
        "type": "table",
        "tool_id": tool_id,
        "table_data": request.table_data,
        "is_table": True,
        "rows_count": len(request.table_data.get('rows', []))
    }
    app_state.document_chunks.append(new_chunk)
    
    app_state.save_to_disk()
    return {"status": "success", "chunk_id": new_chunk["chunk_id"]}


@app.get("/api/tools/{tool_id}/table-entry/{source}")
async def get_table_entry(tool_id: str, source: str):
    """Get a table entry from knowledge base"""
    if tool_id not in app_state.tools:
        raise HTTPException(404, "Tool not found")
    
    from urllib.parse import unquote
    source = unquote(source)
    
    # Find chunk with this source
    for chunk in app_state.document_chunks:
        if chunk.get('tool_id') == tool_id and chunk.get('source') == source:
            return {
                "source": source,
                "content": chunk.get('text', ''),
                "table_data": chunk.get('table_data', {'headers': [], 'rows': []}),
                "is_table": chunk.get('is_table', False)
            }
    
    raise HTTPException(404, "Table entry not found")


@app.put("/api/tools/{tool_id}/table-entry/{source}")
async def update_table_entry(tool_id: str, source: str, request: Dict[str, Any]):
    """Update a table entry in knowledge base"""
    if tool_id not in app_state.tools:
        raise HTTPException(404, "Tool not found")
    
    from urllib.parse import unquote
    source = unquote(source)
    
    # Find and update chunk
    for i, chunk in enumerate(app_state.document_chunks):
        if chunk.get('tool_id') == tool_id and chunk.get('source') == source:
            app_state.document_chunks[i]['text'] = request.get('content', chunk.get('text', ''))
            app_state.document_chunks[i]['table_data'] = request.get('table_data', chunk.get('table_data', {}))
            app_state.document_chunks[i]['rows_count'] = len(request.get('table_data', {}).get('rows', []))
            app_state.save_to_disk()
            return {"status": "success"}
    
    raise HTTPException(404, "Table entry not found")



@app.delete("/api/tools/{tool_id}")
async def delete_tool(tool_id: str, current_user: User = Depends(get_current_user_optional)):
    if tool_id not in app_state.tools:
        raise HTTPException(404, "Tool not found")
    
    tool = app_state.tools[tool_id]
    user_id = str(current_user.id) if current_user else None
    # Get user's groups dynamically
    user_group_ids = get_user_group_ids(user_id) if user_id else []
    
    # Check if user can delete this tool
    if not check_tool_access(tool, user_id, user_group_ids, 'delete'):
        raise HTTPException(403, "You don't have permission to delete this tool")
    
    # Check if tool is being used by any agent
    agents_using_tool = []
    
    # Check in-memory agents
    for agent in app_state.agents.values():
        agent_tool_ids = agent.tool_ids if isinstance(agent.tool_ids, list) else []
        if tool_id in agent_tool_ids:
            agents_using_tool.append({
                "id": agent.id,
                "name": agent.name
            })
    
    # Also check database agents
    try:
        from database.services import AgentService
        db_agents = AgentService.get_all_agents()
        for agent in db_agents:
            agent_tool_ids = agent.tool_ids if isinstance(agent.tool_ids, list) else []
            if tool_id in agent_tool_ids:
                # Avoid duplicates
                if not any(a["id"] == agent.id for a in agents_using_tool):
                    agents_using_tool.append({
                        "id": agent.id,
                        "name": agent.name
                    })
    except Exception as e:
        print(f"⚠️  [DATABASE] Failed to check agents: {e}")
    
    # If tool is being used, don't allow deletion
    if agents_using_tool:
        agent_names = ", ".join([a["name"] for a in agents_using_tool])
        raise HTTPException(
            status_code=400, 
            detail=f"Cannot delete tool: It is currently being used by {len(agents_using_tool)} agent(s): {agent_names}. Please remove the tool from these agents first."
        )
    
    del app_state.tools[tool_id]
    
    # Delete from database
    try:
        from database.services import ToolService
        ToolService.delete_tool(tool_id)
    except Exception as e:
        print(f"⚠️  [DATABASE] Failed to delete tool: {e}")
    
    doc_ids = [d.id for d in app_state.documents.values() if d.tool_id == tool_id]
    for doc_id in doc_ids:
        del app_state.documents[doc_id]
    page_ids = [p.id for p in app_state.scraped_pages.values() if p.tool_id == tool_id]
    for page_id in page_ids:
        del app_state.scraped_pages[page_id]
    app_state.document_chunks = [c for c in app_state.document_chunks if c.get('tool_id') != tool_id]
    
    app_state.save_to_disk()
    return {"status": "success"}


@app.post("/api/tools/cleanup-duplicates")
async def cleanup_duplicate_tools():
    """Remove duplicate tools, keeping only the first occurrence of each name"""
    seen_names = {}
    duplicates_removed = []
    tools_to_delete = []
    
    # Make a copy of items to iterate safely
    tools_list = list(app_state.tools.items())
    
    # Find duplicates (keep first, mark rest for deletion)
    for tool_id, tool in tools_list:
        name_lower = tool.name.lower().strip()
        if name_lower in seen_names:
            tools_to_delete.append(tool_id)
            duplicates_removed.append({"id": tool_id, "name": tool.name})
        else:
            seen_names[name_lower] = tool_id
    
    print(f"[Cleanup] Found {len(tools_to_delete)} duplicates to remove")
    
    # Delete duplicates
    for tool_id in tools_to_delete:
        if tool_id in app_state.tools:
            del app_state.tools[tool_id]
            print(f"[Cleanup] Deleted tool: {tool_id}")
        
        # Also clean up related data
        doc_ids = [d.id for d in list(app_state.documents.values()) if d.tool_id == tool_id]
        for doc_id in doc_ids:
            if doc_id in app_state.documents:
                del app_state.documents[doc_id]
        
        page_ids = [p.id for p in list(app_state.scraped_pages.values()) if p.tool_id == tool_id]
        for page_id in page_ids:
            if page_id in app_state.scraped_pages:
                del app_state.scraped_pages[page_id]
        
        app_state.document_chunks = [c for c in app_state.document_chunks if c.get('tool_id') != tool_id]
        
        for agent in app_state.agents.values():
            if tool_id in agent.tool_ids:
                agent.tool_ids.remove(tool_id)
    
    app_state.save_to_disk()
    
    print(f"[Cleanup] Done. Removed {len(duplicates_removed)}, remaining {len(app_state.tools)}")
    
    return {
        "status": "success",
        "duplicates_removed": len(duplicates_removed),
        "removed_tools": duplicates_removed,
        "remaining_tools": len(app_state.tools)
    }


@app.delete("/api/tools/all")
async def delete_all_tools():
    """Delete all tools - useful for cleanup"""
    count = len(app_state.tools)
    
    # Clear all tools
    app_state.tools.clear()
    
    # Clear related data
    app_state.documents.clear()
    app_state.scraped_pages.clear()
    app_state.document_chunks.clear()
    
    # Remove tool references from agents
    for agent in app_state.agents.values():
        agent.tool_ids = []
    
    app_state.save_to_disk()
    
    return {
        "status": "success",
        "deleted_count": count
    }


@app.post("/api/tools/parse-openapi")
async def parse_openapi_spec(file: UploadFile = File(...)):
    content = await file.read()
    content_str = content.decode('utf-8')
    file_type = file.filename.split('.')[-1].lower() if file.filename else 'json'
    result = OpenAPIParser.parse(content_str, file_type)
    if 'error' in result:
        raise HTTPException(400, f"Failed to parse: {result['error']}")
    return result


@app.post("/api/tools/test-api")
async def test_api(request: APITestRequest):
    try:
        url = request.base_url.rstrip('/')
        path = request.endpoint_path
        for key, value in request.parameters.items():
            path = path.replace(f"{{{key}}}", str(value))
        url = f"{url}/{path.lstrip('/')}" if path else url
        headers = dict(request.headers)
        if request.auth_type == 'bearer' and request.auth_value:
            headers['Authorization'] = f'Bearer {request.auth_value}'
        elif request.auth_type == 'api_key' and request.auth_value:
            if request.api_key_location == 'header':
                headers[request.api_key_name] = request.auth_value
        elif request.auth_type == 'basic' and request.auth_value:
            import base64
            headers['Authorization'] = f'Basic {base64.b64encode(request.auth_value.encode()).decode()}'
        query_params = {}
        if request.auth_type == 'api_key' and request.api_key_location == 'query':
            query_params[request.api_key_name] = request.auth_value
        for key, value in request.parameters.items():
            if f"{{{key}}}" not in request.endpoint_path:
                query_params[key] = value
        async with httpx.AsyncClient(timeout=30.0) as client:
            response = await client.request(request.http_method.upper(), url, headers=headers, params=query_params if query_params else None)
            try:
                data = response.json()
            except:
                data = response.text
            return {"success": 200 <= response.status_code < 300, "status_code": response.status_code, "url": str(response.url), "data": data, "headers": dict(response.headers)}
    except Exception as e:
        return {"success": False, "error": str(e)}


class CodeGenerationRequest(BaseModel):
    description: str
    inputs: str = ""
    output_format: str = "json"
    language: str = "python"

@app.post("/api/tools/generate-code")
async def generate_tool_code(request: CodeGenerationRequest):
    """Use LLM to generate tool code based on description"""
    try:
        # Build the prompt for code generation
        lang_templates = {
            "python": '''def run(input_data: dict) -> dict:
    """
    Tool: {name}
    {description}
    
    Args:
        input_data: Dictionary with keys: {inputs}
    
    Returns:
        {output_format} result
    """
    # TODO: Implement tool logic
    pass''',
            "javascript": '''async function run(inputData) {{
    /**
     * Tool: {name}
     * {description}
     * 
     * @param {{Object}} inputData - {inputs}
     * @returns {{{output_format}}}
     */
    // TODO: Implement tool logic
}}
module.exports = {{ run }};''',
            "typescript": '''interface InputData {{
    // Define input structure based on: {inputs}
}}

async function run(inputData: InputData): Promise<{output_format}> {{
    /**
     * Tool: {name}
     * {description}
     */
    // TODO: Implement tool logic
}}
export {{ run }};'''
        }
        
        system_prompt = f"""You are an expert {request.language} developer. Generate production-ready code for a tool based on the user's description.

REQUIREMENTS:
1. The code MUST be complete and functional
2. Include proper error handling with try/catch
3. Add helpful comments
4. Use best practices for {request.language}
5. Handle edge cases
6. Return results in {request.output_format} format

INPUT PARAMETERS: {request.inputs or 'Will be passed as input_data dict/object'}

OUTPUT FORMAT: {request.output_format}

Generate ONLY the code, no explanations. The code should be ready to execute."""

        user_prompt = f"""Generate {request.language} code for this tool:

DESCRIPTION: {request.description}

INPUT PARAMETERS: {request.inputs or 'Flexible - determine from description'}
OUTPUT FORMAT: {request.output_format}

Requirements:
- Complete, working code
- Proper error handling
- Clear function signature
- Helpful comments
- Production-ready quality"""

        # Call LLM to generate code
        openai_key = os.getenv("OPENAI_API_KEY")
        if not openai_key:
            return {"error": "OpenAI API key not configured"}
        
        async with httpx.AsyncClient(timeout=60.0) as client:
            response = await client.post(
                "https://api.openai.com/v1/chat/completions",
                headers={
                    "Authorization": f"Bearer {openai_key}",
                    "Content-Type": "application/json"
                },
                json={
                    "model": "gpt-4o",
                    "messages": [
                        {"role": "system", "content": system_prompt},
                        {"role": "user", "content": user_prompt}
                    ],
                    "max_tokens": 2000,
                    "temperature": 0.3
                }
            )
            
            result = response.json()
            
            if "choices" in result and len(result["choices"]) > 0:
                code = result["choices"][0]["message"]["content"]
                
                # Clean up the code (remove markdown code blocks if present)
                code = code.strip()
                if code.startswith("```"):
                    lines = code.split("\n")
                    # Remove first and last lines (```language and ```)
                    if lines[0].startswith("```"):
                        lines = lines[1:]
                    if lines and lines[-1].strip() == "```":
                        lines = lines[:-1]
                    code = "\n".join(lines)
                
                return {"code": code, "language": request.language}
            else:
                return {"error": "Failed to generate code", "details": result}
                
    except Exception as e:
        print(f"[Code Generation] Error: {e}")
        return {"error": str(e)}


class ToolTestRequest(BaseModel):
    tool_type: str
    config: Dict[str, Any] = {}

@app.post("/api/tools/test")
async def test_tool_connection(request: ToolTestRequest):
    """Test tool connection before creating"""
    try:
        tool_type = request.tool_type
        config = request.config
        
        if tool_type == 'database':
            # Test database connection
            db_type = config.get('db_type', 'postgresql')
            host = config.get('host', 'localhost')
            port = config.get('port', 5432)
            # Simulate connection test
            return {"success": True, "message": f"Connection to {db_type}://{host}:{port} simulated (actual test requires driver)"}
        
        elif tool_type == 'email':
            provider = config.get('provider', 'smtp')
            if provider == 'smtp':
                host = config.get('smtp_host', '')
                return {"success": True, "message": f"SMTP server {host} ready to test"}
            return {"success": True, "message": f"{provider} configuration looks valid"}
        
        elif tool_type == 'webhook':
            url = config.get('url', '')
            if url:
                async with httpx.AsyncClient(timeout=10.0) as client:
                    response = await client.head(url)
                    return {"success": True, "message": f"Webhook URL reachable (status: {response.status_code})"}
            return {"success": False, "error": "No URL provided"}
        
        elif tool_type == 'websearch':
            provider = config.get('provider', 'tavily')
            api_key = config.get('api_key', '')
            if api_key:
                return {"success": True, "message": f"{provider} API key provided"}
            return {"success": False, "error": "API key required"}
        
        elif tool_type == 'knowledge':
            # Test KB configuration
            return {"success": True, "message": "Knowledge Base configuration valid"}
        
        else:
            return {"success": True, "message": f"Tool type '{tool_type}' ready to create"}
            
    except Exception as e:
        return {"success": False, "error": str(e)}


# Document Endpoints
@app.post("/api/tools/{tool_id}/documents")
async def upload_document(tool_id: str, file: UploadFile = File(...)):
    if tool_id not in app_state.tools:
        raise HTTPException(404, "Tool not found")
    tool = app_state.tools[tool_id]
    if tool.type not in ['document', 'knowledge']:
        raise HTTPException(400, "Tool type does not support documents")
    filename = file.filename
    file_ext = filename.split('.')[-1].lower() if '.' in filename else 'txt'
    content = await file.read()
    file_hash = hashlib.md5(content).hexdigest()[:8]
    stored_filename = f"{file_hash}_{filename}"
    upload_dir = os.environ.get("UPLOAD_PATH", "data/uploads")
    os.makedirs(upload_dir, exist_ok=True)
    file_path = os.path.join(upload_dir, stored_filename)
    with open(file_path, 'wb') as f:
        f.write(content)
    doc = Document(tool_id=tool_id, filename=stored_filename, original_name=filename, file_type=file_ext, file_size=len(content), status="processing")
    app_state.documents[doc.id] = doc
    app_state.save_to_disk()
    asyncio.create_task(process_document(doc.id, file_path))
    return {"status": "success", "document_id": doc.id, "document": doc.dict()}


async def process_document(doc_id: str, file_path: str):
    try:
        doc = app_state.documents.get(doc_id)
        if not doc:
            return
        text = await DocumentProcessor.extract_text(file_path, doc.file_type)
        doc.content = text
        if text.startswith("[") and text.endswith("]"):
            doc.status = "error"
            doc.error_message = text
            app_state.save_to_disk()
            return
        chunks = DocumentProcessor.chunk_text(text)
        doc.chunks = chunks
        for chunk in chunks:
            app_state.document_chunks.append({"tool_id": doc.tool_id, "doc_id": doc.id, "chunk_id": chunk['id'], "text": chunk['text'], "source": doc.original_name, "type": "document"})
        doc.status = "ready"
        app_state.save_to_disk()
    except Exception as e:
        doc = app_state.documents.get(doc_id)
        if doc:
            doc.status = "error"
            doc.error_message = str(e)
            app_state.save_to_disk()


@app.delete("/api/documents/{doc_id}")
async def delete_document(doc_id: str):
    if doc_id not in app_state.documents:
        raise HTTPException(404, "Document not found")
    doc = app_state.documents[doc_id]
    upload_dir = os.environ.get("UPLOAD_PATH", "data/uploads")
    file_path = os.path.join(upload_dir, doc.filename)
    if os.path.exists(file_path):
        os.remove(file_path)
    app_state.document_chunks = [c for c in app_state.document_chunks if c.get('doc_id') != doc_id]
    del app_state.documents[doc_id]
    app_state.save_to_disk()
    return {"status": "success"}


@app.get("/api/documents/{doc_id}/content")
async def get_document_content(doc_id: str):
    """Get document content and metadata"""
    if doc_id not in app_state.documents:
        raise HTTPException(404, "Document not found")
    
    doc = app_state.documents[doc_id]
    
    # Get chunks for this document
    chunks = [c for c in app_state.document_chunks if c.get('doc_id') == doc_id]
    
    # Try to read actual file content
    content = ""
    upload_dir = os.environ.get("UPLOAD_PATH", "data/uploads")
    file_path = os.path.join(upload_dir, doc.filename)
    
    if os.path.exists(file_path):
        try:
            if doc.file_type in ['txt', 'md', 'csv']:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read()
            elif doc.file_type == 'pdf':
                # Try to extract text from PDF
                try:
                    import fitz
                    pdf_doc = fitz.open(file_path)
                    content = ""
                    for page in pdf_doc:
                        content += page.get_text() + "\n"
                    pdf_doc.close()
                except:
                    content = "\n\n".join([c.get('text', '') for c in chunks])
            elif doc.file_type in ['docx', 'doc']:
                try:
                    from docx import Document as DocxDocument
                    docx_doc = DocxDocument(file_path)
                    content = "\n".join([p.text for p in docx_doc.paragraphs])
                except:
                    content = "\n\n".join([c.get('text', '') for c in chunks])
            else:
                content = "\n\n".join([c.get('text', '') for c in chunks])
        except Exception as e:
            content = f"Error reading file: {str(e)}\n\nChunks:\n" + "\n\n".join([c.get('text', '') for c in chunks])
    else:
        content = "\n\n".join([c.get('text', '') for c in chunks])
    
    return {
        "id": doc_id,
        "name": doc.original_name,
        "file_type": doc.file_type,
        "file_size": doc.file_size,
        "content": content,
        "chunks_count": len(chunks),
        "status": doc.status
    }


@app.get("/api/documents/{doc_id}/download")
async def download_document(doc_id: str):
    """Download the original document file"""
    if doc_id not in app_state.documents:
        raise HTTPException(404, "Document not found")
    
    doc = app_state.documents[doc_id]
    upload_dir = os.environ.get("UPLOAD_PATH", "data/uploads")
    file_path = os.path.join(upload_dir, doc.filename)
    
    if not os.path.exists(file_path):
        raise HTTPException(404, "File not found on disk")
    
    return FileResponse(
        file_path, 
        filename=doc.original_name,
        media_type='application/octet-stream'
    )


@app.get("/api/scraped-pages/{page_id}")
async def get_scraped_page(page_id: str):
    """Get scraped page content"""
    if page_id not in app_state.scraped_pages:
        raise HTTPException(404, "Page not found")
    
    page = app_state.scraped_pages[page_id]
    chunks = [c for c in app_state.document_chunks if c.get('page_id') == page_id]
    
    return {
        "id": page_id,
        "url": page.url,
        "title": page.title,
        "content": page.content,
        "chunks": chunks,
        "status": page.status
    }


class KnowledgeSearchRequest(BaseModel):
    query: str
    top_k: int = 5
    use_ai: bool = False  # If True, use LLM to generate answer


@app.post("/api/tools/{tool_id}/ask")
async def ask_knowledge_base(tool_id: str, request: KnowledgeSearchRequest):
    """Ask a question to the knowledge base and get an AI-generated answer"""
    if tool_id not in app_state.tools:
        raise HTTPException(404, "Tool not found")
    
    tool = app_state.tools[tool_id]
    
    # Get all chunks for this tool
    chunks = [c for c in app_state.document_chunks if c.get('tool_id') == tool_id]
    
    if not chunks:
        return {
            "answer": "I don't have any information in this knowledge base yet. Please add some documents or data first.",
            "sources": [],
            "confidence": 0
        }
    
    import re
    import math
    
    # Search for relevant chunks using TF-IDF
    query_lower = request.query.lower().strip()
    stop_words = {'the', 'a', 'an', 'is', 'are', 'was', 'were', 'be', 'been', 'being', 
                  'have', 'has', 'had', 'do', 'does', 'did', 'will', 'would', 'could', 
                  'should', 'may', 'might', 'what', 'which', 'who', 'where', 'when', 'why', 'how',
                  'to', 'of', 'in', 'for', 'on', 'with', 'at', 'by', 'from', 'as',
                  'and', 'but', 'if', 'or', 'this', 'that', 'it', 'its', 'i', 'me', 'my', 'you', 'your'}
    
    query_words = [w for w in re.findall(r'\b\w+\b', query_lower) if w not in stop_words and len(w) > 1]
    if not query_words:
        query_words = re.findall(r'\b\w+\b', query_lower)
    
    # Find relevant chunks
    doc_count = len(chunks)
    idf = {}
    for word in query_words:
        docs_with_word = sum(1 for c in chunks if word in c.get('text', '').lower())
        idf[word] = math.log((doc_count + 1) / (docs_with_word + 1)) + 1
    
    scored_chunks = []
    for chunk in chunks:
        text = chunk.get('text', '').lower()
        text_words = re.findall(r'\b\w+\b', text)
        text_word_count = len(text_words) if text_words else 1
        
        score = 0
        matched_words = []
        
        for word in query_words:
            tf = text.count(word) / text_word_count
            word_score = tf * idf.get(word, 1)
            if word in text:
                matched_words.append(word)
                score += word_score
        
        coverage = len(matched_words) / len(query_words) if query_words else 0
        score *= (1 + coverage)
        
        if score > 0 and matched_words:
            scored_chunks.append({
                "text": chunk.get('text', ''),
                "source": chunk.get('source', 'Unknown'),
                "score": score,
                "coverage": coverage
            })
    
    scored_chunks.sort(key=lambda x: -x['score'])
    relevant_chunks = scored_chunks[:request.top_k]
    
    if not relevant_chunks:
        return {
            "answer": f"I couldn't find any information related to '{request.query}' in this knowledge base.",
            "sources": [],
            "confidence": 0
        }
    
    # Prepare context for LLM
    context = "\n\n".join([
        f"[Source: {c['source']}]\n{c['text']}" 
        for c in relevant_chunks
    ])
    
    # Get tool config for tone
    tone = tool.config.get('tone', 'professional') if tool.config else 'professional'
    tone_instruction = {
        'professional': 'Respond in a professional, clear, and concise manner.',
        'friendly': 'Respond in a warm, friendly, and conversational tone.',
        'formal': 'Respond in a formal and structured manner.',
        'casual': 'Respond in a casual and relaxed tone.'
    }.get(tone, 'Respond in a professional manner.')
    
    # Generate answer using LLM
    try:
        llm_provider = app_state.get_llm_provider()
        
        messages = [
            {
                "role": "system",
                "content": f"""You are a helpful assistant answering questions based on the provided knowledge base.

{tone_instruction}

Rules:
1. ONLY answer based on the information provided in the context below
2. If the answer is not in the context, say you don't have that information
3. Be concise and direct - answer the question without unnecessary details
4. If asked about a specific person/item, only provide information about that specific entity
5. Do NOT list all data - only what's relevant to the question

Context from Knowledge Base:
{context}"""
            },
            {
                "role": "user", 
                "content": request.query
            }
        ]
        
        answer = await llm_provider.generate(messages)
        
        # Calculate confidence based on best match score
        max_score = relevant_chunks[0]['score'] if relevant_chunks else 0
        confidence = min(max_score * 100, 100)
        
        return {
            "answer": answer,
            "sources": [{"source": c['source'], "score": round(c['score'] * 100, 1)} for c in relevant_chunks[:3]],
            "confidence": round(confidence, 1)
        }
        
    except Exception as e:
        # Fallback to simple answer if LLM fails
        best_chunk = relevant_chunks[0] if relevant_chunks else None
        if best_chunk:
            return {
                "answer": f"Based on the knowledge base:\n\n{best_chunk['text'][:500]}...",
                "sources": [{"source": best_chunk['source'], "score": round(best_chunk['score'] * 100, 1)}],
                "confidence": 50,
                "error": str(e)
            }
        return {
            "answer": "Sorry, I couldn't process your question. Please try again.",
            "sources": [],
            "confidence": 0,
            "error": str(e)
        }


@app.post("/api/tools/{tool_id}/search")
async def search_tool_knowledge(tool_id: str, request: KnowledgeSearchRequest):
    """Search within a tool's knowledge base using improved keyword matching"""
    if tool_id not in app_state.tools:
        raise HTTPException(404, "Tool not found")
    
    # Get all chunks for this tool
    chunks = [c for c in app_state.document_chunks if c.get('tool_id') == tool_id]
    
    if not chunks:
        return {"results": [], "message": "No documents indexed for this tool"}
    
    import re
    import math
    
    # Tokenize query
    query_lower = request.query.lower().strip()
    # Remove common stop words for better matching
    stop_words = {'the', 'a', 'an', 'is', 'are', 'was', 'were', 'be', 'been', 'being', 
                  'have', 'has', 'had', 'do', 'does', 'did', 'will', 'would', 'could', 
                  'should', 'may', 'might', 'must', 'shall', 'can', 'need', 'dare',
                  'to', 'of', 'in', 'for', 'on', 'with', 'at', 'by', 'from', 'as',
                  'into', 'through', 'during', 'before', 'after', 'above', 'below',
                  'between', 'under', 'again', 'further', 'then', 'once', 'here',
                  'there', 'when', 'where', 'why', 'how', 'all', 'each', 'few',
                  'more', 'most', 'other', 'some', 'such', 'no', 'nor', 'not',
                  'only', 'own', 'same', 'so', 'than', 'too', 'very', 'just',
                  'and', 'but', 'if', 'or', 'because', 'until', 'while', 'what',
                  'which', 'who', 'whom', 'this', 'that', 'these', 'those', 'am',
                  'i', 'me', 'my', 'myself', 'we', 'our', 'ours', 'ourselves',
                  'you', 'your', 'yours', 'yourself', 'yourselves', 'he', 'him',
                  'his', 'himself', 'she', 'her', 'hers', 'herself', 'it', 'its',
                  'itself', 'they', 'them', 'their', 'theirs', 'themselves'}
    
    # Extract meaningful words from query
    query_words = [w for w in re.findall(r'\b\w+\b', query_lower) if w not in stop_words and len(w) > 1]
    
    if not query_words:
        # If all words were stop words, use original query words
        query_words = re.findall(r'\b\w+\b', query_lower)
    
    if not query_words:
        return {"results": [], "message": "Please enter a valid search query"}
    
    # Calculate IDF (Inverse Document Frequency) for query terms
    doc_count = len(chunks)
    idf = {}
    for word in query_words:
        docs_with_word = sum(1 for c in chunks if word in c.get('text', '').lower())
        idf[word] = math.log((doc_count + 1) / (docs_with_word + 1)) + 1  # Smoothed IDF
    
    scored_chunks = []
    for chunk in chunks:
        text = chunk.get('text', '').lower()
        text_words = re.findall(r'\b\w+\b', text)
        text_word_count = len(text_words) if text_words else 1
        
        # Calculate TF-IDF score
        score = 0
        matched_words = []
        
        for word in query_words:
            # Term Frequency (normalized by document length)
            tf = text.count(word) / text_word_count
            # TF-IDF score
            word_score = tf * idf.get(word, 1)
            if word in text:
                matched_words.append(word)
                score += word_score
        
        # Bonus for phrase match (consecutive words)
        if len(query_words) > 1:
            phrase = ' '.join(query_words)
            if phrase in text:
                score *= 2  # Double score for exact phrase match
        
        # Bonus for matching more unique query terms
        coverage = len(matched_words) / len(query_words) if query_words else 0
        score *= (1 + coverage)  # Boost based on coverage
        
        # Only include results with meaningful matches
        if score > 0 and len(matched_words) > 0:
            # Minimum threshold: at least 30% of query words should match
            if coverage >= 0.3 or len(query_words) == 1:
                scored_chunks.append({
                    "text": chunk.get('text', ''),
                    "source": chunk.get('source', 'Unknown'),
                    "score": min(score, 1.0),  # Cap at 1.0
                    "matched_words": matched_words,
                    "coverage": coverage
                })
    
    # Sort by score and return top_k
    scored_chunks.sort(key=lambda x: (-x['score'], -x['coverage']))
    
    # Return top results
    results = scored_chunks[:request.top_k]
    
    # Clean up results (remove internal fields)
    for r in results:
        r.pop('matched_words', None)
        r.pop('coverage', None)
    
    return {"results": results}


class ApiTestRequest(BaseModel):
    parameters: Dict[str, Any] = {}
    apply_transform: bool = False
    transform_instructions: str = ""


class ExtractParamRequest(BaseModel):
    param_name: str
    file_name: str
    file_content: str
    file_type: str = ""


class GenerateParamRequest(BaseModel):
    param_name: str
    prompt: str = ""


@app.post("/api/tools/{tool_id}/analyze-params")
async def analyze_api_params(tool_id: str):
    """Use LLM to analyze API parameters and suggest best input methods - Pure LLM"""
    if tool_id not in app_state.tools:
        raise HTTPException(404, "Tool not found")
    
    tool = app_state.tools[tool_id]
    api_config = tool.api_config
    
    if not api_config or not api_config.input_parameters:
        return {"parameters": []}
    
    try:
        llm_provider = app_state.get_llm_provider()
        
        # Build description of the API and its parameters
        params_info = []
        for p in api_config.input_parameters:
            params_info.append({
                "name": p.name,
                "type": p.data_type or "string",
                "description": p.description or "",
                "required": p.required
            })
        
        messages = [
            {
                "role": "system",
                "content": """You are an API usability expert. Analyze the API parameters and determine the best input methods for non-technical users.

For EACH parameter, analyze the context and determine:
1. Does it make sense to allow file upload? (e.g., data that could come from documents, receipts, spreadsheets)
2. What file types would be relevant based on the parameter's purpose?
3. What's a contextually appropriate placeholder/hint?
4. Should it be multi-line input based on expected data complexity?
5. Any helpful tips for the user?

Return a JSON array with this structure for each parameter:
{
    "name": "parameter_name",
    "supports_file_upload": true/false,
    "file_upload_label": "Upload [contextual label]" (if supports_file_upload),
    "file_upload_hint": "Contextual hint for what files to upload" (if supports_file_upload),
    "accepted_files": ".pdf,.jpg,.png,.csv,.xlsx" (if supports_file_upload),
    "accepted_file_types": "Human readable file types" (if supports_file_upload),
    "suggested_rows": 2-5 (based on expected data complexity),
    "placeholder": "Contextually appropriate placeholder",
    "ai_prompt_hint": "Hint for AI generation prompt",
    "tip": "Helpful tip for the user" (optional)
}

Return ONLY the JSON array, no explanations or markdown."""
            },
            {
                "role": "user",
                "content": f"""API: {tool.name}
Description: {tool.description or 'No description'}
Endpoint: {api_config.endpoint_path}
Method: {api_config.http_method}

Parameters:
{json.dumps(params_info, indent=2)}

Analyze each parameter based on its context and return the JSON array:"""
            }
        ]
        
        response = await llm_provider.generate(messages)
        
        # Parse JSON response
        response = response.strip()
        if response.startswith('```json'):
            response = response[7:]
        if response.startswith('```'):
            response = response[3:]
        if response.endswith('```'):
            response = response[:-3]
        
        params = json.loads(response.strip())
        return {"parameters": params}
        
    except Exception as e:
        return {"error": str(e), "parameters": []}


@app.post("/api/tools/{tool_id}/extract-param")
async def extract_param_from_file(tool_id: str, request: ExtractParamRequest):
    """Extract parameter data from uploaded file using LLM"""
    if tool_id not in app_state.tools:
        raise HTTPException(404, "Tool not found")
    
    tool = app_state.tools[tool_id]
    api_config = tool.api_config
    
    # Find the parameter definition
    param_def = None
    if api_config and api_config.input_parameters:
        for p in api_config.input_parameters:
            if p.name == request.param_name:
                param_def = p
                break
    
    try:
        llm_provider = app_state.get_llm_provider()
        
        # Build context about what we need to extract
        param_desc = param_def.description if param_def else request.param_name
        param_type = (param_def.data_type if param_def else 'string') or 'string'
        
        messages = [
            {
                "role": "system",
                "content": f"""You are a data extraction assistant. Extract the relevant data from the provided file content.

Parameter to extract: {request.param_name}
Description: {param_desc}
Expected type: {param_type}
File name: {request.file_name}

Rules:
1. Extract ONLY the data relevant to the parameter
2. If it's a list/array, return as JSON array
3. If it's expenses/items, extract each item with relevant fields (description, amount, date, category, etc.)
4. Return ONLY the extracted data, no explanations
5. Format as valid JSON if the parameter expects structured data"""
            },
            {
                "role": "user",
                "content": f"File content:\n\n{request.file_content[:10000]}"  # Limit content size
            }
        ]
        
        extracted = await llm_provider.generate(messages)
        
        # Try to parse as JSON if it looks like JSON
        try:
            if extracted.strip().startswith('[') or extracted.strip().startswith('{'):
                extracted = json.loads(extracted)
        except:
            pass
        
        return {"extracted_data": extracted, "param_name": request.param_name}
        
    except Exception as e:
        return {"error": str(e), "extracted_data": None}


@app.post("/api/tools/{tool_id}/generate-param")
async def generate_param_data(tool_id: str, request: GenerateParamRequest):
    """Generate sample parameter data using LLM - Pure AI-based, no hardcoding"""
    if tool_id not in app_state.tools:
        raise HTTPException(404, "Tool not found")
    
    tool = app_state.tools[tool_id]
    api_config = tool.api_config
    
    # Find the parameter definition
    param_def = None
    if api_config and api_config.input_parameters:
        for p in api_config.input_parameters:
            if p.name == request.param_name:
                param_def = p
                break
    
    # Build context
    param_name = request.param_name
    param_desc = param_def.description if param_def else param_name
    param_type = (param_def.data_type if param_def else 'string') or 'string'
    
    # Get API context
    api_name = tool.name
    api_desc = tool.description or ''
    endpoint = api_config.endpoint_path if api_config else ''
    
    try:
        llm_provider = app_state.get_llm_provider()
        
        messages = [
            {
                "role": "system",
                "content": f"""You are a smart data generator for API testing. Generate realistic sample data based on context.

API Context:
- API Name: {api_name}
- API Description: {api_desc}
- Endpoint: {endpoint}

Parameter to generate:
- Name: {param_name}
- Description: {param_desc}
- Type: {param_type}

User's additional request: {request.prompt or 'Generate appropriate sample data'}

Instructions:
1. Analyze the parameter name and description to understand what data is needed
2. Generate realistic, contextually appropriate data
3. For lists/arrays (like expenses, items, records), generate 2-3 realistic entries with appropriate fields
4. For IDs, generate a realistic format based on context (e.g., EMP001 for employee_id)
5. For dates, use recent dates in YYYY-MM-DD format
6. For amounts/prices, use realistic numbers
7. Return ONLY the data, no explanations or markdown
8. For complex data, return valid JSON"""
            },
            {
                "role": "user",
                "content": f"Generate realistic sample data for the '{param_name}' parameter"
            }
        ]
        
        generated = await llm_provider.generate(messages)
        
        # Clean up response
        generated = generated.strip()
        if generated.startswith('```json'):
            generated = generated[7:]
        if generated.startswith('```'):
            generated = generated[3:]
        if generated.endswith('```'):
            generated = generated[:-3]
        generated = generated.strip()
        
        # Try to parse as JSON
        try:
            if generated.startswith('[') or generated.startswith('{'):
                generated = json.loads(generated)
        except:
            pass
        
        return {"generated_data": generated, "param_name": request.param_name}
        
    except Exception as e:
        return {"error": str(e), "generated_data": None}


@app.post("/api/tools/{tool_id}/auto-fill-test")
async def auto_fill_test_params(tool_id: str):
    """Auto-fill all test parameters with AI-generated sample data - Pure LLM"""
    if tool_id not in app_state.tools:
        raise HTTPException(404, "Tool not found")
    
    tool = app_state.tools[tool_id]
    api_config = tool.api_config
    
    if not api_config or not api_config.input_parameters:
        return {"parameters": {}}
    
    try:
        llm_provider = app_state.get_llm_provider()
        
        # Build description of all parameters
        params_desc = []
        for p in api_config.input_parameters:
            params_desc.append(f"- {p.name} ({p.data_type or 'string'}): {p.description or 'No description'}")
        
        messages = [
            {
                "role": "system",
                "content": f"""You are a smart API test data generator. Generate realistic sample data for ALL parameters.

API Context:
- Name: {tool.name}
- Description: {tool.description or 'No description'}
- Endpoint: {api_config.endpoint_path}
- Method: {api_config.http_method}

Parameters to fill:
{chr(10).join(params_desc)}

Instructions:
1. Analyze each parameter's name and description to understand the context
2. Generate realistic, contextually appropriate data for EACH parameter
3. For IDs, infer the format from context (employee_id → EMP001, order_id → ORD001, etc.)
4. For lists/arrays, generate 2-3 items with appropriate fields based on context
5. For dates, use recent dates in ISO format
6. For amounts/prices, use realistic numbers
7. Return a JSON object with parameter names as keys
8. Return ONLY valid JSON, no explanations or markdown"""
            },
            {
                "role": "user",
                "content": "Generate appropriate test data for all parameters based on the API context"
            }
        ]
        
        generated = await llm_provider.generate(messages)
        
        # Parse JSON response
        generated = generated.strip()
        if generated.startswith('```json'):
            generated = generated[7:]
        if generated.startswith('```'):
            generated = generated[3:]
        if generated.endswith('```'):
            generated = generated[:-3]
        
        params = json.loads(generated.strip())
        return {"parameters": params}
        
    except Exception as e:
        return {"error": str(e), "parameters": {}}


@app.post("/api/tools/{tool_id}/test")
async def test_tool(tool_id: str, request: ApiTestRequest):
    """Test any tool type"""
    if tool_id not in app_state.tools:
        raise HTTPException(404, "Tool not found")
    
    tool = app_state.tools[tool_id]
    
    # Helper function to transform output using LLM
    async def transform_output(original_response: Any, instructions: str) -> Any:
        """Transform API response using LLM based on natural language instructions"""
        try:
            llm_provider = app_state.get_llm_provider()
            
            response_str = json.dumps(original_response, indent=2) if not isinstance(original_response, str) else original_response
            
            messages = [
                {
                    "role": "system",
                    "content": """You are a data transformation assistant. Transform the API response according to the user's instructions.

Rules:
1. Follow the transformation instructions as best as possible
2. If the requested data exists in the response, extract and format it nicely
3. If the requested data does NOT exist in the response, explain what data IS available instead
4. Be helpful - never return empty {} or [] or null
5. For calculations (sum, total, count), perform them if the data exists
6. If user asks for items/list but response has none, show what relevant info IS available
7. Keep the response clean and informative"""
                },
                {
                    "role": "user",
                    "content": f"""API Response:
{response_str}

User's Transformation Request: {instructions}

Transform the response. If the exact data requested isn't available, explain what IS available and show it:"""
                }
            ]
            
            transformed = await llm_provider.generate(messages)
            
            # Clean up response
            transformed = transformed.strip()
            
            # Remove markdown code blocks if present
            if transformed.startswith('```json'):
                transformed = transformed[7:]
            if transformed.startswith('```'):
                transformed = transformed[3:]
            if transformed.endswith('```'):
                transformed = transformed[:-3]
            transformed = transformed.strip()
            
            # Try to parse as JSON if it looks like JSON
            try:
                if transformed.startswith('{') or transformed.startswith('['):
                    parsed = json.loads(transformed)
                    # Don't return empty objects/arrays
                    if parsed == {} or parsed == []:
                        return "No matching data found in the response. The response contains: " + response_str
                    return parsed
            except:
                pass
            
            return transformed
            
        except Exception as e:
            print(f"Transform error: {e}")
            return f"Transformation error: {str(e)}"
    
    # Check if tool has mock response configured (works for any API tool)
    if tool.config and tool.config.get('mock_response'):
        # Get base mock response or generate intelligent one
        base_mock = tool.config.get('mock_response', {})
        
        # Generate intelligent mock response based on input parameters
        async def generate_smart_mock_response(api_tool, input_params: dict, base_response: Any) -> Any:
            """Generate a smart mock response that reflects the input parameters"""
            try:
                llm_provider = app_state.get_llm_provider()
                
                api_config = api_tool.api_config
                base_response_str = json.dumps(base_response, indent=2) if base_response else "{}"
                input_params_str = json.dumps(input_params, indent=2) if input_params else "{}"
                
                messages = [
                    {
                        "role": "system",
                        "content": """You are an API response simulator. Generate a realistic mock response for an API call.

Your task:
1. Analyze the API details and input parameters
2. Generate a realistic response that REFLECTS the input data
3. If the user submitted expenses, include those expenses in the response
4. If the user submitted an employee_id, reference it in the response
5. Include realistic computed values (totals, counts, etc.)
6. Keep the response structure similar to the base mock if provided
7. Return ONLY valid JSON, no explanations or markdown"""
                    },
                    {
                        "role": "user",
                        "content": f"""API Details:
- Name: {api_tool.name}
- Description: {api_tool.description or 'No description'}
- Endpoint: {api_config.endpoint_path if api_config else 'N/A'}
- Method: {api_config.http_method if api_config else 'POST'}

Input Parameters Sent:
{input_params_str}

Base Mock Response (for structure reference):
{base_response_str}

Generate a realistic mock response that includes/reflects the input data. For example:
- If expenses were submitted, include them in the response with a calculated total
- If an employee_id was sent, include it in the response
- Add realistic metadata (timestamps, IDs, status)

Return the mock response as valid JSON:"""
                    }
                ]
                
                response = await llm_provider.generate(messages)
                
                # Clean up response
                response = response.strip()
                if response.startswith('```json'):
                    response = response[7:]
                if response.startswith('```'):
                    response = response[3:]
                if response.endswith('```'):
                    response = response[:-3]
                response = response.strip()
                
                return json.loads(response)
                
            except Exception as e:
                print(f"Smart mock generation failed: {e}")
                # Return base mock with input params merged
                if isinstance(base_response, dict):
                    return {**base_response, "submitted_data": input_params}
                return base_response or {"status": "success", "data": input_params}
        
        # Generate smart mock response
        original_response = await generate_smart_mock_response(tool, request.parameters, base_mock)
        
        # Apply transformation if requested
        if request.apply_transform and request.transform_instructions:
            transformed = await transform_output(original_response, request.transform_instructions)
            if transformed is not None:
                return {
                    "status": "success",
                    "mode": "demo",
                    "tool_type": tool.type,
                    "original_response": original_response,
                    "transformed_response": transformed,
                    "transform_applied": True
                }
        
        return {
            "status": "success",
            "mode": "demo",
            "tool_type": tool.type,
            "response": original_response
        }
    
    # Handle different tool types
    if tool.type == 'api':
        # API tool - make actual call
        api_config = tool.api_config
        if not api_config:
            return {"status": "error", "error": "API not configured"}
        
        try:
            url = api_config.base_url.rstrip('/') + '/' + api_config.endpoint_path.lstrip('/')
            
            for key, value in request.parameters.items():
                url = url.replace('{' + key + '}', str(value))
            
            headers = {}
            if api_config.auth_type == 'bearer' and api_config.auth_config:
                headers['Authorization'] = f"Bearer {api_config.auth_config.get('token', '')}"
            elif api_config.auth_type == 'api_key' and api_config.auth_config:
                key_name = api_config.auth_config.get('key_name', 'X-API-Key')
                headers[key_name] = api_config.auth_config.get('key_value', '')
            
            async with httpx.AsyncClient(timeout=30.0) as client:
                if api_config.http_method == 'GET':
                    response = await client.get(url, headers=headers)
                elif api_config.http_method == 'POST':
                    response = await client.post(url, headers=headers, json=request.parameters)
                elif api_config.http_method == 'PUT':
                    response = await client.put(url, headers=headers, json=request.parameters)
                elif api_config.http_method == 'DELETE':
                    response = await client.delete(url, headers=headers)
                else:
                    response = await client.get(url, headers=headers)
                
                try:
                    original_response = response.json()
                except:
                    original_response = response.text
                
                # Apply transformation if requested
                if request.apply_transform and request.transform_instructions:
                    transformed = await transform_output(original_response, request.transform_instructions)
                    if transformed is not None:
                        return {
                            "status": "success",
                            "mode": "live",
                            "tool_type": "api",
                            "status_code": response.status_code,
                            "original_response": original_response,
                            "transformed_response": transformed,
                            "transform_applied": True
                        }
                
                return {
                    "status": "success",
                    "mode": "live",
                    "tool_type": "api",
                    "status_code": response.status_code,
                    "response": original_response
                }
        except Exception as e:
            return {"status": "error", "tool_type": "api", "error": str(e)}
    
    elif tool.type in ['document', 'knowledge']:
        # Knowledge base - test search
        query = request.parameters.get('query', '')
        if not query:
            return {"status": "error", "error": "Please provide a 'query' parameter"}
        
        chunks = [c for c in app_state.document_chunks if c.get('tool_id') == tool_id]
        if not chunks:
            return {"status": "success", "tool_type": tool.type, "message": "No documents indexed", "results": []}
        
        query_lower = query.lower()
        query_words = query_lower.split()
        
        scored = []
        for chunk in chunks:
            text = chunk.get('text', '').lower()
            score = sum(1 for word in query_words if word in text)
            if score > 0:
                scored.append({"text": chunk.get('text', '')[:500], "source": chunk.get('source', ''), "score": score/len(query_words)})
        
        scored.sort(key=lambda x: x['score'], reverse=True)
        return {"status": "success", "tool_type": tool.type, "results": scored[:5]}
    
    elif tool.type == 'database':
        # Database - test query
        query = request.parameters.get('query', request.parameters.get('sql', ''))
        return {
            "status": "success",
            "mode": "simulated",
            "tool_type": "database",
            "message": "Database query simulated",
            "query": query,
            "sample_response": {
                "rows": [
                    {"id": 1, "name": "Sample Row 1"},
                    {"id": 2, "name": "Sample Row 2"}
                ],
                "row_count": 2
            }
        }
    
    elif tool.type == 'email':
        # Get email parameters
        to_email = request.parameters.get('to', request.parameters.get('recipient', request.parameters.get('email', '')))
        subject = request.parameters.get('subject', 'No Subject')
        body = request.parameters.get('body', request.parameters.get('message', request.parameters.get('content', '')))
        html = request.parameters.get('html', False)
        
        if not to_email:
            return {
                "status": "error",
                "tool_type": "email",
                "message": "No recipient email provided. Please specify 'to' parameter."
            }
        
        if not body:
            return {
                "status": "error",
                "tool_type": "email", 
                "message": "No email body provided. Please specify 'body' or 'message' parameter."
            }
        
        # Get provider config
        config = tool.config or {}
        provider = config.get('provider', 'smtp')
        
        try:
            if provider == 'google':
                # Gmail via OAuth
                access_token = config.get('access_token')
                refresh_token = config.get('refresh_token')
                
                if not access_token and not refresh_token:
                    return {
                        "status": "error",
                        "tool_type": "email",
                        "message": "Gmail not configured. Please reconnect your Gmail account."
                    }
                
                # Try to refresh token if we have refresh_token
                if refresh_token:
                    new_access_token = await EmailService.refresh_google_token(refresh_token)
                    if new_access_token:
                        access_token = new_access_token
                        # Update tool config with new access token
                        tool.config['access_token'] = new_access_token
                
                if not access_token:
                    return {
                        "status": "error",
                        "tool_type": "email",
                        "message": "Could not get access token. Please reconnect your Gmail account."
                    }
                
                # Send via Gmail API
                result = await EmailService.send_gmail(
                    access_token=access_token,
                    to=to_email,
                    subject=subject,
                    body=body,
                    html=html
                )
                
                if result.get('success'):
                    return {
                        "status": "success",
                        "tool_type": "email",
                        "message": f"Email sent successfully to {to_email}",
                        "message_id": result.get('message_id'),
                        "from": config.get('email', 'Gmail')
                    }
                else:
                    return {
                        "status": "error",
                        "tool_type": "email",
                        "message": f"Failed to send email: {result.get('error')}"
                    }
            
            elif provider == 'sendgrid':
                # SendGrid API
                api_key = config.get('apiKey', config.get('api_key'))
                from_email = config.get('fromEmail', config.get('from_email', config.get('email')))
                
                if not api_key:
                    return {
                        "status": "error",
                        "tool_type": "email",
                        "message": "SendGrid API key not configured."
                    }
                
                result = await EmailService.send_sendgrid(
                    api_key=api_key,
                    to=to_email,
                    subject=subject,
                    body=body,
                    from_email=from_email,
                    html=html
                )
                
                if result.get('success'):
                    return {
                        "status": "success",
                        "tool_type": "email",
                        "message": f"Email sent successfully to {to_email} via SendGrid",
                        "from": from_email
                    }
                else:
                    return {
                        "status": "error",
                        "tool_type": "email",
                        "message": f"Failed to send email: {result.get('error')}"
                    }
            
            elif provider in ['smtp', 'outlook', 'microsoft']:
                # SMTP
                host = config.get('smtp_host', config.get('host', 'smtp.gmail.com'))
                port = int(config.get('smtp_port', config.get('port', 587)))
                username = config.get('smtp_user', config.get('username', config.get('email', '')))
                password = config.get('smtp_pass', config.get('password', config.get('app_password', '')))
                from_email = config.get('from_email', username)
                
                if provider == 'outlook' or provider == 'microsoft':
                    host = 'smtp.office365.com'
                    port = 587
                
                result = await EmailService.send_smtp(
                    host=host,
                    port=port,
                    username=username,
                    password=password,
                    to=to_email,
                    subject=subject,
                    body=body,
                    from_email=from_email,
                    html=html
                )
                
                if result.get('success'):
                    return {
                        "status": "success",
                        "tool_type": "email",
                        "message": f"Email sent successfully to {to_email}",
                        "from": from_email
                    }
                else:
                    return {
                        "status": "error",
                        "tool_type": "email",
                        "message": f"Failed to send email: {result.get('error')}"
                    }
            
            else:
                return {
                    "status": "error",
                    "tool_type": "email",
                    "message": f"Unknown email provider: {provider}"
                }
                
        except Exception as e:
            return {
                "status": "error",
                "tool_type": "email",
                "message": f"Error sending email: {str(e)}"
            }
    
    elif tool.type == 'slack':
        return {
            "status": "success",
            "mode": "simulated",
            "tool_type": "slack",
            "message": "Slack connection test successful",
            "config_status": "verified" if tool.config else "not_configured"
        }
    
    elif tool.type == 'calendar':
        return {
            "status": "success",
            "mode": "simulated",
            "tool_type": "calendar",
            "message": "Calendar connection test successful",
            "sample_events": [
                {"title": "Team Meeting", "date": "2024-12-28", "time": "10:00"},
                {"title": "Project Review", "date": "2024-12-29", "time": "14:00"}
            ]
        }
    
    elif tool.type == 'websearch':
        query = request.parameters.get('query', 'test')
        return {
            "status": "success",
            "mode": "simulated",
            "tool_type": "websearch",
            "query": query,
            "sample_results": [
                {"title": f"Search result for: {query}", "url": "https://example.com/1", "snippet": "Sample search result..."},
                {"title": f"Another result for: {query}", "url": "https://example.com/2", "snippet": "More content..."}
            ]
        }
    
    else:
        # Generic tool - just echo config and params
        return {
            "status": "success",
            "mode": "echo",
            "tool_type": tool.type,
            "tool_name": tool.name,
            "config": tool.config,
            "test_input": request.parameters,
            "message": f"Tool '{tool.type}' test completed. Configure mock_response in config for custom test responses."
        }


# ============================================================================
# EMAIL SENDING ENDPOINT
# ============================================================================

class SendEmailRequest(BaseModel):
    to: str
    subject: str
    body: str
    html: bool = False


@app.post("/api/tools/{tool_id}/send-email")
async def send_email_via_tool(tool_id: str, request: SendEmailRequest):
    """
    Send an email using a configured email tool.
    
    This endpoint allows sending actual emails via Gmail OAuth, SendGrid, or SMTP.
    """
    print(f"\n{'='*50}")
    print(f"📧 SEND EMAIL REQUEST")
    print(f"{'='*50}")
    print(f"   Tool ID: {tool_id}")
    print(f"   To: {request.to}")
    print(f"   Subject: {request.subject}")
    
    if tool_id not in app_state.tools:
        print(f"   ❌ Tool not found!")
        raise HTTPException(404, "Email tool not found")
    
    tool = app_state.tools[tool_id]
    print(f"   Tool Name: {tool.name}")
    print(f"   Tool Type: {tool.type}")
    
    if tool.type != 'email':
        print(f"   ❌ Wrong tool type!")
        raise HTTPException(400, f"Tool is not an email tool (type: {tool.type})")
    
    config = tool.config or {}
    provider = config.get('provider', '')
    print(f"   Provider: {provider}")
    print(f"   Config keys: {list(config.keys())}")
    
    if not provider:
        print(f"   ❌ No provider configured!")
        raise HTTPException(400, "Email provider not configured. Please reconnect your email account.")
    
    try:
        if provider == 'google':
            # Gmail via OAuth
            access_token = config.get('access_token')
            refresh_token = config.get('refresh_token')
            
            print(f"   Has access_token: {bool(access_token)}")
            print(f"   Has refresh_token: {bool(refresh_token)}")
            
            if not access_token and not refresh_token:
                print(f"   ❌ No tokens available!")
                raise HTTPException(400, "Gmail not configured. Please reconnect your Gmail account.")
            
            # Try to refresh token if we have refresh_token
            if refresh_token:
                new_access_token = await EmailService.refresh_google_token(refresh_token)
                if new_access_token:
                    access_token = new_access_token
                    # Update tool config with new access token
                    tool.config['access_token'] = new_access_token
                    app_state.save_to_disk()
            
            if not access_token:
                raise HTTPException(400, "Could not get access token. Please reconnect your Gmail account.")
            
            # Send via Gmail API
            result = await EmailService.send_gmail(
                access_token=access_token,
                to=request.to,
                subject=request.subject,
                body=request.body,
                html=request.html
            )
            
            if result.get('success'):
                return {
                    "status": "success",
                    "message": f"Email sent successfully to {request.to}",
                    "message_id": result.get('message_id'),
                    "from": config.get('email', 'Gmail')
                }
            else:
                raise HTTPException(500, f"Failed to send email: {result.get('error')}")
        
        elif provider == 'sendgrid':
            # SendGrid API
            api_key = config.get('apiKey', config.get('api_key'))
            from_email = config.get('fromEmail', config.get('from_email', config.get('email')))
            
            if not api_key:
                raise HTTPException(400, "SendGrid API key not configured.")
            
            result = await EmailService.send_sendgrid(
                api_key=api_key,
                to=request.to,
                subject=request.subject,
                body=request.body,
                from_email=from_email,
                html=request.html
            )
            
            if result.get('success'):
                return {
                    "status": "success",
                    "message": f"Email sent successfully to {request.to} via SendGrid",
                    "from": from_email
                }
            else:
                raise HTTPException(500, f"Failed to send email: {result.get('error')}")
        
        elif provider in ['smtp', 'outlook', 'microsoft']:
            # SMTP
            host = config.get('smtp_host', config.get('host', 'smtp.gmail.com'))
            port = int(config.get('smtp_port', config.get('port', 587)))
            username = config.get('smtp_user', config.get('username', config.get('email', '')))
            password = config.get('smtp_pass', config.get('password', config.get('app_password', '')))
            from_email = config.get('from_email', username)
            
            if provider == 'outlook' or provider == 'microsoft':
                host = 'smtp.office365.com'
                port = 587
            
            result = await EmailService.send_smtp(
                host=host,
                port=port,
                username=username,
                password=password,
                to=request.to,
                subject=request.subject,
                body=request.body,
                from_email=from_email,
                html=request.html
            )
            
            if result.get('success'):
                return {
                    "status": "success",
                    "message": f"Email sent successfully to {request.to}",
                    "from": from_email
                }
            else:
                raise HTTPException(500, f"Failed to send email: {result.get('error')}")
        
        else:
            raise HTTPException(400, f"Unknown email provider: {provider}")
            
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, f"Error sending email: {str(e)}")


class UpdateToolRequest(BaseModel):
    name: Optional[str] = None
    description: Optional[str] = None
    api_config: Optional[Dict[str, Any]] = None
    config: Optional[Dict[str, Any]] = None
    is_active: Optional[bool] = None
    # Access Control
    access_type: Optional[str] = None
    allowed_user_ids: Optional[List[str]] = None
    allowed_group_ids: Optional[List[str]] = None
    can_edit_user_ids: Optional[List[str]] = None
    can_delete_user_ids: Optional[List[str]] = None
    can_execute_user_ids: Optional[List[str]] = None




@app.put("/api/tools/{tool_id}")
async def update_tool(tool_id: str, request: UpdateToolRequest, current_user: User = Depends(get_current_user_optional)):
    """Update tool configuration with automatic re-processing"""
    if tool_id not in app_state.tools:
        raise HTTPException(404, "Tool not found")
    
    tool = app_state.tools[tool_id]
    user_id = str(current_user.id) if current_user else None
    # Get user's groups dynamically
    user_group_ids = get_user_group_ids(user_id) if user_id else []
    
    # Check if user can edit this tool
    if not check_tool_access(tool, user_id, user_group_ids, 'edit'):
        raise HTTPException(403, "You don't have permission to edit this tool")
    
    old_config = tool.config.copy() if tool.config else {}
    
    # Update basic fields
    if request.name is not None:
        tool.name = request.name
    if request.description is not None:
        tool.description = request.description
    if request.is_active is not None:
        tool.is_active = request.is_active
    if request.api_config is not None:
        tool.api_config = APIConfig(**request.api_config)
    
    # Update access control fields (ONLY owner can change these)
    is_owner = tool.owner_id == user_id
    
    if is_owner:
        if request.access_type is not None:
            tool.access_type = request.access_type
        if request.allowed_user_ids is not None:
            tool.allowed_user_ids = request.allowed_user_ids
        if request.allowed_group_ids is not None:
            tool.allowed_group_ids = request.allowed_group_ids
        if request.can_edit_user_ids is not None:
            tool.can_edit_user_ids = request.can_edit_user_ids
        if request.can_delete_user_ids is not None:
            tool.can_delete_user_ids = request.can_delete_user_ids
        if request.can_execute_user_ids is not None:
            tool.can_execute_user_ids = request.can_execute_user_ids
    
    reprocess_result = None
    reprocess_action = None
    
    if request.config is not None:
        new_config = request.config
        
        # Update config first
        if tool.config:
            tool.config.update(new_config)
        else:
            tool.config = new_config
        
        # Auto re-process based on tool type and what changed
        try:
            if tool.type == 'website':
                # Re-scrape if URL changed
                old_url = old_config.get('url', '')
                new_url = new_config.get('url')
                if new_url is not None and old_url != new_url:
                    reprocess_action = 'rescrape'
                    # Clear old scraped pages for this tool
                    pages_to_delete = [pid for pid, p in app_state.scraped_pages.items() if p.tool_id == tool_id]
                    for pid in pages_to_delete:
                        app_state.document_chunks = [c for c in app_state.document_chunks if c.get('page_id') != pid]
                        del app_state.scraped_pages[pid]
                    # Scrape new URL
                    scraper = WebsiteScraper(new_url, new_config.get('max_pages', 10))
                    pages = await scraper.scrape(new_config.get('recursive', False))
                    saved_count = 0
                    for page_data in pages:
                        page = ScrapedPage(tool_id=tool_id, url=page_data['url'], title=page_data['title'], content=page_data['content'], status="processing")
                        chunks = DocumentProcessor.chunk_text(page_data['content'], 
                            chunk_size=new_config.get('chunk_size', 1000),
                            overlap=new_config.get('overlap', 200))
                        page.chunks = chunks
                        for chunk in chunks:
                            app_state.document_chunks.append({"tool_id": tool_id, "page_id": page.id, "chunk_id": chunk['id'], "text": chunk['text'], "source": page_data['title'] or page_data['url'], "type": "website"})
                        page.status = "ready"
                        app_state.scraped_pages[page.id] = page
                        saved_count += 1
                    reprocess_result = {"pages_scraped": saved_count}
                    
            elif tool.type in ['document', 'knowledge']:
                # Re-index if chunk settings changed
                old_chunk = old_config.get('chunk_size', 1000)
                old_overlap = old_config.get('overlap', 200)
                new_chunk = new_config.get('chunk_size')
                new_overlap = new_config.get('overlap')
                if (new_chunk is not None and old_chunk != new_chunk) or (new_overlap is not None and old_overlap != new_overlap):
                    reprocess_action = 'reindex'
                    # Re-chunk all documents for this tool
                    docs_reindexed = 0
                    for doc_id, doc in app_state.documents.items():
                        if doc.tool_id == tool_id and doc.content:
                            # Remove old chunks
                            app_state.document_chunks = [c for c in app_state.document_chunks if c.get('doc_id') != doc_id]
                            # Create new chunks
                            chunks = DocumentProcessor.chunk_text(doc.content, chunk_size=new_chunk, overlap=new_overlap)
                            for chunk in chunks:
                                app_state.document_chunks.append({"tool_id": tool_id, "doc_id": doc_id, "chunk_id": chunk['id'], "text": chunk['text'], "source": doc.filename, "type": "document"})
                            docs_reindexed += 1
                    reprocess_result = {"documents_reindexed": docs_reindexed}
                    
            elif tool.type == 'database':
                # Test connection if connection string changed
                old_conn = old_config.get('connection_string', '')
                new_conn = new_config.get('connection_string')
                if new_conn is not None and old_conn != new_conn:
                    reprocess_action = 'test_connection'
                    reprocess_result = {"status": "connection_updated", "message": "Database connection configured"}
                    
            elif tool.type == 'api':
                # Test if API config changed
                old_url = old_config.get('base_url', '')
                new_url = new_config.get('base_url')
                if new_url is not None and old_url != new_url:
                    reprocess_action = 'test_connection'
                    reprocess_result = {"status": "api_updated", "message": "API endpoint configured"}
                    
            elif tool.type == 'email':
                # Test email config
                old_provider = old_config.get('provider', '')
                new_provider = new_config.get('provider', old_provider)
                old_key = old_config.get('api_key', '')
                new_key = new_config.get('api_key', old_key)
                if old_provider != new_provider or old_key != new_key:
                    reprocess_action = 'test_connection'
                    reprocess_result = {"status": "email_updated", "message": "Email configuration updated"}
                    
            elif tool.type == 'webhook':
                # Test webhook URL
                old_url = old_config.get('url', '')
                new_url = new_config.get('url')
                if new_url is not None and old_url != new_url:
                    reprocess_action = 'test_connection'
                    reprocess_result = {"status": "webhook_updated", "message": "Webhook URL configured"}
                    
            elif tool.type == 'slack':
                # Test slack token
                old_token = old_config.get('bot_token', '')
                new_token = new_config.get('bot_token', old_token)
                if old_token != new_token and new_token:
                    reprocess_action = 'test_connection'
                    reprocess_result = {"status": "slack_updated", "message": "Slack bot token configured"}
                    
            elif tool.type == 'websearch':
                # Update search config
                old_key = old_config.get('api_key', '')
                new_key = new_config.get('api_key', old_key)
                if old_key != new_key and new_key:
                    reprocess_action = 'test_connection'
                    reprocess_result = {"status": "websearch_updated", "message": "Web search API key configured"}
                    
        except Exception as e:
            reprocess_result = {"error": str(e)}
    
    app_state.save_to_disk()
    
    # Also save to database
    try:
        from database.services import ToolService
        tool_dict = tool.dict()
        # Include access control fields explicitly
        tool_dict['access_type'] = tool.access_type
        tool_dict['allowed_user_ids'] = tool.allowed_user_ids or []
        tool_dict['allowed_group_ids'] = tool.allowed_group_ids or []
        tool_dict['can_edit_user_ids'] = tool.can_edit_user_ids or []
        tool_dict['can_delete_user_ids'] = tool.can_delete_user_ids or []
        tool_dict['can_execute_user_ids'] = tool.can_execute_user_ids or []
        tool_dict['owner_id'] = tool.owner_id
        
        ToolService.update_tool(tool_id, tool_dict, "org_default", user_id or "system")
        print(f"✅ [DATABASE] Tool '{tool.name}' updated with access control: allowed_users={tool.allowed_user_ids}, allowed_groups={tool.allowed_group_ids}")
    except Exception as e:
        print(f"⚠️ [DATABASE] Failed to update tool in database: {e}")
    
    response = {
        "status": "success", 
        "tool": tool.dict()
    }
    
    if reprocess_action:
        response["reprocess_action"] = reprocess_action
        response["reprocess_result"] = reprocess_result
        
    return response

@app.post("/api/tools/{tool_id}/scrape")
async def scrape_website(tool_id: str, request: ScrapeRequest):
    if tool_id not in app_state.tools:
        raise HTTPException(404, "Tool not found")
    tool = app_state.tools[tool_id]
    if tool.type not in ['website', 'knowledge']:
        raise HTTPException(400, "Tool type does not support web scraping")
    scraper = WebsiteScraper(request.url, request.max_pages)
    pages = await scraper.scrape(request.recursive)
    saved_pages = []
    for page_data in pages:
        page = ScrapedPage(tool_id=tool_id, url=page_data['url'], title=page_data['title'], content=page_data['content'], status="processing")
        chunks = DocumentProcessor.chunk_text(page_data['content'])
        page.chunks = chunks
        for chunk in chunks:
            app_state.document_chunks.append({"tool_id": tool_id, "page_id": page.id, "chunk_id": chunk['id'], "text": chunk['text'], "source": page_data['title'] or page_data['url'], "type": "website"})
        page.status = "ready"
        app_state.scraped_pages[page.id] = page
        saved_pages.append(page.dict())
    app_state.save_to_disk()
    return {"status": "success", "pages_scraped": len(saved_pages), "pages": saved_pages}


@app.delete("/api/scraped-pages/{page_id}")
async def delete_scraped_page(page_id: str):
    if page_id not in app_state.scraped_pages:
        raise HTTPException(404, "Page not found")
    app_state.document_chunks = [c for c in app_state.document_chunks if c.get('page_id') != page_id]
    del app_state.scraped_pages[page_id]
    app_state.save_to_disk()
    return {"status": "success"}


@app.get("/api/scraped-pages/{page_id}/data")
async def get_scraped_page_data(page_id: str):
    """Get detailed scraped page content including tables and chunks."""
    if page_id not in app_state.scraped_pages:
        raise HTTPException(404, "Page not found")
    
    page = app_state.scraped_pages[page_id]
    
    # Get chunks for this page
    chunks = [c for c in app_state.document_chunks if c.get('page_id') == page_id]
    chunk_texts = [c.get('text', c.get('content', '')) for c in chunks]
    
    # Extract tables from content
    tables = []
    content = page.content or ''
    
    # Find [TABLE]...[/TABLE] blocks
    import re
    table_matches = re.findall(r'\[TABLE\](.*?)\[\/TABLE\]', content, re.DOTALL | re.IGNORECASE)
    tables.extend(table_matches)
    
    return {
        "id": page.id,
        "url": page.url,
        "title": page.title,
        "content": page.content,
        "chunks": chunk_texts,
        "tables": tables,
        "status": page.status,
        "created_at": page.created_at,
        "content_length": len(page.content) if page.content else 0,
        "chunk_count": len(chunks)
    }


# ============================================================================
# AGENT MEMORY ENDPOINTS
# ============================================================================

@app.get("/api/agents/{agent_id}/memory")
async def get_agent_memory(agent_id: str):
    """Get agent memory (conversation summaries)"""
    if agent_id not in app_state.agents:
        raise HTTPException(404, "Agent not found")
    agent = app_state.agents[agent_id]
    return {
        "memory_enabled": agent.memory_enabled,
        "memories": agent.memory,
        "count": len(agent.memory)
    }


@app.put("/api/agents/{agent_id}/memory/toggle")
async def toggle_agent_memory(agent_id: str, enabled: bool = True):
    """Enable or disable memory for an agent"""
    if agent_id not in app_state.agents:
        raise HTTPException(404, "Agent not found")
    agent = app_state.agents[agent_id]
    agent.memory_enabled = enabled
    agent.updated_at = datetime.utcnow().isoformat()
    app_state.save_to_disk()
    return {"status": "success", "memory_enabled": agent.memory_enabled}


@app.delete("/api/agents/{agent_id}/memory")
async def clear_agent_memory(agent_id: str):
    """Clear all memory for an agent"""
    if agent_id not in app_state.agents:
        raise HTTPException(404, "Agent not found")
    agent = app_state.agents[agent_id]
    agent.memory = []
    agent.updated_at = datetime.utcnow().isoformat()
    app_state.save_to_disk()
    return {"status": "success", "message": "Memory cleared"}


# Chat Endpoints
@app.get("/api/agents/{agent_id}/conversations")
async def list_conversations(agent_id: str, current_user: User = Depends(get_current_user)):
    """
    List conversations for a specific agent.
    PRIVACY: Each user can ONLY see their OWN conversations!
    Even the agent owner cannot see other users' conversations.
    """
    user_id = str(current_user.id) if current_user else None
    
    if not user_id:
        return {"conversations": []}  # No user = no conversations
    
    # Try database first
    try:
        from database.services import ConversationService
        convs = ConversationService.get_all_conversations()
        # Filter by agent_id AND user_id (PRIVACY!)
        agent_convs = [
            c for c in convs 
            if c.get('agent_id') == agent_id and c.get('user_id') == user_id
        ]
        sorted_convs = sorted(agent_convs, key=lambda x: x.get("updated_at") or x.get("created_at") or "", reverse=True)
        
        # Log all conversation titles from database
        print(f"📋 [LIST CONVERSATIONS] Agent: {agent_id[:8]}..., User: {user_id[:8]}...")
        print(f"   📊 Total conversations found: {len(sorted_convs)}")
        print(f"   📚 All titles in database:")
        for i, conv in enumerate(sorted_convs):
            print(f"      {i+1}. '{conv.get('title', 'No title')}' (id: {conv.get('id', 'N/A')[:8]}...)")
        
        return {"conversations": sorted_convs}
    except Exception as e:
        print(f"⚠️  Database error, falling back to memory: {e}")
    
    # Fallback to memory - also filter by user_id
    if agent_id not in app_state.agents:
        raise HTTPException(404, "Agent not found")
    convs = [
        {"id": c.id, "title": c.title, "messages_count": len(c.messages), "created_at": c.created_at, "updated_at": c.updated_at} 
        for c in app_state.conversations.values() 
        if c.agent_id == agent_id and getattr(c, 'user_id', None) == user_id
    ]
    return {"conversations": sorted(convs, key=lambda x: x["updated_at"], reverse=True)}


@app.get("/api/conversations/{conversation_id}")
async def get_conversation(conversation_id: str, current_user: User = Depends(get_current_user)):
    """
    Get a specific conversation.
    PRIVACY: User can only access their OWN conversations!
    """
    user_id = str(current_user.id) if current_user else None
    
    # Try database first
    try:
        from database.services import ConversationService
        conv = ConversationService.get_conversation_by_id(conversation_id)
        if conv:
            # PRIVACY CHECK: Verify user owns this conversation
            if conv.get('user_id') and conv.get('user_id') != user_id:
                raise HTTPException(403, "You don't have access to this conversation")
            return conv
    except HTTPException:
        raise
    except Exception as e:
        print(f"⚠️  Database error, falling back to memory: {e}")
    
    # Fallback to memory
    if conversation_id not in app_state.conversations:
        raise HTTPException(404, "Conversation not found")
    
    conv = app_state.conversations[conversation_id]
    # PRIVACY CHECK for memory-based conversations
    conv_user_id = getattr(conv, 'user_id', None) or (conv.access_cache.user_id if hasattr(conv, 'access_cache') and conv.access_cache else None)
    if conv_user_id and conv_user_id != user_id:
        raise HTTPException(403, "You don't have access to this conversation")
    
    return conv.dict()


@app.delete("/api/conversations/{conversation_id}")
async def delete_conversation(conversation_id: str, current_user: User = Depends(get_current_user)):
    """
    Delete a conversation.
    PRIVACY: User can only delete their OWN conversations!
    """
    user_id = str(current_user.id) if current_user else None
    print(f"🗑️  [DELETE] Deleting conversation: {conversation_id} by user: {user_id[:8] if user_id else 'None'}...")
    
    # Try database first
    try:
        from database.services import ConversationService
        # First, verify ownership
        conv = ConversationService.get_conversation_by_id(conversation_id)
        if conv and conv.get('user_id') and conv.get('user_id') != user_id:
            raise HTTPException(403, "You can only delete your own conversations")
        
        result = ConversationService.delete_conversation(conversation_id)
        print(f"🗑️  [DELETE] Database result: {result}")
        if result:
            # Also remove from memory if exists
            if conversation_id in app_state.conversations:
                del app_state.conversations[conversation_id]
            return {"status": "success"}
    except HTTPException:
        raise
    except Exception as e:
        print(f"⚠️  Database error, falling back to memory: {e}")
    
    # Fallback to memory - check ownership
    if conversation_id not in app_state.conversations:
        raise HTTPException(404, "Conversation not found")
    
    conv = app_state.conversations[conversation_id]
    conv_user_id = getattr(conv, 'user_id', None) or (conv.access_cache.user_id if hasattr(conv, 'access_cache') and conv.access_cache else None)
    if conv_user_id and conv_user_id != user_id:
        raise HTTPException(403, "You can only delete your own conversations")
    
    del app_state.conversations[conversation_id]
    app_state.save_to_disk()
    return {"status": "success"}


# ============================================================================
# TEST CHAT - Agent testing environment
# ============================================================================

@app.post("/api/agents/{agent_id}/test-chat")
async def test_chat(agent_id: str, request: ChatRequest, current_user: User = Depends(get_current_user)):
    """
    Test Chat endpoint - for testing agents before deployment.
    Uses the same system as production but in a test environment.
    Requires: owner OR delegated admin with at least one edit permission.
    """
    if not current_user:
        raise HTTPException(401, "Authentication required")
    
    user_id = str(current_user.id)
    org_id = current_user.org_id if current_user else "org_default"
    user_group_ids = get_user_group_ids(user_id) if user_id else []
    
    # Check if user has permission to test this agent
    is_owner = False
    has_edit_permission = False
    
    # Try to get agent from in-memory state first
    agent = app_state.agents.get(agent_id)
    
    # If not found, try to load from database
    if not agent:
        try:
            from database.services import AgentService
            from core.agents.models import AgentData, AgentPersonality, AgentGuardrails, TaskDefinition
            
            # Get current user for org_id
            org_id = current_user.org_id if current_user else "org_default"
            
            agent_dict = AgentService.get_agent_by_id(agent_id, org_id)
            if agent_dict:
                # Convert to AgentData
                agent = AgentData(
                    id=agent_dict['id'],
                    name=agent_dict['name'],
                    icon=agent_dict.get('icon', '🤖'),
                    goal=agent_dict.get('goal', ''),
                    description=agent_dict.get('description', ''),
                    model_id=agent_dict.get('model_id', 'gpt-4o'),
                    personality=AgentPersonality(**agent_dict.get('personality', {})) if agent_dict.get('personality') else AgentPersonality(),
                    guardrails=AgentGuardrails(**agent_dict.get('guardrails', {})) if agent_dict.get('guardrails') else AgentGuardrails(),
                    tasks=[TaskDefinition(**t) for t in agent_dict.get('tasks', [])],
                    tool_ids=agent_dict.get('tool_ids', []),
                    memory=agent_dict.get('memory', []),
                    memory_enabled=agent_dict.get('memory_enabled', True),
                    status=agent_dict.get('status', 'draft'),
                    is_active=agent_dict.get('is_active', False),
                    created_at=agent_dict.get('created_at', datetime.utcnow().isoformat()),
                    updated_at=agent_dict.get('updated_at', datetime.utcnow().isoformat()),
                )
                # Add to in-memory state
                app_state.agents[agent.id] = agent
        except Exception as e:
            print(f"⚠️  Error loading agent from database: {e}")
            import traceback
            traceback.print_exc()
    
    if not agent:
        raise HTTPException(404, "Agent not found")
    
    # Check permissions - owner or delegated admin with edit permissions can test
    try:
        from database.services import AgentService
        agent_dict = AgentService.get_agent_by_id(agent_id, org_id)
        if agent_dict:
            owner_id = str(agent_dict.get('owner_id')) if agent_dict.get('owner_id') else None
            created_by = str(agent_dict.get('created_by')) if agent_dict.get('created_by') else None
            is_owner = (user_id == owner_id) or (user_id == created_by)
        
        if not is_owner and ACCESS_CONTROL_AVAILABLE and AccessControlService:
            from api.modules.access_control.service import get_session, normalize_org_id
            from database.models.agent_access import AgentAccessPolicy
            import json
            
            norm_org_id = normalize_org_id(org_id)
            
            with get_session() as session:
                admin_policy = session.query(AgentAccessPolicy).filter(
                    AgentAccessPolicy.agent_id == agent_id,
                    AgentAccessPolicy.org_id == norm_org_id,
                    AgentAccessPolicy.access_type == 'agent_admin',
                    AgentAccessPolicy.is_active == True
                ).first()
                
                if admin_policy:
                    user_is_admin = user_id in (admin_policy.user_ids or [])
                    group_is_admin = any(g in (admin_policy.group_ids or []) for g in user_group_ids)
                    
                    if user_is_admin or group_is_admin:
                        # Get their permissions
                        if admin_policy.description:
                            try:
                                admin_config = json.loads(admin_policy.description)
                                permissions = []
                                
                                if user_id in admin_config:
                                    entity_config = admin_config[user_id]
                                    if isinstance(entity_config, dict):
                                        permissions = entity_config.get('permissions', [])
                                    elif isinstance(entity_config, list):
                                        permissions = entity_config
                                else:
                                    for group_id in user_group_ids:
                                        if group_id in admin_config:
                                            entity_config = admin_config[group_id]
                                            if isinstance(entity_config, dict):
                                                permissions = entity_config.get('permissions', [])
                                            elif isinstance(entity_config, list):
                                                permissions = entity_config
                                            break
                                
                                # Check if user has any edit permission (not just publish)
                                edit_perms = ['full_admin', 'edit_basic_info', 'edit_personality', 
                                             'edit_model', 'edit_guardrails', 'manage_tasks', 
                                             'manage_tools', 'manage_knowledge']
                                has_edit_permission = any(p in permissions for p in edit_perms)
                                
                            except Exception as e:
                                print(f"⚠️  Error parsing test-chat permissions: {e}")
                                has_edit_permission = True  # Legacy - allow
                        else:
                            has_edit_permission = True  # No description = legacy full_admin
                    else:
                        has_edit_permission = False  # Not an admin
                else:
                    has_edit_permission = False  # No admin policy
        else:
            has_edit_permission = is_owner  # If no access control, owner can test
            
    except Exception as e:
        print(f"⚠️  Error checking test-chat permissions: {e}")
        has_edit_permission = True  # Fallback to allow
    
    if not is_owner and not has_edit_permission:
        raise HTTPException(403, "You need edit permissions to test this agent. Publish-only admins cannot test.")
    
    # Get or create conversation (prefix with demo_ to separate from real chats)
    
    if request.conversation_id and request.conversation_id in app_state.conversations:
        conversation = app_state.conversations[request.conversation_id]
    else:
        title = f"[TEST] {request.message[:40]}..." if len(request.message) > 40 else f"[TEST] {request.message}"
        conversation = Conversation(agent_id=agent_id, user_id=user_id, title=title)
        app_state.conversations[conversation.id] = conversation
        
        # Save to database (marked as TEST conversation - won't show in production chat)
        try:
            from database.services import ConversationService
            ConversationService.create_conversation({
                'id': conversation.id,
                'agent_id': agent_id,
                'user_id': user_id,
                'title': title,
                'is_test': True
            }, org_id, user_id)
        except Exception as e:
            print(f"⚠️  Failed to save conversation to DB: {e}")
    
    # Add user message
    user_msg = ConversationMessage(role="user", content=request.message)
    conversation.messages.append(user_msg)
    
    # Save user message to database
    try:
        from database.services import ConversationService
        ConversationService.add_message(conversation.id, {
            'id': user_msg.id,
            'role': 'user',
            'content': request.message
        }, org_id, user_id)
    except Exception as e:
        print(f"⚠️  Failed to save message to DB: {e}")
    
    # Process with demo agent (dynamic mockup tools)
    result = await process_test_agent_chat(agent, request.message, conversation, timezone=request.timezone)
    print(f"📤 Test chat result: content={len(result.get('content', ''))} chars")
    if not result.get("content"):
        print(f"⚠️  Empty response from LLM! Full result: {result}")
    
    # Add assistant response
    assistant_msg = ConversationMessage(
        role="assistant", 
        content=result["content"], 
        sources=result.get("sources", [])
    )
    conversation.messages.append(assistant_msg)
    conversation.updated_at = datetime.utcnow().isoformat()
    
    # Save assistant message to database
    try:
        from database.services import ConversationService
        ConversationService.add_message(conversation.id, {
            'id': assistant_msg.id,
            'role': 'assistant',
            'content': result["content"],
            'sources': result.get("sources", [])
        }, org_id, user_id)
    except Exception as e:
        print(f"⚠️  Failed to save assistant message to DB: {e}")
    
    app_state.save_to_disk()
    
    return {
        "response": result["content"],
        "conversation_id": conversation.id,
        "sources": result.get("sources", []),
        "formatted": True,
        "demo_mode": True,
        "tool_calls": result.get("tool_calls", [])
    }


@app.post("/api/agents/{agent_id}/test-chat-with-files")
async def test_chat_with_files(
    agent_id: str,
    message: str = Form(""),
    conversation_id: Optional[str] = Form(None),
    timezone: Optional[str] = Form(None),
    files: List[UploadFile] = File(default=[]),
    current_user: User = Depends(get_current_user)
):
    """
    Chat with file attachments - OCR and file processing!
    """
    # Try to get agent from in-memory state first
    agent = app_state.agents.get(agent_id)
    
    # If not found, try to load from database
    if not agent:
        try:
            from database.services import AgentService
            from core.agents.models import AgentData, AgentPersonality, AgentGuardrails, TaskDefinition
            
            # Get current user for org_id
            org_id = current_user.org_id if current_user else "org_default"
            
            agent_dict = AgentService.get_agent_by_id(agent_id, org_id)
            if agent_dict:
                # Convert to AgentData
                agent = AgentData(
                    id=agent_dict['id'],
                    name=agent_dict['name'],
                    icon=agent_dict.get('icon', '🤖'),
                    goal=agent_dict.get('goal', ''),
                    description=agent_dict.get('description', ''),
                    model_id=agent_dict.get('model_id', 'gpt-4o'),
                    personality=AgentPersonality(**agent_dict.get('personality', {})) if agent_dict.get('personality') else AgentPersonality(),
                    guardrails=AgentGuardrails(**agent_dict.get('guardrails', {})) if agent_dict.get('guardrails') else AgentGuardrails(),
                    tasks=[TaskDefinition(**t) for t in agent_dict.get('tasks', [])],
                    tool_ids=agent_dict.get('tool_ids', []),
                    memory=agent_dict.get('memory', []),
                    memory_enabled=agent_dict.get('memory_enabled', True),
                    status=agent_dict.get('status', 'draft'),
                    is_active=agent_dict.get('is_active', False),
                    created_at=agent_dict.get('created_at', datetime.utcnow().isoformat()),
                    updated_at=agent_dict.get('updated_at', datetime.utcnow().isoformat()),
                )
                # Add to in-memory state
                app_state.agents[agent.id] = agent
        except Exception as e:
            print(f"⚠️  Error loading agent from database: {e}")
            import traceback
            traceback.print_exc()
    
    if not agent:
        raise HTTPException(404, "Agent not found")
    
    # Get or create conversation
    org_id = current_user.org_id if current_user else "org_default"
    user_id = str(current_user.id) if current_user else "system"
    
    is_new_conversation = False
    if conversation_id and conversation_id in app_state.conversations:
        conversation = app_state.conversations[conversation_id]
    else:
        is_new_conversation = True
        # Start with temporary title - LLM will update it
        title = "[TEST] New conversation"
        conversation = Conversation(agent_id=agent_id, user_id=user_id, title=title)
        app_state.conversations[conversation.id] = conversation
        
        # Save to database with temporary title
        try:
            from database.services import ConversationService
            ConversationService.create_conversation({
                'id': conversation.id,
                'agent_id': agent_id,
                'user_id': user_id,
                'title': title,
                'is_test': True
            }, org_id, user_id)
        except Exception as e:
            print(f"⚠️  Failed to save conversation to DB: {e}")
        
        # Schedule LLM to generate smart title in background
        if message:
            from api.modules.conversations import ConversationTitleService
            
            async def update_test_title():
                new_title = await ConversationTitleService.generate_title(
                    first_message=message,
                    agent_name=agent.name if agent else None,
                    model_id=agent.model_id if agent else None
                )
                from database.services import ConversationService
                ConversationService.update_title(conversation.id, f"[TEST] {new_title}")
            
            asyncio.create_task(update_test_title())
    
    # Save uploaded files and prepare attachments info
    attachments = []
    upload_dir = os.environ.get("UPLOAD_PATH", "data/uploads")
    os.makedirs(upload_dir, exist_ok=True)
    
    for file in files:
        # Save file
        file_id = f"{uuid.uuid4().hex[:8]}_{file.filename}"
        file_path = os.path.join(upload_dir, file_id)
        
        content = await file.read()
        with open(file_path, "wb") as f:
            f.write(content)
        
        attachments.append({
            "name": file.filename,
            "path": file_path,
            "type": file.content_type or "",
            "size": len(content)
        })
    
    # Build user message with attachment info
    display_message = message
    if attachments:
        attachment_names = [a["name"] for a in attachments]
        display_message += f"\n\n[Attached: {', '.join(attachment_names)}]"
    
    # Add user message
    user_msg = ConversationMessage(role="user", content=display_message)
    conversation.messages.append(user_msg)
    
    # Save user message to database
    try:
        from database.services import ConversationService
        ConversationService.add_message(conversation.id, {
            'id': user_msg.id,
            'role': 'user',
            'content': display_message
        }, org_id, user_id)
    except Exception as e:
        print(f"⚠️  Failed to save message to DB: {e}")
    
    # Process with demo agent - includes REAL OCR/file processing
    result = await process_test_agent_chat(agent, message, conversation, attachments, timezone=timezone)
    
    # Add assistant response
    assistant_msg = ConversationMessage(
        role="assistant",
        content=result["content"],
        sources=result.get("sources", [])
    )
    conversation.messages.append(assistant_msg)
    conversation.updated_at = datetime.utcnow().isoformat()
    
    # Save assistant message to database
    try:
        from database.services import ConversationService
        ConversationService.add_message(conversation.id, {
            'id': assistant_msg.id,
            'role': 'assistant',
            'content': result["content"],
            'sources': result.get("sources", [])
        }, org_id, user_id)
    except Exception as e:
        print(f"⚠️  Failed to save assistant message to DB: {e}")
    
    app_state.save_to_disk()
    
    return {
        "response": result["content"],
        "conversation_id": conversation.id,
        "sources": result.get("sources", []),
        "formatted": True,
        "demo_mode": True,
        "tool_calls": result.get("tool_calls", []),
        "real_extractions": result.get("real_extractions", 0),
        "files_processed": len(attachments)
    }


# ============================================================================
# CHAT SESSION INITIALIZATION - Load permissions ONCE when chat starts
# ============================================================================

class StartChatRequest(BaseModel):
    """Request to start a new chat session"""
    pass


class ChatSessionResponse(BaseModel):
    """Response with session info and pre-loaded permissions"""
    conversation_id: str
    agent_id: str
    agent_name: str
    user_name: str
    user_role: str
    user_groups: List[str]
    accessible_tasks: List[str]
    denied_tasks: List[str]
    has_full_access: bool
    message: str


@app.post("/api/agents/{agent_id}/start-chat")
async def start_chat_session(agent_id: str, current_user: User = Depends(get_current_user)):
    """
    Start a new chat session and load ALL permissions ONCE.
    
    This is the BEST PRACTICE approach:
    1. Load user info (name, role, groups)
    2. Check agent access
    3. Get list of allowed/denied tasks
    4. Cache everything in the conversation
    5. Return session info to frontend
    
    Frontend can use this to:
    - Show personalized greeting
    - Display available features
    - Know what to expect before user types
    """
    # Always log start-chat for debugging (critical path)
    print(f"🚀 [START-CHAT] Request received: agent_id={agent_id}, current_user={current_user.email if current_user else 'NONE'}")
    
    if not current_user:
        print(f"❌ [START-CHAT] No current_user - authentication missing!")
        raise HTTPException(401, "Authentication required")
    
    if agent_id not in app_state.agents:
        print(f"❌ [START-CHAT] Agent not found: {agent_id}")
        print(f"   Available agents: {list(app_state.agents.keys())[:5]}...")
        raise HTTPException(404, "Agent not found")
    
    agent = app_state.agents[agent_id]
    org_id = current_user.org_id if current_user else "org_default"
    user_id = str(current_user.id) if current_user else "system"
    
    # ========================================================================
    # LOAD USER CONTEXT
    # ========================================================================
    print(f"🚀 [START-CHAT] Starting chat session for user: {current_user.email if current_user else 'ANONYMOUS'}")
    print(f"   [START-CHAT] current_user attributes: first_name={getattr(current_user, 'first_name', None)}, last_name={getattr(current_user, 'last_name', None)}, display_name={getattr(current_user, 'display_name', None)}")
    
    user_name = ""
    if hasattr(current_user, 'first_name') and current_user.first_name:
        user_name = current_user.first_name
        if hasattr(current_user, 'last_name') and current_user.last_name:
            user_name += f" {current_user.last_name}"
    elif hasattr(current_user, 'display_name') and current_user.display_name:
        user_name = current_user.display_name
    elif hasattr(current_user, 'email') and current_user.email:
        user_name = current_user.email.split('@')[0].replace('.', ' ').title()
    
    print(f"   [START-CHAT] Final user_name: '{user_name}'")
    
    # Get user's groups
    user_groups = []
    if hasattr(current_user, 'group_ids') and current_user.group_ids:
        try:
            if SECURITY_AVAILABLE and security_state:
                for gid in current_user.group_ids:
                    group = security_state.groups.get(gid)
                    if group:
                        user_groups.append(group.name)
        except:
            pass
    
    # Get user's role
    user_role = ""
    if hasattr(current_user, 'role_ids') and current_user.role_ids:
        try:
            if SECURITY_AVAILABLE and security_state:
                for rid in current_user.role_ids:
                    role = security_state.roles.get(rid)
                    if role:
                        user_role = role.name
                        break
        except:
            pass
    
    # ========================================================================
    # LOAD ACCESS PERMISSIONS (ONCE!)
    # ========================================================================
    accessible_task_names = [task.name for task in agent.tasks]
    denied_task_names = []
    accessible_tool_ids = [tool_id for tool_id in agent.tool_ids]
    denied_tool_ids = []
    has_full_access = True
    
    # Check if user is the OWNER - owners always have full access
    is_owner = False
    agent_owner_id = None
    
    # Try to get owner from in-memory agent
    if hasattr(agent, 'owner_id') and agent.owner_id:
        agent_owner_id = str(agent.owner_id)
    elif hasattr(agent, 'created_by') and agent.created_by:
        agent_owner_id = str(agent.created_by)
    
    # If not in memory, try to get from database
    if not agent_owner_id:
        try:
            from database.services import AgentService
            db_agent = AgentService.get_agent_by_id(agent_id, org_id)
            if db_agent:
                agent_owner_id = str(db_agent.get('owner_id') or db_agent.get('created_by') or '')
                print(f"📋 [START-CHAT] Got owner_id from DB: {agent_owner_id[:8] if agent_owner_id else 'None'}...")
        except Exception as e:
            print(f"⚠️ [START-CHAT] Failed to get owner from DB: {e}")
    
    if agent_owner_id:
        is_owner = agent_owner_id == user_id
        print(f"🔍 [START-CHAT] Owner check: agent_owner={agent_owner_id[:8]}..., user={user_id[:8]}..., is_owner={is_owner}")
    
    # Check if user is a DELEGATED ADMIN - they also have full access
    is_delegated_admin = False
    if not is_owner and ACCESS_CONTROL_AVAILABLE and AccessControlService and current_user:
        try:
            user_role_ids = getattr(current_user, 'role_ids', []) or []
            user_group_ids = get_user_group_ids(user_id) if user_id else []
            
            perm_result = AccessControlService.check_agent_permission(
                user_id=user_id,
                user_role_ids=user_role_ids,
                user_group_ids=user_group_ids,
                agent_id=agent_id,
                org_id=org_id,
                permission='any_admin'
            )
            is_delegated_admin = perm_result.get('has_permission', False)
        except Exception as e:
            print(f"⚠️ [START-CHAT] Delegated admin check failed: {e}")
    
    if is_owner:
        print(f"👑 [START-CHAT] User {user_id[:8]}... is the OWNER - full access granted")
        has_full_access = True
    elif is_delegated_admin:
        print(f"🔑 [START-CHAT] User {user_id[:8]}... is a DELEGATED ADMIN")
        # Check if admin has task restrictions
        try:
            user_group_ids = get_user_group_ids(user_id) if user_id else []
            admin_restrictions = AccessControlService.get_admin_chat_restrictions(
                user_id=user_id,
                user_group_ids=user_group_ids,
                agent_id=agent_id,
                org_id=org_id
            )
            admin_denied_tasks = admin_restrictions.get('denied_task_names', [])
            if admin_denied_tasks:
                denied_task_names = admin_denied_tasks
                accessible_task_names = [t.name for t in agent.tasks if t.name not in denied_task_names]
                has_full_access = False
                print(f"   [START-CHAT] Admin has task restrictions: {denied_task_names}")
            else:
                has_full_access = True
                print(f"   [START-CHAT] Admin has NO task restrictions - full access")
        except Exception as e:
            print(f"⚠️ [START-CHAT] Failed to get admin restrictions: {e}")
            has_full_access = True
    elif ACCESS_CONTROL_AVAILABLE and AccessControlService and current_user:
        try:
            user_role_ids = getattr(current_user, 'role_ids', []) or []
            user_group_ids = get_user_group_ids(user_id) if user_id else []
            
            access_result = AccessControlService.check_user_access(
                user_id=user_id,
                user_role_ids=user_role_ids,
                user_group_ids=user_group_ids,
                agent_id=agent_id,
                org_id=org_id
            )
            
            if not access_result.has_access:
                raise HTTPException(
                    status_code=403,
                    detail={
                        "error": "access_denied",
                        "message": "This assistant is not available for your account.",
                        "user_friendly": True
                    }
                )
            
            # Get denied tasks (by name)
            denied_task_names = access_result.denied_tasks or []
            accessible_task_names = [t.name for t in agent.tasks if t.name not in denied_task_names]
            
            # Get denied tools
            denied_tool_ids = access_result.denied_tools or []
            accessible_tool_ids = [tid for tid in agent.tool_ids if tid not in denied_tool_ids]
            
            has_full_access = len(denied_task_names) == 0 and len(denied_tool_ids) == 0
            
            print(f"🔐 [START-CHAT] Loaded permissions for {user_name}:")
            print(f"   Accessible tasks: {accessible_task_names}")
            print(f"   Denied tasks: {denied_task_names}")
            
        except HTTPException:
            raise
        except Exception as e:
            print(f"⚠️ Access control check failed: {e}")
    
    # ========================================================================
    # CREATE CONVERSATION WITH CACHED PERMISSIONS
    # ========================================================================
    access_cache = ConversationAccessCache(
        user_id=user_id,
        user_name=user_name,
        user_role=user_role,
        user_groups=user_groups,
        accessible_task_names=accessible_task_names,
        denied_task_names=denied_task_names,
        accessible_tool_ids=accessible_tool_ids,
        denied_tool_ids=denied_tool_ids
    )
    
    conversation = Conversation(
        agent_id=agent_id,
        user_id=user_id,  # PRIVACY: Track who owns this conversation
        title=f"Chat with {agent.name}",
        access_cache=access_cache
    )
    app_state.conversations[conversation.id] = conversation
    
    # Save to database
    try:
        from database.services import ConversationService
        ConversationService.create_conversation({
            'id': conversation.id,
            'agent_id': agent_id,
            'user_id': user_id,  # PRIVACY: Track ownership
            'title': conversation.title
        }, org_id, user_id)
    except Exception as e:
        print(f"⚠️ Failed to save conversation to DB: {e}")
    
    # Build personalized welcome message
    welcome = f"Hi {user_name}! " if user_name else "Hello! "
    welcome += f"I'm {agent.name}. "
    
    if has_full_access:
        welcome += "How can I help you today?"
    else:
        welcome += f"I can help you with: {', '.join(accessible_task_names[:3])}..."
    
    print(f"✅ [START-CHAT] Returning session:")
    print(f"   user_name: '{user_name}'")
    print(f"   user_role: '{user_role}'")
    print(f"   user_groups: {user_groups}")
    print(f"   accessible_tasks: {accessible_task_names}")
    print(f"   denied_tasks: {denied_task_names}")
    print(f"   welcome: '{welcome}'")
    
    return ChatSessionResponse(
        conversation_id=conversation.id,
        agent_id=agent_id,
        agent_name=agent.name,
        user_name=user_name,
        user_role=user_role,
        user_groups=user_groups,
        accessible_tasks=accessible_task_names,
        denied_tasks=denied_task_names,
        has_full_access=has_full_access,
        message=welcome
    )


@app.post("/api/agents/{agent_id}/chat")
async def chat(agent_id: str, request: ChatRequest, current_user: User = Depends(get_current_user)):
    if agent_id not in app_state.agents:
        raise HTTPException(404, "Agent not found")
    agent = app_state.agents[agent_id]
    
    org_id = current_user.org_id if current_user else "org_default"
    user_id = str(current_user.id) if current_user else "system"
    
    # Debug: Log user info
    print(f"💬 [CHAT] Agent: {agent.name}, User: {current_user.email if current_user else 'ANONYMOUS'}, user_id: {user_id}")
    print(f"   [CHAT] current_user is None? {current_user is None}, ACCESS_CONTROL_AVAILABLE: {ACCESS_CONTROL_AVAILABLE}")
    
    # ========================================================================
    # ACCESS CONTROL - Use cached permissions if available (BEST PRACTICE)
    # ========================================================================
    access_result = None
    use_cached = False
    
    # Check if we have an existing conversation with cached permissions
    if request.conversation_id and request.conversation_id in app_state.conversations:
        conversation = app_state.conversations[request.conversation_id]
        
        # Use cached permissions if available and from same user
        if conversation.access_cache and conversation.access_cache.user_id == user_id:
            use_cached = True
            print(f"🔐 Using CACHED permissions for user {user_id[:8]}...")
            print(f"   Cached denied tasks: {conversation.access_cache.denied_task_names}")
            
            # Create access_result from cache
            from api.modules.access_control.schemas import AccessCheckResult
            access_result = AccessCheckResult(
                has_access=True,
                denied_tasks=conversation.access_cache.denied_task_names,
                denied_tools=conversation.access_cache.denied_tool_ids
            )
    else:
        conversation = None
    
    # Check if user is the OWNER - owners always have full access
    is_owner = False
    agent_owner_id = None
    
    # Try to get owner from in-memory agent
    if hasattr(agent, 'owner_id') and agent.owner_id:
        agent_owner_id = str(agent.owner_id)
    elif hasattr(agent, 'created_by') and agent.created_by:
        agent_owner_id = str(agent.created_by)
    
    # If not in memory, try to get from database
    if not agent_owner_id:
        try:
            from database.services import AgentService
            db_agent = AgentService.get_agent_by_id(agent_id, org_id)
            if db_agent:
                agent_owner_id = str(db_agent.get('owner_id') or db_agent.get('created_by') or '')
        except Exception as e:
            print(f"⚠️ [CHAT] Failed to get owner from DB: {e}")
    
    if agent_owner_id:
        is_owner = agent_owner_id == user_id
    
    # Check if user is a DELEGATED ADMIN - they also have full access
    is_delegated_admin = False
    if not is_owner and ACCESS_CONTROL_AVAILABLE and AccessControlService and current_user:
        try:
            user_role_ids = getattr(current_user, 'role_ids', []) or []
            user_group_ids = get_user_group_ids(user_id) if user_id else []
            
            perm_result = AccessControlService.check_agent_permission(
                user_id=user_id,
                user_role_ids=user_role_ids,
                user_group_ids=user_group_ids,
                agent_id=agent_id,
                org_id=org_id,
                permission='any_admin'
            )
            is_delegated_admin = perm_result.get('has_permission', False)
        except Exception as e:
            print(f"⚠️ [CHAT] Delegated admin check failed: {e}")
    
    # If no cache, check permissions (first message or new conversation)
    print(f"   [CHAT] Checking access control: use_cached={use_cached}, is_owner={is_owner}, is_delegated_admin={is_delegated_admin}")
    
    if is_owner:
        print(f"👑 [CHAT] User {user_id[:8]}... is the OWNER - full access granted")
        access_result = None  # Owner has full access, no restrictions
    elif is_delegated_admin:
        print(f"🔑 [CHAT] User {user_id[:8]}... is a DELEGATED ADMIN")
        # Check if admin has task restrictions
        try:
            user_group_ids = get_user_group_ids(user_id) if user_id else []
            admin_restrictions = AccessControlService.get_admin_chat_restrictions(
                user_id=user_id,
                user_group_ids=user_group_ids,
                agent_id=agent_id,
                org_id=org_id
            )
            admin_denied_tasks = admin_restrictions.get('denied_task_names', [])
            if admin_denied_tasks:
                from api.modules.access_control.schemas import AccessCheckResult
                access_result = AccessCheckResult(
                    has_access=True,
                    denied_tasks=admin_denied_tasks,
                    denied_tools=[]
                )
                print(f"   [CHAT] Admin has task restrictions: {admin_denied_tasks}")
            else:
                access_result = None  # No restrictions
                print(f"   [CHAT] Admin has NO task restrictions - full access")
        except Exception as e:
            print(f"⚠️ [CHAT] Failed to get admin restrictions: {e}")
            access_result = None
    elif not use_cached and ACCESS_CONTROL_AVAILABLE and AccessControlService and current_user:
        try:
            # Get user's roles and groups
            user_role_ids = getattr(current_user, 'role_ids', []) or []
            user_group_ids = get_user_group_ids(user_id) if user_id else []
            print(f"   [CHAT] User roles: {user_role_ids}, groups: {user_group_ids}")
            
            # Check access
            access_result = AccessControlService.check_user_access(
                user_id=user_id,
                user_role_ids=user_role_ids,
                user_group_ids=user_group_ids,
                agent_id=agent_id,
                org_id=org_id
            )
            
            print(f"🔐 FRESH Access Control Check for user {user_id} on agent {agent_id}:")
            print(f"   Has Access: {access_result.has_access}")
            print(f"   Denied Tasks: {access_result.denied_tasks}")
            print(f"   Denied Tools: {access_result.denied_tools}")
            
            # Level 1: Check if user can access the agent at all
            if not access_result.has_access:
                raise HTTPException(
                    status_code=403,
                    detail={
                        "error": "access_denied",
                        "message": access_result.reason or "This assistant is not available for your account. Please contact your administrator if you need access.",
                        "user_friendly": True
                    }
                )
        except HTTPException:
            raise
        except Exception as e:
            print(f"⚠️ Access control check failed (allowing by default): {e}")
            access_result = None
    
    # Create new conversation if needed
    is_new_conversation = False
    if not conversation:
        is_new_conversation = True
        # Start with temporary title - LLM will update it
        title = "New conversation"
        conversation = Conversation(agent_id=agent_id, user_id=user_id, title=title)
        
        # Cache permissions in new conversation
        if access_result:
            # Build user context for cache
            user_name = ""
            if hasattr(current_user, 'first_name') and current_user.first_name:
                user_name = f"{current_user.first_name} {getattr(current_user, 'last_name', '') or ''}".strip()
            elif hasattr(current_user, 'email'):
                user_name = current_user.email.split('@')[0].replace('.', ' ').title()
            
            conversation.access_cache = ConversationAccessCache(
                user_id=user_id,
                user_name=user_name,
                denied_task_names=access_result.denied_tasks or [],
                denied_tool_ids=access_result.denied_tools or []
            )
            print(f"📦 Cached permissions for future messages")
        
        app_state.conversations[conversation.id] = conversation
        
        # Save to database with temporary title
        try:
            from database.services import ConversationService
            ConversationService.create_conversation({
                'id': conversation.id,
                'agent_id': agent_id,
                'user_id': user_id,  # PRIVACY: Track ownership
                'title': title
            }, org_id, user_id)
        except Exception as e:
            print(f"⚠️  Failed to save conversation to DB: {e}")
        
        # Schedule LLM to generate smart title in background
        from api.modules.conversations import ConversationTitleService
        asyncio.create_task(
            ConversationTitleService.generate_and_update_title(
                conversation_id=conversation.id,
                first_message=request.message,
                agent_name=agent.name,
                model_id=agent.model_id
            )
        )
    
    user_msg = ConversationMessage(role="user", content=request.message)
    conversation.messages.append(user_msg)
    
    # Save user message to database
    try:
        from database.services import ConversationService
        ConversationService.add_message(conversation.id, {
            'id': user_msg.id,
            'role': 'user',
            'content': request.message
        }, org_id, user_id)
    except Exception as e:
        print(f"⚠️  Failed to save message to DB: {e}")
    
    # Pass access control result and user to process_agent_chat for task/tool filtering
    result = await process_agent_chat(
        agent, 
        request.message, 
        conversation, 
        timezone=request.timezone,
        access_control=access_result,  # Pass the access control result
        current_user=current_user  # Pass user for personalization
    )
    print(f"📤 Chat result: content={len(result.get('content', ''))} chars, sources={len(result.get('sources', []))}")
    if not result.get("content"):
        print(f"⚠️  Empty response from LLM!")
    assistant_msg = ConversationMessage(role="assistant", content=result["content"], sources=result["sources"])
    conversation.messages.append(assistant_msg)
    conversation.updated_at = datetime.utcnow().isoformat()
    
    # Save assistant message to database
    try:
        from database.services import ConversationService
        ConversationService.add_message(conversation.id, {
            'id': assistant_msg.id,
            'role': 'assistant',
            'content': result["content"],
            'sources': result.get("sources", [])
        }, org_id, user_id)
    except Exception as e:
        print(f"⚠️  Failed to save assistant message to DB: {e}")
    
    app_state.save_to_disk()
    
    # Update agent memory in background (every 5 messages)
    if len(conversation.messages) % 5 == 0:
        asyncio.create_task(update_agent_memory(agent, conversation))
    
    return ChatResponse(response=result["content"], conversation_id=conversation.id, sources=result["sources"], formatted=True)


# ============================================================================
# STREAMING CHAT - Real-time thinking & reasoning display (like ChatGPT/Cursor)
# ============================================================================

class StreamingChatRequest(BaseModel):
    message: str
    conversation_id: Optional[str] = None
    timezone: Optional[str] = None


@app.post("/api/agents/{agent_id}/chat/stream")
async def chat_stream(agent_id: str, request: StreamingChatRequest, current_user: User = Depends(get_current_user)):
    """
    Streaming chat endpoint with real-time thinking/reasoning display.
    
    Uses Server-Sent Events (SSE) to stream:
    - thinking: Agent is processing/analyzing
    - tool_call: Agent is calling a tool
    - tool_result: Tool returned a result
    - content: Response text chunks
    - sources: Knowledge base sources used
    - done: Streaming complete
    - error: An error occurred
    """
    if agent_id not in app_state.agents:
        raise HTTPException(404, "Agent not found")
    
    agent = app_state.agents[agent_id]
    org_id = current_user.org_id if current_user else "org_default"
    user_id = str(current_user.id) if current_user else "system"
    
    async def event_generator():
        """Generate SSE events for streaming response"""
        try:
            # Detect user's language from message or agent settings
            user_lang = 'en'  # default
            g = agent.guardrails
            if hasattr(g, 'language'):
                lang_setting = g.language
                if lang_setting == 'ar':
                    user_lang = 'ar'
                elif lang_setting == 'user':
                    # Detect from message - simple Arabic detection
                    if any('\u0600' <= c <= '\u06FF' for c in request.message):
                        user_lang = 'ar'
            else:
                # Auto-detect from message
                if any('\u0600' <= c <= '\u06FF' for c in request.message):
                    user_lang = 'ar'
            
            # Thinking messages - some static, main one generated by LLM
            is_ar = user_lang == 'ar'
            msgs = {
                'access_denied': 'عذراً، ليس لديك صلاحية للوصول لهذا المساعد.' if is_ar else 'Sorry, you do not have access to this assistant.',
                'action_failed': 'تعذر إكمال الإجراء' if is_ar else 'Could not complete action',
                'searching': 'جاري البحث في المصادر...' if is_ar else 'Searching sources...',
                'found_sources': 'وجدت {} مصادر مفيدة' if is_ar else 'Found {} helpful sources',
                'preparing': 'جاري تحضير الرد...' if is_ar else 'Preparing response...',
                'processing': 'جاري معالجة النتائج...' if is_ar else 'Processing results...',
                'writing': 'جاري كتابة الرد...' if is_ar else 'Writing response...',
                'sending_email': 'جاري إرسال البريد...' if is_ar else 'Sending email...',
                'connecting': 'جاري الاتصال...' if is_ar else 'Connecting...',
                'triggering': 'جاري التنفيذ...' if is_ar else 'Executing...',
                'slack': 'جاري إرسال الرسالة...' if is_ar else 'Sending message...',
                'websearch': 'جاري البحث في الإنترنت...' if is_ar else 'Searching the web...',
                'database': 'جاري البحث في البيانات...' if is_ar else 'Looking up data...',
                'action': 'جاري التنفيذ...' if is_ar else 'Taking action...',
                'email_sent': 'تم الإرسال' if is_ar else 'Sent',
                'slack_sent': 'تم الإرسال' if is_ar else 'Sent',
                'search_done': 'تم العثور على النتائج' if is_ar else 'Found results',
                'data_done': 'تم' if is_ar else 'Done',
                'service_done': 'تم' if is_ar else 'Done',
                'action_done': 'تم' if is_ar else 'Done',
            }
            
            # Will send thinking after LLM generates it
            await asyncio.sleep(0.01)
            
            # ========================================================================
            # ACCESS CONTROL CHECK
            # ========================================================================
            access_result = None
            is_owner = False
            agent_owner_id = None
            
            # Check ownership
            if hasattr(agent, 'owner_id') and agent.owner_id:
                agent_owner_id = str(agent.owner_id)
            elif hasattr(agent, 'created_by') and agent.created_by:
                agent_owner_id = str(agent.created_by)
            
            if not agent_owner_id:
                try:
                    from database.services import AgentService
                    db_agent = AgentService.get_agent_by_id(agent_id, org_id)
                    if db_agent:
                        agent_owner_id = str(db_agent.get('owner_id') or db_agent.get('created_by') or '')
                except Exception:
                    pass
            
            if agent_owner_id:
                is_owner = agent_owner_id == user_id
            
            # Check permissions if not owner
            if not is_owner and ACCESS_CONTROL_AVAILABLE and AccessControlService and current_user:
                try:
                    user_role_ids = getattr(current_user, 'role_ids', []) or []
                    user_group_ids = get_user_group_ids(user_id) if user_id else []
                    
                    access_result = AccessControlService.check_user_access(
                        user_id=user_id,
                        user_role_ids=user_role_ids,
                        user_group_ids=user_group_ids,
                        agent_id=agent_id,
                        org_id=org_id
                    )
                    
                    if not access_result.has_access:
                        yield f"data: {json.dumps({'type': 'error', 'content': msgs['access_denied']})}\n\n"
                        yield f"data: {json.dumps({'type': 'done'})}\n\n"
                        return
                except Exception as e:
                    print(f"⚠️ Access control check failed: {e}")
            
            # ========================================================================
            # CONVERSATION HANDLING
            # ========================================================================
            
            is_new_conversation = False
            is_first_message = False
            print(f"🔍 [STREAM] Checking conversation: request.conversation_id={request.conversation_id}")
            if request.conversation_id and request.conversation_id in app_state.conversations:
                conversation = app_state.conversations[request.conversation_id]
                print(f"🔍 [STREAM] Found existing conversation in memory: {conversation.id[:8]}")
                # Check if this is the first message (no messages yet)
                if len(conversation.messages) == 0:
                    is_first_message = True
                    print(f"🔍 [STREAM] This is the FIRST MESSAGE in conversation - will generate title")
            else:
                is_new_conversation = True
                is_first_message = True
                print(f"🔍 [STREAM] Creating NEW conversation (is_new={is_new_conversation})")
                # Start with temporary title - LLM will update it
                title = "New conversation"
                
                conversation = Conversation(agent_id=agent_id, user_id=user_id, title=title)
                app_state.conversations[conversation.id] = conversation
                
                # Save to database with temporary title
                try:
                    from database.services import ConversationService
                    ConversationService.create_conversation({
                        'id': conversation.id,
                        'agent_id': agent_id,
                        'user_id': user_id,
                        'title': title
                    }, org_id, user_id)
                except Exception as e:
                    print(f"⚠️ Failed to save conversation to DB: {e}")
                
            # Schedule LLM to generate smart title for FIRST message
            # (whether new conversation or existing conversation with no messages yet)
            if is_first_message:
                from api.modules.conversations import ConversationTitleService
                print(f"🏷️ [STREAM] Scheduling title generation for conversation {conversation.id[:8]}...")
                try:
                    task = asyncio.create_task(
                        ConversationTitleService.generate_and_update_title(
                            conversation_id=conversation.id,
                            first_message=request.message,
                            agent_name=agent.name,
                            model_id=agent.model_id
                        )
                    )
                    print(f"🏷️ [STREAM] Task created: {task}")
                except Exception as title_err:
                    print(f"❌ [STREAM] Failed to create title task: {title_err}")
            
            # Send conversation ID
            yield f"data: {json.dumps({'type': 'conversation_id', 'content': conversation.id})}\n\n"
            
            # Add user message
            user_msg = ConversationMessage(role="user", content=request.message)
            conversation.messages.append(user_msg)
            
            # Save user message
            try:
                from database.services import ConversationService
                ConversationService.add_message(conversation.id, {
                    'id': user_msg.id,
                    'role': 'user',
                    'content': request.message
                }, org_id, user_id)
            except Exception:
                pass
            
            # ========================================================================
            # KNOWLEDGE BASE SEARCH
            # ========================================================================
            agent_tools = get_agent_tools(agent)
            
            # Filter tools if access control is active
            if access_result and hasattr(access_result, 'denied_tools') and access_result.denied_tools:
                agent_tools = [t for t in agent_tools if t.id not in access_result.denied_tools]
            
            tool_ids = [t.id for t in agent_tools]
            
            search_results = search_documents(request.message, tool_ids, top_k=5)
            context = ""
            sources = []
            if search_results:
                context = "\n\nRELEVANT INFORMATION FROM KNOWLEDGE BASE:\n"
                for i, result in enumerate(search_results):
                    context += f"\n[Source {i+1}: {result['source']}]\n{result['text'][:800]}\n"
                    sources.append({
                        "source": result['source'], 
                        "type": result['type'], 
                        "relevance": round(result['score'] * 100)
                    })
            
            # ========================================================================
            # BUILD SYSTEM PROMPT
            # ========================================================================
            
            # Get denied tasks
            denied_task_names = []
            if access_result and hasattr(access_result, 'denied_tasks') and access_result.denied_tasks:
                denied_task_names = access_result.denied_tasks
            
            g = agent.guardrails
            p = agent.personality
            date_info = get_current_datetime_for_user(request.timezone)
            
            # Build action tools
            action_tools = [t for t in agent_tools if t.type in ['email', 'api', 'webhook', 'slack', 'websearch', 'database']]
            
            # Filter tools based on denied tasks
            if denied_task_names:
                denied_keywords = []
                for task_name in denied_task_names:
                    words = task_name.lower().replace('-', ' ').replace('_', ' ').split()
                    denied_keywords.extend([w for w in words if len(w) > 3])
                
                filtered_tools = []
                for tool in action_tools:
                    tool_name_lower = (tool.name or '').lower()
                    tool_desc_lower = (tool.description or '').lower()
                    is_denied = any(kw in tool_name_lower or kw in tool_desc_lower for kw in denied_keywords)
                    if not is_denied:
                        filtered_tools.append(tool)
                action_tools = filtered_tools
            
            tool_definitions = build_tool_definitions(action_tools) if action_tools else []
            
            # Notify about capabilities - Business friendly (skip to reduce messages)
            
            # Build accessible tasks
            accessible_tasks = [task for task in agent.tasks if task.name not in denied_task_names]
            
            # User context
            user_name = ""
            if current_user:
                if hasattr(current_user, 'first_name') and current_user.first_name:
                    user_name = current_user.first_name
                elif hasattr(current_user, 'email'):
                    user_name = current_user.email.split('@')[0].replace('.', ' ').title()
            
            # Build system prompt with instruction enforcement
            creativity_desc = "(Factual only)" if p.creativity <= 3 else "(Creative)" if p.creativity >= 7 else "(Balanced)"
            length_desc = "(Brief)" if p.length <= 3 else "(Detailed)" if p.length >= 7 else "(Moderate)"
            
            # Determine language for enforcement
            enforce_lang = 'ar' if user_lang == 'ar' else 'en'
            
            # Build base system prompt
            system_prompt = f"""You are {agent.name}.
{f'The current user is {user_name}.' if user_name else ''}

{date_info}

=== GOAL ===
{agent.goal}

=== PERSONALITY ===
• Tone: {p.tone}
• Voice: {p.voice}
• Creativity: {p.creativity}/10 {creativity_desc}
• Response Length: {p.length}/10 {length_desc}

=== LANGUAGE ===
{get_language_instruction(g.language)}

"""
            # Add instruction enforcement section
            if INSTRUCTION_ENFORCER_AVAILABLE and InstructionEnforcer:
                # Get enforcement phrases
                phrases = InstructionEnforcer.ENFORCEMENT_PHRASES.get(enforce_lang, InstructionEnforcer.ENFORCEMENT_PHRASES['en'])
                
                system_prompt += f"""
=== ⚠️ INSTRUCTION ENFORCEMENT ===
{phrases['critical']}
{phrases['mandatory']}
{phrases['no_skip']}

"""
            
            system_prompt += "=== TASKS ==="
            
            # Build tasks with enforced instructions
            for task in accessible_tasks:
                system_prompt += f"\n\n### TASK: {task.name}\n{task.description}"
                if task.instructions:
                    if INSTRUCTION_ENFORCER_AVAILABLE and InstructionEnforcer:
                        # Use enforcer formatting
                        system_prompt += InstructionEnforcer.format_instructions_for_prompt(
                            [inst.text for inst in task.instructions],
                            task.name,
                            enforce_lang
                        )
                    else:
                        # Fallback formatting with emphasis
                        system_prompt += "\n**⚠️ MANDATORY STEPS (Execute ALL in order):**"
                        for i, inst in enumerate(task.instructions, 1):
                            system_prompt += f"\n  [{i}] ✓ {inst.text}"
            
            # Add verification reminder
            if INSTRUCTION_ENFORCER_AVAILABLE and InstructionEnforcer:
                phrases = InstructionEnforcer.ENFORCEMENT_PHRASES.get(enforce_lang, InstructionEnforcer.ENFORCEMENT_PHRASES['en'])
                system_prompt += f"""

=== BEFORE YOU RESPOND ===
{phrases['verify']}
{phrases['confirm']}
"""
            
            system_prompt += context
            
            # ========================================================================
            # GENERATE THINKING FIRST (quick LLM call)
            # ========================================================================
            thinking_system = f"""You are a helpful assistant. Generate 2-3 SHORT thinking lines (max 8 words each) that show:
1. What you understood from the user's question
2. How you'll help them

Respond ONLY with the thinking lines, one per line. No other text.
Use the SAME language as the user's message.

Examples:
If user asks "show me employees details" respond:
User wants to see employee information
Searching employee records
Will display the details found

If user asks "مين مدير فادي" respond:
المستخدم يسأل عن مدير فادي
أبحث في بيانات الموظفين
سأعرض اسم المدير"""

            try:
                thinking_result = await call_llm_with_tools(
                    [
                        {"role": "system", "content": thinking_system},
                        {"role": "user", "content": request.message}
                    ],
                    [],  # no tools
                    agent.model_id
                )
                thinking_text = thinking_result.get("content", "").strip()
                if thinking_text:
                    thinking_lines = [line.strip() for line in thinking_text.split('\n') if line.strip()]
                    for line in thinking_lines[:3]:  # Max 3 lines
                        if len(line) < 80:  # Skip if too long
                            yield f"data: {json.dumps({'type': 'thinking', 'content': line})}\n\n"
                            await asyncio.sleep(0.4)
            except Exception as e:
                print(f"Thinking generation failed: {e}")
                # Continue without thinking
            
            # ========================================================================
            # LLM CALL WITH TOOL EXECUTION
            # ========================================================================
            messages = [{"role": "system", "content": system_prompt}]
            
            # Add conversation history
            for msg in conversation.messages[-10:]:  # Last 10 messages
                role = msg.role if msg.role in ['user', 'assistant'] else 'user'
                messages.append({"role": role, "content": msg.content})
            
            # Execute tools if needed (agentic loop)
            max_iterations = 5
            final_content = ""
            tool_calls_made = []
            
            for iteration in range(max_iterations):
                if iteration > 0:
                    yield f"data: {json.dumps({'type': 'thinking', 'content': msgs['processing']})}\n\n"
                
                # Call LLM using the existing function (messages, tools, model_id)
                llm_result = await call_llm_with_tools(
                    messages, 
                    tool_definitions if action_tools else [],
                    agent.model_id
                )
                
                content = llm_result.get("content", "")
                tool_calls = llm_result.get("tool_calls", [])
                
                if tool_calls:
                    # Process tool calls
                    for tc in tool_calls:
                        tool_name = tc.get("name", tc.get("function", {}).get("name", "unknown"))
                        tool_args_str = tc.get("arguments", tc.get("function", {}).get("arguments", "{}"))
                        
                        # Parse arguments
                        try:
                            tool_args = json.loads(tool_args_str) if isinstance(tool_args_str, str) else tool_args_str
                        except:
                            tool_args = {}
                        
                        # Find tool info from tool_definitions
                        tool_id = None
                        tool_type = None
                        for td in tool_definitions:
                            if td.get('function', {}).get('name') == tool_name:
                                tool_id = td.get('_tool_id')
                                tool_type = td.get('_tool_type')
                                break
                        
                        # Business-friendly action message based on tool type (in user's language)
                        action_msg = msgs['action']
                        if tool_type == 'email':
                            action_msg = msgs['sending_email']
                        elif tool_type == 'api':
                            action_msg = msgs['connecting']
                        elif tool_type == 'webhook':
                            action_msg = msgs['triggering']
                        elif tool_type == 'slack':
                            action_msg = msgs['slack']
                        elif tool_type == 'websearch':
                            action_msg = msgs['websearch']
                        elif tool_type == 'database':
                            action_msg = msgs['database']
                        
                        yield f"data: {json.dumps({'type': 'tool_call', 'content': action_msg, 'tool': tool_name})}\n\n"
                        
                        # Execute the tool using existing execute_tool function
                        try:
                            if tool_id and tool_type:
                                tool_result = await execute_tool(tool_id, tool_type, tool_args)
                            else:
                                tool_result = {"success": False, "error": f"Action not available"}
                            
                            success = tool_result.get("success", False)
                            result_str = json.dumps(tool_result) if isinstance(tool_result, dict) else str(tool_result)
                            
                            tool_calls_made.append({
                                "tool": tool_name,
                                "arguments": tool_args,
                                "result": result_str[:500],
                                "success": success
                            })
                            
                            # Business-friendly result message (in user's language)
                            if success:
                                if tool_type == 'email':
                                    status_msg = msgs['email_sent']
                                elif tool_type == 'slack':
                                    status_msg = msgs['slack_sent']
                                elif tool_type == 'websearch':
                                    status_msg = msgs['search_done']
                                elif tool_type == 'database':
                                    status_msg = msgs['data_done']
                                elif tool_type == 'api':
                                    status_msg = msgs['service_done']
                                else:
                                    status_msg = msgs['action_done']
                            else:
                                status_msg = msgs['action_failed']
                            
                            yield f"data: {json.dumps({'type': 'tool_result', 'content': status_msg, 'tool': tool_name, 'success': success})}\n\n"
                            
                            # Add to messages for next iteration - ensure proper format
                            tool_call_formatted = {
                                "id": tc.get("id", f"call_{tool_name}"),
                                "type": "function",
                                "function": {
                                    "name": tool_name,
                                    "arguments": json.dumps(tool_args) if isinstance(tool_args, dict) else tool_args
                                }
                            }
                            messages.append({
                                "role": "assistant",
                                "content": content or "",
                                "tool_calls": [tool_call_formatted]
                            })
                            messages.append({
                                "role": "tool",
                                "tool_call_id": tool_call_formatted["id"],
                                "content": result_str
                            })
                        except Exception as e:
                            tool_calls_made.append({
                                "tool": tool_name,
                                "arguments": tool_args,
                                "error": str(e),
                                "success": False
                            })
                            yield f"data: {json.dumps({'type': 'tool_result', 'content': msgs['action_failed'], 'tool': tool_name, 'success': False})}\n\n"
                            
                            # Add error to messages with proper format
                            tool_call_formatted = {
                                "id": tc.get("id", f"call_{tool_name}"),
                                "type": "function",
                                "function": {
                                    "name": tool_name,
                                    "arguments": json.dumps(tool_args) if isinstance(tool_args, dict) else tool_args
                                }
                            }
                            messages.append({
                                "role": "assistant",
                                "content": content or "",
                                "tool_calls": [tool_call_formatted]
                            })
                            messages.append({
                                "role": "tool",
                                "tool_call_id": tool_call_formatted["id"],
                                "content": f"Error: {str(e)}"
                            })
                else:
                    # No tool calls, we have the final response
                    final_content = content
                    break
            
            if not final_content:
                final_content = "I apologize, but I couldn't complete the task within the allowed iterations."
            
            # ========================================================================
            # STREAM RESPONSE CONTENT
            # ========================================================================
            # Stream content in chunks for better UX
            chunk_size = 50
            for i in range(0, len(final_content), chunk_size):
                chunk = final_content[i:i+chunk_size]
                yield f"data: {json.dumps({'type': 'content', 'content': chunk})}\n\n"
                await asyncio.sleep(0.02)  # Small delay for streaming effect
            
            # Send sources if any
            if sources:
                yield f"data: {json.dumps({'type': 'sources', 'content': sources})}\n\n"
            
            # ========================================================================
            # SAVE ASSISTANT MESSAGE
            # ========================================================================
            assistant_msg = ConversationMessage(role="assistant", content=final_content, sources=sources)
            conversation.messages.append(assistant_msg)
            conversation.updated_at = datetime.utcnow().isoformat()
            
            try:
                from database.services import ConversationService
                ConversationService.add_message(conversation.id, {
                    'id': assistant_msg.id,
                    'role': 'assistant',
                    'content': final_content,
                    'sources': sources
                }, org_id, user_id)
            except Exception:
                pass
            
            app_state.save_to_disk()
            
            # Send completion
            yield f"data: {json.dumps({'type': 'done', 'tool_calls': tool_calls_made})}\n\n"
            
        except Exception as e:
            print(f"❌ Streaming error: {e}")
            import traceback
            traceback.print_exc()
            yield f"data: {json.dumps({'type': 'error', 'content': str(e)})}\n\n"
            yield f"data: {json.dumps({'type': 'done'})}\n\n"
    
    return StreamingResponse(
        event_generator(),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "Connection": "keep-alive",
            "X-Accel-Buffering": "no"  # Disable nginx buffering
        }
    )


@app.post("/api/agents/{agent_id}/chat-with-files")
async def chat_with_files(
    agent_id: str,
    message: str = Form(""),
    conversation_id: Optional[str] = Form(None),
    timezone: Optional[str] = Form(None),
    files: List[UploadFile] = File(default=[]),
    current_user: User = Depends(get_current_user)
):
    """Chat with agent including file attachments."""
    if agent_id not in app_state.agents:
        raise HTTPException(404, "Agent not found")
    
    agent = app_state.agents[agent_id]
    
    org_id = current_user.org_id if current_user else "org_default"
    user_id = str(current_user.id) if current_user else "system"
    
    # ========================================================================
    # ACCESS CONTROL ENFORCEMENT (3-Level Check)
    # ========================================================================
    access_result = None
    if ACCESS_CONTROL_AVAILABLE and AccessControlService and current_user:
        try:
            user_role_ids = getattr(current_user, 'role_ids', []) or []
            user_group_ids = get_user_group_ids(user_id) if user_id else []
            
            access_result = AccessControlService.check_user_access(
                user_id=user_id,
                user_role_ids=user_role_ids,
                user_group_ids=user_group_ids,
                agent_id=agent_id,
                org_id=org_id
            )
            
            if not access_result.has_access:
                raise HTTPException(
                    status_code=403,
                    detail={
                        "error": "access_denied",
                        "message": access_result.reason or "This assistant is not available for your account. Please contact your administrator if you need access.",
                        "user_friendly": True
                    }
                )
        except HTTPException:
            raise
        except Exception as e:
            print(f"⚠️ Access control check failed (allowing by default): {e}")
            access_result = None
    
    # Get or create conversation
    is_new_conversation = False
    if conversation_id and conversation_id in app_state.conversations:
        conversation = app_state.conversations[conversation_id]
    else:
        is_new_conversation = True
        # Start with temporary title - LLM will update it
        title = "New conversation"
        conversation = Conversation(agent_id=agent_id, user_id=user_id, title=title)
        app_state.conversations[conversation.id] = conversation
        
        # Save to database with temporary title
        try:
            from database.services import ConversationService
            ConversationService.create_conversation({
                'id': conversation.id,
                'agent_id': agent_id,
                'user_id': user_id,  # PRIVACY: Track ownership
                'title': title
            }, org_id, user_id)
        except Exception as e:
            print(f"⚠️  Failed to save conversation to DB: {e}")
        
        # Schedule LLM to generate smart title in background
        if message:
            from api.modules.conversations import ConversationTitleService
            asyncio.create_task(
                ConversationTitleService.generate_and_update_title(
                    conversation_id=conversation.id,
                    first_message=message,
                    agent_name=agent.name if agent else None,
                    model_id=agent.model_id if agent else None
                )
            )
    
    # Process uploaded files and extract text
    file_contents = []
    for file in files:
        try:
            content = await file.read()
            file_text = ""
            filename = file.filename.lower()
            
            if filename.endswith('.txt') or filename.endswith('.md'):
                file_text = content.decode('utf-8', errors='ignore')
            elif filename.endswith('.csv'):
                file_text = content.decode('utf-8', errors='ignore')
            elif filename.endswith('.pdf'):
                # Try to extract text from PDF
                try:
                    import fitz  # PyMuPDF
                    pdf_doc = fitz.open(stream=content, filetype="pdf")
                    for page in pdf_doc:
                        file_text += page.get_text() + "\n"
                    pdf_doc.close()
                except:
                    file_text = "[PDF file - text extraction failed]"
            elif filename.endswith('.docx'):
                try:
                    from docx import Document as DocxDocument
                    import io
                    doc = DocxDocument(io.BytesIO(content))
                    file_text = "\n".join([para.text for para in doc.paragraphs])
                except:
                    file_text = "[DOCX file - text extraction failed]"
            elif filename.endswith('.xlsx') or filename.endswith('.xls'):
                try:
                    import openpyxl
                    import io
                    wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True)
                    for sheet in wb.worksheets:
                        for row in sheet.iter_rows(values_only=True):
                            row_text = " | ".join([str(cell) if cell else "" for cell in row])
                            if row_text.strip():
                                file_text += row_text + "\n"
                except:
                    file_text = "[Excel file - text extraction failed]"
            elif filename.endswith('.pptx') or filename.endswith('.ppt'):
                try:
                    from pptx import Presentation
                    import io
                    prs = Presentation(io.BytesIO(content))
                    slide_num = 0
                    for slide in prs.slides:
                        slide_num += 1
                        file_text += f"\n--- Slide {slide_num} ---\n"
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text:
                                file_text += shape.text + "\n"
                            # Extract table data if present
                            if shape.has_table:
                                table = shape.table
                                for row in table.rows:
                                    row_text = " | ".join([cell.text for cell in row.cells])
                                    file_text += row_text + "\n"
                except:
                    file_text = "[PowerPoint file - text extraction failed]"
            elif filename.endswith(('.png', '.jpg', '.jpeg', '.gif', '.webp', '.bmp')):
                # Handle image files - store as base64 for vision API
                import base64
                image_base64 = base64.b64encode(content).decode('utf-8')
                # Determine media type
                if filename.endswith('.png'):
                    media_type = 'image/png'
                elif filename.endswith(('.jpg', '.jpeg')):
                    media_type = 'image/jpeg'
                elif filename.endswith('.gif'):
                    media_type = 'image/gif'
                elif filename.endswith('.webp'):
                    media_type = 'image/webp'
                else:
                    media_type = 'image/png'
                
                file_contents.append({
                    "filename": file.filename,
                    "content": "[Image file - will be analyzed with vision]",
                    "is_image": True,
                    "image_data": image_base64,
                    "media_type": media_type
                })
                continue  # Skip the normal file_contents append below
            elif filename.endswith('.svg'):
                # SVG is text-based, extract the content
                file_text = content.decode('utf-8', errors='ignore')
            else:
                file_text = f"[Unsupported file type: {filename}]"
            
            if file_text:
                file_contents.append({
                    "filename": file.filename,
                    "content": file_text[:10000]  # Limit content size
                })
        except Exception as e:
            file_contents.append({
                "filename": file.filename,
                "content": f"[Error reading file: {str(e)}]"
            })
    
    # Build enhanced message with file contents
    enhanced_message = message
    image_files = [fc for fc in file_contents if fc.get('is_image')]
    text_files = [fc for fc in file_contents if not fc.get('is_image')]
    
    if text_files:
        enhanced_message += "\n\n--- ATTACHED FILES ---\n"
        for fc in text_files:
            enhanced_message += f"\n📄 **{fc['filename']}**:\n```\n{fc['content']}\n```\n"
    
    # Add user message to conversation
    display_message = message
    if files:
        display_message += f"\n\n[Attached: {', '.join([f.filename for f in files])}]"
    user_msg = ConversationMessage(role="user", content=display_message)
    conversation.messages.append(user_msg)
    
    # Handle images with vision API
    if image_files:
        # Use OpenAI vision to analyze images
        try:
            from openai import OpenAI
            api_key = os.getenv("OPENAI_API_KEY")
            if not api_key:
                result = {"content": "Error: OpenAI API key not configured", "sources": []}
            else:
                user_text = enhanced_message or "Please analyze these images."
                content_parts = [{"type": "text", "text": user_text}]
                
                for img in image_files:
                    content_parts.append({
                        "type": "image_url",
                        "image_url": {
                            "url": f"data:{img['media_type']};base64,{img['image_data']}",
                            "detail": "high"
                        }
                    })
                
                # Build system prompt for agent with FULL configuration
                p = agent.personality
                g = agent.guardrails
                
                # Personality descriptions
                creativity_desc = "(Factual only)" if p.creativity <= 3 else "(Creative)" if p.creativity >= 7 else "(Balanced)"
                length_desc = "(Brief)" if p.length <= 3 else "(Detailed)" if p.length >= 7 else "(Moderate)"
                formality_desc = "(Casual)" if p.formality <= 3 else "(Formal)" if p.formality >= 7 else "(Professional)"
                empathy_desc = "(Direct)" if p.empathy <= 3 else "(Empathetic)" if p.empathy >= 7 else "(Supportive)"
                
                # Current date with user's timezone
                date_info = get_current_datetime_for_user(timezone)
                
                # Get language instruction - put it prominently at the top
                lang_instruction = get_language_instruction(g.language)
                
                system_prompt = f"""You are {agent.name}.

=== LANGUAGE ===
{lang_instruction}

{date_info}
=== GOAL ===
{agent.goal}

=== PERSONALITY ===
• Tone: {p.tone}
• Voice: {p.voice}
• Creativity: {p.creativity}/10 {creativity_desc}
• Response Length: {p.length}/10 {length_desc}
• Formality: {p.formality}/10 {formality_desc}
• Empathy: {p.empathy}/10 {empathy_desc}

=== TASKS ==="""
                # Add instruction enforcement
                if INSTRUCTION_ENFORCER_AVAILABLE and InstructionEnforcer:
                    enforce_lang = 'ar' if g.language == 'ar' else 'en'
                    phrases = InstructionEnforcer.ENFORCEMENT_PHRASES.get(enforce_lang, InstructionEnforcer.ENFORCEMENT_PHRASES['en'])
                    system_prompt += f"""

=== ⚠️ INSTRUCTION ENFORCEMENT ===
{phrases['critical']}
{phrases['mandatory']}
"""
                
                for task in agent.tasks:
                    system_prompt += f"\n### TASK: {task.name}\n{task.description}"
                    if task.instructions:
                        if INSTRUCTION_ENFORCER_AVAILABLE and InstructionEnforcer:
                            enforce_lang = 'ar' if g.language == 'ar' else 'en'
                            inst_texts = [inst.text if hasattr(inst, 'text') else str(inst) for inst in task.instructions]
                            system_prompt += InstructionEnforcer.format_instructions_for_prompt(inst_texts, task.name, enforce_lang)
                        else:
                            system_prompt += "\n**⚠️ MANDATORY STEPS:**"
                            for inst in task.instructions:
                                inst_text = inst.text if hasattr(inst, 'text') else str(inst)
                                system_prompt += f"\n  ✓ {inst_text}"
                
                # Add guardrails
                if g.anti_hallucination:
                    system_prompt += "\n\n=== ACCURACY ===\nUse facts from knowledge base. Don't make up data."
                
                # Response length
                length_map = {'short': '1-2 paragraphs MAX', 'medium': '2-3 paragraphs', 'long': 'Detailed', 'unlimited': 'Thorough'}
                system_prompt += f"\n\n**Response Length:** {length_map.get(g.max_length, '2-3 paragraphs')}"
                
                system_prompt += "\n\nAnalyze the images based on your tasks and respond according to your settings."
                
                vision_messages = [
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": content_parts}
                ]
                
                # Call OpenAI with vision
                client = OpenAI(api_key=api_key)
                response = client.chat.completions.create(
                    model="gpt-4o",  # Use gpt-4o for vision
                    messages=vision_messages,
                    max_tokens=4096
                )
                
                result = {"content": response.choices[0].message.content, "sources": []}
        except Exception as e:
            result = {"content": f"Error analyzing images: {str(e)}", "sources": []}
    else:
        # Process chat with enhanced message (no images) - with access control and user context
        result = await process_agent_chat(agent, enhanced_message, conversation, timezone=timezone, access_control=access_result, current_user=current_user)
    
    # Add assistant response
    assistant_msg = ConversationMessage(role="assistant", content=result["content"], sources=result["sources"])
    conversation.messages.append(assistant_msg)
    conversation.updated_at = datetime.utcnow().isoformat()
    app_state.save_to_disk()
    
    return ChatResponse(
        response=result["content"],
        conversation_id=conversation.id,
        sources=result["sources"],
        formatted=True
    )


@app.get("/ui", response_class=HTMLResponse)
@app.get("/ui/", response_class=HTMLResponse)
@app.get("/ui/{path:path}", response_class=HTMLResponse)
async def serve_ui(path: str = ""):
    ui_file = "ui/index.html"
    if os.path.exists(ui_file):
        with open(ui_file) as f:
            return f.read()
    return "<html><body><h1>UI not found</h1></body></html>"


# ============================================================================
# END USER CHAT PORTAL - Separate interface for end users
# ============================================================================
@app.get("/chat", response_class=HTMLResponse)
@app.get("/chat/", response_class=HTMLResponse)
@app.get("/chat/{path:path}", response_class=HTMLResponse)
async def serve_chat_portal(path: str = ""):
    """End User Chat Portal - Modern chat interface for using agents"""
    chat_file = "ui/chat.html"
    if os.path.exists(chat_file):
        with open(chat_file) as f:
            return f.read()
    return "<html><body><h1>Chat Portal not found</h1></body></html>"


# ============================================================================
# LAB PORTAL - Test Data Generator
# ============================================================================
@app.get("/lab", response_class=HTMLResponse)
@app.get("/lab/", response_class=HTMLResponse)
async def serve_lab_portal():
    """Lab Portal - Generate test APIs, documents, and images"""
    lab_file = "ui/lab.html"
    if os.path.exists(lab_file):
        with open(lab_file) as f:
            return f.read()
    return "<html><body><h1>Lab Portal not found</h1></body></html>"


@app.get("/frontend", response_class=HTMLResponse)
@app.get("/frontend/", response_class=HTMLResponse)
@app.get("/frontend/{path:path}", response_class=HTMLResponse)
async def serve_frontend(path: str = ""):
    ui_file = "ui/index.html"
    if os.path.exists(ui_file):
        with open(ui_file) as f:
            return f.read()
    return "<html><body><h1>Frontend not found</h1></body></html>"


@app.get("/admin", response_class=HTMLResponse)
@app.get("/admin/", response_class=HTMLResponse)
async def serve_admin():
    # Try ui folder first, then root
    for path in ["ui/admin.html", "admin.html"]:
        if os.path.exists(path):
            with open(path) as f:
                return f.read()
    return "<html><body><h1>Admin not found</h1></body></html>"


@app.get("/monitor", response_class=HTMLResponse)
@app.get("/monitor/", response_class=HTMLResponse)
async def serve_monitor():
    # Try ui folder first, then root
    for path in ["ui/monitor.html", "monitor.html"]:
        if os.path.exists(path):
            with open(path) as f:
                return f.read()
    return "<html><body><h1>Monitor not found</h1></body></html>"


# ============================================================================
# DEMO LAB - Mock API & Document Generator
# ============================================================================

class DemoItem(BaseModel):
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    type: str  # api, document, image
    name: str
    description: str = ""
    icon: str = "📄"
    url: str = ""
    config: Dict[str, Any] = {}
    created_at: str = Field(default_factory=lambda: datetime.utcnow().isoformat())

class DemoGenerateRequest(BaseModel):
    message: str
    conversation: List[Dict[str, str]] = []
    mode_hint: Optional[str] = None  # api, document, or image

class DemoCreateToolRequest(BaseModel):
    demo_id: str

# Store for demo items
demo_items: Dict[str, DemoItem] = {}
mock_api_data: Dict[str, Dict[str, Any]] = {}
mock_api_llm_cache: Dict[str, Dict[str, Any]] = {}  # Cache LLM responses

@app.get("/api/demo/items")
async def get_demo_items():
    """Get all generated demo items"""
    items = list(demo_items.values())
    return {
        "apis": [i.dict() for i in items if i.type == "api"],
        "docs": [i.dict() for i in items if i.type == "document"],
        "images": [i.dict() for i in items if i.type == "image"]
    }


@app.put("/api/demo/items/{item_id}")
async def update_demo_item(item_id: str, request: dict):
    """Update a demo item's name and description"""
    if item_id not in demo_items:
        raise HTTPException(404, "Item not found")
    
    item = demo_items[item_id]
    if "name" in request:
        item.name = request["name"]
    if "description" in request:
        item.description = request["description"]
    
    demo_items[item_id] = item
    return {"success": True, "item": item.dict()}


@app.delete("/api/demo/items/{item_id}")
async def delete_demo_item(item_id: str, type: str = None):
    """Delete a demo item"""
    if item_id not in demo_items:
        raise HTTPException(404, "Item not found")
    
    item = demo_items[item_id]
    
    # Try to delete associated file if it exists
    if hasattr(item, 'config') and item.config:
        filepath = item.config.get('filepath')
        if filepath and os.path.exists(filepath):
            try:
                os.remove(filepath)
                print(f"[Demo Lab] Deleted file: {filepath}")
            except Exception as e:
                print(f"[Demo Lab] Could not delete file: {e}")
    
    del demo_items[item_id]
    return {"success": True}


@app.delete("/api/demo/items/clear")
async def clear_all_demo_items():
    """Clear all demo items"""
    global demo_items
    
    # Try to delete associated files
    for item in demo_items.values():
        if hasattr(item, 'config') and item.config:
            filepath = item.config.get('filepath')
            if filepath and os.path.exists(filepath):
                try:
                    os.remove(filepath)
                except Exception:
                    pass
    
    demo_items.clear()
    return {"success": True}

@app.post("/api/demo/generate")
async def generate_demo(request: DemoGenerateRequest):
    """Use LLM to understand request and generate mock API, document, or image"""
    message = request.message.lower()
    original_message = request.message
    mode_hint = request.mode_hint  # User's selected mode (api, document, image)
    
    # ============================================================
    # CHECK FOR IMAGE REQUESTS FIRST (highest priority)
    # ============================================================
    is_image_request = any(word in message for word in ["image", "picture", "photo", "screenshot", "visual", "graphic", "ocr"])
    
    # If mode_hint is image OR message contains image keywords, treat as image request
    if mode_hint == "image" or is_image_request:
        print(f"[Demo Lab] Image generation detected for: {message}")
        
        # Determine image type
        if "invoice" in message:
            image_type = "invoice"
            image_name = "Invoice Image"
        elif "receipt" in message:
            image_type = "receipt"
            image_name = "Receipt Image"
        elif "logo" in message:
            image_type = "logo"
            image_name = "Company Logo"
        elif "chart" in message or "graph" in message:
            image_type = "chart"
            image_name = "Business Chart"
        elif "screenshot" in message or "ui" in message:
            image_type = "screenshot"
            image_name = "UI Screenshot"
        else:
            image_type = "general"
            image_name = "Generated Image"
        
        # Generate image
        image_result = await generate_demo_image(image_name, original_message, image_type)
        
        if image_result.get("success"):
            item = DemoItem(
                id=str(uuid.uuid4()),
                name=image_name,
                description=original_message[:100],
                type="image",
                url=image_result.get("url", ""),
                config={"image_type": image_type, "filepath": image_result.get("filepath")}
            )
            demo_items[item.id] = item
            
            return {
                "response": f"🖼️ Here's your **{image_name}**! I've generated a {'clear, OCR-friendly' if image_type in ['invoice', 'receipt', 'document'] else 'professional'} image based on your request.",
                "generated": {"type": "image", "id": item.id, "url": image_result.get("url", ""), "name": image_name}
            }
        else:
            return {
                "response": f"❌ Image generation failed: {image_result.get('error', 'Unknown error')}",
                "generated": None
            }
    
    # ============================================================
    # DIRECT DOCUMENT GENERATION (Skip LLM for simple document requests)
    # ============================================================
    is_doc_request = any(word in message for word in ["document", "pdf", "word", "excel", "docx", "xlsx"])
    is_policy = any(word in message for word in ["policy", "hr policy", "employee", "handbook"])
    is_invoice = any(word in message for word in ["invoice", "bill", "receipt"]) and not is_image_request
    is_report = any(word in message for word in ["report", "analysis", "summary"])
    
    # If mode_hint is document OR message contains document keywords
    if mode_hint == "document" or is_doc_request or is_policy or is_invoice or is_report:
        print(f"[Demo Lab] Direct document generation for: {message}")
        
        # Determine document type
        if "excel" in message or "xlsx" in message or "spreadsheet" in message:
            doc_type = "excel"
        elif "word" in message or "docx" in message:
            doc_type = "word"
        else:
            doc_type = "pdf"  # Default to PDF
        
        # Determine content type
        if is_invoice:
            content_type = "invoice"
            doc_name = "Sales Invoice"
        elif is_policy or "hr" in message:
            content_type = "policy"
            doc_name = "HR Policy Document"
        elif is_report:
            content_type = "report"
            doc_name = "Business Report"
        else:
            content_type = "document"
            doc_name = "Generated Document"
        
        # Extract custom name if provided
        if "for" in message:
            parts = message.split("for")
            if len(parts) > 1:
                doc_name = parts[1].strip().title()[:50]
        
        print(f"[Demo Lab] Creating {doc_type} document, content: {content_type}, name: {doc_name}")
        
        try:
            # Generate the document directly
            doc_result = await generate_demo_document(doc_name, doc_type, content_type, {}, original_message)
            print(f"[Demo Lab] Document result: {doc_result}")
            
            if doc_result and doc_result.get("url"):
                item = DemoItem(
                    type="document",
                    name=doc_name,
                    description=original_message,
                    icon={"word": "📘", "excel": "📗", "pdf": "📕"}.get(doc_type, "📄"),
                    url=doc_result["url"],
                    config={"doc_type": doc_type, "content_type": content_type}
                )
                demo_items[item.id] = item
                
                return {
                    "response": f"""✅ Here's your document! I've created **{doc_name}**.

📄 **Format:** {doc_type.upper()}
📋 **Type:** {content_type.title()}

Click the Download button below to save it.""",
                    "generated": {
                        "type": "document",
                        "id": item.id,
                        "url": doc_result["url"],
                        "name": doc_name
                    }
                }
            else:
                error_msg = doc_result.get("error", "Unknown error") if doc_result else "Generation failed"
                return {
                    "response": f"""❌ Sorry, I couldn't generate the document.

**Error:** {error_msg}

Please check server logs or try a different format.""",
                    "generated": None
                }
        except Exception as e:
            print(f"[Demo Lab] Direct document generation error: {e}")
            import traceback
            traceback.print_exc()
            return {
                "response": f"❌ Error generating document: {str(e)}",
                "generated": None
            }
    
    # ============================================================
    # LLM-BASED GENERATION (for APIs and complex requests)
    # ============================================================
    
    # Build conversation for LLM
    system_prompt = """You are a Demo Lab assistant that INSTANTLY creates mock APIs, sample documents, and images for demos.

CRITICAL RULES:
1. NEVER ask questions or request more information
2. ALWAYS generate everything immediately with realistic sample data
3. Make smart assumptions based on the request
4. Use industry-standard data patterns
5. For documents: ALWAYS include doc_type, content_type, and sample_data with FULL content
6. IMPORTANT: Do NOT use Unicode characters like bullet points (•), smart quotes (" " ' '), or special dashes (– —). Use ASCII only: - for bullets, " for quotes, - for dashes.

DETECT REQUEST TYPE (check in this ORDER - first match wins):
1. If user mentions "image", "picture", "photo", "screenshot", "visual", "graphic", "OCR" -> type: "image"
   Examples: "invoice image", "receipt photo", "document image for OCR" = ALL are type "image"
2. If user mentions "API", "endpoint", "REST", "mock API" -> type: "api"
3. If user mentions "document", "PDF", "Word", "Excel", "policy", "report" -> type: "document"

IMPORTANT: "invoice image" or "receipt image" = type "image" (an IMAGE of an invoice)
           "invoice document" or "invoice PDF" = type "document" (a PDF document)

ALWAYS respond with this JSON format:
{
    "type": "api" | "document" | "image",
    "name": "Short descriptive name",
    "description": "Brief description",
    "response_text": "Here's your [item]! I've created [description].",
    "details": { ... }
}

FOR IMAGES - details MUST include:
{
    "image_type": "invoice" | "receipt" | "document" | "chart" | "logo" | "screenshot" | "general",
    "image_description": "Detailed description of what the image should show",
    "purpose": "OCR testing" | "demo" | "mockup" | "general"
}

IMAGE EXAMPLE - "Generate invoice image for OCR":
{
    "type": "image",
    "name": "Sample Invoice Image",
    "description": "Clear invoice image optimized for OCR testing",
    "response_text": "Here's your invoice image! I've created a clear, high-contrast invoice perfect for OCR testing.",
    "details": {
        "image_type": "invoice",
        "image_description": "Professional invoice with clear text: Company logo at top, invoice number INV-2024-001, date, billing details, itemized list with quantities and prices, subtotal, tax, and total. High contrast black text on white background for optimal OCR accuracy.",
        "purpose": "OCR testing"
    }
}

FOR DOCUMENTS - details MUST include:
{
    "doc_type": "pdf" | "word" | "excel" | "text",  // REQUIRED!
    "content_type": "policy" | "invoice" | "report" | "letter" | "contract",
    "sample_data": {
        "title": "Document Title",
        "company": "Company Name",
        "date": "2024-12-25",
        "sections": {
            "Section 1 Title": "Full paragraph content here...",
            "Section 2 Title": "More detailed content..."
        }
    }
}

DOCUMENT EXAMPLES:

If user asks for "HR Policy PDF":
{
    "type": "document",
    "name": "HR Policy Document",
    "description": "Comprehensive HR policy for tech company",
    "response_text": "Here's your HR Policy document!",
    "details": {
        "doc_type": "pdf",
        "content_type": "policy",
        "sample_data": {
            "title": "Employee Handbook & HR Policies",
            "company": "TechCorp Solutions",
            "effective_date": "January 1, 2024",
            "version": "2.0",
            "approved_by": "Sarah Johnson, CHRO",
            "sections": {
                "1. Introduction": "Welcome to TechCorp Solutions. This handbook outlines our policies and procedures...",
                "2. Employment Policies": "2.1 Equal Opportunity: TechCorp is an equal opportunity employer...",
                "3. Work Schedule": "Standard work hours are 9:00 AM to 6:00 PM, Monday through Friday...",
                "4. Leave Policy": "Employees are entitled to 21 days of annual leave...",
                "5. Code of Conduct": "All employees are expected to maintain professional behavior..."
            }
        }
    }
}

If user asks for "Invoice Excel":
{
    "type": "document",
    "name": "Sales Invoice",
    "description": "Professional invoice template",
    "response_text": "Here's your Invoice spreadsheet!",
    "details": {
        "doc_type": "excel",
        "content_type": "invoice",
        "sample_data": {
            "invoice_number": "INV-2024-001234",
            "company": "Acme Corp",
            "customer": "John Smith",
            "items": [
                {"description": "Product A", "qty": 2, "price": 150.00},
                {"description": "Service B", "qty": 1, "price": 500.00}
            ],
            "subtotal": 800.00,
            "tax": 64.00,
            "total": 864.00
        }
    }
}

FOR APIs - details MUST include endpoints array with sample_response

REMEMBER: 
- Generate IMMEDIATELY with FULL content
- Never ask questions
- Always include doc_type for documents
- Include realistic, detailed sample_data"""
    
    try:
        # Call LLM
        llm_provider = app_state.get_llm_provider()
        
        messages = [{"role": "system", "content": system_prompt}]
        for msg in request.conversation[-6:]:  # Last 6 messages for context
            messages.append(msg)
        messages.append({"role": "user", "content": request.message})
        
        response = await llm_provider.generate(messages)
        
        # Try to parse JSON from response
        try:
            # Find JSON in response
            import re
            json_match = re.search(r'\{[\s\S]*\}', response)
            if json_match:
                data = json.loads(json_match.group())
            else:
                # LLM didn't return JSON, provide helpful response
                return {
                    "response": response,
                    "generated": None
                }
        except json.JSONDecodeError:
            return {
                "response": response,
                "generated": None
            }
        
        item_type = data.get("type", "api")
        name = data.get("name", "Demo Item")
        description = data.get("description", "")
        response_text = data.get("response_text", "I've created that for you!")
        details = data.get("details", {})
        
        generated_item = None
        
        if item_type == "api":
            # Create mock API
            base_name = details.get("base_name", name.lower().replace(" ", "_"))
            endpoints = details.get("endpoints", [])
            
            # Store mock data
            mock_api_data[base_name] = {
                "endpoints": endpoints,
                "name": name,
                "description": description
            }
            
            # Create demo item
            item = DemoItem(
                type="api",
                name=name,
                description=description,
                icon="🔗",
                url=f"/demo-api/{base_name}",
                config={"endpoints": endpoints, "base_name": base_name}
            )
            demo_items[item.id] = item
            generated_item = {"type": "api", "id": item.id, "url": item.url, "name": name}
            
            # Format endpoint list for response
            endpoint_list = "\n".join([f"  {e.get('method', 'GET')} {e.get('path', '/')} - {e.get('description', '')}" for e in endpoints])
            response_text = f"""I've created the **{name}** mock API with the following endpoints:

```
{endpoint_list}
```

🔗 **Base URL:** `{item.url}`

You can test it now or create an API Tool to use with your agents."""
        
        elif item_type == "document":
            # Generate document
            import random as rand_mod
            doc_type = details.get("doc_type", "pdf")  # Default to PDF
            content_type = details.get("content_type", "document")
            sample_data = details.get("sample_data", {})
            
            # If no sample_data from LLM, generate some based on content_type
            if not sample_data:
                if "policy" in content_type.lower() or "hr" in name.lower():
                    sample_data = {
                        "sections": {
                            "1. Introduction": "This policy establishes guidelines and procedures for all employees.",
                            "2. Scope": "This policy applies to all full-time and part-time employees.",
                            "3. Policy Statement": "The company is committed to maintaining professional standards.",
                            "4. Procedures": "All employees must follow the established procedures.",
                            "5. Compliance": "Non-compliance may result in disciplinary action."
                        }
                    }
                elif "invoice" in content_type.lower():
                    sample_data = {
                        "invoice_number": f"INV-{rand_mod.randint(10000, 99999)}",
                        "date": datetime.utcnow().strftime("%B %d, %Y"),
                        "customer": "Sample Customer",
                        "items": [
                            {"description": "Service Fee", "amount": 500},
                            {"description": "Consulting", "amount": 1500}
                        ],
                        "total": 2000
                    }
                else:
                    sample_data = {
                        "title": name,
                        "content": description or "Sample document content generated by Demo Lab."
                    }
            
            print(f"[Demo Lab] Generating document: {name}, type: {doc_type}, content: {content_type}")
            
            # Generate the document
            try:
                doc_result = await generate_demo_document(name, doc_type, content_type, sample_data, description)
                print(f"[Demo Lab] Document result: {doc_result}")
            except Exception as doc_err:
                print(f"[Demo Lab] Document generation error: {doc_err}")
                doc_result = None
            
            if doc_result and doc_result.get("url"):
                item = DemoItem(
                    type="document",
                    name=name,
                    description=description,
                    icon={"word": "📘", "excel": "📗", "pdf": "📕", "ppt": "📙", "text": "📄"}.get(doc_type, "📄"),
                    url=doc_result["url"],
                    config={"doc_type": doc_type, "content_type": content_type}
                )
                demo_items[item.id] = item
                generated_item = {"type": "document", "id": item.id, "url": doc_result["url"], "name": name}
                
                response_text = f"""✅ Here's your document! I've created **{name}**.

📄 **Format:** {doc_type.upper()}
📋 **Type:** {content_type}

Click the Download button below to save it."""
            else:
                # Document generation failed - still create item but with error
                error_msg = doc_result.get("error", "Unknown error") if doc_result else "Generation failed"
                print(f"[Demo Lab] Document generation failed: {error_msg}")
                
                response_text = f"""❌ Sorry, I couldn't generate the document.

**Requested:** {name} ({doc_type.upper()})
**Error:** {error_msg}

This might be due to missing libraries (fpdf2, python-docx, openpyxl). Please try:
- A different format (word, excel, pdf, text)
- Or check server logs for errors"""
                generated_item = None
        
        elif item_type == "image":
            # Generate image using OpenAI DALL-E
            image_desc = details.get("description", description)
            image_type = details.get("image_type", "general")
            
            # Try to generate the image
            image_result = await generate_demo_image(name, image_desc, image_type)
            
            if image_result and image_result.get("success"):
                item = DemoItem(
                    type="image",
                    name=name,
                    description=description,
                    icon="🖼️",
                    url=image_result.get("url", ""),
                    config={"image_description": image_desc, "image_type": image_type}
                )
                demo_items[item.id] = item
                
                response_text = f"""I've generated the **{name}** image!

🖼️ **Type:** {image_type.title()}
📝 **Description:** {image_desc[:100]}...

You can download it now."""
                generated_item = {"type": "image", "id": item.id, "url": image_result.get("url", ""), "name": name}
            else:
                # Image generation failed - provide helpful message
                error_msg = image_result.get("error", "Unknown error") if image_result else "No image generator configured"
                
                item = DemoItem(
                    type="image",
                    name=name,
                    description=description,
                    icon="🖼️",
                    url="",
                    config={"image_description": image_desc, "image_type": image_type, "pending": True}
                )
                demo_items[item.id] = item
                
                response_text = f"""⚠️ Image generation requires OpenAI API with DALL-E access.

**What I would create:** {name}
📝 **Description:** {image_desc}

**To enable image generation:**
1. Go to ⚙️ Settings
2. Ensure OpenAI API key is configured
3. The API key needs DALL-E access

Alternatively, I can create:
- 📄 A document with the specifications
- 🔗 A mock API that returns image metadata"""
                generated_item = {"type": "image", "id": item.id, "name": name, "pending": True}
        
        return {
            "response": response_text,
            "generated": generated_item
        }
        
    except Exception as e:
        print(f"Demo generation error: {e}")
        return {
            "response": f"I encountered an error: {str(e)}. Please try again or be more specific about what you need.",
            "generated": None
        }


async def generate_demo_image(name: str, description: str, image_type: str) -> Dict:
    """Generate image - uses programmatic generation for documents/OCR, DALL-E for creative images"""
    try:
        output_dir = os.environ.get("UPLOAD_PATH", "data/uploads")
        os.makedirs(output_dir, exist_ok=True)
        
        # Check if this is a document-type image that needs clear text (for OCR)
        ocr_types = ["invoice", "receipt", "document", "form", "certificate", "letter"]
        is_ocr_image = image_type.lower() in ocr_types or "ocr" in description.lower()
        
        if is_ocr_image:
            # Generate programmatically for clear, readable text
            return await generate_programmatic_document_image(name, description, image_type, output_dir)
        else:
            # Use DALL-E for creative/visual images
            return await generate_dalle_image(name, description, image_type, output_dir)
            
    except Exception as e:
        print(f"[Demo Lab] Image generation error: {e}")
        import traceback
        traceback.print_exc()
        return {"success": False, "error": str(e)}


async def generate_programmatic_document_image(name: str, description: str, image_type: str, output_dir: str) -> Dict:
    """Generate document images programmatically for clear text (OCR-friendly)"""
    try:
        from PIL import Image, ImageDraw, ImageFont
        import random
        from datetime import datetime, timedelta
        
        print(f"[Demo Lab] Generating programmatic {image_type} image for OCR")
        
        # Create high-resolution image for clarity
        width, height = 1200, 1600
        img = Image.new('RGB', (width, height), color='white')
        draw = ImageDraw.Draw(img)
        
        # Try to load fonts (fallback to default if not available)
        try:
            font_large = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 36)
            font_medium = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 24)
            font_normal = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 18)
            font_small = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 14)
        except:
            try:
                font_large = ImageFont.truetype("/usr/share/fonts/TTF/DejaVuSans-Bold.ttf", 36)
                font_medium = ImageFont.truetype("/usr/share/fonts/TTF/DejaVuSans-Bold.ttf", 24)
                font_normal = ImageFont.truetype("/usr/share/fonts/TTF/DejaVuSans.ttf", 18)
                font_small = ImageFont.truetype("/usr/share/fonts/TTF/DejaVuSans.ttf", 14)
            except:
                font_large = ImageFont.load_default()
                font_medium = font_large
                font_normal = font_large
                font_small = font_large
        
        # Generate realistic data
        companies = ["Apex Technologies Inc.", "GlobalTech Solutions", "Premier Services LLC", "Quantum Dynamics Corp.", "Nexus Industries"]
        names = ["John Smith", "Sarah Johnson", "Michael Chen", "Emily Davis", "Robert Wilson"]
        
        company = random.choice(companies)
        customer = random.choice(names)
        invoice_num = f"INV-{datetime.now().year}-{random.randint(10000, 99999)}"
        date_str = datetime.now().strftime("%B %d, %Y")
        due_date = (datetime.now() + timedelta(days=30)).strftime("%B %d, %Y")
        
        # Generate items
        items = [
            ("Professional Consulting Services - Phase 1", random.randint(10, 40), random.uniform(150, 300)),
            ("Software Development - Custom Module", random.randint(20, 80), random.uniform(100, 200)),
            ("Project Management Services", random.randint(15, 30), random.uniform(125, 175)),
            ("Technical Documentation", random.randint(5, 15), random.uniform(75, 125)),
            ("Quality Assurance Testing", random.randint(10, 25), random.uniform(85, 150)),
        ]
        selected_items = random.sample(items, random.randint(3, 5))
        
        y = 50
        
        # Company Header
        draw.text((50, y), company, fill='#003366', font=font_large)
        y += 50
        
        draw.text((50, y), f"123 Business Park Drive, Suite {random.randint(100, 999)}", fill='#333333', font=font_small)
        y += 20
        draw.text((50, y), f"San Francisco, CA {random.randint(94100, 94199)}", fill='#333333', font=font_small)
        y += 20
        draw.text((50, y), f"Phone: (415) 555-{random.randint(1000, 9999)} | Email: billing@{company.lower().replace(' ', '').replace('.', '').replace(',', '')[:10]}.com", fill='#333333', font=font_small)
        y += 50
        
        # Invoice Title
        draw.text((width - 250, 50), "INVOICE", fill='#666666', font=font_large)
        draw.text((width - 250, 100), f"#{invoice_num}", fill='#333333', font=font_normal)
        draw.text((width - 250, 130), f"Date: {date_str}", fill='#333333', font=font_small)
        draw.text((width - 250, 150), f"Due: {due_date}", fill='#333333', font=font_small)
        
        # Separator line
        draw.line((50, y, width - 50, y), fill='#003366', width=2)
        y += 30
        
        # Bill To
        draw.text((50, y), "BILL TO:", fill='#003366', font=font_medium)
        y += 35
        draw.text((50, y), customer, fill='#333333', font=font_normal)
        y += 25
        draw.text((50, y), f"{random.randint(100, 9999)} {random.choice(['Oak Street', 'Main Avenue', 'Park Boulevard'])}", fill='#333333', font=font_small)
        y += 20
        draw.text((50, y), f"{random.choice(['Los Angeles', 'Chicago', 'Houston', 'Phoenix'])}, {random.choice(['CA', 'IL', 'TX', 'AZ'])} {random.randint(10000, 99999)}", fill='#333333', font=font_small)
        y += 50
        
        # Table Header
        draw.rectangle((50, y, width - 50, y + 35), fill='#003366')
        draw.text((60, y + 8), "Description", fill='white', font=font_normal)
        draw.text((650, y + 8), "Qty/Hours", fill='white', font=font_normal)
        draw.text((800, y + 8), "Rate", fill='white', font=font_normal)
        draw.text((950, y + 8), "Amount", fill='white', font=font_normal)
        y += 35
        
        # Table Items
        subtotal = 0
        for desc, qty, rate in selected_items:
            amount = qty * rate
            subtotal += amount
            
            # Alternate row colors
            if selected_items.index((desc, qty, rate)) % 2 == 0:
                draw.rectangle((50, y, width - 50, y + 30), fill='#f5f5f5')
            
            draw.text((60, y + 6), desc[:50], fill='#333333', font=font_small)
            draw.text((680, y + 6), str(qty), fill='#333333', font=font_small)
            draw.text((800, y + 6), f"${rate:.2f}", fill='#333333', font=font_small)
            draw.text((950, y + 6), f"${amount:,.2f}", fill='#333333', font=font_small)
            y += 30
        
        # Separator
        draw.line((700, y + 10, width - 50, y + 10), fill='#cccccc', width=1)
        y += 25
        
        # Totals
        tax_rate = 0.0825
        tax = subtotal * tax_rate
        total = subtotal + tax
        
        draw.text((800, y), "Subtotal:", fill='#333333', font=font_normal)
        draw.text((950, y), f"${subtotal:,.2f}", fill='#333333', font=font_normal)
        y += 30
        
        draw.text((800, y), f"Tax ({tax_rate*100:.2f}%):", fill='#333333', font=font_normal)
        draw.text((950, y), f"${tax:,.2f}", fill='#333333', font=font_normal)
        y += 30
        
        draw.rectangle((780, y - 5, width - 50, y + 30), fill='#003366')
        draw.text((800, y + 3), "TOTAL DUE:", fill='white', font=font_medium)
        draw.text((950, y + 3), f"${total:,.2f}", fill='white', font=font_medium)
        y += 60
        
        # Payment Info
        draw.text((50, y), "Payment Information:", fill='#003366', font=font_medium)
        y += 30
        draw.text((50, y), f"Bank: First National Bank | Account: ****{random.randint(1000, 9999)} | Routing: {random.randint(100000000, 999999999)}", fill='#333333', font=font_small)
        y += 25
        draw.text((50, y), "Payment Terms: Net 30 | Please include invoice number with payment", fill='#333333', font=font_small)
        y += 50
        
        # Footer
        draw.line((50, height - 80, width - 50, height - 80), fill='#cccccc', width=1)
        draw.text((width // 2 - 150, height - 60), "Thank you for your business!", fill='#666666', font=font_normal)
        draw.text((width // 2 - 200, height - 35), f"Questions? Contact us at billing@{company.lower().replace(' ', '')[:10]}.com", fill='#999999', font=font_small)
        
        # Save image
        filename = f"demo_{name.lower().replace(' ', '_')}_{uuid.uuid4().hex[:6]}.png"
        filepath = os.path.join(output_dir, filename)
        img.save(filepath, 'PNG', quality=95)
        
        print(f"[Demo Lab] Programmatic invoice image saved: {filepath}")
        
        return {
            "success": True,
            "url": f"/api/demo/files/{filename}",
            "filepath": filepath
        }
        
    except ImportError as ie:
        print(f"[Demo Lab] PIL not available: {ie}, falling back to DALL-E")
        return await generate_dalle_image(name, description, image_type, output_dir)
    except Exception as e:
        print(f"[Demo Lab] Programmatic image error: {e}")
        import traceback
        traceback.print_exc()
        return {"success": False, "error": str(e)}


async def generate_dalle_image(name: str, description: str, image_type: str, output_dir: str) -> Dict:
    """Generate image using OpenAI DALL-E for creative/visual images"""
    try:
        import httpx
        
        # Get OpenAI API key from settings
        openai_key = os.environ.get("OPENAI_API_KEY") or app_state.settings.get("openai_api_key")
        
        if not openai_key:
            return {"success": False, "error": "OpenAI API key not configured"}
        
        # Build a detailed prompt based on image type
        prompt_prefix = ""
        if image_type == "chart":
            prompt_prefix = "A professional business chart or graph showing "
        elif image_type == "logo":
            prompt_prefix = "A modern, professional company logo for "
        elif image_type == "product":
            prompt_prefix = "A professional product photo showing "
        elif image_type == "screenshot":
            prompt_prefix = "A clean UI screenshot showing "
        elif image_type == "diagram":
            prompt_prefix = "A professional diagram showing "
        elif image_type == "infographic":
            prompt_prefix = "A modern infographic showing "
        else:
            prompt_prefix = "A professional image showing "
        
        full_prompt = f"{prompt_prefix}{description}. High quality, professional, clean design, suitable for business presentations."
        
        # Call OpenAI DALL-E API
        async with httpx.AsyncClient(timeout=60.0) as client:
            response = await client.post(
                "https://api.openai.com/v1/images/generations",
                headers={
                    "Authorization": f"Bearer {openai_key}",
                    "Content-Type": "application/json"
                },
                json={
                    "model": "dall-e-3",
                    "prompt": full_prompt[:4000],  # DALL-E 3 limit
                    "n": 1,
                    "size": "1024x1024",
                    "quality": "standard",
                    "response_format": "url"
                }
            )
            
            if response.status_code != 200:
                error_data = response.json()
                return {"success": False, "error": error_data.get("error", {}).get("message", "API error")}
            
            data = response.json()
            image_url = data.get("data", [{}])[0].get("url", "")
            
            if not image_url:
                return {"success": False, "error": "No image URL returned"}
            
            # Download and save the image locally
            output_dir = os.environ.get("UPLOAD_PATH", "data/uploads")
            os.makedirs(output_dir, exist_ok=True)
            
            img_response = await client.get(image_url)
            if img_response.status_code == 200:
                filename = f"demo_{name.lower().replace(' ', '_')}_{uuid.uuid4().hex[:6]}.png"
                filepath = os.path.join(output_dir, filename)
                
                with open(filepath, 'wb') as f:
                    f.write(img_response.content)
                
                return {
                    "success": True,
                    "url": f"/api/demo/files/{filename}",
                    "filepath": filepath,
                    "original_url": image_url
                }
            else:
                # Return the OpenAI URL directly if download fails
                return {
                    "success": True,
                    "url": image_url,
                    "external": True
                }
    
    except httpx.TimeoutException:
        return {"success": False, "error": "Image generation timed out. Please try again."}
    except Exception as e:
        print(f"Image generation error: {e}")
        return {"success": False, "error": str(e)}


async def generate_document_content_with_llm(name: str, content_type: str, description: str) -> Optional[Dict]:
    """Use LLM to generate professional, relevant document content based on best practices"""
    try:
        llm_provider = app_state.get_llm_provider()
        if not llm_provider:
            print("[Demo Lab] No LLM provider available for document content generation")
            return None
        
        content_lower = content_type.lower() if content_type else ""
        
        # Build specialized prompts based on document type
        if "invoice" in content_lower:
            system_prompt = """You are creating a REALISTIC invoice for a demo/showcase. Generate professional, believable invoice data.

CRITICAL: Analyze the request context to determine the appropriate industry and items:
- "consulting invoice" -> Professional services, hourly rates ($150-300/hr), project phases
- "software invoice" -> Licenses, subscriptions, implementation fees, support packages
- "construction invoice" -> Materials, labor hours, equipment rental, permits
- "medical invoice" -> Procedures, consultations, lab tests, medications
- "restaurant invoice" -> Food items, beverages, service charges, gratuity
- "freelance invoice" -> Project deliverables, revisions, rush fees
- "retail invoice" -> Products with SKUs, quantities, discounts

Generate items that MATCH the context. If someone asks for a "consulting invoice", don't include physical products.

OUTPUT FORMAT (JSON only):
{
    "company": "Realistic company name matching the industry",
    "company_address": "123 Business District\\nSuite 400\\nSan Francisco, CA 94102",
    "company_phone": "+1 (415) 555-0123",
    "company_email": "billing@companyname.com",
    "invoice_number": "INV-2024-78432",
    "invoice_date": "December 15, 2024",
    "due_date": "January 14, 2025",
    "customer_name": "Realistic customer/client name",
    "customer_address": "Customer address\\nCity, State ZIP",
    "customer_email": "contact@customer.com",
    "items": [
        {"description": "Detailed, specific item description", "qty": 1, "price": 1500.00},
        {"description": "Another realistic item", "qty": 5, "price": 200.00}
    ],
    "subtotal": 2500.00,
    "tax_rate": "8.5%",
    "tax": 212.50,
    "total": 2712.50,
    "payment_terms": "Net 30 - 2% discount if paid within 10 days",
    "notes": "Professional thank you note with payment instructions"
}

RULES:
- Generate 4-8 line items that are SPECIFIC to the industry/context
- Use realistic pricing for the industry (not round numbers like $100, use $97.50, $1,247.00)
- Calculate subtotal, tax, and total CORRECTLY
- Item descriptions should be detailed: "Strategic Planning Workshop (8 hours)" not just "Consulting"
- Include realistic company and customer names
- ONLY output valid JSON
- Do NOT use Unicode characters"""""

        elif "policy" in content_lower or "hr" in content_lower:
            system_prompt = """You are an expert HR policy writer creating REAL, DETAILED policy documents for a demo/showcase.

CRITICAL: This is for a DEMO - generate REALISTIC, SPECIFIC content as if this were an actual company policy.
DO NOT write generic descriptions or placeholders. Write ACTUAL policy content with real numbers, dates, and procedures.

EXAMPLES OF WHAT TO INCLUDE:

For Vacation/PTO Policy:
- "Full-time employees accrue 15 days of paid vacation annually, increasing to 20 days after 3 years of service"
- "Vacation requests must be submitted at least 2 weeks in advance via the HR portal"
- "Maximum carryover of 5 unused days to the following year"

For Sick Leave:
- "Employees receive 10 paid sick days per calendar year"
- "A doctor's note is required for absences exceeding 3 consecutive days"

For Remote Work:
- "Employees may work remotely up to 3 days per week with manager approval"
- "Core hours are 10 AM - 3 PM in your local timezone"
- "Company provides $500 annual stipend for home office equipment"

For Performance Reviews:
- "Annual reviews conducted in Q4 with mid-year check-ins in June"
- "Rating scale: 1 (Needs Improvement) to 5 (Exceptional)"
- "Bonus eligibility requires minimum rating of 3"

OUTPUT FORMAT (JSON only):
{
    "company": "Generate a realistic company name based on context",
    "document_title": "Specific Policy Title",
    "effective_date": "January 1, 2024",
    "version": "2.1",
    "approved_by": "Sarah Mitchell, Chief Human Resources Officer",
    "sections": {
        "1. Purpose": "2-3 paragraphs with specific purpose...",
        "2. Scope": "Specific details on who this applies to...",
        "3. Policy Details": "THE MAIN CONTENT - multiple paragraphs with SPECIFIC numbers, timeframes, procedures...",
        "4. Eligibility": "Who qualifies and requirements...",
        "5. Procedures": "Step-by-step process with specific details...",
        "6. Manager Responsibilities": "What managers must do...",
        "7. Employee Responsibilities": "What employees must do...",
        "8. Exceptions": "How to request exceptions...",
        "9. Violations": "Consequences with specific examples...",
        "10. Contact Information": "HR contact with email and extension..."
    }
}

RULES:
- NEVER write "This section covers..." or "This policy addresses..." - write the ACTUAL content
- Include SPECIFIC numbers: days, percentages, dollar amounts, timeframes
- Include SPECIFIC procedures: "Submit Form HR-101 to your manager..."
- Include SPECIFIC dates and deadlines
- Make it feel like a REAL policy from a REAL company
- Each section should be 3-5 detailed paragraphs minimum
- Generate realistic names for approvers and contacts
- ONLY output valid JSON, no explanations
- Do NOT use Unicode characters - use "-" for bullets, regular quotes only"""

        elif "report" in content_lower or "analysis" in content_lower:
            system_prompt = """You are creating a REALISTIC business report for a demo/showcase. Generate professional, data-driven content.

CRITICAL: This is for a DEMO - include REAL-LOOKING data, statistics, and insights. Not placeholders!

Analyze the request to determine report type:
- "Sales Report" -> Revenue figures, growth rates, top products, regional breakdown
- "Financial Report" -> P&L data, cash flow, margins, YoY comparisons
- "Marketing Report" -> Campaign metrics, ROI, conversion rates, channel performance
- "HR Report" -> Headcount, turnover rates, hiring metrics, satisfaction scores
- "Project Report" -> Timeline, milestones, budget vs actual, risks
- "Performance Report" -> KPIs, targets vs actuals, trends, forecasts

EXAMPLE OF GOOD CONTENT:
"Q3 revenue reached $4.2M, representing a 23% increase over Q2. The Enterprise segment drove 67% of new bookings, with average deal size increasing from $45K to $58K. Customer acquisition cost decreased by 12% to $1,247 per customer, while lifetime value increased to $18,500."

OUTPUT FORMAT (JSON only):
{
    "company": "Realistic company name",
    "document_title": "Q3 2024 Sales Performance Report",
    "report_date": "October 15, 2024",
    "prepared_by": "Jennifer Chen, Director of Analytics",
    "prepared_for": "Executive Leadership Team",
    "sections": {
        "Executive Summary": "3-4 paragraphs with KEY NUMBERS and findings...",
        "Introduction": "Context and objectives of this report...",
        "Methodology": "Data sources and analysis approach...",
        "Key Findings": "5-7 major discoveries with SPECIFIC data points...",
        "Detailed Analysis": "Deep dive with percentages, comparisons, trends...",
        "Regional/Segment Breakdown": "Performance by region, product, or segment...",
        "Challenges and Risks": "Issues identified with impact assessment...",
        "Recommendations": "3-5 SPECIFIC, actionable recommendations...",
        "Implementation Roadmap": "Timeline and resources needed...",
        "Conclusion": "Summary and next steps with deadlines..."
    },
    "key_metrics": [
        {"metric": "Total Revenue", "value": "$4.2M", "trend": "up", "change": "+23%", "insight": "Exceeded target by $400K"},
        {"metric": "Customer Acquisition Cost", "value": "$1,247", "trend": "down", "change": "-12%", "insight": "Improved marketing efficiency"},
        {"metric": "Net Promoter Score", "value": "67", "trend": "up", "change": "+8 pts", "insight": "Product improvements driving satisfaction"}
    ]
}

RULES:
- Include SPECIFIC numbers throughout: $4.2M, 23%, 1,247 customers
- Use realistic business metrics and industry benchmarks
- Each section should be 2-4 substantial paragraphs
- Data should be internally consistent (totals should add up)
- Include YoY, QoQ, or MoM comparisons where relevant
- ONLY output valid JSON
- Do NOT use Unicode characters"""

        elif "contract" in content_lower or "agreement" in content_lower:
            system_prompt = """You are a legal document specialist creating professional contracts.

OUTPUT FORMAT (JSON only):
{
    "document_title": "Master Services Agreement",
    "effective_date": "January 1, 2025",
    "parties": {
        "party_a": {"name": "Nexus Technologies Inc.", "address": "500 Innovation Drive, Suite 300\\nAustin, TX 78701", "represented_by": "Michael Torres, Chief Executive Officer"},
        "party_b": {"name": "GlobalTech Solutions LLC", "address": "1200 Corporate Center\\nDenver, CO 80202", "represented_by": "Amanda Foster, VP of Operations"}
    },
    "sections": {
        "1. Definitions": "Define ALL key terms with specifics: 'Services' means..., 'Deliverables' means..., 'Confidential Information' includes...",
        "2. Scope of Work": "DETAILED description of services, deliverables, timelines, and acceptance criteria",
        "3. Term and Termination": "Initial term of 24 months, auto-renewal provisions, 90-day termination notice, termination for cause conditions",
        "4. Compensation": "SPECIFIC amounts: $150,000 base fee, payment schedule, expense reimbursement up to $10,000/month",
        "5. Confidentiality": "5-year obligation, specific exclusions, return of materials within 30 days",
        "6. Intellectual Property": "Work product ownership, license grants, pre-existing IP carve-outs",
        "7. Warranties": "Specific performance warranties, compliance warranties, 30-day cure period",
        "8. Limitation of Liability": "Cap at 12 months of fees paid, exclusions for gross negligence",
        "9. Indemnification": "Mutual indemnification, IP infringement coverage, third-party claims",
        "10. Insurance": "Required coverage: $2M general liability, $5M professional liability, $1M cyber",
        "11. Dispute Resolution": "30-day negotiation, mediation in Austin TX, binding arbitration under AAA rules",
        "12. General Provisions": "Governing law (Texas), assignment restrictions, force majeure, entire agreement",
        "13. Signatures": "Signature blocks with printed name, title, date lines"
    }
}

RULES:
- Use ACTUAL legal language, not descriptions of what should be there
- Include SPECIFIC dollar amounts, timeframes, and obligations
- Reference specific laws or standards where appropriate
- Make it read like a REAL contract that could be signed
- Each section should be legally substantive
- ONLY output valid JSON
- Do NOT use Unicode characters"""

        else:
            # Generic document - but still make it realistic!
            system_prompt = """You are creating a REALISTIC professional document for a demo/showcase.

CRITICAL: Generate ACTUAL content, not placeholders or descriptions. This should look like a real document.

Analyze the request to determine document type and generate appropriate content:
- Proposals: Include pricing, timelines, deliverables, terms
- Memos: Include clear purpose, background, action items with deadlines
- Procedures: Include step-by-step instructions with specifics
- Guidelines: Include do's and don'ts with examples
- Plans: Include objectives, milestones, resources, metrics

OUTPUT FORMAT (JSON only):
{
    "company": "Realistic company name",
    "document_title": "Specific, descriptive title",
    "date": "December 20, 2024",
    "author": "David Park, Senior Manager",
    "sections": {
        "Purpose/Overview": "Clear statement of document purpose with context...",
        "Background": "Relevant history and context...",
        "Main Content Section 1": "Detailed, specific content...",
        "Main Content Section 2": "More detailed content...",
        "Implementation/Next Steps": "Specific action items with owners and deadlines...",
        "Appendix/References": "Supporting information..."
    }
}

RULES:
- Create 5-10 relevant sections based on document type
- Each section should have REAL, detailed content (2-4 paragraphs minimum)
- Include specific names, dates, numbers, and details
- Make it feel like an actual business document
- ONLY output valid JSON
- Do NOT use Unicode characters"""""

        user_prompt = f"""Create a REALISTIC, DETAILED document for a DEMO/SHOWCASE:

Document Name: {name}
Document Type: {content_type}
Context: {description}

IMPORTANT INSTRUCTIONS:
1. This is for a DEMO - generate content that looks 100% REAL and PROFESSIONAL
2. Include SPECIFIC numbers, dates, percentages, dollar amounts
3. Use realistic company names, employee names, and contact info
4. DO NOT write descriptions of what content should be - write the ACTUAL content
5. Each section should be substantial (3-5 paragraphs with real details)
6. Make it indistinguishable from a real business document

Return ONLY valid JSON with complete, realistic content."""

        messages = [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt}
        ]
        
        print(f"[Demo Lab] Generating document content with LLM for: {name} ({content_type})")
        # Increased max_tokens for more detailed content
        response = await llm_provider.generate(messages, max_tokens=4000)
        
        # Parse JSON from response
        import re
        # Remove markdown code blocks if present
        response = re.sub(r'```json\s*', '', response)
        response = re.sub(r'```\s*', '', response)
        
        json_match = re.search(r'\{[\s\S]*\}', response)
        if json_match:
            try:
                data = json.loads(json_match.group())
                print(f"[Demo Lab] LLM generated document content successfully")
                return data
            except json.JSONDecodeError as e:
                print(f"[Demo Lab] JSON parse error: {e}")
                return None
        
        print("[Demo Lab] No valid JSON found in LLM response")
        return None
        
    except Exception as e:
        print(f"[Demo Lab] LLM document content generation error: {e}")
        return None


async def generate_demo_document(name: str, doc_type: str, content_type: str, sample_data: Dict, description: str) -> Optional[Dict]:
    """Generate professional, realistic sample documents using LLM for content"""
    import random
    from datetime import datetime, timedelta
    
    try:
        output_dir = os.environ.get("UPLOAD_PATH", "data/uploads")
        print(f"[Demo Lab] generate_demo_document called: name={name}, doc_type={doc_type}, content_type={content_type}")
        print(f"[Demo Lab] Output directory: {output_dir}")
        os.makedirs(output_dir, exist_ok=True)
        
        filename = f"demo_{name.lower().replace(' ', '_')}_{uuid.uuid4().hex[:6]}"
        content_lower = content_type.lower() if content_type else ""
        print(f"[Demo Lab] Filename base: {filename}, content_lower: {content_lower}")
        
        # ============================================================
        # ALWAYS USE LLM-POWERED CONTENT GENERATION
        # ============================================================
        print("[Demo Lab] Using LLM to generate professional document content...")
        llm_content = await generate_document_content_with_llm(name, content_type, description)
        if llm_content:
            sample_data = llm_content
            print("[Demo Lab] ✅ Using LLM-generated content")
        else:
            print("[Demo Lab] ⚠️ LLM content generation failed, falling back to templates")
        
        # ============================================================
        # FALLBACK: REALISTIC DATA GENERATORS (only if LLM fails)
        # ============================================================
        
        first_names = ["Ahmed", "Sarah", "Mohamed", "John", "Emily", "Michael", "Jessica", "David", "Emma", "James"]
        last_names = ["Hassan", "Ali", "Smith", "Johnson", "Williams", "Brown", "Jones", "Garcia", "Miller"]
        companies = ["TechCorp Solutions", "GlobalBank International", "Sunrise Electronics", "Blue Ocean Trading", "Apex Industries", "Pinnacle Holdings"]
        
        def random_date(days_back=30):
            return (datetime.now() - timedelta(days=random.randint(0, days_back))).strftime("%B %d, %Y")
        
        def random_phone():
            return f"+1 ({random.randint(200,999)}) {random.randint(100,999)}-{random.randint(1000,9999)}"
        
        def random_email(name):
            return f"{name.lower().replace(' ', '.')}@{random.choice(['company.com', 'corp.com', 'business.org'])}"
        
        # Only use fallback templates if sample_data is still empty
        if not sample_data or len(sample_data) == 0:
            print("[Demo Lab] Using fallback template data...")
            # ============================================================
            # INVOICE CONTENT (Fallback)
            # ============================================================
            if "invoice" in content_lower:
                company_name = random.choice(companies)
                customer_name = f"{random.choice(first_names)} {random.choice(last_names)}"
                invoice_num = f"INV-{datetime.now().year}-{random.randint(10000, 99999)}"
                
                items = [
                    {"description": "Apple MacBook Pro 14\" M3 Pro", "qty": 1, "price": 1999.00},
                    {"description": "Apple Magic Keyboard", "qty": 2, "price": 99.00},
                    {"description": "Dell 27\" 4K Monitor", "qty": 1, "price": 449.00},
                    {"description": "Logitech MX Master 3S Mouse", "qty": 2, "price": 99.99},
                    {"description": "Professional Setup & Configuration", "qty": 1, "price": 150.00},
                ]
                selected_items = random.sample(items, random.randint(2, 4))
                subtotal = sum(i["qty"] * i["price"] for i in selected_items)
                tax = subtotal * 0.08
                total = subtotal + tax
                
                invoice_data = {
                    "company": company_name,
                    "company_address": f"{random.randint(100, 9999)} Business Park Drive\nSuite {random.randint(100, 999)}\nNew York, NY {random.randint(10000, 10999)}",
                    "company_phone": random_phone(),
                    "company_email": f"billing@{company_name.lower().replace(' ', '')}.com",
                    "invoice_number": invoice_num,
                    "invoice_date": random_date(5),
                    "due_date": random_date(-30),
                    "customer_name": customer_name,
                    "customer_address": f"{random.randint(100, 9999)} {random.choice(['Main St', 'Oak Ave', 'Park Blvd'])}\n{random.choice(['Los Angeles', 'Chicago', 'Houston'])}, {random.choice(['CA', 'IL', 'TX'])} {random.randint(10000, 99999)}",
                    "customer_email": random_email(customer_name),
                    "items": selected_items,
                    "subtotal": subtotal,
                    "tax_rate": "8%",
                    "tax": tax,
                    "total": total,
                    "payment_terms": "Net 30",
                    "notes": "Thank you for your business!"
                }
                sample_data = invoice_data
            
            # ============================================================
            # HR POLICY CONTENT
            # ============================================================
            elif "policy" in content_lower or "hr" in content_lower:
                company_name = random.choice(companies)
                policy_data = {
                    "company": company_name,
                    "document_title": name,
                    "effective_date": random_date(60),
                    "version": f"{random.randint(1,3)}.{random.randint(0,9)}",
                    "approved_by": f"{random.choice(first_names)} {random.choice(last_names)}, Chief Human Resources Officer",
                    "sections": {
                        "1. Purpose": f"This policy establishes guidelines for {name.lower()} at {company_name}. It applies to all employees, contractors, and temporary staff across all departments and locations.",
                        
                        "2. Scope": f"This policy applies to:\n- All full-time and part-time employees\n- Contractors and consultants\n- Temporary and seasonal workers\n- Interns and trainees\n- All company locations worldwide",
                        
                        "3. Policy Statement": f"{company_name} is committed to maintaining a professional, productive, and respectful workplace. We believe that clear policies help create an environment where all team members can thrive and contribute to our shared success.\n\nAll employees are expected to:\n- Act with integrity and professionalism\n- Treat colleagues with respect and dignity\n- Comply with all applicable laws and regulations\n- Report any violations or concerns promptly",
                        
                        "4. Procedures": "4.1 Implementation\nDepartment managers are responsible for ensuring all team members understand and comply with this policy.\n\n4.2 Training\nAll new employees will receive policy training during onboarding. Annual refresher training is mandatory for all staff.\n\n4.3 Documentation\nAll policy-related documentation must be maintained for a minimum of seven (7) years.",
                        
                        "5. Responsibilities": f"5.1 Human Resources\n- Maintain and update this policy\n- Provide training and guidance\n- Investigate reported violations\n\n5.2 Managers\n- Communicate policy to team members\n- Ensure compliance within their departments\n- Report violations to HR\n\n5.3 Employees\n- Read and understand this policy\n- Comply with all requirements\n- Report violations or concerns",
                        
                        "6. Violations": "Violations of this policy may result in disciplinary action, up to and including termination of employment. The severity of disciplinary action will depend on the nature and circumstances of the violation.",
                        
                        "7. Questions": f"For questions about this policy, please contact:\n\nHuman Resources Department\nEmail: hr@{company_name.lower().replace(' ', '')}.com\nPhone: {random_phone()}\nInternal Extension: {random.randint(1000, 9999)}"
                    }
                }
                sample_data = policy_data
            
            # ============================================================
            # EMPLOYEE REPORT CONTENT
            # ============================================================
            elif "employee" in content_lower or "report" in content_lower:
                employees = []
                departments = ["Engineering", "Sales", "Marketing", "Finance", "Human Resources", "Operations"]
                for i in range(10):
                    fname, lname = random.choice(first_names), random.choice(last_names)
                    employees.append({
                        "id": f"EMP-{random.randint(10000, 99999)}",
                        "name": f"{fname} {lname}",
                        "department": random.choice(departments),
                        "position": random.choice(["Manager", "Senior Analyst", "Specialist", "Coordinator", "Director"]),
                        "hire_date": random_date(1000),
                        "salary": f"${random.randint(55, 180),000}",
                        "status": "Active"
                    })
                sample_data = {"employees": employees, "total": len(employees), "report_date": random_date(0)}
        
        # ============================================================
        # GENERATE WORD DOCUMENT
        # ============================================================
        if doc_type == "word":
            try:
                from docx import Document
                from docx.shared import Inches, Pt, RGBColor
                from docx.enum.text import WD_ALIGN_PARAGRAPH
                from docx.enum.table import WD_TABLE_ALIGNMENT
                
                doc = Document()
                
                # Set default font
                style = doc.styles['Normal']
                style.font.name = 'Calibri'
                style.font.size = Pt(11)
                
                if "invoice" in content_lower and sample_data:
                    # === PROFESSIONAL INVOICE ===
                    # Company Header
                    header = doc.add_paragraph()
                    header_run = header.add_run(sample_data.get("company", "Company Name"))
                    header_run.bold = True
                    header_run.font.size = Pt(24)
                    header_run.font.color.rgb = RGBColor(0, 51, 102)
                    header.alignment = WD_ALIGN_PARAGRAPH.LEFT
                    
                    # Company Address
                    addr = doc.add_paragraph(sample_data.get("company_address", ""))
                    addr.add_run(f"\nPhone: {sample_data.get('company_phone', '')}")
                    addr.add_run(f"\nEmail: {sample_data.get('company_email', '')}")
                    
                    # Invoice Title
                    doc.add_paragraph()
                    title = doc.add_paragraph()
                    title_run = title.add_run("INVOICE")
                    title_run.bold = True
                    title_run.font.size = Pt(28)
                    title_run.font.color.rgb = RGBColor(128, 128, 128)
                    title.alignment = WD_ALIGN_PARAGRAPH.RIGHT
                    
                    # Invoice Details Table
                    details_table = doc.add_table(rows=3, cols=2)
                    details_table.alignment = WD_TABLE_ALIGNMENT.RIGHT
                    details_table.cell(0, 0).text = "Invoice Number:"
                    details_table.cell(0, 1).text = sample_data.get("invoice_number", "")
                    details_table.cell(1, 0).text = "Invoice Date:"
                    details_table.cell(1, 1).text = sample_data.get("invoice_date", "")
                    details_table.cell(2, 0).text = "Due Date:"
                    details_table.cell(2, 1).text = sample_data.get("due_date", "")
                    
                    doc.add_paragraph()
                    
                    # Bill To
                    bill_to = doc.add_paragraph()
                    bill_to_run = bill_to.add_run("BILL TO:")
                    bill_to_run.bold = True
                    bill_to_run.font.size = Pt(10)
                    
                    customer = doc.add_paragraph()
                    cust_name = customer.add_run(sample_data.get("customer_name", ""))
                    cust_name.bold = True
                    customer.add_run(f"\n{sample_data.get('customer_address', '')}")
                    customer.add_run(f"\nEmail: {sample_data.get('customer_email', '')}")
                    
                    doc.add_paragraph()
                    
                    # Items Table
                    items = sample_data.get("items", [])
                    table = doc.add_table(rows=1 + len(items) + 4, cols=4)
                    table.style = 'Table Grid'
                    
                    # Header row
                    header_cells = table.rows[0].cells
                    headers = ["Description", "Quantity", "Unit Price", "Amount"]
                    for i, h in enumerate(headers):
                        header_cells[i].text = h
                        for paragraph in header_cells[i].paragraphs:
                            for run in paragraph.runs:
                                run.bold = True
                    
                    # Item rows
                    for idx, item in enumerate(items):
                        row = table.rows[idx + 1].cells
                        row[0].text = item["description"]
                        row[1].text = str(item["qty"])
                        row[2].text = f"${item['price']:,.2f}"
                        row[3].text = f"${item['qty'] * item['price']:,.2f}"
                    
                    # Totals
                    subtotal_row = len(items) + 1
                    table.rows[subtotal_row].cells[2].text = "Subtotal:"
                    table.rows[subtotal_row].cells[3].text = f"${sample_data.get('subtotal', 0):,.2f}"
                    
                    table.rows[subtotal_row + 1].cells[2].text = f"Tax ({sample_data.get('tax_rate', '8%')}):"
                    table.rows[subtotal_row + 1].cells[3].text = f"${sample_data.get('tax', 0):,.2f}"
                    
                    table.rows[subtotal_row + 2].cells[2].text = "TOTAL:"
                    total_cell = table.rows[subtotal_row + 2].cells[3]
                    total_cell.text = f"${sample_data.get('total', 0):,.2f}"
                    for paragraph in total_cell.paragraphs:
                        for run in paragraph.runs:
                            run.bold = True
                    
                    # Payment Terms & Notes
                    doc.add_paragraph()
                    terms = doc.add_paragraph()
                    terms_run = terms.add_run(f"Payment Terms: {sample_data.get('payment_terms', 'Net 30')}")
                    terms_run.bold = True
                    
                    doc.add_paragraph()
                    notes = doc.add_paragraph(sample_data.get("notes", "Thank you for your business!"))
                    notes.alignment = WD_ALIGN_PARAGRAPH.CENTER
                
                elif "policy" in content_lower and sample_data:
                    # === PROFESSIONAL HR POLICY ===
                    # Document Header
                    header = doc.add_paragraph()
                    header_run = header.add_run(sample_data.get("company", "Company"))
                    header_run.bold = True
                    header_run.font.size = Pt(14)
                    header.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    
                    # Document Title
                    title = doc.add_heading(sample_data.get("document_title", name), 0)
                    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    
                    # Document Info Table
                    info_table = doc.add_table(rows=3, cols=2)
                    info_table.cell(0, 0).text = "Effective Date:"
                    info_table.cell(0, 1).text = sample_data.get("effective_date", "")
                    info_table.cell(1, 0).text = "Version:"
                    info_table.cell(1, 1).text = sample_data.get("version", "1.0")
                    info_table.cell(2, 0).text = "Approved By:"
                    info_table.cell(2, 1).text = sample_data.get("approved_by", "")
                    
                    doc.add_paragraph()
                    doc.add_paragraph("_" * 60)
                    doc.add_paragraph()
                    
                    # Policy Sections
                    sections = sample_data.get("sections", {})
                    for section_title, section_content in sections.items():
                        heading = doc.add_heading(section_title, level=1)
                        content_para = doc.add_paragraph(section_content)
                        doc.add_paragraph()
                    
                    # Footer
                    doc.add_paragraph("_" * 60)
                    footer = doc.add_paragraph()
                    footer_text = footer.add_run(f"\nThis document is the property of {sample_data.get('company', 'the Company')}. Unauthorized distribution is prohibited.")
                    footer_text.italic = True
                    footer_text.font.size = Pt(9)
                    footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
                
                else:
                    # === GENERAL PROFESSIONAL DOCUMENT ===
                    # Title
                    title = doc.add_heading(name, 0)
                    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    
                    # Metadata
                    meta = doc.add_paragraph()
                    meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    meta.add_run(f"Generated: {datetime.utcnow().strftime('%B %d, %Y')}\n").italic = True
                    meta.add_run(f"Document ID: DOC-{uuid.uuid4().hex[:8].upper()}").italic = True
                    
                    doc.add_paragraph()
                    
                    if description:
                        intro = doc.add_heading("Executive Summary", level=1)
                        doc.add_paragraph(description)
                        doc.add_paragraph()
                    
                    if sample_data:
                        for key, value in sample_data.items():
                            if key not in ["company", "document_title"]:
                                doc.add_heading(str(key).replace("_", " ").title(), level=1)
                                if isinstance(value, list):
                                    for item in value:
                                        if isinstance(item, dict):
                                            for k, v in item.items():
                                                p = doc.add_paragraph()
                                                p.add_run(f"{k}: ").bold = True
                                                p.add_run(str(v))
                                            doc.add_paragraph()
                                        else:
                                            doc.add_paragraph(f"• {item}")
                                elif isinstance(value, dict):
                                    for k, v in value.items():
                                        p = doc.add_paragraph()
                                        p.add_run(f"{k}: ").bold = True
                                        p.add_run(str(v))
                                else:
                                    doc.add_paragraph(str(value))
                                doc.add_paragraph()
                
                filepath = os.path.join(output_dir, f"{filename}.docx")
                doc.save(filepath)
                return {"url": f"/api/demo/files/{filename}.docx", "filepath": filepath}
                
            except ImportError:
                return await generate_demo_document(name, "text", content_type, sample_data, description)
        
        # ============================================================
        # GENERATE EXCEL SPREADSHEET
        # ============================================================
        elif doc_type == "excel":
            try:
                from openpyxl import Workbook
                from openpyxl.styles import Font, Alignment, Border, Side, PatternFill
                from openpyxl.utils import get_column_letter
                
                wb = Workbook()
                ws = wb.active
                ws.title = name[:30]
                
                # Styles
                header_font = Font(bold=True, size=12, color="FFFFFF")
                header_fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
                border = Border(
                    left=Side(style='thin'),
                    right=Side(style='thin'),
                    top=Side(style='thin'),
                    bottom=Side(style='thin')
                )
                
                if "invoice" in content_lower and sample_data:
                    # Invoice Spreadsheet
                    ws['A1'] = sample_data.get("company", "Company")
                    ws['A1'].font = Font(bold=True, size=18)
                    ws.merge_cells('A1:D1')
                    
                    ws['A3'] = "INVOICE"
                    ws['A3'].font = Font(bold=True, size=14)
                    
                    ws['A5'] = "Invoice #:"
                    ws['B5'] = sample_data.get("invoice_number", "")
                    ws['A6'] = "Date:"
                    ws['B6'] = sample_data.get("invoice_date", "")
                    ws['A7'] = "Due Date:"
                    ws['B7'] = sample_data.get("due_date", "")
                    
                    ws['A9'] = "Bill To:"
                    ws['A9'].font = Font(bold=True)
                    ws['A10'] = sample_data.get("customer_name", "")
                    
                    # Items header
                    headers = ["Description", "Qty", "Unit Price", "Amount"]
                    for col, header in enumerate(headers, 1):
                        cell = ws.cell(row=12, column=col, value=header)
                        cell.font = header_font
                        cell.fill = header_fill
                        cell.alignment = Alignment(horizontal='center')
                        cell.border = border
                    
                    # Items
                    items = sample_data.get("items", [])
                    for idx, item in enumerate(items, 13):
                        ws.cell(row=idx, column=1, value=item["description"]).border = border
                        ws.cell(row=idx, column=2, value=item["qty"]).border = border
                        ws.cell(row=idx, column=3, value=f"${item['price']:,.2f}").border = border
                        ws.cell(row=idx, column=4, value=f"${item['qty'] * item['price']:,.2f}").border = border
                    
                    # Totals
                    total_row = 13 + len(items) + 1
                    ws.cell(row=total_row, column=3, value="Subtotal:").font = Font(bold=True)
                    ws.cell(row=total_row, column=4, value=f"${sample_data.get('subtotal', 0):,.2f}")
                    ws.cell(row=total_row + 1, column=3, value="Tax (8%):").font = Font(bold=True)
                    ws.cell(row=total_row + 1, column=4, value=f"${sample_data.get('tax', 0):,.2f}")
                    ws.cell(row=total_row + 2, column=3, value="TOTAL:").font = Font(bold=True, size=12)
                    ws.cell(row=total_row + 2, column=4, value=f"${sample_data.get('total', 0):,.2f}").font = Font(bold=True, size=12)
                    
                    # Adjust column widths
                    ws.column_dimensions['A'].width = 40
                    ws.column_dimensions['B'].width = 10
                    ws.column_dimensions['C'].width = 15
                    ws.column_dimensions['D'].width = 15
                
                elif "employee" in content_lower and sample_data:
                    # Employee Report Spreadsheet
                    ws['A1'] = "Employee Report"
                    ws['A1'].font = Font(bold=True, size=16)
                    ws.merge_cells('A1:G1')
                    
                    ws['A2'] = f"Generated: {sample_data.get('report_date', datetime.now().strftime('%B %d, %Y'))}"
                    
                    # Headers
                    headers = ["Employee ID", "Name", "Department", "Position", "Hire Date", "Salary", "Status"]
                    for col, header in enumerate(headers, 1):
                        cell = ws.cell(row=4, column=col, value=header)
                        cell.font = header_font
                        cell.fill = header_fill
                        cell.alignment = Alignment(horizontal='center')
                        cell.border = border
                    
                    # Employee data
                    employees = sample_data.get("employees", [])
                    for idx, emp in enumerate(employees, 5):
                        ws.cell(row=idx, column=1, value=emp.get("id", "")).border = border
                        ws.cell(row=idx, column=2, value=emp.get("name", "")).border = border
                        ws.cell(row=idx, column=3, value=emp.get("department", "")).border = border
                        ws.cell(row=idx, column=4, value=emp.get("position", "")).border = border
                        ws.cell(row=idx, column=5, value=emp.get("hire_date", "")).border = border
                        ws.cell(row=idx, column=6, value=emp.get("salary", "")).border = border
                        ws.cell(row=idx, column=7, value=emp.get("status", "")).border = border
                    
                    # Adjust widths
                    widths = [15, 25, 15, 20, 15, 12, 10]
                    for i, w in enumerate(widths, 1):
                        ws.column_dimensions[get_column_letter(i)].width = w
                
                else:
                    # Generic Spreadsheet
                    ws['A1'] = name
                    ws['A1'].font = Font(bold=True, size=14)
                    ws['A2'] = f"Generated: {datetime.now().strftime('%B %d, %Y')}"
                    
                    if sample_data:
                        row = 4
                        for key, value in sample_data.items():
                            ws.cell(row=row, column=1, value=str(key).title())
                            ws.cell(row=row, column=1).font = Font(bold=True)
                            ws.cell(row=row, column=2, value=str(value) if not isinstance(value, (list, dict)) else json.dumps(value))
                            row += 1
                
                filepath = os.path.join(output_dir, f"{filename}.xlsx")
                wb.save(filepath)
                return {"url": f"/api/demo/files/{filename}.xlsx", "filepath": filepath}
                
            except ImportError:
                return await generate_demo_document(name, "text", content_type, sample_data, description)
        
        # ============================================================
        # GENERATE PDF
        # ============================================================
        elif doc_type == "pdf":
            print(f"[Demo Lab] Generating PDF document")
            try:
                from fpdf import FPDF
                print(f"[Demo Lab] fpdf imported successfully")
                
                # Helper function to clean text for PDF (replace Unicode with ASCII)
                def clean_text(text):
                    if not isinstance(text, str):
                        text = str(text)
                    # Replace Unicode characters with ASCII equivalents
                    replacements = {
                        '•': '-',
                        '–': '-',
                        '—': '-',
                        '"': '"',
                        '"': '"',
                        ''': "'",
                        ''': "'",
                        '…': '...',
                        '→': '->',
                        '←': '<-',
                        '✓': '[x]',
                        '✗': '[ ]',
                        '★': '*',
                        '☆': '*',
                        '\u2022': '-',  # bullet point
                        '\u2013': '-',  # en dash
                        '\u2014': '-',  # em dash
                        '\u201c': '"',  # left double quote
                        '\u201d': '"',  # right double quote
                        '\u2018': "'",  # left single quote
                        '\u2019': "'",  # right single quote
                    }
                    for old, new in replacements.items():
                        text = text.replace(old, new)
                    # Remove any remaining non-latin1 characters
                    text = text.encode('latin-1', errors='replace').decode('latin-1')
                    return text
                
                class PDF(FPDF):
                    def header(self):
                        if hasattr(self, 'doc_company') and self.doc_company:
                            self.set_font('Helvetica', 'B', 10)
                            self.cell(0, 10, self.doc_company, 0, 1, 'R')
                    
                    def footer(self):
                        self.set_y(-15)
                        self.set_font('Helvetica', 'I', 8)
                        self.cell(0, 10, f'Page {self.page_no()}', 0, 0, 'C')
                
                pdf = PDF()
                pdf.doc_company = clean_text(sample_data.get("company", "")) if sample_data else ""
                pdf.add_page()
                print(f"[Demo Lab] PDF initialized, content_lower={content_lower}, sample_data keys={list(sample_data.keys()) if sample_data else 'None'}")
                
                if "invoice" in content_lower and sample_data:
                    # Professional Invoice PDF
                    pdf.set_font("Helvetica", "B", 24)
                    pdf.set_text_color(0, 51, 102)
                    pdf.cell(0, 15, clean_text(sample_data.get("company", "Company")), ln=True)
                    
                    pdf.set_font("Helvetica", "", 10)
                    pdf.set_text_color(0, 0, 0)
                    pdf.multi_cell(0, 5, clean_text(sample_data.get("company_address", "")))
                    pdf.cell(0, 5, clean_text(f"Phone: {sample_data.get('company_phone', '')}"), ln=True)
                    pdf.cell(0, 5, clean_text(f"Email: {sample_data.get('company_email', '')}"), ln=True)
                    
                    pdf.ln(10)
                    pdf.set_font("Helvetica", "B", 28)
                    pdf.set_text_color(128, 128, 128)
                    pdf.cell(0, 15, "INVOICE", ln=True, align="R")
                    
                    pdf.set_text_color(0, 0, 0)
                    pdf.set_font("Helvetica", "", 10)
                    pdf.cell(0, 6, clean_text(f"Invoice #: {sample_data.get('invoice_number', '')}"), ln=True, align="R")
                    pdf.cell(0, 6, clean_text(f"Date: {sample_data.get('invoice_date', '')}"), ln=True, align="R")
                    pdf.cell(0, 6, clean_text(f"Due: {sample_data.get('due_date', '')}"), ln=True, align="R")
                    
                    pdf.ln(10)
                    pdf.set_font("Helvetica", "B", 10)
                    pdf.cell(0, 6, "BILL TO:", ln=True)
                    pdf.set_font("Helvetica", "B", 11)
                    pdf.cell(0, 6, clean_text(sample_data.get("customer_name", "")), ln=True)
                    pdf.set_font("Helvetica", "", 10)
                    pdf.multi_cell(0, 5, clean_text(sample_data.get("customer_address", "")))
                    
                    pdf.ln(10)
                    
                    # Items table
                    pdf.set_fill_color(0, 51, 102)
                    pdf.set_text_color(255, 255, 255)
                    pdf.set_font("Helvetica", "B", 10)
                    pdf.cell(90, 8, "Description", 1, 0, "L", True)
                    pdf.cell(20, 8, "Qty", 1, 0, "C", True)
                    pdf.cell(35, 8, "Unit Price", 1, 0, "R", True)
                    pdf.cell(35, 8, "Amount", 1, 1, "R", True)
                    
                    pdf.set_text_color(0, 0, 0)
                    pdf.set_font("Helvetica", "", 10)
                    
                    items = sample_data.get("items", [])
                    for item in items:
                        pdf.cell(90, 7, clean_text(item["description"][:45]), 1, 0, "L")
                        pdf.cell(20, 7, str(item["qty"]), 1, 0, "C")
                        pdf.cell(35, 7, f"${item['price']:,.2f}", 1, 0, "R")
                        pdf.cell(35, 7, f"${item['qty'] * item['price']:,.2f}", 1, 1, "R")
                    
                    pdf.ln(5)
                    pdf.set_font("Helvetica", "", 10)
                    pdf.cell(145, 7, "Subtotal:", 0, 0, "R")
                    pdf.cell(35, 7, f"${sample_data.get('subtotal', 0):,.2f}", 0, 1, "R")
                    pdf.cell(145, 7, f"Tax ({sample_data.get('tax_rate', '8%')}):", 0, 0, "R")
                    pdf.cell(35, 7, f"${sample_data.get('tax', 0):,.2f}", 0, 1, "R")
                    
                    pdf.set_font("Helvetica", "B", 12)
                    pdf.cell(145, 10, "TOTAL:", 0, 0, "R")
                    pdf.cell(35, 10, f"${sample_data.get('total', 0):,.2f}", 0, 1, "R")
                    
                    pdf.ln(15)
                    pdf.set_font("Helvetica", "B", 10)
                    pdf.cell(0, 6, clean_text(f"Payment Terms: {sample_data.get('payment_terms', 'Net 30')}"), ln=True)
                    pdf.ln(5)
                    pdf.set_font("Helvetica", "I", 10)
                    pdf.cell(0, 6, clean_text(sample_data.get("notes", "Thank you for your business!")), ln=True, align="C")
                
                elif "policy" in content_lower and sample_data:
                    # Professional Policy PDF
                    pdf.set_font("Helvetica", "B", 12)
                    pdf.cell(0, 10, clean_text(sample_data.get("company", "")), ln=True, align="C")
                    
                    pdf.set_font("Helvetica", "B", 20)
                    pdf.cell(0, 15, clean_text(sample_data.get("document_title", name)), ln=True, align="C")
                    
                    pdf.set_font("Helvetica", "", 10)
                    pdf.cell(0, 6, clean_text(f"Effective Date: {sample_data.get('effective_date', '')}"), ln=True, align="C")
                    pdf.cell(0, 6, clean_text(f"Version: {sample_data.get('version', '1.0')}"), ln=True, align="C")
                    
                    pdf.ln(5)
                    pdf.cell(0, 0, "", "T", 1)
                    pdf.ln(10)
                    
                    sections = sample_data.get("sections", {})
                    for section_title, section_content in sections.items():
                        pdf.set_font("Helvetica", "B", 12)
                        pdf.set_text_color(0, 51, 102)
                        pdf.cell(0, 8, clean_text(section_title), ln=True)
                        
                        pdf.set_font("Helvetica", "", 10)
                        pdf.set_text_color(0, 0, 0)
                        pdf.multi_cell(0, 5, clean_text(section_content))
                        pdf.ln(5)
                    
                    pdf.ln(10)
                    pdf.cell(0, 0, "", "T", 1)
                    pdf.ln(5)
                    pdf.set_font("Helvetica", "I", 8)
                    pdf.multi_cell(0, 4, clean_text(f"This document is the property of {sample_data.get('company', 'the Company')}. Unauthorized distribution is prohibited."), align="C")
                
                else:
                    # General PDF
                    pdf.set_font("Helvetica", "B", 20)
                    pdf.cell(0, 15, clean_text(name), ln=True, align="C")
                    
                    pdf.set_font("Helvetica", "I", 10)
                    pdf.cell(0, 8, f"Generated: {datetime.now().strftime('%B %d, %Y')}", ln=True, align="C")
                    
                    if description:
                        pdf.ln(10)
                        pdf.set_font("Helvetica", "B", 14)
                        pdf.cell(0, 8, "Overview", ln=True)
                        pdf.set_font("Helvetica", "", 11)
                        pdf.multi_cell(0, 6, clean_text(description))
                    
                    if sample_data:
                        for key, value in sample_data.items():
                            pdf.ln(5)
                            pdf.set_font("Helvetica", "B", 12)
                            pdf.cell(0, 8, clean_text(str(key).replace("_", " ").title()), ln=True)
                            pdf.set_font("Helvetica", "", 10)
                            if isinstance(value, (list, dict)):
                                pdf.multi_cell(0, 5, clean_text(json.dumps(value, indent=2)))
                            else:
                                pdf.multi_cell(0, 5, clean_text(str(value)))
                
                filepath = os.path.join(output_dir, f"{filename}.pdf")
                print(f"[Demo Lab] Saving PDF to: {filepath}")
                pdf.output(filepath)
                print(f"[Demo Lab] PDF saved successfully!")
                return {"url": f"/api/demo/files/{filename}.pdf", "filepath": filepath}
                
            except ImportError as ie:
                print(f"[Demo Lab] PDF ImportError: {ie} - falling back to text")
                return await generate_demo_document(name, "text", content_type, sample_data, description)
            except Exception as pdf_err:
                print(f"[Demo Lab] PDF generation error: {pdf_err}")
                import traceback
                traceback.print_exc()
                return {"error": str(pdf_err)}
        
        # ============================================================
        # TEXT FILE (Fallback)
        # ============================================================
        elif doc_type == "text":
            filepath = os.path.join(output_dir, f"{filename}.txt")
            
            content = f"""{'=' * 60}
{name.upper()}
{'=' * 60}

Generated: {datetime.now().strftime('%B %d, %Y')}
Document ID: DOC-{uuid.uuid4().hex[:8].upper()}

{'─' * 60}

{description if description else ''}

"""
            if sample_data:
                for key, value in sample_data.items():
                    content += f"\n{str(key).upper().replace('_', ' ')}\n{'-' * 40}\n"
                    if isinstance(value, dict):
                        for k, v in value.items():
                            content += f"  {k}: {v}\n"
                    elif isinstance(value, list):
                        for item in value:
                            if isinstance(item, dict):
                                for k, v in item.items():
                                    content += f"  {k}: {v}\n"
                                content += "\n"
                            else:
                                content += f"  - {item}\n"
                    else:
                        content += f"  {value}\n"
                    content += "\n"
            
            content += f"""
{'=' * 60}
End of Document
{'=' * 60}
"""
            
            with open(filepath, 'w') as f:
                f.write(content)
            
            return {"url": f"/api/demo/files/{filename}.txt", "filepath": filepath}
        
        # If we get here, no document type matched - generate text as fallback
        print(f"[Demo Lab] Unknown doc_type '{doc_type}', falling back to text")
        filepath = os.path.join(output_dir, f"{filename}.txt")
        content = f"# {name}\n\n{description}\n\n{json.dumps(sample_data, indent=2)}"
        with open(filepath, 'w') as f:
            f.write(content)
        return {"url": f"/api/demo/files/{filename}.txt", "filepath": filepath}
        
    except Exception as e:
        print(f"Document generation error: {e}")
        import traceback
        traceback.print_exc()
        return {"error": str(e)}

@app.get("/api/demo/files/{filename}")
async def serve_demo_file(filename: str):
    """Serve generated demo files"""
    output_dir = os.environ.get("UPLOAD_PATH", "data/uploads")
    filepath = os.path.join(output_dir, filename)
    
    if os.path.exists(filepath):
        return FileResponse(filepath, filename=filename)
    raise HTTPException(404, "File not found")


@app.get("/api/demo/test-document")
async def test_document_generation():
    """Test endpoint to verify document generation works"""
    try:
        output_dir = os.environ.get("UPLOAD_PATH", "data/uploads")
        
        # Test 1: Check output directory
        os.makedirs(output_dir, exist_ok=True)
        dir_exists = os.path.exists(output_dir)
        dir_writable = os.access(output_dir, os.W_OK)
        
        # Test 2: Check libraries
        libs = {}
        try:
            from fpdf import FPDF
            libs["fpdf2"] = "✅ Installed"
        except ImportError:
            libs["fpdf2"] = "❌ Missing"
        
        try:
            from docx import Document
            libs["python-docx"] = "✅ Installed"
        except ImportError:
            libs["python-docx"] = "❌ Missing"
        
        try:
            from openpyxl import Workbook
            libs["openpyxl"] = "✅ Installed"
        except ImportError:
            libs["openpyxl"] = "❌ Missing"
        
        # Test 3: Try generating a test PDF
        test_result = None
        try:
            from fpdf import FPDF
            pdf = FPDF()
            pdf.add_page()
            pdf.set_font("Arial", "B", 16)
            pdf.cell(0, 10, "Test Document", ln=True)
            
            test_filename = f"test_{uuid.uuid4().hex[:6]}.pdf"
            test_filepath = os.path.join(output_dir, test_filename)
            pdf.output(test_filepath)
            
            if os.path.exists(test_filepath):
                test_result = {"success": True, "url": f"/api/demo/files/{test_filename}"}
                # Clean up
                # os.remove(test_filepath)
            else:
                test_result = {"success": False, "error": "File not created"}
        except Exception as e:
            test_result = {"success": False, "error": str(e)}
        
        return {
            "output_dir": output_dir,
            "dir_exists": dir_exists,
            "dir_writable": dir_writable,
            "libraries": libs,
            "test_pdf": test_result
        }
    except Exception as e:
        return {"error": str(e)}

@app.post("/api/demo/create-tool")
async def create_tool_from_demo(request: DemoCreateToolRequest, current_user: User = Depends(get_current_user_optional)):
    """Create an API tool from a generated mock API"""
    demo_id = request.demo_id
    
    if demo_id not in demo_items:
        return {"success": False, "error": "Demo item not found"}
    
    item = demo_items[demo_id]
    
    if item.type != "api":
        return {"success": False, "error": "Only API demos can be converted to tools"}
    
    # Get owner_id from current user
    owner_id = str(current_user.id) if current_user else "system"
    print(f"🔧 [CREATE_TOOL_FROM_DEMO] Creating tool '{item.name}' with owner_id='{owner_id}', user={current_user.email if current_user else 'None'}")
    
    # Get base URL from environment or default
    base_url = os.environ.get("PUBLIC_URL", "http://localhost:8000")
    
    # Create the API tool
    api_config = APIEndpointConfig(
        base_url=f"{base_url}{item.url}",
        http_method="GET",
        endpoint_path="/",
        auth_type="none"
    )
    
    # Demo Lab tools default to PUBLIC access so everyone can use them
    tool = ToolConfiguration(
        type="api",
        name=f"Demo: {item.name}",
        description=f"Mock API - {item.description}",
        api_config=api_config,
        config=item.config,
        # Access Control - Demo tools are PUBLIC by default
        owner_id=owner_id,
        access_type="public"  # Demo tools should be accessible to everyone
    )
    print(f"   ✅ Tool created in memory: id={tool.id}, owner_id={tool.owner_id}, access_type={tool.access_type}")
    
    app_state.tools[tool.id] = tool
    app_state.save_to_disk()
    
    return {"success": True, "tool_id": tool.id, "tool_name": tool.name, "access_type": "public"}


# Dynamic Mock API Routes with Auto-Generated Data
@app.get("/demo-api")
async def list_mock_apis():
    """List available mock API patterns and examples"""
    return {
        "message": "🧪 Smart Mock API Generator - LLM Powered!",
        "description": "Ask for ANY API and get realistic data. The LLM understands your context and generates appropriate data.",
        "how_it_works": [
            "1. Describe what you need in the API name",
            "2. LLM analyzes your request and generates realistic data",
            "3. Data includes real company names, proper formats, regional context"
        ],
        "examples": [
            {"category": "🏢 Company Master Data", "endpoints": [
                {"url": "/demo-api/egyptian-retail-companies", "description": "Egyptian retail companies with full details"},
                {"url": "/demo-api/uae-real-estate-companies", "description": "UAE real estate developers"},
                {"url": "/demo-api/saudi-telecom-companies", "description": "Saudi telecom providers"},
                {"url": "/demo-api/us-tech-startups", "description": "US tech startup directory"},
            ]},
            {"category": "🏦 Banking & Finance", "endpoints": [
                {"url": "/demo-api/banking/accounts", "description": "Banking accounts list"},
                {"url": "/demo-api/egypt-banks", "description": "Egyptian banks master data"},
                {"url": "/demo-api/investment-portfolio", "description": "Investment portfolio data"},
            ]},
            {"category": "🛒 Retail & E-commerce", "endpoints": [
                {"url": "/demo-api/electronics-store-products", "description": "Electronics store inventory"},
                {"url": "/demo-api/fashion-brands-catalog", "description": "Fashion brands and products"},
                {"url": "/demo-api/grocery-store-inventory", "description": "Grocery items with prices"},
            ]},
            {"category": "✈️ Travel & Hospitality", "endpoints": [
                {"url": "/demo-api/egypt-airlines", "description": "Egyptian airlines data"},
                {"url": "/demo-api/dubai-hotels", "description": "Dubai hotels directory"},
                {"url": "/demo-api/car-rental-companies", "description": "Car rental services"},
            ]},
            {"category": "🏥 Healthcare", "endpoints": [
                {"url": "/demo-api/cairo-hospitals", "description": "Cairo hospitals directory"},
                {"url": "/demo-api/pharmacy-chains", "description": "Pharmacy chain data"},
            ]},
            {"category": "🍔 Food & Restaurants", "endpoints": [
                {"url": "/demo-api/riyadh-restaurants", "description": "Restaurants in Riyadh"},
                {"url": "/demo-api/food-delivery-services", "description": "Food delivery platforms"},
            ]},
        ],
        "features": {
            "llm_powered": True,
            "understands_context": True,
            "regional_awareness": ["Egypt", "UAE", "Saudi", "US", "UK", "etc."],
            "industry_awareness": ["Retail", "Banking", "Healthcare", "Tech", "Real Estate", "etc."],
            "generates": ["Companies", "Products", "People", "Transactions", "Orders", "etc."]
        },
        "tips": [
            "Be specific: 'egyptian-retail-companies' is better than just 'companies'",
            "Include region for localized data: 'dubai-restaurants', 'cairo-hospitals'",
            "Include industry for context: 'tech-startups', 'fashion-brands'",
            "The LLM will generate realistic names, numbers, and details!"
        ],
        "supported_methods": ["GET", "POST", "PUT", "DELETE"]
    }

@app.get("/demo-api/_cache")
async def get_mock_api_cache():
    """Get cache statistics"""
    return {
        "cached_endpoints": len(mock_api_llm_cache),
        "endpoints": list(mock_api_llm_cache.keys())
    }

@app.delete("/demo-api/_cache")
async def clear_mock_api_cache():
    """Clear LLM response cache"""
    count = len(mock_api_llm_cache)
    mock_api_llm_cache.clear()
    return {"message": f"Cleared {count} cached responses"}

@app.api_route("/demo-api/{api_name}", methods=["GET", "POST", "PUT", "DELETE"])
@app.api_route("/demo-api/{api_name}/{path:path}", methods=["GET", "POST", "PUT", "DELETE"])
async def mock_api_handler(api_name: str, path: str = "", request: Request = None):
    """Handle requests to mock APIs with auto-generated sample data"""
    
    method = request.method if request else "GET"
    full_path = f"/{path}" if path else "/"
    
    # Check if we have stored mock data from Demo Lab
    if api_name in mock_api_data:
        api_info = mock_api_data[api_name]
        endpoints = api_info.get("endpoints", [])
        
        for endpoint in endpoints:
            ep_method = endpoint.get("method", "GET")
            ep_path = endpoint.get("path", "/")
            
            if ep_method.upper() == method.upper():
                if ep_path == full_path or ep_path.rstrip("/") == full_path.rstrip("/"):
                    return endpoint.get("sample_response", await generate_smart_response(api_name, full_path, method))
                
                if "{" in ep_path:
                    import re
                    pattern = ep_path.replace("{", "(?P<").replace("}", ">[^/]+)")
                    match = re.match(pattern, full_path)
                    if match:
                        params = match.groupdict()
                        return await generate_smart_response(api_name, full_path, method, params)
    
    # Use smart LLM-powered generation OR fallback to pattern matching
    return await generate_smart_response(api_name, full_path, method)


async def generate_smart_response(api_name: str, path: str, method: str, params: dict = None):
    """
    Smart response generator that:
    1. First tries LLM for dynamic, contextual data
    2. Falls back to pattern matching for common cases
    """
    
    # Try LLM-powered generation first
    llm_response = await generate_llm_mock_data(api_name, path, method, params)
    if llm_response:
        return llm_response
    
    # Fallback to pattern-based generation
    return generate_auto_response(api_name, path, method, params)


async def generate_llm_mock_data(api_name: str, path: str, method: str, params: dict = None) -> Optional[Dict]:
    """Use LLM to generate realistic, contextual mock data"""
    try:
        # Check cache first
        cache_key = f"{api_name}:{path}:{method}"
        if cache_key in mock_api_llm_cache:
            cached = mock_api_llm_cache[cache_key].copy()
            cached["_cached"] = True
            return cached
        
        llm_provider = app_state.get_llm_provider()
        if not llm_provider:
            return None
        
        # Build context from API name and path
        context = f"API: {api_name}, Path: {path}, Method: {method}"
        if params:
            context += f", Parameters: {params}"
        
        system_prompt = """You are a Mock API data generator. Generate realistic, professional JSON data based on the API request.

RULES:
1. ONLY output valid JSON - no explanations, no markdown, no code blocks
2. Generate realistic data with real-world names, numbers, dates
3. For companies: use real company names or realistic fictional ones with proper details
4. For people: use realistic names appropriate to the region/context
5. For numbers: use realistic ranges (salaries, prices, etc.)
6. Include relevant fields based on context (e.g., Arabic names for Egyptian data)
7. Generate 5-15 items for lists
8. Use current year (2024-2025) for dates

REGIONAL DATA:
- Egyptian companies: Carrefour Egypt, B.TECH, El Araby, Vodafone Egypt, Orange Egypt, Talaat Moustafa Group, etc.
- UAE companies: Emaar, ADNOC, Emirates NBD, Majid Al Futtaim, etc.
- Saudi companies: SABIC, Saudi Aramco, STC, Al Rajhi Bank, etc.

For company/business master data, include:
- company_id, name, name_ar (Arabic if Middle East), type, industry, sub_industry
- headquarters (city, address, country), founded year
- employees count, annual_revenue, stores/branches count
- contact (phone, email, website)
- registration_number, tax_id
- services (array), brands (array), status
- ownership, stock_symbol (if public)

OUTPUT FORMAT:
For lists: {"companies": [...], "total": X, "metadata": {"region": "...", "industry": "..."}}
For single item: {full object with all details}
For POST: {"id": "...", "status": "created", "message": "..."}

IMPORTANT: Output ONLY the JSON object, nothing else!"""

        user_prompt = f"""Generate mock API response for:
{context}

Return ONLY valid JSON with realistic data."""

        messages = [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt}
        ]
        
        response = await llm_provider.generate(messages, max_tokens=3000)
        
        # Try to parse JSON from response
        import re
        # Remove markdown code blocks if present
        response = re.sub(r'```json\s*', '', response)
        response = re.sub(r'```\s*', '', response)
        
        json_match = re.search(r'\{[\s\S]*\}', response)
        if json_match:
            try:
                data = json.loads(json_match.group())
                # Add metadata
                data["_generated_by"] = "llm"
                data["_api"] = api_name
                data["_path"] = path
                
                # Cache the response (without metadata)
                cache_data = {k: v for k, v in data.items() if not k.startswith("_")}
                mock_api_llm_cache[cache_key] = cache_data
                
                return data
            except json.JSONDecodeError as e:
                print(f"JSON decode error: {e}")
                pass
        
        return None
        
    except Exception as e:
        print(f"LLM mock data generation error: {e}")
        return None


def generate_auto_response(api_name: str, path: str, method: str, params: dict = None):
    """Auto-generate realistic sample data based on API name and path"""
    import random
    from datetime import datetime, timedelta
    
    # ============================================================
    # REALISTIC DATA POOLS
    # ============================================================
    
    # People Names
    first_names = ["Ahmed", "Sarah", "Mohamed", "Fatima", "Omar", "Layla", "Hassan", "Nour", "Ali", "Yasmin", "John", "Emily", "Michael", "Jessica", "David", "Emma", "James", "Olivia", "William", "Sophia"]
    last_names = ["Hassan", "Ali", "Ibrahim", "Ahmed", "Mohamed", "Khalil", "Nasser", "Smith", "Johnson", "Williams", "Brown", "Jones", "Garcia", "Miller", "Davis"]
    
    # Companies
    companies = ["TechCorp Solutions", "GlobalBank International", "MegaRetail Inc", "HealthPlus Medical", "EduLearn Academy", "FinanceHub Capital", "CloudSoft Technologies", "DataDrive Analytics", "Sunrise Electronics", "Blue Ocean Trading"]
    
    # RETAIL PRODUCTS - Realistic categories
    electronics = [
        {"name": "Apple iPhone 15 Pro Max 256GB", "sku": "APPL-IP15PM-256", "price": 1199.00, "category": "Smartphones"},
        {"name": "Samsung Galaxy S24 Ultra 512GB", "sku": "SAMS-S24U-512", "price": 1299.00, "category": "Smartphones"},
        {"name": "Apple MacBook Pro 14\" M3 Pro", "sku": "APPL-MBP14-M3P", "price": 1999.00, "category": "Laptops"},
        {"name": "Dell XPS 15 Intel Core i7 32GB", "sku": "DELL-XPS15-I7", "price": 1549.00, "category": "Laptops"},
        {"name": "Sony WH-1000XM5 Wireless Headphones", "sku": "SONY-WH1000XM5", "price": 349.00, "category": "Audio"},
        {"name": "Apple AirPods Pro 2nd Gen", "sku": "APPL-APP2-WHT", "price": 249.00, "category": "Audio"},
        {"name": "LG C3 65\" OLED 4K Smart TV", "sku": "LG-C3-65-OLED", "price": 1799.00, "category": "TVs"},
        {"name": "Samsung 55\" Neo QLED 4K", "sku": "SAMS-NEO55-4K", "price": 1299.00, "category": "TVs"},
        {"name": "Apple iPad Pro 12.9\" M2 256GB", "sku": "APPL-IPADP-M2", "price": 1099.00, "category": "Tablets"},
        {"name": "Microsoft Surface Pro 9 i7", "sku": "MSFT-SP9-I7", "price": 1599.00, "category": "Tablets"},
        {"name": "PlayStation 5 Console", "sku": "SONY-PS5-STD", "price": 499.00, "category": "Gaming"},
        {"name": "Xbox Series X 1TB", "sku": "MSFT-XBSX-1TB", "price": 499.00, "category": "Gaming"},
        {"name": "Nintendo Switch OLED Model", "sku": "NTDO-SWCH-OLED", "price": 349.00, "category": "Gaming"},
        {"name": "Canon EOS R6 Mark II Body", "sku": "CANN-R6M2-BDY", "price": 2499.00, "category": "Cameras"},
        {"name": "GoPro HERO12 Black", "sku": "GPRO-H12-BLK", "price": 399.00, "category": "Cameras"},
    ]
    
    clothing = [
        {"name": "Nike Air Max 270 Men's Running Shoes", "sku": "NIKE-AM270-BLK", "price": 150.00, "category": "Footwear"},
        {"name": "Adidas Ultraboost 23 Women's", "sku": "ADID-UB23-WHT", "price": 190.00, "category": "Footwear"},
        {"name": "Levi's 501 Original Fit Jeans", "sku": "LEVI-501-BLU", "price": 69.50, "category": "Pants"},
        {"name": "Ralph Lauren Polo Shirt Classic", "sku": "RL-POLO-NAV", "price": 98.00, "category": "Shirts"},
        {"name": "North Face Thermoball Eco Jacket", "sku": "TNF-THRM-BLK", "price": 220.00, "category": "Outerwear"},
        {"name": "Calvin Klein Slim Fit Dress Shirt", "sku": "CK-SLIM-WHT", "price": 79.00, "category": "Shirts"},
        {"name": "Under Armour Tech 2.0 T-Shirt", "sku": "UA-TECH2-GRY", "price": 25.00, "category": "Activewear"},
        {"name": "Zara Oversized Blazer", "sku": "ZARA-BLZR-BLK", "price": 119.00, "category": "Outerwear"},
    ]
    
    home_goods = [
        {"name": "Dyson V15 Detect Cordless Vacuum", "sku": "DYSN-V15-DET", "price": 749.00, "category": "Appliances"},
        {"name": "Ninja Foodi 9-in-1 Pressure Cooker", "sku": "NINJ-FD-9IN1", "price": 179.00, "category": "Kitchen"},
        {"name": "KitchenAid Artisan Stand Mixer", "sku": "KAID-ART-RED", "price": 449.00, "category": "Kitchen"},
        {"name": "iRobot Roomba j7+ Self-Emptying", "sku": "IRBT-J7P-BLK", "price": 799.00, "category": "Appliances"},
        {"name": "Casper Original Mattress Queen", "sku": "CSPR-ORIG-QN", "price": 1095.00, "category": "Bedroom"},
        {"name": "Nespresso Vertuo Next Coffee Machine", "sku": "NESP-VNXT-CHR", "price": 179.00, "category": "Kitchen"},
        {"name": "Philips Hue Starter Kit 4-Pack", "sku": "PHIL-HUE-4PK", "price": 199.00, "category": "Lighting"},
        {"name": "Instant Pot Duo 7-in-1 8Qt", "sku": "INST-DUO-8QT", "price": 89.00, "category": "Kitchen"},
    ]
    
    grocery = [
        {"name": "Organic Whole Milk 1 Gallon", "sku": "ORG-MILK-1GL", "price": 6.99, "category": "Dairy"},
        {"name": "Fresh Atlantic Salmon Fillet 1lb", "sku": "FSH-SALM-1LB", "price": 12.99, "category": "Seafood"},
        {"name": "Organic Free-Range Eggs 12ct", "sku": "ORG-EGGS-12", "price": 5.99, "category": "Dairy"},
        {"name": "Avocados Hass 4-Pack", "sku": "PRD-AVOC-4PK", "price": 4.99, "category": "Produce"},
        {"name": "Grass-Fed Ground Beef 1lb", "sku": "MT-GFBEEF-1LB", "price": 8.99, "category": "Meat"},
        {"name": "Organic Baby Spinach 5oz", "sku": "ORG-SPIN-5OZ", "price": 4.49, "category": "Produce"},
        {"name": "Sourdough Bread Loaf", "sku": "BKY-SRDGH-LF", "price": 5.99, "category": "Bakery"},
        {"name": "Greek Yogurt Variety Pack 12ct", "sku": "DRY-GRYG-12", "price": 9.99, "category": "Dairy"},
    ]
    
    departments = ["Engineering", "Sales", "Marketing", "Human Resources", "Finance", "Operations", "IT", "Legal", "Customer Service", "Research & Development"]
    job_titles = ["Software Engineer", "Sales Manager", "Marketing Specialist", "HR Coordinator", "Financial Analyst", "Operations Manager", "IT Administrator", "Legal Counsel", "Customer Success Rep", "Product Manager", "Senior Developer", "VP of Sales", "Director of Marketing", "Chief Technology Officer"]
    
    currencies = ["USD", "EUR", "GBP", "AED", "SAR"]
    order_statuses = ["pending", "confirmed", "processing", "shipped", "out_for_delivery", "delivered", "cancelled", "returned"]
    payment_methods = ["credit_card", "debit_card", "paypal", "apple_pay", "google_pay", "bank_transfer"]
    
    def random_date(days_back=365):
        d = datetime.now() - timedelta(days=random.randint(0, days_back))
        return d.strftime("%Y-%m-%d")
    
    def random_datetime(days_back=30):
        d = datetime.now() - timedelta(days=random.randint(0, days_back), hours=random.randint(0, 23), minutes=random.randint(0, 59))
        return d.isoformat()
    
    def random_id(prefix=""):
        return f"{prefix}{random.randint(100000, 999999)}"
    
    def random_uuid():
        import uuid
        return str(uuid.uuid4())[:8].upper()
    
    def random_phone():
        return f"+1-{random.randint(200,999)}-{random.randint(100,999)}-{random.randint(1000,9999)}"
    
    def random_email(name):
        domains = ["gmail.com", "outlook.com", "yahoo.com", "company.com", "work.org"]
        return f"{name.lower().replace(' ', '.').replace('-', '')}@{random.choice(domains)}"
    
    def random_address():
        streets = ["Main St", "Oak Avenue", "Park Boulevard", "Cedar Lane", "Maple Drive", "First Street", "Washington Ave", "King Road"]
        cities = ["New York", "Los Angeles", "Chicago", "Houston", "Phoenix", "Philadelphia", "San Antonio", "San Diego", "Dallas", "Austin"]
        states = ["NY", "CA", "IL", "TX", "AZ", "PA", "TX", "CA", "TX", "TX"]
        return {
            "street": f"{random.randint(100, 9999)} {random.choice(streets)}",
            "city": random.choice(cities),
            "state": random.choice(states),
            "zip": f"{random.randint(10000, 99999)}",
            "country": "USA"
        }
    
    api_lower = api_name.lower()
    path_lower = path.lower()
    
    # ============================================================
    # CONTEXT DETECTION - Understand what user REALLY wants
    # ============================================================
    
    # Check if asking for COMPANIES (not products)
    is_company_request = any(x in api_lower or x in path_lower for x in [
        "company", "companies", "business", "businesses", "corporation", "corp",
        "enterprise", "firm", "organization", "org", "vendor", "supplier",
        "master", "directory", "registry", "merchant", "retailer", "operator"
    ])
    
    # Check for Regional context
    is_egyptian = any(x in api_lower for x in ["egypt", "egyptian", "مصر", "مصري", "eg"])
    is_uae = any(x in api_lower for x in ["uae", "emirates", "dubai", "emirati"])
    is_saudi = any(x in api_lower for x in ["saudi", "ksa", "riyadh", "سعودي"])
    
    # Check industry
    is_retail = any(x in api_lower or x in path_lower for x in ["retail", "store", "shop", "market", "supermarket", "hypermarket"])
    is_banking = any(x in api_lower or x in path_lower for x in ["bank", "banking", "finance", "financial"])
    is_tech = any(x in api_lower or x in path_lower for x in ["tech", "technology", "software", "startup", "fintech"])
    
    # ============================================================
    # EGYPTIAN COMPANIES DATA
    # ============================================================
    
    egyptian_retail_companies = [
        {
            "company_id": "EG-RET-001",
            "name": "Carrefour Egypt",
            "name_ar": "كارفور مصر",
            "type": "Hypermarket Chain",
            "industry": "Retail",
            "sub_industry": "Hypermarkets & Supercenters",
            "founded": 2002,
            "headquarters": {
                "city": "Cairo",
                "address": "City Stars Mall, Nasr City",
                "governorate": "Cairo"
            },
            "ownership": "Majid Al Futtaim Group",
            "employees": 8500,
            "stores_count": 38,
            "annual_revenue_egp_million": 12500,
            "website": "www.carrefouregypt.com",
            "contact": {
                "phone": "+20-2-2480-0000",
                "email": "customercare@carrefouregypt.com"
            },
            "registration_number": "CR-12345-2002",
            "tax_id": "100-234-567",
            "services": ["Groceries", "Electronics", "Fashion", "Home & Garden", "Online Delivery"],
            "brands": ["Carrefour", "Carrefour Market", "Carrefour Express"],
            "status": "active"
        },
        {
            "company_id": "EG-RET-002",
            "name": "Spinneys Egypt",
            "name_ar": "سبينيز مصر",
            "type": "Premium Supermarket",
            "industry": "Retail",
            "sub_industry": "Supermarkets",
            "founded": 2007,
            "headquarters": {
                "city": "Cairo",
                "address": "Zamalek, Cairo",
                "governorate": "Cairo"
            },
            "ownership": "Spinneys Dubai LLC",
            "employees": 2200,
            "stores_count": 15,
            "annual_revenue_egp_million": 3200,
            "website": "www.spinneys-egypt.com",
            "contact": {
                "phone": "+20-2-2735-0000",
                "email": "info@spinneys-egypt.com"
            },
            "registration_number": "CR-23456-2007",
            "tax_id": "100-345-678",
            "services": ["Premium Groceries", "Organic Products", "Fresh Bakery", "Deli"],
            "brands": ["Spinneys"],
            "status": "active"
        },
        {
            "company_id": "EG-RET-003",
            "name": "Metro Markets",
            "name_ar": "مترو ماركت",
            "type": "Supermarket Chain",
            "industry": "Retail",
            "sub_industry": "Supermarkets",
            "founded": 1990,
            "headquarters": {
                "city": "Cairo",
                "address": "Heliopolis, Cairo",
                "governorate": "Cairo"
            },
            "ownership": "Metro Group Egypt",
            "employees": 3500,
            "stores_count": 45,
            "annual_revenue_egp_million": 5800,
            "website": "www.metromarkets.com.eg",
            "contact": {
                "phone": "+20-2-2690-0000",
                "email": "info@metromarkets.com.eg"
            },
            "registration_number": "CR-34567-1990",
            "tax_id": "100-456-789",
            "services": ["Groceries", "Fresh Produce", "Household Items"],
            "brands": ["Metro Market"],
            "status": "active"
        },
        {
            "company_id": "EG-RET-004",
            "name": "Seoudi Market",
            "name_ar": "سعودي ماركت",
            "type": "Supermarket Chain",
            "industry": "Retail",
            "sub_industry": "Supermarkets",
            "founded": 1985,
            "headquarters": {
                "city": "Cairo",
                "address": "Dokki, Giza",
                "governorate": "Giza"
            },
            "ownership": "Seoudi Group",
            "employees": 4000,
            "stores_count": 52,
            "annual_revenue_egp_million": 4200,
            "website": "www.seoudimarket.com",
            "contact": {
                "phone": "+20-2-3760-0000",
                "email": "info@seoudimarket.com"
            },
            "registration_number": "CR-45678-1985",
            "tax_id": "100-567-890",
            "services": ["Groceries", "Fresh Meat", "Fresh Vegetables"],
            "brands": ["Seoudi"],
            "status": "active"
        },
        {
            "company_id": "EG-RET-005",
            "name": "Kazyon",
            "name_ar": "كازيون",
            "type": "Discount Store Chain",
            "industry": "Retail",
            "sub_industry": "Discount Stores",
            "founded": 2014,
            "headquarters": {
                "city": "Cairo",
                "address": "New Cairo, Cairo",
                "governorate": "Cairo"
            },
            "ownership": "Kazyon Egypt LLC",
            "employees": 6000,
            "stores_count": 450,
            "annual_revenue_egp_million": 8500,
            "website": "www.kazyon.com.eg",
            "contact": {
                "phone": "+20-2-2537-0000",
                "email": "info@kazyon.com.eg"
            },
            "registration_number": "CR-56789-2014",
            "tax_id": "100-678-901",
            "services": ["Discount Groceries", "Household Essentials", "Personal Care"],
            "brands": ["Kazyon"],
            "status": "active"
        },
        {
            "company_id": "EG-RET-006",
            "name": "B.TECH",
            "name_ar": "بي تك",
            "type": "Electronics Retail Chain",
            "industry": "Retail",
            "sub_industry": "Electronics & Appliances",
            "founded": 1997,
            "headquarters": {
                "city": "Cairo",
                "address": "Smart Village, 6th of October",
                "governorate": "Giza"
            },
            "ownership": "B.TECH Group",
            "employees": 5000,
            "stores_count": 120,
            "annual_revenue_egp_million": 15000,
            "website": "www.btech.com",
            "contact": {
                "phone": "+20-2-3827-0000",
                "email": "info@btech.com"
            },
            "registration_number": "CR-67890-1997",
            "tax_id": "100-789-012",
            "services": ["Electronics", "Home Appliances", "Mobile Phones", "Computers", "Installation Services"],
            "brands": ["B.TECH"],
            "status": "active"
        },
        {
            "company_id": "EG-RET-007",
            "name": "2B Egypt",
            "name_ar": "تو بي",
            "type": "Electronics Retail",
            "industry": "Retail",
            "sub_industry": "Electronics & Appliances",
            "founded": 2016,
            "headquarters": {
                "city": "Cairo",
                "address": "Maadi, Cairo",
                "governorate": "Cairo"
            },
            "ownership": "2B Group",
            "employees": 1500,
            "stores_count": 35,
            "annual_revenue_egp_million": 4500,
            "website": "www.2b.com.eg",
            "contact": {
                "phone": "+20-2-2516-0000",
                "email": "info@2b.com.eg"
            },
            "registration_number": "CR-78901-2016",
            "tax_id": "100-890-123",
            "services": ["Computers", "Laptops", "Accessories", "Gaming"],
            "brands": ["2B"],
            "status": "active"
        },
        {
            "company_id": "EG-RET-008",
            "name": "El Araby Group",
            "name_ar": "مجموعة العربي",
            "type": "Electronics & Appliances",
            "industry": "Retail & Manufacturing",
            "sub_industry": "Consumer Electronics",
            "founded": 1964,
            "headquarters": {
                "city": "Benha",
                "address": "Industrial Zone, Benha",
                "governorate": "Qalyubia"
            },
            "ownership": "El Araby Family",
            "employees": 15000,
            "stores_count": 400,
            "annual_revenue_egp_million": 25000,
            "website": "www.elarabygroup.com",
            "contact": {
                "phone": "+20-13-323-0000",
                "email": "info@elarabygroup.com"
            },
            "registration_number": "CR-89012-1964",
            "tax_id": "100-901-234",
            "services": ["Home Appliances", "Air Conditioners", "TVs", "Washing Machines"],
            "brands": ["TOSHIBA El Araby", "SHARP El Araby", "Tornado"],
            "status": "active"
        },
        {
            "company_id": "EG-RET-009",
            "name": "LC Waikiki Egypt",
            "name_ar": "إل سي وايكيكي مصر",
            "type": "Fashion Retail",
            "industry": "Retail",
            "sub_industry": "Apparel & Fashion",
            "founded": 2011,
            "headquarters": {
                "city": "Cairo",
                "address": "City Stars Mall, Nasr City",
                "governorate": "Cairo"
            },
            "ownership": "LC Waikiki Turkey",
            "employees": 1800,
            "stores_count": 65,
            "annual_revenue_egp_million": 2800,
            "website": "www.lcwaikiki.com.eg",
            "contact": {
                "phone": "+20-2-2480-5000",
                "email": "egypt@lcwaikiki.com"
            },
            "registration_number": "CR-90123-2011",
            "tax_id": "100-012-345",
            "services": ["Men's Fashion", "Women's Fashion", "Kids' Fashion", "Accessories"],
            "brands": ["LC Waikiki"],
            "status": "active"
        },
        {
            "company_id": "EG-RET-010",
            "name": "Tradeline",
            "name_ar": "تريد لاين",
            "type": "Electronics Retail",
            "industry": "Retail",
            "sub_industry": "Electronics & Computers",
            "founded": 1995,
            "headquarters": {
                "city": "Cairo",
                "address": "Dokki, Giza",
                "governorate": "Giza"
            },
            "ownership": "Tradeline Group",
            "employees": 800,
            "stores_count": 25,
            "annual_revenue_egp_million": 3500,
            "website": "www.tradeline.com.eg",
            "contact": {
                "phone": "+20-2-3338-0000",
                "email": "info@tradeline.com.eg"
            },
            "registration_number": "CR-01234-1995",
            "tax_id": "100-123-456",
            "services": ["Computers", "Laptops", "Printers", "IT Solutions"],
            "brands": ["Tradeline"],
            "status": "active"
        },
        {
            "company_id": "EG-RET-011",
            "name": "Ragab Sons (Hyper One)",
            "name_ar": "أولاد رجب (هايبر وان)",
            "type": "Hypermarket Chain",
            "industry": "Retail",
            "sub_industry": "Hypermarkets",
            "founded": 1980,
            "headquarters": {
                "city": "Cairo",
                "address": "6th of October City",
                "governorate": "Giza"
            },
            "ownership": "Ragab Sons Group",
            "employees": 4500,
            "stores_count": 35,
            "annual_revenue_egp_million": 6500,
            "website": "www.hyperone.com.eg",
            "contact": {
                "phone": "+20-2-3837-0000",
                "email": "info@hyperone.com.eg"
            },
            "registration_number": "CR-12345-1980",
            "tax_id": "100-234-567",
            "services": ["Groceries", "Fresh Food", "Household Items", "Electronics"],
            "brands": ["Hyper One", "Ragab Sons"],
            "status": "active"
        },
        {
            "company_id": "EG-RET-012",
            "name": "Fathalla Market",
            "name_ar": "فتح الله ماركت",
            "type": "Supermarket Chain",
            "industry": "Retail",
            "sub_industry": "Supermarkets",
            "founded": 1960,
            "headquarters": {
                "city": "Alexandria",
                "address": "Smouha, Alexandria",
                "governorate": "Alexandria"
            },
            "ownership": "Fathalla Family",
            "employees": 2500,
            "stores_count": 30,
            "annual_revenue_egp_million": 2800,
            "website": "www.fathalla-market.com",
            "contact": {
                "phone": "+20-3-424-0000",
                "email": "info@fathalla-market.com"
            },
            "registration_number": "CR-23456-1960",
            "tax_id": "200-345-678",
            "services": ["Groceries", "Fresh Produce", "Butchery"],
            "brands": ["Fathalla"],
            "status": "active"
        }
    ]
    
    egyptian_banks = [
        {"company_id": "EG-BNK-001", "name": "National Bank of Egypt", "name_ar": "البنك الأهلي المصري", "type": "Public Bank", "founded": 1898, "headquarters": {"city": "Cairo"}, "branches": 520, "employees": 18000, "assets_egp_billion": 4200},
        {"company_id": "EG-BNK-002", "name": "Banque Misr", "name_ar": "بنك مصر", "type": "Public Bank", "founded": 1920, "headquarters": {"city": "Cairo"}, "branches": 690, "employees": 15000, "assets_egp_billion": 2100},
        {"company_id": "EG-BNK-003", "name": "Commercial International Bank (CIB)", "name_ar": "البنك التجاري الدولي", "type": "Private Bank", "founded": 1975, "headquarters": {"city": "Cairo"}, "branches": 210, "employees": 7500, "assets_egp_billion": 850},
        {"company_id": "EG-BNK-004", "name": "QNB Alahli", "name_ar": "بنك قطر الوطني الأهلي", "type": "Private Bank", "founded": 1978, "headquarters": {"city": "Cairo"}, "branches": 230, "employees": 5000, "assets_egp_billion": 680},
    ]
    
    egyptian_tech_companies = [
        {"company_id": "EG-TECH-001", "name": "Fawry", "name_ar": "فوري", "type": "Fintech", "founded": 2008, "headquarters": {"city": "Smart Village"}, "employees": 2500, "services": ["E-payments", "Bill Payment", "Financial Services"]},
        {"company_id": "EG-TECH-002", "name": "Paymob", "name_ar": "باي موب", "type": "Fintech", "founded": 2015, "headquarters": {"city": "Cairo"}, "employees": 350, "services": ["Payment Gateway", "POS", "E-commerce"]},
        {"company_id": "EG-TECH-003", "name": "Instabug", "name_ar": "انستاباج", "type": "Software", "founded": 2012, "headquarters": {"city": "Cairo"}, "employees": 200, "services": ["Bug Reporting", "App Analytics", "User Feedback"]},
    ]
    
    # ============================================================
    # COMPANIES / MASTER DATA APIs - CHECK FIRST!
    # ============================================================
    if is_company_request or "master" in api_lower:
        # Determine which company list to use
        if is_egyptian:
            if is_banking:
                company_list = egyptian_banks
                industry = "Banking"
            elif is_tech:
                company_list = egyptian_tech_companies
                industry = "Technology"
            else:  # Default to retail
                company_list = egyptian_retail_companies
                industry = "Retail"
            region = "Egypt"
        else:
            # Generate generic companies
            company_list = [
                {
                    "company_id": f"COMP-{random.randint(100000, 999999)}",
                    "name": c,
                    "type": "Corporation",
                    "industry": "Retail" if is_retail else "General",
                    "founded": random.randint(1980, 2020),
                    "headquarters": {"city": random.choice(["New York", "London", "Dubai"])},
                    "employees": random.randint(500, 10000),
                    "status": "active"
                } for c in companies
            ]
            region = "Global"
            industry = "Mixed"
        
        # Return single company if ID provided
        if params and ("id" in params or "name" in params):
            company = random.choice(company_list)
            return {
                **company,
                "retrieved_at": random_datetime(0),
                "data_source": "Master Data API"
            }
        
        # Return list of companies
        return {
            "companies": company_list,
            "metadata": {
                "total": len(company_list),
                "page": 1,
                "per_page": len(company_list),
                "region": region,
                "industry": industry,
                "data_type": "Master Data",
                "last_updated": random_date(7)
            }
        }
    
    # ============================================================
    # RETAIL / E-COMMERCE / PRODUCTS APIs (Only if NOT company request)
    # ============================================================
    if any(x in api_lower for x in ["retail", "ecommerce", "shop", "store", "product", "catalog", "inventory"]):
        # Combine all products
        all_products = electronics + clothing + home_goods + grocery
        
        if any(x in path_lower for x in ["product", "item", "catalog"]):
            if params and ("id" in params or "sku" in params):
                prod = random.choice(all_products)
                return {
                    "product_id": params.get("id", random_id("PRD-")),
                    "sku": prod["sku"],
                    "name": prod["name"],
                    "category": prod["category"],
                    "price": prod["price"],
                    "compare_at_price": round(prod["price"] * 1.2, 2),
                    "currency": "USD",
                    "in_stock": random.choice([True, True, True, False]),
                    "quantity_available": random.randint(0, 150),
                    "description": f"High-quality {prod['category'].lower()} product. {prod['name']} offers excellent value and performance.",
                    "brand": prod["name"].split()[0],
                    "rating": round(random.uniform(3.5, 5.0), 1),
                    "reviews_count": random.randint(10, 500),
                    "images": [f"https://images.example.com/{prod['sku'].lower()}-{i}.jpg" for i in range(1, 4)],
                    "created_at": random_date(365),
                    "updated_at": random_datetime(30)
                }
            
            # List of products
            selected = random.sample(all_products, min(10, len(all_products)))
            return {
                "products": [
                    {
                        "product_id": random_id("PRD-"),
                        "sku": p["sku"],
                        "name": p["name"],
                        "category": p["category"],
                        "price": p["price"],
                        "currency": "USD",
                        "in_stock": random.choice([True, True, True, False]),
                        "quantity_available": random.randint(0, 150),
                        "rating": round(random.uniform(3.5, 5.0), 1),
                        "image_url": f"https://images.example.com/{p['sku'].lower()}-1.jpg"
                    } for p in selected
                ],
                "total": random.randint(150, 500),
                "page": 1,
                "per_page": 10
            }
        
        if "order" in path_lower:
            if method == "POST":
                order_items = random.sample(all_products, random.randint(1, 4))
                subtotal = sum(p["price"] for p in order_items)
                return {
                    "order_id": f"ORD-{random_uuid()}",
                    "status": "confirmed",
                    "items": [{"name": p["name"], "sku": p["sku"], "price": p["price"], "quantity": 1} for p in order_items],
                    "subtotal": round(subtotal, 2),
                    "tax": round(subtotal * 0.08, 2),
                    "shipping": 9.99,
                    "total": round(subtotal * 1.08 + 9.99, 2),
                    "currency": "USD",
                    "payment_method": random.choice(payment_methods),
                    "shipping_address": random_address(),
                    "estimated_delivery": random_date(-7),
                    "created_at": random_datetime(0),
                    "message": "Order placed successfully!"
                }
            
            if params and "id" in params:
                order_items = random.sample(all_products, random.randint(1, 4))
                subtotal = sum(p["price"] for p in order_items)
                return {
                    "order_id": params.get("id"),
                    "customer": {
                        "name": f"{random.choice(first_names)} {random.choice(last_names)}",
                        "email": random_email(f"customer{random.randint(1,999)}"),
                        "phone": random_phone()
                    },
                    "items": [{"name": p["name"], "sku": p["sku"], "price": p["price"], "quantity": random.randint(1, 3)} for p in order_items],
                    "subtotal": round(subtotal, 2),
                    "tax": round(subtotal * 0.08, 2),
                    "shipping": 9.99,
                    "total": round(subtotal * 1.08 + 9.99, 2),
                    "status": random.choice(order_statuses),
                    "payment_method": random.choice(payment_methods),
                    "shipping_address": random_address(),
                    "tracking_number": f"TRK{random_uuid()}{random_uuid()}",
                    "created_at": random_datetime(30),
                    "updated_at": random_datetime(5)
                }
            
            return {
                "orders": [
                    {
                        "order_id": f"ORD-{random_uuid()}",
                        "customer_name": f"{random.choice(first_names)} {random.choice(last_names)}",
                        "items_count": random.randint(1, 5),
                        "total": round(random.uniform(50, 500), 2),
                        "status": random.choice(order_statuses),
                        "payment_method": random.choice(payment_methods),
                        "created_at": random_datetime(30)
                    } for _ in range(8)
                ],
                "total": random.randint(500, 2000),
                "page": 1
            }
        
        if "cart" in path_lower:
            cart_items = random.sample(all_products, random.randint(1, 5))
            subtotal = sum(p["price"] * random.randint(1, 2) for p in cart_items)
            return {
                "cart_id": f"CART-{random_uuid()}",
                "items": [
                    {
                        "product_id": random_id("PRD-"),
                        "name": p["name"],
                        "sku": p["sku"],
                        "price": p["price"],
                        "quantity": random.randint(1, 2),
                        "image_url": f"https://images.example.com/{p['sku'].lower()}-1.jpg"
                    } for p in cart_items
                ],
                "items_count": len(cart_items),
                "subtotal": round(subtotal, 2),
                "estimated_tax": round(subtotal * 0.08, 2),
                "estimated_total": round(subtotal * 1.08, 2),
                "currency": "USD"
            }
        
        if "category" in path_lower or "categories" in path_lower:
            categories = ["Electronics", "Smartphones", "Laptops", "TVs", "Audio", "Gaming", "Cameras", "Clothing", "Footwear", "Home & Kitchen", "Appliances", "Grocery", "Beauty", "Sports", "Toys"]
            return {
                "categories": [
                    {
                        "id": random_id("CAT-"),
                        "name": cat,
                        "slug": cat.lower().replace(" ", "-").replace("&", "and"),
                        "products_count": random.randint(20, 200),
                        "image_url": f"https://images.example.com/categories/{cat.lower().replace(' ', '-')}.jpg"
                    } for cat in categories
                ]
            }
        
        # Default - return products
        selected = random.sample(all_products, min(8, len(all_products)))
        return {
            "products": [
                {
                    "product_id": random_id("PRD-"),
                    "sku": p["sku"],
                    "name": p["name"],
                    "category": p["category"],
                    "price": p["price"],
                    "in_stock": True,
                    "rating": round(random.uniform(4.0, 5.0), 1)
                } for p in selected
            ],
            "total": len(all_products),
            "page": 1
        }
    
    # ============================================================
    # BANKING / FINANCE APIs
    # ============================================================
    if any(x in api_lower for x in ["bank", "finance", "account", "payment", "wallet"]):
        if "account" in path_lower:
            if params and "id" in params:
                fname, lname = random.choice(first_names), random.choice(last_names)
                return {
                    "account_id": params.get("id"),
                    "account_number": f"****{random.randint(1000, 9999)}",
                    "routing_number": f"****{random.randint(100, 999)}",
                    "account_type": random.choice(["checking", "savings", "money_market", "business"]),
                    "holder": {
                        "first_name": fname,
                        "last_name": lname,
                        "email": random_email(f"{fname} {lname}")
                    },
                    "balance": {
                        "available": round(random.uniform(1000, 50000), 2),
                        "current": round(random.uniform(1000, 50000), 2),
                        "pending": round(random.uniform(0, 500), 2)
                    },
                    "currency": "USD",
                    "status": "active",
                    "opened_date": random_date(2000),
                    "interest_rate": round(random.uniform(0.01, 0.05), 4),
                    "branch": {
                        "id": f"BR-{random.randint(100, 999)}",
                        "name": f"{random.choice(['Downtown', 'Uptown', 'Midtown', 'Westside', 'Eastside'])} Branch"
                    }
                }
            return {
                "accounts": [
                    {
                        "account_id": random_id("ACC-"),
                        "account_number": f"****{random.randint(1000, 9999)}",
                        "account_type": random.choice(["checking", "savings", "money_market"]),
                        "holder_name": f"{random.choice(first_names)} {random.choice(last_names)}",
                        "balance": round(random.uniform(500, 75000), 2),
                        "currency": "USD",
                        "status": "active"
                    } for _ in range(5)
                ],
                "total": 5
            }
        
        if "transaction" in path_lower or "transfer" in path_lower:
            transaction_types = ["deposit", "withdrawal", "transfer", "payment", "refund", "fee", "interest"]
            merchants = ["Amazon", "Walmart", "Target", "Starbucks", "Shell Gas", "Netflix", "Spotify", "Uber", "DoorDash", "Apple", "Costco", "Whole Foods"]
            
            if method == "POST":
                amount = round(random.uniform(10, 5000), 2)
                return {
                    "transaction_id": f"TXN-{random_uuid()}",
                    "status": "completed",
                    "type": "transfer",
                    "amount": amount,
                    "currency": "USD",
                    "from_account": f"****{random.randint(1000, 9999)}",
                    "to_account": f"****{random.randint(1000, 9999)}",
                    "description": "Transfer completed",
                    "timestamp": random_datetime(0),
                    "confirmation_number": f"CNF{random_uuid()}{random_uuid()}",
                    "message": "Transaction processed successfully"
                }
            
            return {
                "transactions": [
                    {
                        "transaction_id": f"TXN-{random_uuid()}",
                        "type": random.choice(transaction_types),
                        "amount": round(random.uniform(5, 500), 2),
                        "currency": "USD",
                        "merchant": random.choice(merchants) if random.random() > 0.3 else None,
                        "description": f"{random.choice(['Purchase at', 'Payment to', 'Transfer from', 'Deposit from'])} {random.choice(merchants)}",
                        "category": random.choice(["shopping", "food", "transportation", "entertainment", "bills", "income"]),
                        "date": random_datetime(30),
                        "status": "completed",
                        "balance_after": round(random.uniform(1000, 50000), 2)
                    } for _ in range(10)
                ],
                "total": random.randint(100, 500),
                "page": 1
            }
        
        if "balance" in path_lower:
            return {
                "account_id": params.get("id", random_id("ACC-")) if params else random_id("ACC-"),
                "available_balance": round(random.uniform(1000, 50000), 2),
                "current_balance": round(random.uniform(1000, 50000), 2),
                "pending_transactions": round(random.uniform(0, 500), 2),
                "currency": "USD",
                "as_of": random_datetime(0)
            }
    
    # ============================================================
    # HR / EMPLOYEE APIs
    # ============================================================
    if any(x in api_lower for x in ["hr", "employee", "staff", "people", "workforce"]):
        if "employee" in path_lower or "staff" in path_lower:
            if params and "id" in params:
                fname, lname = random.choice(first_names), random.choice(last_names)
                dept = random.choice(departments)
                return {
                    "employee_id": params.get("id"),
                    "first_name": fname,
                    "last_name": lname,
                    "email": random_email(f"{fname}.{lname}"),
                    "phone": random_phone(),
                    "department": dept,
                    "job_title": random.choice(job_titles),
                    "employment_type": random.choice(["full-time", "part-time", "contract"]),
                    "hire_date": random_date(3000),
                    "salary": {
                        "amount": random.randint(45000, 180000),
                        "currency": "USD",
                        "frequency": "annual"
                    },
                    "manager": {
                        "id": random_id("EMP-"),
                        "name": f"{random.choice(first_names)} {random.choice(last_names)}"
                    },
                    "location": {
                        "office": f"{random.choice(['New York', 'San Francisco', 'Chicago', 'Austin', 'Seattle'])} Office",
                        "remote": random.choice([True, False])
                    },
                    "status": "active"
                }
            
            return {
                "employees": [
                    {
                        "employee_id": random_id("EMP-"),
                        "name": f"{random.choice(first_names)} {random.choice(last_names)}",
                        "email": random_email(f"emp{i}"),
                        "department": random.choice(departments),
                        "job_title": random.choice(job_titles),
                        "hire_date": random_date(2000),
                        "status": "active"
                    } for i in range(8)
                ],
                "total": random.randint(100, 500),
                "page": 1
            }
        
        if "department" in path_lower:
            return {
                "departments": [
                    {
                        "id": random_id("DEPT-"),
                        "name": dept,
                        "head": f"{random.choice(first_names)} {random.choice(last_names)}",
                        "employee_count": random.randint(10, 100),
                        "budget": random.randint(500000, 5000000),
                        "location": random.choice(["New York", "San Francisco", "Chicago", "Austin"])
                    } for dept in departments
                ]
            }
        
        if "leave" in path_lower or "vacation" in path_lower or "pto" in path_lower:
            leave_types = ["annual", "sick", "personal", "parental", "bereavement", "jury_duty"]
            return {
                "leave_requests": [
                    {
                        "request_id": random_id("LV-"),
                        "employee": {
                            "id": random_id("EMP-"),
                            "name": f"{random.choice(first_names)} {random.choice(last_names)}"
                        },
                        "type": random.choice(leave_types),
                        "start_date": random_date(30),
                        "end_date": random_date(20),
                        "days_requested": random.randint(1, 14),
                        "status": random.choice(["approved", "pending", "rejected"]),
                        "reason": random.choice(["Family vacation", "Medical appointment", "Personal matters", "Wedding", "Moving"]),
                        "submitted_at": random_datetime(45)
                    } for _ in range(5)
                ]
            }
        
        if "payroll" in path_lower:
            return {
                "payroll": [
                    {
                        "employee_id": random_id("EMP-"),
                        "employee_name": f"{random.choice(first_names)} {random.choice(last_names)}",
                        "pay_period": f"{random_date(30)} to {random_date(15)}",
                        "gross_pay": round(random.uniform(3000, 15000), 2),
                        "deductions": {
                            "federal_tax": round(random.uniform(300, 2000), 2),
                            "state_tax": round(random.uniform(100, 800), 2),
                            "social_security": round(random.uniform(200, 900), 2),
                            "medicare": round(random.uniform(50, 200), 2),
                            "health_insurance": round(random.uniform(100, 500), 2),
                            "401k": round(random.uniform(100, 1000), 2)
                        },
                        "net_pay": round(random.uniform(2000, 10000), 2),
                        "payment_date": random_date(5),
                        "payment_method": random.choice(["direct_deposit", "check"])
                    } for _ in range(6)
                ]
            }
    
    # ============================================================
    # CRM / CUSTOMER APIs
    # ============================================================
    if any(x in api_lower for x in ["crm", "customer", "client", "sales", "lead"]):
        if "customer" in path_lower or "client" in path_lower or "contact" in path_lower:
            if params and "id" in params:
                fname, lname = random.choice(first_names), random.choice(last_names)
                company = random.choice(companies)
                return {
                    "customer_id": params.get("id"),
                    "first_name": fname,
                    "last_name": lname,
                    "email": random_email(f"{fname} {lname}"),
                    "phone": random_phone(),
                    "company": company,
                    "job_title": random.choice(["CEO", "CTO", "VP of Sales", "Director", "Manager", "Buyer"]),
                    "status": random.choice(["lead", "prospect", "customer", "churned"]),
                    "source": random.choice(["website", "referral", "linkedin", "trade_show", "cold_call"]),
                    "lifetime_value": round(random.uniform(1000, 100000), 2),
                    "total_orders": random.randint(1, 50),
                    "address": random_address(),
                    "created_at": random_date(500),
                    "last_contact": random_datetime(30),
                    "assigned_to": f"{random.choice(first_names)} {random.choice(last_names)}",
                    "tags": random.sample(["enterprise", "smb", "startup", "hot_lead", "decision_maker", "vip"], k=random.randint(1, 3))
                }
            
            return {
                "customers": [
                    {
                        "customer_id": random_id("CUST-"),
                        "name": f"{random.choice(first_names)} {random.choice(last_names)}",
                        "email": random_email(f"customer{i}"),
                        "company": random.choice(companies),
                        "status": random.choice(["lead", "prospect", "customer"]),
                        "lifetime_value": round(random.uniform(500, 50000), 2),
                        "last_contact": random_datetime(30)
                    } for i in range(10)
                ],
                "total": random.randint(500, 5000),
                "page": 1
            }
        
        if "deal" in path_lower or "opportunity" in path_lower:
            stages = ["qualification", "discovery", "proposal", "negotiation", "closed_won", "closed_lost"]
            return {
                "deals": [
                    {
                        "deal_id": random_id("DEAL-"),
                        "name": f"{random.choice(companies)} - {random.choice(['Enterprise', 'Professional', 'Starter', 'Custom'])} Plan",
                        "value": round(random.uniform(5000, 250000), 2),
                        "currency": "USD",
                        "stage": random.choice(stages),
                        "probability": random.randint(10, 95),
                        "expected_close_date": random_date(-60),
                        "owner": f"{random.choice(first_names)} {random.choice(last_names)}",
                        "company": random.choice(companies),
                        "contact": f"{random.choice(first_names)} {random.choice(last_names)}",
                        "created_at": random_datetime(90)
                    } for _ in range(8)
                ],
                "pipeline_value": round(random.uniform(500000, 2000000), 2),
                "total": random.randint(50, 200)
            }
    
    # ============================================================
    # HEALTHCARE / MEDICAL APIs
    # ============================================================
    if any(x in api_lower for x in ["health", "medical", "patient", "hospital", "clinic", "doctor"]):
        specialties = ["Cardiology", "Neurology", "Orthopedics", "Pediatrics", "Oncology", "Dermatology", "Internal Medicine", "Family Medicine", "Emergency Medicine", "Radiology"]
        
        if "patient" in path_lower:
            if params and "id" in params:
                fname, lname = random.choice(first_names), random.choice(last_names)
                return {
                    "patient_id": params.get("id"),
                    "mrn": f"MRN-{random.randint(100000, 999999)}",
                    "first_name": fname,
                    "last_name": lname,
                    "date_of_birth": random_date(25000),
                    "gender": random.choice(["male", "female"]),
                    "blood_type": random.choice(["A+", "A-", "B+", "B-", "O+", "O-", "AB+", "AB-"]),
                    "phone": random_phone(),
                    "email": random_email(f"{fname} {lname}"),
                    "address": random_address(),
                    "emergency_contact": {
                        "name": f"{random.choice(first_names)} {random.choice(last_names)}",
                        "relationship": random.choice(["spouse", "parent", "sibling", "child"]),
                        "phone": random_phone()
                    },
                    "insurance": {
                        "provider": random.choice(["Blue Cross", "Aetna", "Cigna", "United Healthcare", "Kaiser"]),
                        "policy_number": f"POL-{random.randint(100000, 999999)}",
                        "group_number": f"GRP-{random.randint(1000, 9999)}"
                    },
                    "allergies": random.sample(["Penicillin", "Aspirin", "Sulfa", "Latex", "Peanuts", "None"], k=random.randint(0, 2)),
                    "conditions": random.sample(["Hypertension", "Diabetes Type 2", "Asthma", "None"], k=random.randint(0, 2)),
                    "last_visit": random_date(30)
                }
            
            return {
                "patients": [
                    {
                        "patient_id": random_id("PAT-"),
                        "name": f"{random.choice(first_names)} {random.choice(last_names)}",
                        "date_of_birth": random_date(25000),
                        "gender": random.choice(["male", "female"]),
                        "phone": random_phone(),
                        "last_visit": random_date(60),
                        "primary_doctor": f"Dr. {random.choice(last_names)}"
                    } for _ in range(8)
                ],
                "total": random.randint(500, 5000)
            }
        
        if "appointment" in path_lower:
            return {
                "appointments": [
                    {
                        "appointment_id": random_id("APT-"),
                        "patient": {
                            "id": random_id("PAT-"),
                            "name": f"{random.choice(first_names)} {random.choice(last_names)}"
                        },
                        "doctor": {
                            "id": random_id("DOC-"),
                            "name": f"Dr. {random.choice(first_names)} {random.choice(last_names)}",
                            "specialty": random.choice(specialties)
                        },
                        "date": random_date(-14),
                        "time": f"{random.randint(8, 17)}:{random.choice(['00', '15', '30', '45'])}",
                        "duration_minutes": random.choice([15, 30, 45, 60]),
                        "type": random.choice(["check-up", "follow-up", "consultation", "procedure", "emergency"]),
                        "status": random.choice(["scheduled", "confirmed", "completed", "cancelled", "no-show"]),
                        "notes": random.choice(["Regular checkup", "Follow-up visit", "New patient consultation", ""])
                    } for _ in range(6)
                ]
            }
    
    # ============================================================
    # GENERIC / DEFAULT RESPONSES
    # ============================================================
    
    if method == "POST":
        return {
            "id": random_id(),
            "status": "created",
            "message": "Resource created successfully",
            "created_at": random_datetime(0)
        }
    
    if method in ["PUT", "PATCH"]:
        return {
            "id": params.get("id", random_id()) if params else random_id(),
            "status": "updated",
            "message": "Resource updated successfully",
            "updated_at": random_datetime(0)
        }
    
    if method == "DELETE":
        return {
            "status": "deleted",
            "message": "Resource deleted successfully",
            "deleted_at": random_datetime(0)
        }
    
    # Generic list with contextual data
    return {
        "items": [
            {
                "id": random_id(),
                "name": f"Sample {api_name.title()} Item {i+1}",
                "description": f"Auto-generated sample data for {api_name}",
                "status": random.choice(["active", "pending", "completed"]),
                "created_at": random_datetime(60),
                "updated_at": random_datetime(7)
            } for i in range(5)
        ],
        "total": random.randint(50, 200),
        "page": 1,
        "message": f"Try more specific paths like /{api_name}/products, /{api_name}/orders, /{api_name}/customers"
    }


if __name__ == "__main__":
    uvicorn.run("main:app", host="0.0.0.0", port=8000, reload=True)
