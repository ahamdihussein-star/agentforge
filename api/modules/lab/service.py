"""
Lab Service - AI-powered test data generation

This service handles:
- Mock API generation with realistic data using AI
- Document generation (Word, PDF, Excel, PPT)
- Image generation for OCR testing
"""

import os
import re
import json
import uuid
import base64
from datetime import datetime
from typing import Dict, Any, List, Optional
from io import BytesIO

# Document generation libraries
try:
    from docx import Document
    from docx.shared import Inches, Pt
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    from openpyxl import Workbook
    from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx.util import Inches as PptxInches, Pt as PptxPt
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from PIL import Image, ImageDraw, ImageFont
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False


class LabService:
    """Service for generating test data using AI"""
    
    # Storage path for generated files
    STORAGE_PATH = os.path.join(os.path.dirname(__file__), '..', '..', '..', 'data', 'lab')
    
    @classmethod
    def _ensure_storage(cls):
        """Ensure storage directory exists"""
        os.makedirs(cls.STORAGE_PATH, exist_ok=True)
    
    @classmethod
    def _get_ai_client(cls):
        """Get the AI client for generation"""
        # Try to import the existing AI setup
        try:
            from api.main import get_llm_client
            return get_llm_client()
        except:
            pass
        
        # Fallback to OpenAI
        try:
            import openai
            api_key = os.getenv('OPENAI_API_KEY')
            if api_key:
                return openai.OpenAI(api_key=api_key)
        except:
            pass
        
        return None
    
    @classmethod
    async def generate_api_data(
        cls,
        name: str,
        description: str,
        record_count: int = 10,
        org_id: str = None,
        user_id: str = None
    ) -> Dict[str, Any]:
        """
        Generate mock API data using AI
        
        The AI generates realistic data based on the description,
        without knowing it's for demo/testing purposes.
        """
        cls._ensure_storage()
        
        item_id = str(uuid.uuid4())
        
        # Build the AI prompt
        prompt = f"""Generate exactly {record_count} records of realistic data based on this description:

{description}

Requirements:
- Return ONLY a valid JSON array
- Each record should have realistic, varied data
- Use proper data types (strings, numbers, dates, etc.)
- Make names, emails, and other identifiers look real
- Ensure data is consistent and makes sense together
- Do not include any explanation, just the JSON array

Example format:
[
  {{"field1": "value1", "field2": "value2"}},
  {{"field1": "value3", "field2": "value4"}}
]"""

        # Try to use AI
        client = cls._get_ai_client()
        data = []
        
        if client:
            try:
                response = client.chat.completions.create(
                    model="gpt-4o-mini",
                    messages=[
                        {"role": "system", "content": "You are a data generation expert. Generate realistic, production-quality data. Return only valid JSON."},
                        {"role": "user", "content": prompt}
                    ],
                    temperature=0.7,
                    max_tokens=4000
                )
                
                content = response.choices[0].message.content.strip()
                
                # Extract JSON from response
                if content.startswith('```'):
                    content = content.split('```')[1]
                    if content.startswith('json'):
                        content = content[4:]
                
                data = json.loads(content)
                
            except Exception as e:
                print(f"AI generation failed: {e}")
                # Fall back to sample data
                data = cls._generate_sample_data(description, record_count)
        else:
            # No AI available, generate sample data
            data = cls._generate_sample_data(description, record_count)
        
        # Infer schema from data
        schema = cls._infer_schema(data)
        
        # URL-friendly slug from name so endpoint can be /api/lab/mock/{api-name}
        slug = (re.sub(r"[^a-z0-9-]+", "-", name.lower()).strip("-") or item_id)[:80]
        
        # AI-generated metadata for agent tool: description and parameters (no hardcoded rules)
        agent_meta = cls._generate_agent_metadata(name, description, client)
        
        # Save to storage (response_schema for APIGenerateResponse compatibility)
        result = {
            "id": item_id,
            "name": name,
            "slug": slug,
            "endpoint": f"/api/lab/mock/{slug}",
            "data": data,
            "response_schema": schema,
            "record_count": len(data),
            "created_at": datetime.utcnow().isoformat(),
            "org_id": org_id,
            "agent_description": agent_meta.get("agent_description"),
            "parameters": agent_meta.get("parameters") or []
        }
        
        filepath = os.path.join(cls.STORAGE_PATH, f"{item_id}.json")
        with open(filepath, 'w') as f:
            json.dump(result, f, indent=2)

        cls._save_mock_api_to_db(result, user_id)

        print(f"💾 [LAB] Saved API: name='{name}', slug='{slug}', file={filepath}")
        return result

    @classmethod
    def _save_mock_api_to_db(cls, result: Dict[str, Any], user_id: Optional[str] = None) -> None:
        """Persist mock API to database (id, slug, data, parameters, etc.)."""
        try:
            from database.base import get_db_session
            from database.models import LabMockAPI

            item_id = result.get("id")
            slug = (result.get("slug") or "")[:120]
            if not item_id or not slug:
                return
            with get_db_session() as db:
                row = LabMockAPI(
                    id=item_id,
                    user_id=user_id,
                    name=result.get("name", ""),
                    slug=slug,
                    description=None,
                    endpoint=result.get("endpoint", f"/api/lab/mock/{slug}"),
                    agent_description=result.get("agent_description"),
                    parameters=result.get("parameters") or [],
                    response_schema=result.get("response_schema") or {},
                    data=result.get("data") or [],
                    record_count=result.get("record_count", 0),
                )
                db.add(row)
        except Exception as e:
            print(f"⚠️ [LAB] Could not save mock API to DB: {e}")
    
    @classmethod
    def _generate_agent_metadata(
        cls, name: str, description: str, client: Any
    ) -> Dict[str, Any]:
        """
        Use AI to generate agent-facing description and input parameters.
        No hardcoded rules; suitable for any enterprise API.
        """
        out = {"agent_description": None, "parameters": []}
        if not client:
            out["agent_description"] = f"Returns data for: {name}. Use when the user asks for this information."
            return out
        prompt = f"""API name: {name}
User description: {description}

Return a JSON object with exactly two keys:
1. "agent_description": One clear sentence for an AI agent: what this API returns and when to use it. Do not mention mock, demo, or test.
2. "parameters": If this API logically accepts parameters (e.g. ID, filter, date range), return an array of objects, each with: "name", "description", "data_type" (string/number/boolean), "required" (boolean), "location" (query or path). If the API does not take parameters, return [].

Return only valid JSON, no markdown or explanation."""
        try:
            response = client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[
                    {"role": "system", "content": "You output only valid JSON. No markdown code blocks."},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.3,
                max_tokens=800
            )
            content = response.choices[0].message.content.strip()
            if content.startswith("```"):
                content = content.split("```")[1]
                if content.startswith("json"):
                    content = content[4:]
            obj = json.loads(content)
            out["agent_description"] = obj.get("agent_description") or out["agent_description"]
            out["parameters"] = obj.get("parameters") if isinstance(obj.get("parameters"), list) else []
        except Exception as e:
            print(f"Agent metadata generation failed: {e}")
            out["agent_description"] = f"Returns data for: {name}. Use when the user asks for this information."
        return out
    
    @classmethod
    def _generate_sample_data(cls, description: str, count: int) -> List[Dict]:
        """Generate realistic fallback data when AI is not available"""
        import random
        from datetime import timedelta

        first_names = ["Ahmed", "Sarah", "Michael", "Emily", "James", "Maria", "David",
                        "Jessica", "Robert", "Laura", "Daniel", "Sophia", "Omar", "Lina",
                        "William", "Olivia", "Hassan", "Nora", "Thomas", "Fatima"]
        last_names = ["Al-Rashid", "Johnson", "Chen", "Garcia", "Müller", "Tanaka",
                       "Petrov", "Williams", "Kim", "Santos", "Ali", "Martin",
                       "Anderson", "Lee", "Hernandez", "Smith", "Patel", "Brown"]
        companies = ["Apex Technologies Inc.", "GlobalTech Solutions", "Premier Services LLC",
                      "Quantum Dynamics Corp.", "Nexus Industries", "Blue Ocean Trading",
                      "Pinnacle Holdings", "Vantage Consulting", "Meridian Partners",
                      "Atlas Logistics", "Silverline Financial", "Nova Enterprises"]
        domains = ["company.com", "corp.com", "group.io", "global.com", "hq.com"]
        departments = ["Engineering", "Finance", "Marketing", "Operations", "Sales",
                        "Human Resources", "Legal", "Product", "Customer Success"]

        def _rand_name():
            return f"{random.choice(first_names)} {random.choice(last_names)}"

        def _rand_email(name: str):
            local = name.lower().replace(" ", ".")
            return f"{local}@{random.choice(domains)}"

        def _rand_date(days_back=365):
            d = datetime.utcnow() - timedelta(days=random.randint(0, days_back))
            return d.strftime("%Y-%m-%d")

        desc_lower = description.lower()

        if any(w in desc_lower for w in ['customer', 'user', 'person', 'employee']):
            results = []
            for i in range(count):
                name = _rand_name()
                results.append({
                    "id": i + 1,
                    "name": name,
                    "email": _rand_email(name),
                    "phone": f"+1-{random.randint(200,999)}-{random.randint(100,999)}-{random.randint(1000,9999)}",
                    "company": random.choice(companies),
                    "department": random.choice(departments),
                    "created_at": _rand_date(730)
                })
            return results

        elif any(w in desc_lower for w in ['product', 'item', 'inventory']):
            product_names = [
                "Wireless Bluetooth Headphones", "Smart Watch Pro", "4K Ultra Monitor",
                "Ergonomic Office Chair", "Portable SSD 1TB", "Mechanical Keyboard",
                "Noise-Canceling Earbuds", "USB-C Hub Adapter", "LED Desk Lamp",
                "Webcam HD 1080p", "Standing Desk Converter", "Laptop Backpack",
                "Wireless Charging Pad", "External Battery Pack", "Graphics Tablet",
                "Smart Thermostat", "Security Camera Kit", "Robot Vacuum Cleaner",
                "Air Purifier HEPA", "Espresso Machine Pro"
            ]
            random.shuffle(product_names)
            return [
                {
                    "id": i + 1,
                    "name": product_names[i % len(product_names)],
                    "sku": f"SKU-{random.randint(10000, 99999)}",
                    "price": round(random.uniform(19.99, 599.99), 2),
                    "quantity": random.randint(0, 250),
                    "category": random.choice(["Electronics", "Office", "Home & Garden", "Audio", "Computing"])
                }
                for i in range(count)
            ]

        elif any(w in desc_lower for w in ['order', 'transaction', 'sale']):
            statuses = ["pending", "processing", "shipped", "delivered", "completed"]
            return [
                {
                    "id": i + 1,
                    "order_number": f"ORD-{random.randint(100000, 999999)}",
                    "customer": _rand_name(),
                    "total": round(random.uniform(45, 2500), 2),
                    "status": random.choice(statuses),
                    "created_at": _rand_date(90)
                }
                for i in range(count)
            ]

        else:
            record_names = [
                "Annual Budget Review", "Q4 Performance Report", "Client Onboarding",
                "Maintenance Schedule", "Compliance Audit", "Partnership Agreement",
                "Training Program", "Market Analysis", "Supply Chain Update",
                "Risk Assessment", "Strategic Plan", "Resource Allocation"
            ]
            random.shuffle(record_names)
            return [
                {
                    "id": i + 1,
                    "name": record_names[i % len(record_names)],
                    "value": random.randint(1, 100),
                    "active": random.choice([True, False]),
                    "created_at": _rand_date(180)
                }
                for i in range(count)
            ]
    
    @classmethod
    def _infer_schema(cls, data: List[Dict]) -> Dict[str, Any]:
        """Infer JSON schema from data"""
        if not data:
            return {"type": "array", "items": {}}
        
        sample = data[0]
        properties = {}
        
        for key, value in sample.items():
            if isinstance(value, bool):
                properties[key] = {"type": "boolean"}
            elif isinstance(value, int):
                properties[key] = {"type": "integer"}
            elif isinstance(value, float):
                properties[key] = {"type": "number"}
            elif isinstance(value, str):
                properties[key] = {"type": "string"}
            elif isinstance(value, list):
                properties[key] = {"type": "array"}
            elif isinstance(value, dict):
                properties[key] = {"type": "object"}
            else:
                properties[key] = {"type": "string"}
        
        return {
            "type": "array",
            "items": {
                "type": "object",
                "properties": properties
            }
        }
    
    @classmethod
    async def generate_document(
        cls,
        name: str,
        description: str,
        format: str = "docx",
        org_id: str = None
    ) -> Dict[str, Any]:
        """
        Generate a professional document using AI.
        Financial documents (invoices, receipts) use a structured JSON pipeline
        with proper table rendering; other documents use markdown content.
        """
        cls._ensure_storage()

        item_id = str(uuid.uuid4())
        filename = f"{item_id}.{format}"
        filepath = os.path.join(cls.STORAGE_PATH, filename)

        is_financial = cls._is_financial_document(description, name)
        size = 0

        if is_financial:
            fin_data = await cls._generate_financial_document_data(name, description)
            if format == "docx" and DOCX_AVAILABLE:
                size = cls._create_financial_docx(filepath, fin_data)
            elif format == "pdf" and PDF_AVAILABLE:
                size = cls._create_financial_pdf(filepath, fin_data)
            elif format == "xlsx" and XLSX_AVAILABLE:
                size = cls._create_financial_xlsx(filepath, fin_data)
            elif format == "pptx" and PPTX_AVAILABLE:
                content = await cls._generate_document_content(description)
                size = cls._create_pptx(filepath, name, content)
            else:
                content = await cls._generate_document_content(description)
                with open(filepath.replace(f'.{format}', '.txt'), 'w') as f:
                    f.write(f"# {name}\n\n{content}")
                filepath = filepath.replace(f'.{format}', '.txt')
                format = 'txt'
                size = os.path.getsize(filepath)
        else:
            content = await cls._generate_document_content(description)
            if format == "docx" and DOCX_AVAILABLE:
                size = cls._create_docx(filepath, name, content)
            elif format == "pdf" and PDF_AVAILABLE:
                size = cls._create_pdf(filepath, name, content)
            elif format == "xlsx" and XLSX_AVAILABLE:
                size = cls._create_xlsx(filepath, name, content)
            elif format == "pptx" and PPTX_AVAILABLE:
                size = cls._create_pptx(filepath, name, content)
            else:
                with open(filepath.replace(f'.{format}', '.txt'), 'w') as f:
                    f.write(f"# {name}\n\n{content}")
                filepath = filepath.replace(f'.{format}', '.txt')
                format = 'txt'
                size = os.path.getsize(filepath)

        result = {
            "id": item_id,
            "name": name,
            "format": format,
            "size": size,
            "download_url": f"/api/lab/download/{item_id}",
            "created_at": datetime.utcnow().isoformat(),
            "org_id": org_id
        }

        with open(os.path.join(cls.STORAGE_PATH, f"{item_id}_meta.json"), 'w') as f:
            json.dump(result, f, indent=2)

        return result
    
    _FINANCIAL_KEYWORDS = ['invoice', 'receipt', 'bill', 'purchase order', 'quotation', 'estimate']

    @classmethod
    def _is_financial_document(cls, description: str, name: str = "") -> bool:
        combined = f"{name} {description}".lower()
        return any(kw in combined for kw in cls._FINANCIAL_KEYWORDS)

    @classmethod
    async def _generate_document_content(cls, description: str) -> str:
        """Generate document content using AI -- always with realistic filled-in data."""
        client = cls._get_ai_client()

        if client:
            try:
                response = client.chat.completions.create(
                    model="gpt-4o-mini",
                    messages=[
                        {"role": "system", "content": (
                            "You are a professional document writer. Produce a COMPLETE, FINISHED, real-looking "
                            "business document with ALL data filled in.\n"
                            "FORMATTING RULES (the markdown is rendered into a real Word/PDF document, so follow exactly):\n"
                            "- Start directly with the body. Do NOT repeat the document title (it is added automatically).\n"
                            "- Use '## Section Heading' for each section, never bold text as a fake heading.\n"
                            "- Present structured / key-value information (passenger details, line items, specs, "
                            "totals, etc.) as proper MARKDOWN TABLES using pipes, e.g.\n"
                            "    | Field | Value |\n    | --- | --- |\n    | Name | Ahmed Hamdy |\n"
                            "  Tables render as real bordered tables and look far more professional than bullet lists.\n"
                            "- Use '- ' only for genuine bullet lists (terms, notes), and **bold** sparingly for emphasis.\n"
                            "- NEVER output horizontal-rule separator lines (---, ***, ___). Use section headings instead.\n"
                            "- NEVER use placeholders like [Company Name] or [Date]; every value must be realistic and specific.\n"
                            "- Match the real-world layout and tone of the requested document type "
                            "(e.g. a ticket, invoice, contract, letter, or report)."
                        )},
                        {"role": "user", "content": f"Write a complete, professional, realistic document with all data filled in for:\n\n{description}"}
                    ],
                    temperature=0.7,
                    max_tokens=3000
                )
                return response.choices[0].message.content
            except Exception as e:
                print(f"AI content generation failed: {e}")

        # Fallback content
        return f"""# {description[:60] if description else 'Document'}

## Overview
{description}

## Details
This document provides a comprehensive overview of the subject matter outlined above.
All information has been compiled from relevant sources and reviewed for accuracy.

## Summary
For further details or inquiries, please contact the responsible department.
"""

    @classmethod
    async def _generate_financial_document_data(cls, name: str, description: str) -> Dict[str, Any]:
        """Generate structured JSON data for financial documents (invoices, receipts, etc.)."""
        client = cls._get_ai_client()

        prompt = f"""Generate realistic financial document data for this request:

Name: {name}
Description: {description}

Return a JSON object with this structure:
{{
    "company": "Realistic company name",
    "company_address": "Full street address\\nCity, State ZIP",
    "company_phone": "+1 (XXX) XXX-XXXX",
    "company_email": "billing@company.com",
    "document_number": "INV-2025-XXXXX",
    "document_date": "February 25, 2025",
    "due_date": "March 27, 2025",
    "customer_name": "Realistic customer name",
    "customer_address": "Full street address\\nCity, State ZIP",
    "customer_email": "contact@customer.com",
    "items": [
        {{"description": "Specific detailed item/service description", "qty": 1, "price": 1500.00}},
        {{"description": "Another specific item", "qty": 5, "price": 200.00}}
    ],
    "subtotal": 2500.00,
    "tax_rate": "8%",
    "tax": 200.00,
    "total": 2700.00,
    "payment_terms": "Net 30",
    "notes": "Thank you for your business!"
}}

RULES:
- Generate 3-6 line items that are SPECIFIC to the context described
- Use realistic pricing (not round numbers -- use $97.50, $1,247.00, etc.)
- Calculate subtotal, tax, and total CORRECTLY
- Item descriptions must be detailed: "Strategic Planning Workshop (8 hours)" not just "Consulting"
- NEVER use placeholders. Every field must have a realistic value.
- Return ONLY valid JSON, no markdown or explanation."""

        if client:
            try:
                response = client.chat.completions.create(
                    model="gpt-4o-mini",
                    messages=[
                        {"role": "system", "content": "You generate realistic financial document data in JSON format. Never use placeholders. Output valid JSON only."},
                        {"role": "user", "content": prompt}
                    ],
                    temperature=0.7,
                    max_tokens=2000
                )
                content = response.choices[0].message.content.strip()
                if '```' in content:
                    content = content.split('```')[1]
                    if content.startswith('json'):
                        content = content[4:]
                return json.loads(content)
            except Exception as e:
                print(f"Financial document AI generation failed: {e}")

        return cls._get_fallback_financial_data(name, description)

    @classmethod
    def _get_fallback_financial_data(cls, name: str, description: str) -> Dict[str, Any]:
        """Realistic fallback data for invoices/receipts when AI is unavailable."""
        import random
        from datetime import timedelta

        first_names = ["Ahmed", "Sarah", "Michael", "Emily", "James", "Maria", "David", "Olivia"]
        last_names = ["Al-Rashid", "Johnson", "Chen", "Garcia", "Müller", "Williams", "Patel", "Santos"]
        companies = ["Apex Technologies Inc.", "Meridian Partners LLC", "Blue Ocean Trading Co.",
                      "Pinnacle Holdings Group", "Vantage Consulting Ltd.", "Silverline Financial Corp."]

        company = random.choice(companies)
        customer = f"{random.choice(first_names)} {random.choice(last_names)}"
        inv_num = f"INV-{datetime.utcnow().year}-{random.randint(10000, 99999)}"
        inv_date = datetime.utcnow().strftime("%B %d, %Y")
        due_date = (datetime.utcnow() + timedelta(days=30)).strftime("%B %d, %Y")

        all_items = [
            {"description": "Professional Consulting Services - Phase 1", "qty": 24, "price": 175.00},
            {"description": "Software License (Enterprise, Annual)", "qty": 1, "price": 4800.00},
            {"description": "Cloud Infrastructure Setup & Configuration", "qty": 1, "price": 3250.00},
            {"description": "Technical Documentation Package", "qty": 1, "price": 1500.00},
            {"description": "Quality Assurance & Testing (40 hours)", "qty": 40, "price": 95.00},
            {"description": "Project Management - Monthly Retainer", "qty": 3, "price": 2200.00},
            {"description": "Data Migration & Integration Services", "qty": 1, "price": 5750.00},
            {"description": "Staff Training Workshop (2 days)", "qty": 1, "price": 3400.00},
        ]
        items = random.sample(all_items, random.randint(3, 5))
        subtotal = sum(i["qty"] * i["price"] for i in items)
        tax = round(subtotal * 0.08, 2)
        total = round(subtotal + tax, 2)

        short = company.split()[0].lower().replace(".", "")
        return {
            "company": company,
            "company_address": f"{random.randint(100, 9999)} Commerce Blvd\nSuite {random.randint(100, 800)}\nNew York, NY 10017",
            "company_phone": f"+1 ({random.randint(200,999)}) {random.randint(100,999)}-{random.randint(1000,9999)}",
            "company_email": f"billing@{short}.com",
            "document_number": inv_num,
            "document_date": inv_date,
            "due_date": due_date,
            "customer_name": customer,
            "customer_address": f"{random.randint(100, 9999)} {random.choice(['Park Ave', 'Oak Street', 'Main Blvd'])}\n{random.choice(['Los Angeles, CA 90001', 'Chicago, IL 60601', 'Houston, TX 77001'])}",
            "customer_email": f"{customer.lower().replace(' ', '.')}@{random.choice(['corp.com', 'company.com', 'group.io'])}",
            "items": items,
            "subtotal": subtotal,
            "tax_rate": "8%",
            "tax": tax,
            "total": total,
            "payment_terms": "Net 30",
            "notes": "Thank you for your business!"
        }

    # ------------------------------------------------------------------
    # Structured document renderers (for financial docs with tables)
    # ------------------------------------------------------------------

    @classmethod
    def _create_financial_docx(cls, filepath: str, data: Dict[str, Any]) -> int:
        """Create a professional Word invoice/receipt with tables."""
        doc = Document()

        style = doc.styles['Normal']
        style.font.name = 'Calibri'
        style.font.size = Pt(11)

        # Company header
        header = doc.add_paragraph()
        run = header.add_run(data.get("company", ""))
        run.bold = True
        run.font.size = Pt(22)

        addr_lines = (data.get("company_address") or "").split('\n')
        for line in addr_lines:
            doc.add_paragraph(line)
        phone = data.get("company_phone", "")
        email = data.get("company_email", "")
        if phone or email:
            doc.add_paragraph(f"Phone: {phone}   |   Email: {email}")

        doc.add_paragraph()

        # Title (INVOICE / RECEIPT)
        title_p = doc.add_paragraph()
        title_p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        title_run = title_p.add_run("INVOICE")
        title_run.bold = True
        title_run.font.size = Pt(26)

        # Metadata
        meta_items = [
            ("Invoice #:", data.get("document_number", "")),
            ("Date:", data.get("document_date", "")),
            ("Due Date:", data.get("due_date", "")),
        ]
        for label, val in meta_items:
            p = doc.add_paragraph()
            p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
            p.add_run(f"{label} ").bold = True
            p.add_run(val)

        doc.add_paragraph()

        # Bill To
        bt = doc.add_paragraph()
        bt.add_run("BILL TO:").bold = True
        cust = doc.add_paragraph()
        cust.add_run(data.get("customer_name", "")).bold = True
        for line in (data.get("customer_address") or "").split('\n'):
            doc.add_paragraph(line)
        cust_email = data.get("customer_email", "")
        if cust_email:
            doc.add_paragraph(f"Email: {cust_email}")

        doc.add_paragraph()

        # Items table
        items = data.get("items", [])
        table = doc.add_table(rows=1 + len(items) + 3, cols=4)
        table.style = 'Table Grid'

        from docx.shared import RGBColor as _RGB
        headers = ["Description", "Quantity", "Unit Price", "Amount"]
        for i, h in enumerate(headers):
            cell = table.rows[0].cells[i]
            cell.text = h
            cls._shade_cell(cell, cls._TABLE_HEADER_FILL)
            for p in cell.paragraphs:
                for r in p.runs:
                    r.bold = True
                    r.font.color.rgb = _RGB(0xFF, 0xFF, 0xFF)

        for idx, item in enumerate(items):
            row = table.rows[idx + 1].cells
            row[0].text = item.get("description", "")
            row[1].text = str(item.get("qty", 0))
            row[2].text = f"${item.get('price', 0):,.2f}"
            row[3].text = f"${item.get('qty', 0) * item.get('price', 0):,.2f}"

        # Totals
        st_row = len(items) + 1
        table.rows[st_row].cells[2].text = "Subtotal:"
        table.rows[st_row].cells[3].text = f"${data.get('subtotal', 0):,.2f}"
        table.rows[st_row + 1].cells[2].text = f"Tax ({data.get('tax_rate', '8%')}):"
        table.rows[st_row + 1].cells[3].text = f"${data.get('tax', 0):,.2f}"
        total_cell = table.rows[st_row + 2].cells[2]
        total_cell.text = "TOTAL:"
        for p in total_cell.paragraphs:
            for r in p.runs:
                r.bold = True
        total_val = table.rows[st_row + 2].cells[3]
        total_val.text = f"${data.get('total', 0):,.2f}"
        for p in total_val.paragraphs:
            for r in p.runs:
                r.bold = True

        doc.add_paragraph()
        terms = doc.add_paragraph()
        terms.add_run(f"Payment Terms: {data.get('payment_terms', 'Net 30')}").bold = True
        doc.add_paragraph()
        notes = doc.add_paragraph(data.get("notes", ""))
        notes.alignment = WD_ALIGN_PARAGRAPH.CENTER

        doc.save(filepath)
        return os.path.getsize(filepath)

    @classmethod
    def _create_financial_pdf(cls, filepath: str, data: Dict[str, Any]) -> int:
        """Create a professional PDF invoice/receipt with tables."""
        doc = SimpleDocTemplate(filepath, pagesize=A4)
        styles = getSampleStyleSheet()
        story = []

        company_style = ParagraphStyle('Company', parent=styles['Heading1'], fontSize=22, textColor=colors.HexColor('#003366'), spaceAfter=6)
        subtitle_style = ParagraphStyle('Sub', parent=styles['Normal'], fontSize=10, textColor=colors.grey)
        title_style = ParagraphStyle('InvTitle', parent=styles['Heading1'], fontSize=28, textColor=colors.grey, alignment=2)
        meta_style = ParagraphStyle('Meta', parent=styles['Normal'], fontSize=10, alignment=2)
        bold_style = ParagraphStyle('Bold', parent=styles['Normal'], fontSize=10, fontName='Helvetica-Bold')

        story.append(Paragraph(data.get("company", ""), company_style))
        for line in (data.get("company_address") or "").split('\n'):
            story.append(Paragraph(line, subtitle_style))
        phone = data.get("company_phone", "")
        email = data.get("company_email", "")
        if phone or email:
            story.append(Paragraph(f"Phone: {phone}  |  Email: {email}", subtitle_style))
        story.append(Spacer(1, 20))

        story.append(Paragraph("INVOICE", title_style))
        story.append(Paragraph(f"Invoice #: {data.get('document_number', '')}", meta_style))
        story.append(Paragraph(f"Date: {data.get('document_date', '')}", meta_style))
        story.append(Paragraph(f"Due: {data.get('due_date', '')}", meta_style))
        story.append(Spacer(1, 20))

        story.append(Paragraph("<b>BILL TO:</b>", bold_style))
        story.append(Paragraph(f"<b>{data.get('customer_name', '')}</b>", styles['Normal']))
        for line in (data.get("customer_address") or "").split('\n'):
            story.append(Paragraph(line, styles['Normal']))
        cust_email = data.get("customer_email", "")
        if cust_email:
            story.append(Paragraph(f"Email: {cust_email}", styles['Normal']))
        story.append(Spacer(1, 20))

        # Items table
        items = data.get("items", [])
        table_data = [["Description", "Qty", "Unit Price", "Amount"]]
        for item in items:
            table_data.append([
                item.get("description", ""),
                str(item.get("qty", 0)),
                f"${item.get('price', 0):,.2f}",
                f"${item.get('qty', 0) * item.get('price', 0):,.2f}",
            ])
        table_data.append(["", "", "Subtotal:", f"${data.get('subtotal', 0):,.2f}"])
        table_data.append(["", "", f"Tax ({data.get('tax_rate', '8%')}):", f"${data.get('tax', 0):,.2f}"])
        table_data.append(["", "", "TOTAL:", f"${data.get('total', 0):,.2f}"])

        t = Table(table_data, colWidths=[250, 50, 80, 80])
        t.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#003366')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 10),
            ('ALIGN', (1, 0), (-1, -1), 'RIGHT'),
            ('GRID', (0, 0), (-1, -4), 0.5, colors.grey),
            ('ROWBACKGROUNDS', (0, 1), (-1, -4), [colors.white, colors.HexColor('#f5f5f5')]),
            ('FONTNAME', (2, -1), (-1, -1), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 1), (-1, -1), 9),
            ('TOPPADDING', (0, 0), (-1, -1), 6),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
        ]))
        story.append(t)
        story.append(Spacer(1, 30))

        story.append(Paragraph(f"<b>Payment Terms:</b> {data.get('payment_terms', 'Net 30')}", styles['Normal']))
        story.append(Spacer(1, 12))
        story.append(Paragraph(data.get("notes", ""), ParagraphStyle('Notes', parent=styles['Normal'], alignment=1, textColor=colors.grey)))

        doc.build(story)
        return os.path.getsize(filepath)

    @classmethod
    def _create_financial_xlsx(cls, filepath: str, data: Dict[str, Any]) -> int:
        """Create a professional Excel invoice/receipt."""
        wb = Workbook()
        ws = wb.active
        ws.title = "Invoice"

        header_font = Font(bold=True, size=12, color="FFFFFF")
        header_fill = PatternFill(start_color="003366", end_color="003366", fill_type="solid")
        bold_font = Font(bold=True)
        thin_border = Border(
            left=Side(style='thin'), right=Side(style='thin'),
            top=Side(style='thin'), bottom=Side(style='thin')
        )

        ws['A1'] = data.get("company", "")
        ws['A1'].font = Font(bold=True, size=18, color="003366")
        ws.merge_cells('A1:D1')

        ws['A2'] = (data.get("company_address") or "").replace('\n', ', ')
        ws['A3'] = f"Phone: {data.get('company_phone', '')}  |  Email: {data.get('company_email', '')}"

        ws['C5'] = "INVOICE"
        ws['C5'].font = Font(bold=True, size=16)
        ws['A6'] = "Invoice #:"
        ws['A6'].font = bold_font
        ws['B6'] = data.get("document_number", "")
        ws['A7'] = "Date:"
        ws['A7'].font = bold_font
        ws['B7'] = data.get("document_date", "")
        ws['A8'] = "Due Date:"
        ws['A8'].font = bold_font
        ws['B8'] = data.get("due_date", "")

        ws['A10'] = "Bill To:"
        ws['A10'].font = bold_font
        ws['A11'] = data.get("customer_name", "")
        ws['A11'].font = bold_font
        ws['A12'] = (data.get("customer_address") or "").replace('\n', ', ')

        headers = ["Description", "Qty", "Unit Price", "Amount"]
        for col, h in enumerate(headers, 1):
            cell = ws.cell(row=14, column=col, value=h)
            cell.font = header_font
            cell.fill = header_fill
            cell.alignment = Alignment(horizontal='center')
            cell.border = thin_border

        items = data.get("items", [])
        for idx, item in enumerate(items):
            r = 15 + idx
            ws.cell(row=r, column=1, value=item.get("description", "")).border = thin_border
            ws.cell(row=r, column=2, value=item.get("qty", 0)).border = thin_border
            ws.cell(row=r, column=3, value=f"${item.get('price', 0):,.2f}").border = thin_border
            ws.cell(row=r, column=4, value=f"${item.get('qty', 0) * item.get('price', 0):,.2f}").border = thin_border

        tr = 15 + len(items) + 1
        ws.cell(row=tr, column=3, value="Subtotal:").font = bold_font
        ws.cell(row=tr, column=4, value=f"${data.get('subtotal', 0):,.2f}")
        ws.cell(row=tr + 1, column=3, value=f"Tax ({data.get('tax_rate', '8%')}):").font = bold_font
        ws.cell(row=tr + 1, column=4, value=f"${data.get('tax', 0):,.2f}")
        ws.cell(row=tr + 2, column=3, value="TOTAL:").font = Font(bold=True, size=12)
        ws.cell(row=tr + 2, column=4, value=f"${data.get('total', 0):,.2f}").font = Font(bold=True, size=12)

        ws.column_dimensions['A'].width = 45
        ws.column_dimensions['B'].width = 10
        ws.column_dimensions['C'].width = 15
        ws.column_dimensions['D'].width = 15

        wb.save(filepath)
        return os.path.getsize(filepath)

    # ------------------------------------------------------------------
    # Markdown-based renderers (for general/non-financial docs)
    # ------------------------------------------------------------------

    # ---- Markdown rendering helpers (shared by docx/pptx) ----
    @staticmethod
    def _md_runs(paragraph, text):
        """Add text to a docx paragraph, rendering **bold**, *italic*, `code`."""
        for tok in re.split(r'(\*\*.+?\*\*|\*.+?\*|`.+?`)', text or ''):
            if not tok:
                continue
            if tok.startswith('**') and tok.endswith('**'):
                paragraph.add_run(tok[2:-2]).bold = True
            elif tok.startswith('*') and tok.endswith('*'):
                paragraph.add_run(tok[1:-1]).italic = True
            elif tok.startswith('`') and tok.endswith('`'):
                paragraph.add_run(tok[1:-1])
            else:
                paragraph.add_run(tok)

    @staticmethod
    def _md_strip(text):
        t = re.sub(r'\*\*(.+?)\*\*', r'\1', text or '')
        t = re.sub(r'`(.+?)`', r'\1', t)
        return t.replace('*', '').strip()

    @staticmethod
    def _is_md_table_sep(line):
        s = (line or '').strip()
        return bool(re.match(r'^\|?\s*:?-{2,}:?\s*(\|\s*:?-{2,}:?\s*)+\|?$', s))

    @staticmethod
    def _md_table_cells(line):
        s = (line or '').strip()
        if s.startswith('|'):
            s = s[1:]
        if s.endswith('|'):
            s = s[:-1]
        return [c.strip() for c in s.split('|')]

    @staticmethod
    def _shade_cell(cell, hex_fill):
        """Apply a solid background fill to a table cell (python-docx has no direct API)."""
        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement
        tcPr = cell._tc.get_or_add_tcPr()
        shd = OxmlElement('w:shd')
        shd.set(qn('w:val'), 'clear')
        shd.set(qn('w:color'), 'auto')
        shd.set(qn('w:fill'), hex_fill)
        tcPr.append(shd)

    # Header fill + zebra stripe used for all generated tables.
    _TABLE_HEADER_FILL = '1F3A5F'
    _TABLE_ZEBRA_FILL = 'EEF2F9'

    @classmethod
    def _docx_table(cls, doc, header, rows):
        from docx.shared import RGBColor
        ncols = max(len(header), max((len(r) for r in rows), default=0)) or 1
        table = doc.add_table(rows=1, cols=ncols)
        try:
            table.style = 'Table Grid'  # clean thin borders; we colour cells ourselves
        except Exception:
            pass
        # Header row: navy fill + white bold text (readable, professional).
        hdr = table.rows[0].cells
        for c in range(ncols):
            hdr[c].text = ''
            cls._shade_cell(hdr[c], cls._TABLE_HEADER_FILL)
            cls._md_runs(hdr[c].paragraphs[0], header[c] if c < len(header) else '')
            for run in hdr[c].paragraphs[0].runs:
                run.bold = True
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        # Body rows: zebra striping for readability.
        for ri, r in enumerate(rows):
            cells = table.add_row().cells
            for c in range(ncols):
                cells[c].text = ''
                if ri % 2 == 1:
                    cls._shade_cell(cells[c], cls._TABLE_ZEBRA_FILL)
                cls._md_runs(cells[c].paragraphs[0], r[c] if c < len(r) else '')
        return table

    @classmethod
    def _create_docx(cls, filepath: str, title: str, content: str) -> int:
        """Create a Word document from markdown content (native tables + bold)."""
        doc = Document()
        doc.add_heading(title, 0)

        lines = content.split('\n')
        i = 0
        while i < len(lines):
            line = lines[i].strip()
            if not line:
                i += 1
                continue

            # Skip markdown horizontal rules (---, ***, ___) — they otherwise print literally.
            if re.match(r'^(-{3,}|\*{3,}|_{3,})$', line):
                i += 1
                continue

            if line.startswith('|') and i + 1 < len(lines) and cls._is_md_table_sep(lines[i + 1]):
                header = cls._md_table_cells(line)
                body = []
                j = i + 2
                while j < len(lines) and lines[j].strip().startswith('|'):
                    body.append(cls._md_table_cells(lines[j]))
                    j += 1
                cls._docx_table(doc, header, body)
                i = j
                continue

            if line.startswith('### '):
                doc.add_heading(line[4:], 3)
            elif line.startswith('## '):
                doc.add_heading(line[3:], 2)
            elif line.startswith('# '):
                doc.add_heading(line[2:], 1)
            elif line.startswith('- ') or line.startswith('* '):
                cls._md_runs(doc.add_paragraph(style='List Bullet'), line[2:])
            elif re.match(r'^\d+\.\s', line):
                cls._md_runs(doc.add_paragraph(style='List Number'), re.sub(r'^\d+\.\s', '', line))
            else:
                cls._md_runs(doc.add_paragraph(), line)
            i += 1

        doc.save(filepath)
        return os.path.getsize(filepath)
    
    @staticmethod
    def _md_to_rl(text):
        """Convert inline markdown to ReportLab markup (escapes &<> then **bold**, *italic*, `code`)."""
        import html as _html
        t = _html.escape(text or '')
        t = re.sub(r'\*\*(.+?)\*\*', r'<b>\1</b>', t)
        t = re.sub(r'`(.+?)`', r'<font face="Courier">\1</font>', t)
        t = re.sub(r'(?<!\*)\*(?!\*)([^*]+?)\*(?!\*)', r'<i>\1</i>', t)
        return t

    @classmethod
    def _create_pdf(cls, filepath: str, title: str, content: str) -> int:
        """Create a PDF document from markdown (renders bold, bullets and tables)."""
        from reportlab.platypus import Table, TableStyle
        from reportlab.lib import colors
        doc = SimpleDocTemplate(filepath, pagesize=A4)
        styles = getSampleStyleSheet()
        story = []

        title_style = ParagraphStyle('CustomTitle', parent=styles['Heading1'], fontSize=24, spaceAfter=24)
        story.append(Paragraph(cls._md_to_rl(title), title_style))
        story.append(Spacer(1, 12))

        lines = content.split('\n')
        i = 0
        while i < len(lines):
            line = lines[i].strip()
            if not line:
                story.append(Spacer(1, 8)); i += 1; continue
            # Skip horizontal rules
            if re.match(r'^(-{3,}|\*{3,}|_{3,})$', line):
                i += 1; continue
            # Markdown table -> ReportLab Table
            if line.startswith('|') and i + 1 < len(lines) and cls._is_md_table_sep(lines[i + 1]):
                header = cls._md_table_cells(line)
                rows = []
                j = i + 2
                while j < len(lines) and lines[j].strip().startswith('|'):
                    rows.append(cls._md_table_cells(lines[j])); j += 1
                cell = ParagraphStyle('cell', parent=styles['Normal'], fontSize=9, leading=12)
                data = [[Paragraph(cls._md_to_rl(c), cell) for c in header]] + \
                       [[Paragraph(cls._md_to_rl(c), cell) for c in r] for r in rows]
                tbl = Table(data, hAlign='LEFT')
                tbl.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#1f3a5f')),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor('#cbd5e1')),
                    ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.HexColor('#f3f6fb')]),
                    ('VALIGN', (0, 0), (-1, -1), 'TOP'),
                    ('LEFTPADDING', (0, 0), (-1, -1), 6), ('RIGHTPADDING', (0, 0), (-1, -1), 6),
                    ('TOPPADDING', (0, 0), (-1, -1), 4), ('BOTTOMPADDING', (0, 0), (-1, -1), 4),
                ]))
                story.append(tbl); story.append(Spacer(1, 10)); i = j; continue

            if line.startswith('### '):
                story.append(Paragraph(cls._md_to_rl(line[4:]), styles['Heading3']))
            elif line.startswith('## '):
                story.append(Paragraph(cls._md_to_rl(line[3:]), styles['Heading2']))
            elif line.startswith('# '):
                story.append(Paragraph(cls._md_to_rl(line[2:]), styles['Heading1']))
            elif line.startswith('- ') or line.startswith('* '):
                story.append(Paragraph(cls._md_to_rl(line[2:]), styles['Normal'], bulletText='•'))
            elif re.match(r'^\d+\.\s', line):
                story.append(Paragraph(cls._md_to_rl(re.sub(r'^\d+\.\s', '', line)), styles['Normal'], bulletText='•'))
            else:
                story.append(Paragraph(cls._md_to_rl(line), styles['Normal']))
            i += 1

        doc.build(story)
        return os.path.getsize(filepath)
    
    @classmethod
    def _create_xlsx(cls, filepath: str, title: str, content: str) -> int:
        """Create an Excel document"""
        wb = Workbook()
        ws = wb.active
        ws.title = title[:31]  # Excel sheet names limited to 31 chars
        
        # Style
        header_font = Font(bold=True, size=14)
        header_fill = PatternFill(start_color="366092", end_color="366092", fill_type="solid")
        
        # Title
        ws['A1'] = title
        ws['A1'].font = Font(bold=True, size=18)
        ws.merge_cells('A1:E1')
        
        # Parse content and try to create table-like structure
        row = 3
        for line in content.split('\n'):
            line = line.strip()
            if not line:
                continue
            
            if line.startswith('|'):
                # Table row
                cells = [c.strip() for c in line.split('|')[1:-1]]
                for col, cell in enumerate(cells, 1):
                    ws.cell(row=row, column=col, value=cell)
                row += 1
            elif line.startswith('#'):
                # Header
                ws.cell(row=row, column=1, value=line.lstrip('#').strip())
                ws.cell(row=row, column=1).font = header_font
                row += 1
            else:
                ws.cell(row=row, column=1, value=line)
                row += 1
        
        wb.save(filepath)
        return os.path.getsize(filepath)
    
    @staticmethod
    def _pptx_runs(paragraph, text):
        for tok in re.split(r'(\*\*.+?\*\*|\*.+?\*|`.+?`)', text or ''):
            if not tok:
                continue
            run = paragraph.add_run()
            if tok.startswith('**') and tok.endswith('**'):
                run.text = tok[2:-2]; run.font.bold = True
            elif tok.startswith('*') and tok.endswith('*'):
                run.text = tok[1:-1]; run.font.italic = True
            elif tok.startswith('`') and tok.endswith('`'):
                run.text = tok[1:-1]
            else:
                run.text = tok

    @classmethod
    def _pptx_table_slide(cls, prs, title, header, rows):
        try:
            slide = prs.slides.add_slide(prs.slide_layouts[5])
        except Exception:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
        try:
            if slide.shapes.title:
                slide.shapes.title.text = cls._md_strip(title) or 'Details'
        except Exception:
            pass
        ncols = max(len(header), max((len(r) for r in rows), default=1)) or 1
        nrows = 1 + len(rows)
        left = PptxInches(0.5); top = PptxInches(1.6)
        width = PptxInches(9); height = PptxInches(min(0.4 * nrows + 0.3, 5))
        tbl = slide.shapes.add_table(nrows, ncols, left, top, width, height).table
        for c in range(ncols):
            tbl.cell(0, c).text = cls._md_strip(header[c]) if c < len(header) else ''
        for ri, r in enumerate(rows, start=1):
            for c in range(ncols):
                tbl.cell(ri, c).text = cls._md_strip(r[c]) if c < len(r) else ''
        return tbl

    @classmethod
    def _create_pptx(cls, filepath: str, title: str, content: str) -> int:
        """Create a PowerPoint from markdown (native tables + bold bullets)."""
        prs = Presentation()

        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
        try:
            slide.placeholders[1].text = datetime.utcnow().strftime("%B %Y")
        except Exception:
            pass

        bullet_layout = prs.slide_layouts[1]
        state = {'slide': None, 'points': []}

        def flush():
            if state['slide'] and state['points']:
                tf = state['slide'].shapes.placeholders[1].text_frame
                tf.clear()
                for idx, point in enumerate(state['points']):
                    p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                    cls._pptx_runs(p, point)
            state['points'] = []

        def ensure_slide():
            if not state['slide']:
                state['slide'] = prs.slides.add_slide(bullet_layout)
                state['slide'].shapes.title.text = title

        lines = content.split('\n')
        i = 0
        while i < len(lines):
            line = lines[i].strip()
            if not line:
                i += 1
                continue
            if line.startswith('|') and i + 1 < len(lines) and cls._is_md_table_sep(lines[i + 1]):
                flush()
                header = cls._md_table_cells(line)
                body = []
                j = i + 2
                while j < len(lines) and lines[j].strip().startswith('|'):
                    body.append(cls._md_table_cells(lines[j]))
                    j += 1
                ttl = state['slide'].shapes.title.text if state['slide'] else title
                cls._pptx_table_slide(prs, ttl, header, body)
                i = j
                continue
            if line.startswith('# ') or line.startswith('## '):
                flush()
                state['slide'] = prs.slides.add_slide(bullet_layout)
                state['slide'].shapes.title.text = cls._md_strip(line.lstrip('#').strip())
            elif line.startswith('### '):
                ensure_slide()
                state['points'].append(cls._md_strip(line[4:]))
            elif line.startswith('- ') or line.startswith('* '):
                ensure_slide()
                state['points'].append(line[2:])
            else:
                ensure_slide()
                state['points'].append(line)
            i += 1

        flush()
        prs.save(filepath)
        return os.path.getsize(filepath)
    
    @classmethod
    async def generate_image(
        cls,
        name: str,
        description: str,
        document_type: str = "invoice",
        format: str = "png",
        org_id: str = None
    ) -> Dict[str, Any]:
        """
        Generate a document image for OCR testing
        """
        cls._ensure_storage()
        
        item_id = str(uuid.uuid4())
        filename = f"{item_id}.{format}"
        filepath = os.path.join(cls.STORAGE_PATH, filename)
        
        # Generate image
        if PIL_AVAILABLE:
            size = await cls._create_document_image(filepath, name, description, document_type, format)
        else:
            # Create placeholder
            size = 0
        # NOTE: Do not inline base64 in the JSON response by default.
        # The UI can display images via `image_url`, which is faster and cacheable.
        base64_data = None
        
        result = {
            "id": item_id,
            "name": name,
            "format": format,
            "document_type": document_type,
            "size": size,
            "download_url": f"/api/lab/download/{item_id}",
            "image_url": f"/api/lab/image/{item_id}",
            "base64": base64_data,
            "created_at": datetime.utcnow().isoformat(),
            "org_id": org_id
        }
        
        # Save metadata
        with open(os.path.join(cls.STORAGE_PATH, f"{item_id}_meta.json"), 'w') as f:
            json.dump({k: v for k, v in result.items() if k != 'base64'}, f, indent=2)
        
        return result
    
    @classmethod
    async def _create_document_image(
        cls,
        filepath: str,
        name: str,
        description: str,
        document_type: str,
        format: str
    ) -> int:
        """Create a document image using PIL with AI-generated content"""
        
        # Get structured content from AI based on description
        doc_data = await cls._generate_structured_image_content(description, document_type)
        
        # Image dimensions (A4-like ratio)
        width = 800
        height = 1100
        
        # Create image
        img = Image.new('RGB', (width, height), color='white')
        draw = ImageDraw.Draw(img)
        
        # Try to load fonts
        try:
            font_large = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 24)
            font_medium = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 16)
            font_small = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 12)
            font_title = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 28)
        except:
            font_large = ImageFont.load_default()
            font_medium = font_large
            font_small = font_large
            font_title = font_large
        
        margin = 50
        y = 40
        
        # Use AI-generated data
        title = doc_data.get('title', name.upper())
        subtitle = doc_data.get('subtitle', '')
        header_color = doc_data.get('header_color', '#1e3a5f')
        sections = doc_data.get('sections', [])
        items = doc_data.get('items', [])
        footer = doc_data.get('footer', '')
        
        # Draw header
        draw.rectangle([0, 0, width, 120], fill=header_color)
        draw.text((margin, 35), title[:30], font=font_title, fill='white')
        if subtitle:
            draw.text((margin, 75), subtitle[:50], font=font_small, fill='#cccccc')
        
        # Draw reference number and date on right
        ref_num = doc_data.get('reference', f"#{cls._random_number(6)}")
        draw.text((width - 200, 35), ref_num, font=font_medium, fill='white')
        draw.text((width - 200, 60), datetime.now().strftime("%Y-%m-%d"), font=font_small, fill='#cccccc')
        
        y = 150
        
        # Draw sections (From, To, etc.)
        for section in sections[:4]:
            section_title = section.get('title', '')
            section_lines = section.get('lines', [])
            
            draw.text((margin, y), section_title, font=font_medium, fill='black')
            y += 25
            
            for line in section_lines[:3]:
                draw.text((margin + 10, y), str(line)[:60], font=font_small, fill='#666666')
                y += 20
            
            y += 15
        
        # Draw items table if present
        columns = doc_data.get('columns', [])
        if items or columns:
            # Ensure we have columns
            if not columns and items:
                # Try to get columns from first item if it's a dict
                if isinstance(items[0], dict):
                    columns = list(items[0].keys())[:4]
                else:
                    columns = [f"Column {i+1}" for i in range(min(4, len(items[0]) if items else 4))]
            
            num_cols = len(columns) if columns else 4
            col_width = (width - 2 * margin) // num_cols
            
            # Table header
            draw.rectangle([margin, y, width - margin, y + 30], fill='#f0f0f0')
            
            for i, col in enumerate(columns[:4]):
                draw.text((margin + 10 + i * col_width, y + 8), str(col)[:18], font=font_small, fill='black')
            
            y += 35
            
            # Table rows
            calculated_total = 0
            for item in items[:8]:
                # Handle both list and dict formats
                if isinstance(item, list):
                    values = item[:4]
                elif isinstance(item, dict):
                    values = list(item.values())[:4]
                else:
                    values = [str(item)]
                
                for i, val in enumerate(values):
                    text = str(val)[:20]
                    draw.text((margin + 10 + i * col_width, y), text, font=font_small, fill='#333333')
                    
                    # Try to extract numeric for total (last column usually)
                    if i == len(values) - 1:
                        try:
                            num_val = float(str(val).replace('$', '').replace(',', ''))
                            calculated_total += num_val
                        except:
                            pass
                
                draw.line([(margin, y + 22), (width - margin, y + 22)], fill='#eeeeee')
                y += 28
            
            # Total - use provided total or calculated
            total_text = doc_data.get('total', '')
            if total_text:
                y += 15
                draw.rectangle([width - 280, y, width - margin, y + 40], fill=header_color)
                draw.text((width - 270, y + 10), f"TOTAL: {total_text}", font=font_medium, fill='white')
                y += 50
            elif calculated_total > 0:
                y += 15
                draw.rectangle([width - 280, y, width - margin, y + 40], fill=header_color)
                draw.text((width - 270, y + 10), f"TOTAL: ${calculated_total:,.2f}", font=font_medium, fill='white')
                y += 50
        
        # Draw additional content/notes
        notes = doc_data.get('notes', [])
        if notes:
            y += 20
            draw.text((margin, y), "Notes:", font=font_medium, fill='black')
            y += 25
            for note in notes[:5]:
                draw.text((margin + 10, y), f"• {str(note)[:70]}", font=font_small, fill='#666666')
                y += 20
        
        # Footer
        draw.line([(0, height - 60), (width, height - 60)], fill='#eeeeee')
        footer_text = footer or doc_data.get('company', doc_data.get('subtitle', ''))
        draw.text((margin, height - 45), footer_text[:60], font=font_small, fill='#999999')
        
        # Save
        if format == 'jpg':
            img = img.convert('RGB')
            # Web-friendly JPEG: smaller + progressive
            img.save(
                filepath,
                'JPEG',
                quality=85,
                optimize=True,
                progressive=True,
            )
        else:
            # Web-friendly PNG: max compression + optimization.
            # Keep RGB first for best OCR; only quantize if still too large.
            try:
                from io import BytesIO
                buf = BytesIO()
                img.save(buf, 'PNG', optimize=True, compress_level=9)
                data = buf.getvalue()
                # If still large, quantize without dithering to avoid introducing OCR artifacts.
                if len(data) > 350_000:
                    q = img.quantize(colors=256, dither=Image.Dither.NONE)
                    q.save(filepath, 'PNG', optimize=True, compress_level=9)
                else:
                    with open(filepath, 'wb') as f:
                        f.write(data)
            except Exception:
                img.save(filepath, 'PNG', optimize=True, compress_level=9)
        
        return os.path.getsize(filepath)
    
    @classmethod
    async def _generate_structured_image_content(cls, description: str, document_type: str) -> Dict[str, Any]:
        """Generate structured content for image using AI"""
        client = cls._get_ai_client()
        
        prompt = f"""Generate realistic document data for this request:

Description: {description}
Document Type: {document_type}

Return a JSON object with this EXACT structure:
{{
    "title": "Main document title matching the description (e.g., 'AIRLINE TICKET', 'BOARDING PASS', 'FLIGHT INVOICE')",
    "subtitle": "Company or organization name",
    "reference": "Unique reference number (realistic format)",
    "header_color": "Hex color matching the brand (e.g., '#c60c30' for Emirates, '#00205b' for Delta)",
    "sections": [
        {{"title": "Section Title", "lines": ["Detail line 1", "Detail line 2", "Detail line 3"]}}
    ],
    "columns": ["Column Header 1", "Column Header 2", "Column Header 3", "Column Header 4"],
    "items": [
        ["Row 1 Col 1", "Row 1 Col 2", "Row 1 Col 3", "Row 1 Col 4"],
        ["Row 2 Col 1", "Row 2 Col 2", "Row 2 Col 3", "Row 2 Col 4"]
    ],
    "total": "Total amount if applicable (e.g., '$1,250.00')",
    "notes": ["Note 1", "Note 2"],
    "footer": "Footer text"
}}

IMPORTANT:
- "columns" is a list of column header NAMES (strings)
- "items" is a list of ROWS, each row is a list of VALUES (not objects)
- Make all data realistic and contextual to the description
- For airline tickets: include passenger name, flight number, route, date/time, seat, class
- For invoices: include line items with quantities and prices

Return ONLY the JSON, no explanation or markdown."""

        if client:
            try:
                response = client.chat.completions.create(
                    model="gpt-4o-mini",
                    messages=[
                        {"role": "system", "content": "You generate realistic document data in JSON format. Be accurate and detailed."},
                        {"role": "user", "content": prompt}
                    ],
                    temperature=0.7,
                    max_tokens=1500
                )
                
                content = response.choices[0].message.content.strip()
                
                # Extract JSON from response
                if '```' in content:
                    content = content.split('```')[1]
                    if content.startswith('json'):
                        content = content[4:]
                
                data = json.loads(content)
                return data
                
            except Exception as e:
                print(f"AI structured content generation failed: {e}")
        
        # Fallback - basic structure based on document type
        return cls._get_fallback_structure(description, document_type)
    
    @classmethod
    def _get_fallback_structure(cls, description: str, document_type: str) -> Dict[str, Any]:
        """Fallback structure when AI is not available"""
        import random
        desc_lower = description.lower()

        first_names = ["Ahmed", "Sarah", "Michael", "Emily", "James", "Maria", "David", "Olivia"]
        last_names = ["Al-Rashid", "Johnson", "Chen", "Garcia", "Müller", "Williams", "Patel", "Santos"]
        companies = ["Apex Technologies Inc.", "Meridian Partners", "Blue Ocean Trading",
                      "Pinnacle Holdings", "Vantage Consulting", "Silverline Financial"]

        def _name():
            return f"{random.choice(first_names)} {random.choice(last_names)}"

        def _company():
            return random.choice(companies)

        if 'airline' in desc_lower or 'flight' in desc_lower or 'ticket' in desc_lower:
            passenger = _name()
            return {
                "title": "AIRLINE TICKET",
                "subtitle": "E-Ticket Receipt",
                "reference": f"TKT-{cls._random_number(10)}",
                "header_color": "#0066cc",
                "sections": [
                    {"title": "Passenger", "lines": [passenger, f"Passport: {chr(random.randint(65,90))}{chr(random.randint(65,90))}{cls._random_number(7)}"]},
                    {"title": "Flight Details", "lines": [f"Flight: EK {random.randint(100,999)}", f"Date: {datetime.now().strftime('%Y-%m-%d')}", "Route: DXB \u2192 LHR"]}
                ],
                "columns": ["Segment", "Class", "Status", "Fare"],
                "items": [
                    ["DXB - LHR", "Business", "Confirmed", "$1,850.00"],
                    ["Taxes & Fees", "-", "-", "$124.50"]
                ],
                "total": "$1,974.50",
                "notes": ["Check-in opens 24 hours before departure", "Baggage allowance: 2 x 32kg"],
                "footer": "Thank you for flying with us"
            }
        elif 'receipt' in desc_lower:
            store = _company()
            return {
                "title": "RECEIPT",
                "subtitle": store,
                "reference": f"RCP-{cls._random_number(8)}",
                "header_color": "#2d2d2d",
                "sections": [
                    {"title": "Store", "lines": [store, f"{random.randint(100,999)} Market Street, Suite {random.randint(100,400)}"]}
                ],
                "columns": ["Item", "Qty", "Price", "Total"],
                "items": [
                    ["Wireless Keyboard", "1", "$79.99", "$79.99"],
                    ["USB-C Cable (2m)", "3", "$12.50", "$37.50"],
                    ["Monitor Stand", "1", "$45.00", "$45.00"]
                ],
                "total": "$162.49",
                "notes": ["Returns accepted within 30 days with receipt"],
                "footer": store
            }
        elif 'invoice' in desc_lower:
            company = _company()
            customer = _name()
            return {
                "title": "INVOICE",
                "subtitle": company,
                "reference": f"INV-{cls._random_number(8)}",
                "header_color": "#1e3a5f",
                "sections": [
                    {"title": "Bill To", "lines": [customer, f"{random.randint(100,9999)} Commerce Blvd, Suite {random.randint(100,800)}", "New York, NY 10017"]}
                ],
                "columns": ["Description", "Qty", "Unit Price", "Amount"],
                "items": [
                    ["Consulting Services - Phase 1", "40", "$150.00", "$6,000.00"],
                    ["Software License (Annual)", "1", "$2,400.00", "$2,400.00"],
                    ["Technical Support Plan", "12", "$200.00", "$2,400.00"]
                ],
                "total": "$10,800.00",
                "notes": ["Payment due within 30 days", "Wire transfer or ACH preferred"],
                "footer": company
            }
        else:
            company = _company()
            return {
                "title": document_type.upper() if document_type else "DOCUMENT",
                "subtitle": company,
                "reference": f"DOC-{cls._random_number(6)}",
                "header_color": "#1e3a5f",
                "sections": [
                    {"title": "Details", "lines": [description[:100] if description else "Refer to the attached documentation for full details."]}
                ],
                "columns": [],
                "items": [],
                "total": "",
                "notes": [],
                "footer": company
            }
    
    @classmethod
    async def _generate_image_content(cls, description: str, document_type: str) -> str:
        """Generate content for image using AI"""
        client = cls._get_ai_client()
        
        if client:
            try:
                response = client.chat.completions.create(
                    model="gpt-4o-mini",
                    messages=[
                        {"role": "system", "content": f"Generate realistic content for a {document_type} document. Be brief and use real-looking data."},
                        {"role": "user", "content": description}
                    ],
                    temperature=0.7,
                    max_tokens=500
                )
                return response.choices[0].message.content
            except:
                pass
        
        return description
    
    @classmethod
    def _random_number(cls, digits: int, as_int: bool = False):
        """Generate a random number string"""
        import random
        num = ''.join([str(random.randint(0, 9)) for _ in range(digits)])
        return int(num) if as_int else num
    
    @classmethod
    def get_file(cls, item_id: str) -> Optional[str]:
        """Get the file path for a generated item"""
        cls._ensure_storage()
        
        # Look for the file
        for ext in ['json', 'docx', 'pdf', 'xlsx', 'pptx', 'txt', 'png', 'jpg']:
            filepath = os.path.join(cls.STORAGE_PATH, f"{item_id}.{ext}")
            if os.path.exists(filepath):
                return filepath
        
        return None
    
    @classmethod
    def get_mock_data(cls, item_id: str) -> Optional[List[Dict]]:
        """Get mock API data by item id (UUID) or by slug. DB first, then file."""
        # 1) Try DB by id (UUID)
        if len(item_id) == 36 and item_id.count("-") == 4:
            data = cls._get_mock_data_from_db_by_id(item_id)
            if data is not None:
                return data
        # 2) Try DB by slug
        data = cls._get_mock_data_from_db_by_slug(item_id)
        if data is not None:
            return data
        # 3) Fallback: file by id
        if len(item_id) == 36 and item_id.count("-") == 4:
            filepath = os.path.join(cls.STORAGE_PATH, f"{item_id}.json")
            if os.path.exists(filepath):
                with open(filepath, 'r') as f:
                    obj = json.load(f)
                    return obj.get("data", [])
            return None
        # 4) Fallback: file by slug
        return cls.get_mock_data_by_slug(item_id)

    @classmethod
    def get_mock_api_record(cls, item_id: str) -> Optional[Dict[str, Any]]:
        """Return mock API record (user_id, id) from DB for access check. None if file-only or not found."""
        if not item_id or item_id in ("null", "undefined", "None"):
            return None
        try:
            from database.base import get_db_session
            from database.models import LabMockAPI
            with get_db_session() as db:
                # Try by id (UUID)
                if len(item_id) == 36 and item_id.count("-") == 4:
                    row = db.query(LabMockAPI).filter(LabMockAPI.id == item_id).first()
                    if row:
                        return {"user_id": str(row.user_id) if row.user_id else None, "id": str(row.id)}
                # Try by slug
                row = (
                    db.query(LabMockAPI)
                    .filter(LabMockAPI.slug == item_id)
                    .order_by(LabMockAPI.created_at.desc())
                    .first()
                )
                if row:
                    return {"user_id": str(row.user_id) if row.user_id else None, "id": str(row.id)}
        except Exception as e:
            print(f"⚠️ [LAB] get_mock_api_record failed: {e}")
        return None

    @classmethod
    def _get_mock_data_from_db_by_id(cls, item_id: str) -> Optional[List[Dict]]:
        """Load mock API data from DB by primary key id."""
        try:
            from database.base import get_db_session
            from database.models import LabMockAPI
            with get_db_session() as db:
                row = db.query(LabMockAPI).filter(LabMockAPI.id == item_id).first()
                if row and row.data is not None:
                    return list(row.data) if isinstance(row.data, list) else []
        except Exception as e:
            print(f"⚠️ [LAB] DB get by id failed: {e}")
        return None

    @classmethod
    def _get_mock_data_from_db_by_slug(cls, slug: str) -> Optional[List[Dict]]:
        """Load mock API data from DB by slug (most recent when multiple)."""
        if not slug or slug in ("null", "undefined", "None"):
            return None
        try:
            from database.base import get_db_session
            from database.models import LabMockAPI
            with get_db_session() as db:
                row = (
                    db.query(LabMockAPI)
                    .filter(LabMockAPI.slug == slug)
                    .order_by(LabMockAPI.created_at.desc())
                    .first()
                )
                if row and row.data is not None:
                    return list(row.data) if isinstance(row.data, list) else []
        except Exception as e:
            print(f"⚠️ [LAB] DB get by slug failed: {e}")
        return None

    @classmethod
    def get_mock_data_by_slug(cls, slug: str) -> Optional[List[Dict]]:
        """Get mock API data by URL slug from DB first, then from files."""
        if not slug or slug in ("null", "undefined", "None"):
            return None
        data = cls._get_mock_data_from_db_by_slug(slug)
        if data is not None:
            return data
        cls._ensure_storage()
        for f in os.listdir(cls.STORAGE_PATH):
            if not f.endswith(".json"):
                continue
            filepath = os.path.join(cls.STORAGE_PATH, f)
            try:
                with open(filepath, 'r') as file:
                    obj = json.load(file)
                    if obj.get("slug") == slug:
                        return obj.get("data", [])
            except Exception:
                continue
        return None
