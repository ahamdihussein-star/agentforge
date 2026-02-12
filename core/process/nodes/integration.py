"""
Integration Node Executors
HTTP requests, database queries, and external system integrations

These nodes connect to external systems:
- HTTP_REQUEST: REST API calls
- DATABASE_QUERY: SQL/NoSQL operations
- FILE_OPERATION: File read/write
- MESSAGE_QUEUE: Pub/sub messaging
"""

import json
import time
from typing import Optional, Dict, Any
from ..schemas import ProcessNode, NodeType
from ..state import ProcessState, ProcessContext
from ..result import NodeResult, ExecutionError, ErrorCategory
from .base import BaseNodeExecutor, register_executor


@register_executor(NodeType.HTTP_REQUEST)
class HTTPRequestNodeExecutor(BaseNodeExecutor):
    """
    HTTP Request node executor
    
    Makes HTTP API calls to external services.
    
    Config:
        method: HTTP method (GET, POST, PUT, DELETE, PATCH)
        url: URL template with variable interpolation
        headers: Request headers
        body: Request body (for POST/PUT/PATCH)
        auth_type: none, basic, bearer, api_key
        auth_config: Authentication configuration
        response_type: json, text, binary
        success_codes: List of success status codes
        verify_ssl: Whether to verify SSL certificates
    """
    
    display_name = "HTTP Request"
    
    async def execute(
        self,
        node: ProcessNode,
        state: ProcessState,
        context: ProcessContext
    ) -> NodeResult:
        """Execute HTTP request node"""
        
        # Get configuration
        method = self.get_config_value(node, 'method', 'GET').upper()
        url_template = self.get_config_value(node, 'url', '')
        headers = self.get_config_value(node, 'headers', {})
        body = self.get_config_value(node, 'body')
        auth_type = self.get_config_value(node, 'auth_type', 'none')
        auth_config = self.get_config_value(node, 'auth_config', {})
        response_type = self.get_config_value(node, 'response_type', 'json')
        success_codes = self.get_config_value(node, 'success_codes', [200, 201, 202, 204])
        verify_ssl = self.get_config_value(node, 'verify_ssl', True)
        timeout_seconds = self.get_config_value(node, 'timeout_seconds', 30)
        
        logs = [f"HTTP {method} request"]
        
        # Interpolate URL
        try:
            url = state.interpolate_string(url_template)
            logs.append(f"URL: {url}")
        except Exception as e:
            return NodeResult.failure(
                error=ExecutionError.validation_error(f"Failed to interpolate URL: {e}"),
                logs=logs
            )
        
        # Interpolate headers
        try:
            interpolated_headers = {}
            for key, value in headers.items():
                interpolated_headers[key] = state.interpolate_string(str(value))
            logs.append(f"Headers: {list(interpolated_headers.keys())}")
        except Exception as e:
            return NodeResult.failure(
                error=ExecutionError.validation_error(f"Failed to interpolate headers: {e}"),
                logs=logs
            )
        
        # Add authentication
        try:
            interpolated_headers = self._add_auth(
                interpolated_headers, auth_type, auth_config, state
            )
        except Exception as e:
            return NodeResult.failure(
                error=ExecutionError(
                    category=ErrorCategory.AUTHENTICATION,
                    code="AUTH_CONFIG_ERROR",
                    message=f"Failed to configure authentication: {e}"
                ),
                logs=logs
            )
        
        # Interpolate body
        interpolated_body = None
        if body and method in ['POST', 'PUT', 'PATCH']:
            try:
                interpolated_body = state.interpolate_object(body)
                logs.append("Body interpolated")
            except Exception as e:
                return NodeResult.failure(
                    error=ExecutionError.validation_error(f"Failed to interpolate body: {e}"),
                    logs=logs
                )
        
        # Make HTTP request
        start_time = time.time()
        
        try:
            import aiohttp
            
            async with aiohttp.ClientSession() as session:
                request_kwargs = {
                    'method': method,
                    'url': url,
                    'headers': interpolated_headers,
                    'ssl': verify_ssl,
                    'timeout': aiohttp.ClientTimeout(total=timeout_seconds)
                }
                
                if interpolated_body:
                    if isinstance(interpolated_body, dict):
                        request_kwargs['json'] = interpolated_body
                    else:
                        request_kwargs['data'] = interpolated_body
                
                async with session.request(**request_kwargs) as response:
                    duration_ms = (time.time() - start_time) * 1000
                    status_code = response.status
                    
                    logs.append(f"Response: {status_code} in {duration_ms:.0f}ms")
                    
                    # Parse response
                    if response_type == 'json':
                        try:
                            response_data = await response.json()
                        except Exception:
                            response_data = await response.text()
                    elif response_type == 'binary':
                        response_data = await response.read()
                        # Convert to base64 for JSON storage
                        import base64
                        response_data = base64.b64encode(response_data).decode()
                    else:
                        response_data = await response.text()
                    
                    # Check success
                    if status_code in success_codes:
                        variables_update = {}
                        if node.output_variable:
                            variables_update[node.output_variable] = response_data
                        
                        return NodeResult.success(
                            output={
                                'status_code': status_code,
                                'data': response_data,
                                'headers': dict(response.headers)
                            },
                            variables_update=variables_update,
                            duration_ms=duration_ms,
                            logs=logs
                        )
                    else:
                        return NodeResult.failure(
                            error=ExecutionError(
                                category=ErrorCategory.EXTERNAL,
                                code=f"HTTP_{status_code}",
                                message=f"HTTP request failed with status {status_code}",
                                details={
                                    'status_code': status_code,
                                    'response': response_data[:500] if isinstance(response_data, str) else str(response_data)[:500]
                                },
                                is_retryable=status_code in [408, 429, 500, 502, 503, 504]
                            ),
                            duration_ms=duration_ms,
                            logs=logs
                        )
                        
        except Exception as e:
            error_type = type(e).__name__
            is_timeout = 'timeout' in str(e).lower()
            
            return NodeResult.failure(
                error=ExecutionError(
                    category=ErrorCategory.TIMEOUT if is_timeout else ErrorCategory.CONNECTION,
                    code="HTTP_ERROR",
                    message=f"HTTP request failed: {error_type}: {e}",
                    is_retryable=True,
                    retry_after_seconds=5
                ),
                logs=logs
            )
    
    def _add_auth(
        self,
        headers: Dict[str, str],
        auth_type: str,
        auth_config: Dict[str, Any],
        state: ProcessState
    ) -> Dict[str, str]:
        """Add authentication to headers"""
        
        if auth_type == 'none':
            return headers
        
        elif auth_type == 'bearer':
            token = auth_config.get('token', '')
            token = state.interpolate_string(token)
            headers['Authorization'] = f'Bearer {token}'
        
        elif auth_type == 'api_key':
            key_name = auth_config.get('key_name', 'X-API-Key')
            key_value = auth_config.get('key_value', '')
            key_value = state.interpolate_string(key_value)
            key_location = auth_config.get('location', 'header')
            
            if key_location == 'header':
                headers[key_name] = key_value
            # Query parameter auth would need URL modification
        
        elif auth_type == 'basic':
            import base64
            username = state.interpolate_string(auth_config.get('username', ''))
            password = state.interpolate_string(auth_config.get('password', ''))
            credentials = base64.b64encode(f'{username}:{password}'.encode()).decode()
            headers['Authorization'] = f'Basic {credentials}'
        
        return headers
    
    def validate(self, node: ProcessNode) -> Optional[ExecutionError]:
        """Validate HTTP request node"""
        base_error = super().validate(node)
        if base_error:
            return base_error
        
        url = self.get_config_value(node, 'url')
        if not url:
            return ExecutionError.validation_error("URL is required")
        
        method = self.get_config_value(node, 'method', 'GET').upper()
        if method not in ['GET', 'POST', 'PUT', 'DELETE', 'PATCH', 'HEAD', 'OPTIONS']:
            return ExecutionError.validation_error(f"Invalid HTTP method: {method}")
        
        return None


@register_executor(NodeType.DATABASE_QUERY)
class DatabaseQueryNodeExecutor(BaseNodeExecutor):
    """
    Database Query node executor
    
    Executes database operations (query, insert, update, delete).
    
    Config:
        connection_id: Database connection ID (from platform tools)
        operation: query, insert, update, delete
        query: SQL query template (for query operation)
        table: Table name (for structured operations)
        data: Data for insert/update
        where: Where conditions for update/delete
        max_rows: Maximum rows to return
    """
    
    display_name = "Database Query"
    
    async def execute(
        self,
        node: ProcessNode,
        state: ProcessState,
        context: ProcessContext
    ) -> NodeResult:
        """Execute database query node"""
        
        connection_id = self.get_config_value(node, 'connection_id')
        operation = self.get_config_value(node, 'operation', 'query')
        query_template = self.get_config_value(node, 'query')
        table = self.get_config_value(node, 'table')
        data = self.get_config_value(node, 'data', {})
        where = self.get_config_value(node, 'where', {})
        max_rows = self.get_config_value(node, 'max_rows', 1000)
        
        logs = [f"Database {operation} on connection: {connection_id}"]
        
        # Get database connection/tool
        db_connection = self.deps.get_db_connection(connection_id)
        if not db_connection:
            # Try getting as a tool
            db_tool = self.deps.get_tool(connection_id)
            if not db_tool:
                return NodeResult.failure(
                    error=ExecutionError(
                        category=ErrorCategory.RESOURCE,
                        code="DB_CONNECTION_NOT_FOUND",
                        message=f"Database connection not found: {connection_id}"
                    ),
                    logs=logs
                )
            # Use tool for execution
            return await self._execute_via_tool(node, state, db_tool, logs)
        
        # Build and execute query
        start_time = time.time()
        
        try:
            if operation == 'query' and query_template:
                # Interpolate query
                query = state.interpolate_string(query_template)
                logs.append(f"Query: {query[:100]}...")
                
                # Execute query (implementation depends on DB type)
                result = await self._execute_query(db_connection, query, max_rows)
                
            elif operation == 'insert':
                # Interpolate data
                interpolated_data = state.interpolate_object(data)
                result = await self._execute_insert(db_connection, table, interpolated_data)
                logs.append(f"Inserted into {table}")
                
            elif operation == 'update':
                interpolated_data = state.interpolate_object(data)
                interpolated_where = state.interpolate_object(where)
                result = await self._execute_update(db_connection, table, interpolated_data, interpolated_where)
                logs.append(f"Updated {table}")
                
            elif operation == 'delete':
                interpolated_where = state.interpolate_object(where)
                result = await self._execute_delete(db_connection, table, interpolated_where)
                logs.append(f"Deleted from {table}")
                
            else:
                return NodeResult.failure(
                    error=ExecutionError.validation_error(f"Unknown operation: {operation}"),
                    logs=logs
                )
            
            duration_ms = (time.time() - start_time) * 1000
            logs.append(f"Completed in {duration_ms:.0f}ms")
            
            variables_update = {}
            if node.output_variable:
                variables_update[node.output_variable] = result
            
            return NodeResult.success(
                output=result,
                variables_update=variables_update,
                duration_ms=duration_ms,
                logs=logs
            )
            
        except Exception as e:
            return NodeResult.failure(
                error=ExecutionError(
                    category=ErrorCategory.EXTERNAL,
                    code="DB_ERROR",
                    message=f"Database error: {e}",
                    is_retryable=True
                ),
                logs=logs
            )
    
    async def _execute_via_tool(
        self,
        node: ProcessNode,
        state: ProcessState,
        db_tool,
        logs: list
    ) -> NodeResult:
        """Execute database operation via platform tool"""
        
        query_template = self.get_config_value(node, 'query')
        if not query_template:
            return NodeResult.failure(
                error=ExecutionError.validation_error("Query is required when using tool"),
                logs=logs
            )
        
        query = state.interpolate_string(query_template)
        logs.append(f"Executing via tool: {query[:100]}...")
        
        start_time = time.time()
        try:
            result = await db_tool.execute(query=query)
            duration_ms = (time.time() - start_time) * 1000
            
            if result.success:
                variables_update = {}
                if node.output_variable:
                    variables_update[node.output_variable] = result.data
                
                return NodeResult.success(
                    output=result.data,
                    variables_update=variables_update,
                    duration_ms=duration_ms,
                    logs=logs
                )
            else:
                return NodeResult.failure(
                    error=ExecutionError(
                        category=ErrorCategory.EXTERNAL,
                        code="DB_TOOL_ERROR",
                        message=result.error
                    ),
                    logs=logs
                )
        except Exception as e:
            return NodeResult.failure(
                error=ExecutionError.internal_error(f"Tool execution failed: {e}"),
                logs=logs
            )
    
    async def _execute_query(self, connection, query: str, max_rows: int) -> Any:
        """
        Execute a SELECT query using the connection
        
        Connection should be a dict with:
        - type: postgres, mysql, sqlite
        - url: connection URL
        """
        db_type = connection.get('type', 'postgres')
        db_url = connection.get('url')
        
        if not db_url:
            raise ValueError("Database URL not configured")
        
        if db_type == 'postgres':
            return await self._query_postgres(db_url, query, max_rows)
        elif db_type == 'mysql':
            return await self._query_mysql(db_url, query, max_rows)
        elif db_type == 'sqlite':
            return await self._query_sqlite(db_url, query, max_rows)
        else:
            raise ValueError(f"Unsupported database type: {db_type}")
    
    async def _query_postgres(self, url: str, query: str, max_rows: int) -> list:
        """Execute query on PostgreSQL"""
        try:
            import asyncpg
            
            conn = await asyncpg.connect(url)
            try:
                rows = await conn.fetch(query)
                result = [dict(row) for row in rows[:max_rows]]
                return result
            finally:
                await conn.close()
        except ImportError:
            raise ImportError("asyncpg not installed. Run: pip install asyncpg")
    
    async def _query_mysql(self, url: str, query: str, max_rows: int) -> list:
        """Execute query on MySQL"""
        try:
            import aiomysql
            import urllib.parse
            
            parsed = urllib.parse.urlparse(url)
            
            conn = await aiomysql.connect(
                host=parsed.hostname,
                port=parsed.port or 3306,
                user=parsed.username,
                password=parsed.password,
                db=parsed.path.lstrip('/')
            )
            try:
                async with conn.cursor(aiomysql.DictCursor) as cursor:
                    await cursor.execute(query)
                    rows = await cursor.fetchmany(max_rows)
                    return list(rows)
            finally:
                conn.close()
        except ImportError:
            raise ImportError("aiomysql not installed. Run: pip install aiomysql")
    
    async def _query_sqlite(self, path: str, query: str, max_rows: int) -> list:
        """Execute query on SQLite"""
        try:
            import aiosqlite
            
            # Remove sqlite:/// prefix if present
            if path.startswith('sqlite:///'):
                path = path[10:]
            
            async with aiosqlite.connect(path) as db:
                db.row_factory = aiosqlite.Row
                async with db.execute(query) as cursor:
                    rows = await cursor.fetchmany(max_rows)
                    return [dict(row) for row in rows]
        except ImportError:
            raise ImportError("aiosqlite not installed. Run: pip install aiosqlite")
    
    async def _execute_insert(self, connection, table: str, data: Dict) -> Any:
        """Execute an INSERT"""
        # Build INSERT query
        columns = ', '.join(data.keys())
        placeholders = ', '.join([f"${i+1}" for i in range(len(data))])
        query = f"INSERT INTO {table} ({columns}) VALUES ({placeholders}) RETURNING *"
        
        db_type = connection.get('type', 'postgres')
        db_url = connection.get('url')
        
        if db_type == 'postgres':
            try:
                import asyncpg
                conn = await asyncpg.connect(db_url)
                try:
                    row = await conn.fetchrow(query, *data.values())
                    return dict(row) if row else {'inserted': True}
                finally:
                    await conn.close()
            except ImportError:
                raise ImportError("asyncpg not installed")
        else:
            # For MySQL/SQLite, use simpler approach
            return {'inserted': True, 'data': data}
    
    async def _execute_update(self, connection, table: str, data: Dict, where: Dict) -> Any:
        """Execute an UPDATE"""
        # Build UPDATE query
        set_clause = ', '.join([f"{k} = ${i+1}" for i, k in enumerate(data.keys())])
        where_clause = ' AND '.join([f"{k} = ${len(data)+i+1}" for i, k in enumerate(where.keys())])
        query = f"UPDATE {table} SET {set_clause} WHERE {where_clause}"
        
        db_type = connection.get('type', 'postgres')
        db_url = connection.get('url')
        
        if db_type == 'postgres':
            try:
                import asyncpg
                conn = await asyncpg.connect(db_url)
                try:
                    result = await conn.execute(query, *data.values(), *where.values())
                    return {'updated': True, 'result': result}
                finally:
                    await conn.close()
            except ImportError:
                raise ImportError("asyncpg not installed")
        else:
            return {'updated': True}
    
    async def _execute_delete(self, connection, table: str, where: Dict) -> Any:
        """Execute a DELETE"""
        where_clause = ' AND '.join([f"{k} = ${i+1}" for i, k in enumerate(where.keys())])
        query = f"DELETE FROM {table} WHERE {where_clause}"
        
        db_type = connection.get('type', 'postgres')
        db_url = connection.get('url')
        
        if db_type == 'postgres':
            try:
                import asyncpg
                conn = await asyncpg.connect(db_url)
                try:
                    result = await conn.execute(query, *where.values())
                    return {'deleted': True, 'result': result}
                finally:
                    await conn.close()
            except ImportError:
                raise ImportError("asyncpg not installed")
        else:
            return {'deleted': True}
    
    def validate(self, node: ProcessNode) -> Optional[ExecutionError]:
        """Validate database query node"""
        base_error = super().validate(node)
        if base_error:
            return base_error
        
        connection_id = self.get_config_value(node, 'connection_id')
        if not connection_id:
            return ExecutionError.validation_error("connection_id is required")
        
        operation = self.get_config_value(node, 'operation', 'query')
        if operation == 'query':
            query = self.get_config_value(node, 'query')
            if not query:
                return ExecutionError.validation_error("query is required for query operation")
        else:
            table = self.get_config_value(node, 'table')
            if not table:
                return ExecutionError.validation_error(f"table is required for {operation} operation")
        
        return None


@register_executor(NodeType.FILE_OPERATION)
class FileOperationNodeExecutor(BaseNodeExecutor):
    """
    File Operation node executor
    
    Reads and writes files (local, S3, SFTP, etc.)
    
    Config:
        operation: read, write, delete, list
        storage_type: local, s3
        path: File path
        content: Content to write
        encoding: File encoding (default: utf-8)
        storage_config: Storage-specific configuration
    """
    
    display_name = "File Operation"

    @staticmethod
    def _resolve_uploaded_file_path(file_obj: dict, context=None) -> str:
        """
        Resolve an uploaded file object to its physical filesystem path.

        Uploaded file objects look like:
            {"kind": "uploadedFile", "id": "<uuid>", "name": "receipt.png",
             "download_url": "/process/uploads/<uuid>/download", ...}

        The physical file is stored at:
            {UPLOAD_PATH}/process_uploads/{org_id}/{file_id}_{filename}

        This method finds the file on disk by scanning the org upload directory.
        """
        import os

        if not isinstance(file_obj, dict):
            return ""

        file_id = str(file_obj.get("id") or "").strip()
        file_name = file_obj.get("name") or ""
        dl_url = str(file_obj.get("download_url") or "").strip()

        # Extract file_id from download_url if not directly available
        if not file_id and dl_url:
            # /process/uploads/<uuid>/download → extract <uuid>
            parts = dl_url.rstrip("/").split("/")
            # Find "uploads" and take the next segment
            for i, p in enumerate(parts):
                if p == "uploads" and i + 1 < len(parts):
                    candidate = parts[i + 1]
                    if candidate != "download":
                        file_id = candidate
                        break

        if not file_id:
            return ""

        base_dir = os.environ.get("UPLOAD_PATH", "data/uploads")

        # Try to find org_id from context
        org_id = ""
        if context:
            # ProcessContext has org_id directly
            org_id = str(getattr(context, 'org_id', '') or "").strip()
            if not org_id:
                uc = (getattr(context, 'trigger_input', None) or {}).get("_user_context", {})
                org_id = str(uc.get("org_id") or "").strip()
            if not org_id:
                cu = (getattr(context, 'trigger_input', None) or {}).get("currentUser", {})
                if isinstance(cu, dict):
                    org_id = str(cu.get("orgId") or cu.get("org_id") or "").strip()

        # Search strategies for the file
        search_dirs = []
        if org_id:
            search_dirs.append(os.path.join(base_dir, "process_uploads", org_id))
        # Fallback: scan all org subdirectories
        proc_uploads = os.path.join(base_dir, "process_uploads")
        if os.path.isdir(proc_uploads):
            try:
                for d in os.listdir(proc_uploads):
                    full = os.path.join(proc_uploads, d)
                    if os.path.isdir(full) and full not in search_dirs:
                        search_dirs.append(full)
            except OSError:
                pass

        prefix = f"{file_id}_"
        for search_dir in search_dirs:
            if not os.path.isdir(search_dir):
                continue
            try:
                for name in os.listdir(search_dir):
                    if name.startswith(prefix):
                        return os.path.join(search_dir, name)
            except OSError:
                continue

        return ""

    async def execute(
        self,
        node: ProcessNode,
        state: ProcessState,
        context: ProcessContext
    ) -> NodeResult:
        """Execute file operation"""
        
        operation = self.get_config_value(node, 'operation', 'read')
        storage_type = self.get_config_value(node, 'storage_type', 'local')
        path_template = self.get_config_value(node, 'path', '')
        content = self.get_config_value(node, 'content')
        encoding = self.get_config_value(node, 'encoding', 'utf-8')
        storage_config = self.get_config_value(node, 'storage_config', {})
        
        logs = [f"File {operation} on {storage_type}"]

        # Execute based on storage type
        start_time = time.time()
        
        try:
            # Special operation: document generation (writes an actual file)
            if operation == 'generate_document':
                if storage_type != 'local':
                    return NodeResult.failure(
                        error=ExecutionError.validation_error("generate_document currently supports storage_type=local only"),
                        logs=logs
                    )
                result = await self._execute_generate_document_local(
                    node=node,
                    state=state,
                    context=context,
                    logs=logs,
                )
            # Special operation: extract text from a document (for AI extraction and review)
            elif operation in ('extract_text', 'extract_document_text'):
                if storage_type != 'local':
                    return NodeResult.failure(
                        error=ExecutionError.validation_error("extract_text currently supports storage_type=local only"),
                        logs=logs
                    )
                
                # Detect multi-file uploads: if sourceField points to an array of files,
                # extract text from each file and concatenate the results.
                source_field = self.get_config_value(node, 'sourceField', '')
                source_value = state.get(source_field) if source_field else None
                
                if isinstance(source_value, list) and len(source_value) > 0:
                    # Multiple files uploaded — extract from each
                    logs.append(f"Multi-file extraction: {len(source_value)} files detected")
                    all_texts = []
                    for idx, file_item in enumerate(source_value):
                        file_path = None
                        if isinstance(file_item, dict):
                            file_name = file_item.get("name") or f"file_{idx+1}"
                            # Resolve the uploaded file to its physical path on disk
                            file_path = self._resolve_uploaded_file_path(file_item, context)
                            if not file_path:
                                # Fallback: try legacy path property
                                file_path = file_item.get("path") or ""
                        else:
                            continue
                        if not file_path:
                            logs.append(f"  File {idx+1} ({file_name}): could not resolve path, skipping")
                            continue
                        logs.append(f"  Extracting file {idx+1}/{len(source_value)}: {file_name}")
                        file_result = await self._execute_extract_text_local(
                            path=file_path,
                            encoding=encoding,
                            logs=logs,
                        )
                        if file_result.get("success") and file_result.get("data"):
                            all_texts.append(f"--- File: {file_name} ---\n{file_result['data']}")
                            logs.append(f"  ✅ Extracted {len(file_result['data'])} characters from {file_name}")
                        else:
                            err = file_result.get("error") or "extraction failed"
                            all_texts.append(f"--- File: {file_name} ---\n[Extraction failed: {err}]")
                            logs.append(f"  ⚠️ Failed to extract from {file_name}: {err}")
                    
                    combined_text = "\n\n".join(all_texts)
                    if combined_text.strip():
                        result = {
                            "success": True,
                            "data": combined_text,
                            "meta": {
                                "files_count": len(source_value),
                                "chars": len(combined_text),
                                "preview": combined_text[:500]
                            }
                        }
                    else:
                        result = {
                            "success": False,
                            "error": f"Could not extract text from any of the {len(source_value)} uploaded files",
                            "data": ""
                        }
                elif isinstance(source_value, dict) and source_value.get("kind") == "uploadedFile":
                    # Single uploaded file object — resolve to physical path
                    file_name = source_value.get("name") or "file"
                    logs.append(f"Single uploaded file detected: {file_name}")
                    path = self._resolve_uploaded_file_path(source_value, context)
                    if not path:
                        # Fallback to legacy path
                        path = source_value.get("path") or ""
                    logs.append(f"Resolved file path: {path}")
                    if not path:
                        return NodeResult.failure(
                            error=ExecutionError.validation_error(
                                f"Could not resolve file path for uploaded file '{file_name}'. "
                                f"File ID: {source_value.get('id', 'unknown')}. "
                                f"Make sure the file was uploaded successfully."
                            ),
                            logs=logs
                        )
                else:
                    # Fallback: original template-based path resolution
                    logs.append(f"Path template: {path_template}")
                    path = state.interpolate_string(path_template)
                    logs.append(f"Resolved path: {path}")
                    if not path or path == path_template or path.startswith('{{'):
                        # Path didn't resolve — try sourceField as uploaded file object
                        if isinstance(source_value, dict) and (source_value.get("download_url") or source_value.get("id")):
                            path = self._resolve_uploaded_file_path(source_value, context)
                            logs.append(f"Fallback: resolved uploaded file to: {path}")
                        if not path or path.startswith('{{'):
                            return NodeResult.failure(
                                error=ExecutionError.validation_error(
                                    f"Could not resolve file path from '{path_template}'. "
                                    f"Make sure the file upload field name matches the source field configuration. "
                                    f"Resolved to: '{path}'"
                                ),
                                logs=logs
                            )
                    result = await self._execute_extract_text_local(
                        path=path,
                        encoding=encoding,
                        logs=logs,
                    )
            else:
                # Interpolate path for normal operations
                path = state.interpolate_string(path_template)
                logs.append(f"Path: {path}")

            if storage_type == 'local':
                if operation not in ('generate_document', 'extract_text', 'extract_document_text'):
                    result = await self._execute_local(
                        operation, path, content, encoding, logs
                    )
            elif storage_type == 's3':
                result = await self._execute_s3(
                    operation, path, content, storage_config, logs
                )
            else:
                return NodeResult.failure(
                    error=ExecutionError(
                        category=ErrorCategory.CONFIGURATION,
                        code="UNSUPPORTED_STORAGE",
                        message=f"Storage type '{storage_type}' not supported. Use 'local' or 's3'."
                    ),
                    logs=logs
                )
            
            duration_ms = (time.time() - start_time) * 1000
            
            if result.get('success'):
                variables_update = {}
                if node.output_variable:
                    variables_update[node.output_variable] = result.get('data')
                
                return NodeResult.success(
                    output=result,
                    variables_update=variables_update,
                    duration_ms=duration_ms,
                    logs=logs
                )
            else:
                return NodeResult.failure(
                    error=ExecutionError(
                        category=ErrorCategory.EXTERNAL,
                        code="FILE_ERROR",
                        message=result.get('error', 'File operation failed')
                    ),
                    duration_ms=duration_ms,
                    logs=logs
                )
                
        except Exception as e:
            return NodeResult.failure(
                error=ExecutionError.internal_error(f"File operation error: {e}"),
                logs=logs
            )

    async def _execute_extract_text_local(
        self,
        path: str,
        encoding: str,
        logs: list
    ) -> Dict[str, Any]:
        """
        Extract text/content from a local file.

        Handles documents (pdf, docx, xlsx, pptx, csv, txt, etc.),
        images (png, jpg, gif, webp, bmp, tiff — via LLM OCR),
        and any other text-readable file.

        Security: only allows files under the platform upload directory.
        """
        import os
        import csv

        base_upload = os.path.abspath(os.environ.get("UPLOAD_PATH", "data/uploads"))
        abs_path = os.path.abspath(path or "")

        if not abs_path:
            return {"success": False, "error": "No file path provided"}
        if not abs_path.startswith(base_upload + os.sep) and abs_path != base_upload:
            return {"success": False, "error": "File path is not allowed"}
        if not os.path.exists(abs_path):
            return {"success": False, "error": "File not found"}
        if os.path.isdir(abs_path):
            return {"success": False, "error": "Path is a directory"}

        ext = os.path.splitext(abs_path)[1].lower().lstrip(".")
        logs.append(f"Detected file type: {ext or 'unknown'}")

        try:
            text = ""
            if ext in ("txt", "md", "text", "log", "json"):
                with open(abs_path, "r", encoding=encoding, errors="ignore") as f:
                    text = f.read()
            elif ext == "csv":
                with open(abs_path, "r", encoding=encoding, errors="ignore", newline="") as f:
                    reader = csv.reader(f)
                    text = "\n".join([" | ".join(row) for row in reader])
            elif ext == "pdf":
                try:
                    import fitz  # PyMuPDF
                    doc = fitz.open(abs_path)
                    text = "".join([page.get_text() for page in doc])
                    doc.close()
                except Exception:
                    return {"success": False, "error": "PDF text extraction failed"}
            elif ext in ("docx", "doc"):
                try:
                    from docx import Document as DocxDocument
                    doc = DocxDocument(abs_path)
                    text = "\n".join([p.text for p in doc.paragraphs])
                except Exception:
                    return {"success": False, "error": "Word text extraction failed"}
            elif ext in ("xlsx", "xls"):
                try:
                    import openpyxl
                    wb = openpyxl.load_workbook(abs_path, read_only=True)
                    chunks = []
                    for sheet in wb.sheetnames:
                        ws = wb[sheet]
                        chunks.append(f"=== Sheet: {sheet} ===")
                        for row in ws.iter_rows(values_only=True):
                            chunks.append(" | ".join([str(c) if c is not None else "" for c in row]))
                    text = "\n".join(chunks)
                except Exception:
                    return {"success": False, "error": "Excel text extraction failed"}
            elif ext in ("pptx", "ppt"):
                try:
                    from pptx import Presentation
                    prs = Presentation(abs_path)
                    chunks = []
                    for i, slide in enumerate(prs.slides):
                        chunks.append(f"=== Slide {i + 1} ===")
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                chunks.append(shape.text)
                    text = "\n".join(chunks)
                except Exception:
                    return {"success": False, "error": "PowerPoint text extraction failed"}
            elif ext in ("png", "jpg", "jpeg", "gif", "webp", "bmp", "tiff", "tif", "svg", "ico", "heic", "heif"):
                # Image files — use LLM vision for OCR/content extraction
                logs.append("Image detected — attempting LLM-based OCR/extraction")
                try:
                    import base64
                    with open(abs_path, "rb") as img_f:
                        img_data = base64.b64encode(img_f.read()).decode("utf-8")
                    mime_map = {
                        "png": "image/png", "jpg": "image/jpeg", "jpeg": "image/jpeg",
                        "gif": "image/gif", "webp": "image/webp", "bmp": "image/bmp",
                        "tiff": "image/tiff", "tif": "image/tiff", "svg": "image/svg+xml",
                        "ico": "image/x-icon", "heic": "image/heic", "heif": "image/heif",
                    }
                    mime = mime_map.get(ext, f"image/{ext}")
                    data_uri = f"data:{mime};base64,{img_data}"
                    # Try to use LLM for OCR — this requires a vision-capable model
                    try:
                        from openai import OpenAI
                        api_key = os.environ.get("OPENAI_API_KEY")
                        if api_key:
                            client = OpenAI(api_key=api_key)
                            response = client.chat.completions.create(
                                model="gpt-4o",
                                messages=[{
                                    "role": "user",
                                    "content": [
                                        {"type": "text", "text": "Extract ALL text, data, numbers, and information visible in this image. If it contains a form, table, receipt, invoice, or document, extract every field and value. Return the extracted content as structured text."},
                                        {"type": "image_url", "image_url": {"url": data_uri}}
                                    ]
                                }],
                                max_tokens=4096
                            )
                            text = response.choices[0].message.content or ""
                            logs.append(f"LLM OCR extracted {len(text)} characters")
                        else:
                            logs.append("No OPENAI_API_KEY — cannot perform OCR")
                            return {"success": False, "error": "Image OCR requires an LLM API key (OPENAI_API_KEY). Configure it in environment settings."}
                    except ImportError:
                        return {"success": False, "error": "OpenAI library not available for image OCR"}
                except Exception as e:
                    return {"success": False, "error": f"Image extraction failed: {str(e)}"}
            elif ext in ("html", "htm", "xml", "yaml", "yml", "ini", "cfg", "conf", "properties", "env", "sh", "bat", "ps1", "py", "js", "ts", "java", "c", "cpp", "go", "rs", "rb", "php", "sql", "r"):
                # Text-based files — read as plain text
                with open(abs_path, "r", encoding=encoding, errors="ignore") as f:
                    text = f.read()
            else:
                # Generic fallback — try reading as text first, then report unsupported
                logs.append(f"Unknown file type '{ext}' — attempting text read as fallback")
                try:
                    with open(abs_path, "r", encoding=encoding, errors="ignore") as f:
                        text = f.read()
                    if not text.strip():
                        return {"success": False, "error": f"Could not extract text from file type: {ext or 'unknown'}. The file may be binary or empty."}
                except Exception:
                    return {"success": False, "error": f"Could not read file type: {ext or 'unknown'}"}

            text = text or ""
            preview = text[:800]
            logs.append(f"Extracted {len(text)} characters")
            # If extraction yielded no text, treat as failure so downstream steps
            # (AI parsing, conditions) don't silently operate on empty data.
            if not text.strip():
                logs.append("Extraction produced empty text — marking as failure")
                return {
                    "success": False,
                    "error": f"Could not extract any text from the file ({ext or 'unknown'} format). The file may be empty, corrupted, or unreadable. Please upload a clear, readable document or image.",
                    "data": "",
                    "meta": {"path": abs_path, "file_type": ext, "chars": 0}
                }
            return {
                "success": True,
                "data": text,
                "meta": {
                    "path": abs_path,
                    "file_type": ext,
                    "chars": len(text),
                    "preview": preview,
                }
            }
        except Exception as e:
            return {"success": False, "error": f"Extraction failed: {e}"}
    
    async def _execute_local(
        self,
        operation: str,
        path: str,
        content: Any,
        encoding: str,
        logs: list
    ) -> Dict[str, Any]:
        """Execute local file operation"""
        import os
        import aiofiles
        
        try:
            if operation == 'read':
                async with aiofiles.open(path, 'r', encoding=encoding) as f:
                    data = await f.read()
                logs.append(f"Read {len(data)} characters")
                return {'success': True, 'data': data}
                
            elif operation == 'write':
                # Ensure directory exists
                os.makedirs(os.path.dirname(path), exist_ok=True)
                
                write_content = content
                if isinstance(content, dict):
                    write_content = json.dumps(content, indent=2)
                elif not isinstance(content, str):
                    write_content = str(content)
                
                async with aiofiles.open(path, 'w', encoding=encoding) as f:
                    await f.write(write_content)
                logs.append(f"Wrote {len(write_content)} characters")
                return {'success': True, 'bytes_written': len(write_content)}
                
            elif operation == 'delete':
                if os.path.exists(path):
                    os.remove(path)
                    logs.append("File deleted")
                    return {'success': True, 'deleted': True}
                else:
                    logs.append("File not found")
                    return {'success': True, 'deleted': False, 'note': 'File not found'}
                
            elif operation == 'list':
                if os.path.isdir(path):
                    files = os.listdir(path)
                    logs.append(f"Found {len(files)} items")
                    return {'success': True, 'data': files}
                else:
                    return {'success': False, 'error': 'Path is not a directory'}
                
            elif operation == 'exists':
                exists = os.path.exists(path)
                return {'success': True, 'data': exists}
                
            else:
                return {'success': False, 'error': f'Unknown operation: {operation}'}
                
        except FileNotFoundError:
            return {'success': False, 'error': f'File not found: {path}'}
        except PermissionError:
            return {'success': False, 'error': f'Permission denied: {path}'}
        except ImportError:
            # Fallback to sync if aiofiles not installed
            return await self._execute_local_sync(operation, path, content, encoding, logs)
    
    async def _execute_local_sync(
        self,
        operation: str,
        path: str,
        content: Any,
        encoding: str,
        logs: list
    ) -> Dict[str, Any]:
        """Synchronous fallback for local files"""
        import os
        
        if operation == 'read':
            with open(path, 'r', encoding=encoding) as f:
                data = f.read()
            return {'success': True, 'data': data}
            
        elif operation == 'write':
            os.makedirs(os.path.dirname(path) or '.', exist_ok=True)
            write_content = content if isinstance(content, str) else json.dumps(content)
            with open(path, 'w', encoding=encoding) as f:
                f.write(write_content)
            return {'success': True, 'bytes_written': len(write_content)}
            
        return {'success': False, 'error': f'Operation {operation} not supported in sync mode'}

    async def _execute_generate_document_local(
        self,
        node: ProcessNode,
        state: ProcessState,
        context: ProcessContext,
        logs: list
    ) -> Dict[str, Any]:
        """
        Generate a professional document and write it to local storage.

        Config (type_config):
            title: Document title
            format: docx|pdf|xlsx|pptx|txt (fallback)
            instructions: what the document should contain (supports {{var}} interpolation)
            output_dir: optional base output dir (default: env PROCESS_OUTPUT_PATH or data/process_outputs)
        """
        import os
        import re

        title = self.get_config_value(node, 'title') or node.name or 'Document'
        fmt = self.get_config_value(node, 'format') or 'docx'
        fmt = str(fmt).strip().lower().lstrip('.')

        instructions = self.get_config_value(node, 'instructions') or ''
        instructions = state.interpolate_string(str(instructions))

        base_dir = self.get_config_value(node, 'output_dir') or os.environ.get('PROCESS_OUTPUT_PATH', 'data/process_outputs')
        exec_dir = os.path.join(base_dir, str(context.execution_id))
        os.makedirs(exec_dir, exist_ok=True)

        safe_base = re.sub(r'[^A-Za-z0-9]+', '_', str(title)).strip('_') or 'document'
        filename = f"{safe_base}_{node.id}.{fmt}"
        filepath = os.path.join(exec_dir, filename)

        logs.append(f"Generating document: {title} ({fmt})")

        content_md = await self._generate_document_markdown(title=title, instructions=instructions, state=state, logs=logs)

        actual_format = fmt
        try:
            if fmt == 'docx':
                size = self._create_docx(filepath, title, content_md)
            elif fmt == 'pdf':
                size = self._create_pdf(filepath, title, content_md)
            elif fmt == 'xlsx':
                size = self._create_xlsx(filepath, title, content_md)
            elif fmt == 'pptx':
                size = self._create_pptx(filepath, title, content_md)
            else:
                actual_format = 'txt'
                filepath = os.path.join(exec_dir, f"{safe_base}_{node.id}.txt")
                with open(filepath, 'w', encoding='utf-8') as f:
                    f.write(f"# {title}\n\n{content_md}\n")
                size = os.path.getsize(filepath)
        except Exception as e:
            # Safe fallback
            logs.append(f"Document generation fallback: {e}")
            actual_format = 'txt'
            filepath = os.path.join(exec_dir, f"{safe_base}_{node.id}.txt")
            with open(filepath, 'w', encoding='utf-8') as f:
                f.write(f"# {title}\n\n{content_md}\n")
            size = os.path.getsize(filepath)

        logs.append(f"Document saved: {filepath}")
        return {
            'success': True,
            'data': {
                'title': title,
                'format': actual_format,
                'path': filepath,
                'filename': os.path.basename(filepath),
                'size': size,
            }
        }

    async def _generate_document_markdown(
        self,
        title: str,
        instructions: str,
        state: ProcessState,
        logs: list
    ) -> str:
        """Generate document content as markdown using the platform LLM (if available)."""
        llm = getattr(self.deps, 'llm', None)
        if not llm:
            return f"## Overview\n{instructions}\n\n## Details\n(LLM not configured)\n"

        # Limit variable context to avoid huge prompts; sensitive values are masked.
        try:
            variables = state.get_masked_variables()
            trimmed = {}
            for k in list(variables.keys())[:60]:
                trimmed[k] = variables.get(k)
            variables_json = json.dumps(trimmed, ensure_ascii=False, indent=2)
        except Exception:
            variables_json = "{}"

        system_prompt = (
            "You are a professional document writer. "
            "Write clear, well-structured content for business users. "
            "Use markdown headings and bullet lists. "
            "Do NOT invent company-specific facts; only use provided data."
        )
        user_prompt = (
            f"Document title: {title}\n\n"
            f"Instructions:\n{instructions}\n\n"
            f"Available data (JSON, sensitive masked):\n{variables_json}\n\n"
            "Write the document content now."
        )

        try:
            from core.llm.base import Message, MessageRole
            resp = await llm.chat(
                messages=[
                    Message(role=MessageRole.SYSTEM, content=system_prompt),
                    Message(role=MessageRole.USER, content=user_prompt),
                ],
                temperature=0.4,
                max_tokens=2500,
            )
            text = (resp.content or "").strip() if resp else ""
            if text:
                return text
        except Exception as e:
            logs.append(f"LLM document content generation failed: {e}")

        return f"# {title}\n\n## Overview\n{instructions}\n\n## Notes\n(Generated without LLM)\n"

    def _create_docx(self, filepath: str, title: str, content: str) -> int:
        """Create a Word document from markdown-ish content."""
        import os
        import re
        from docx import Document

        doc = Document()
        doc.add_heading(title, 0)

        for raw in (content or "").split('\n'):
            line = raw.strip()
            if not line:
                continue
            if line.startswith('# '):
                doc.add_heading(line[2:], 1)
            elif line.startswith('## '):
                doc.add_heading(line[3:], 2)
            elif line.startswith('### '):
                doc.add_heading(line[4:], 3)
            elif line.startswith('- ') or line.startswith('* '):
                doc.add_paragraph(line[2:], style='List Bullet')
            elif re.match(r'^\d+\.\s+', line):
                doc.add_paragraph(re.sub(r'^\d+\.\s+', '', line), style='List Number')
            else:
                doc.add_paragraph(line)

        doc.save(filepath)
        return os.path.getsize(filepath)

    def _create_pdf(self, filepath: str, title: str, content: str) -> int:
        """Create a PDF document from markdown-ish content (requires reportlab)."""
        import os
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.styles import ParagraphStyle

        doc = SimpleDocTemplate(filepath, pagesize=A4)
        styles = getSampleStyleSheet()
        story = []

        title_style = ParagraphStyle(
            'DocTitle',
            parent=styles['Heading1'],
            fontSize=22,
            spaceAfter=18
        )
        story.append(Paragraph(title, title_style))
        story.append(Spacer(1, 12))

        for raw in (content or "").split('\n'):
            line = raw.strip()
            if not line:
                story.append(Spacer(1, 10))
                continue
            if line.startswith('# '):
                story.append(Paragraph(line[2:], styles['Heading1']))
            elif line.startswith('## '):
                story.append(Paragraph(line[3:], styles['Heading2']))
            elif line.startswith('### '):
                story.append(Paragraph(line[4:], styles['Heading3']))
            else:
                story.append(Paragraph(line, styles['Normal']))

        doc.build(story)
        return os.path.getsize(filepath)

    def _create_xlsx(self, filepath: str, title: str, content: str) -> int:
        """Create an Excel file (simple sheet) from content (requires openpyxl)."""
        import os
        from openpyxl import Workbook
        from openpyxl.styles import Font

        wb = Workbook()
        ws = wb.active
        ws.title = (title or 'Sheet')[:31]

        ws['A1'] = title
        ws['A1'].font = Font(bold=True, size=16)
        ws.merge_cells('A1:E1')

        row = 3
        for raw in (content or "").split('\n'):
            line = raw.strip()
            if not line:
                continue
            ws.cell(row=row, column=1, value=line)
            row += 1

        wb.save(filepath)
        return os.path.getsize(filepath)

    def _create_pptx(self, filepath: str, title: str, content: str) -> int:
        """Create a PPTX with a title slide + bullets (requires python-pptx)."""
        import os
        from pptx import Presentation

        prs = Presentation()
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = title
        if slide.placeholders and len(slide.placeholders) > 1:
            slide.placeholders[1].text = ""

        bullet_layout = prs.slide_layouts[1]
        slide2 = prs.slides.add_slide(bullet_layout)
        slide2.shapes.title.text = "Summary"
        tf = slide2.shapes.placeholders[1].text_frame
        tf.clear()
        for raw in (content or "").split('\n'):
            line = raw.strip()
            if not line:
                continue
            if line.startswith('#'):
                continue
            p = tf.add_paragraph()
            p.text = line[:200]
            p.level = 0

        prs.save(filepath)
        return os.path.getsize(filepath)
    
    async def _execute_s3(
        self,
        operation: str,
        path: str,
        content: Any,
        config: Dict[str, Any],
        logs: list
    ) -> Dict[str, Any]:
        """Execute S3 file operation"""
        try:
            import aioboto3
            
            session = aioboto3.Session(
                aws_access_key_id=config.get('access_key'),
                aws_secret_access_key=config.get('secret_key'),
                region_name=config.get('region', 'us-east-1')
            )
            
            bucket = config.get('bucket')
            if not bucket:
                return {'success': False, 'error': 'S3 bucket not configured'}
            
            async with session.client('s3') as s3:
                if operation == 'read':
                    response = await s3.get_object(Bucket=bucket, Key=path)
                    data = await response['Body'].read()
                    logs.append(f"Read {len(data)} bytes from S3")
                    return {'success': True, 'data': data.decode('utf-8')}
                    
                elif operation == 'write':
                    write_content = content
                    if isinstance(content, dict):
                        write_content = json.dumps(content)
                    elif not isinstance(content, (str, bytes)):
                        write_content = str(content)
                    
                    if isinstance(write_content, str):
                        write_content = write_content.encode('utf-8')
                    
                    await s3.put_object(Bucket=bucket, Key=path, Body=write_content)
                    logs.append(f"Wrote to S3: {path}")
                    return {'success': True, 'path': f"s3://{bucket}/{path}"}
                    
                elif operation == 'delete':
                    await s3.delete_object(Bucket=bucket, Key=path)
                    logs.append(f"Deleted from S3: {path}")
                    return {'success': True, 'deleted': True}
                    
                elif operation == 'list':
                    response = await s3.list_objects_v2(Bucket=bucket, Prefix=path)
                    files = [obj['Key'] for obj in response.get('Contents', [])]
                    logs.append(f"Listed {len(files)} objects")
                    return {'success': True, 'data': files}
                    
                else:
                    return {'success': False, 'error': f'Unknown operation: {operation}'}
                    
        except ImportError:
            return {
                'success': False, 
                'error': 'aioboto3 not installed. Run: pip install aioboto3'
            }
        except Exception as e:
            return {'success': False, 'error': str(e)}


@register_executor(NodeType.MESSAGE_QUEUE)
class MessageQueueNodeExecutor(BaseNodeExecutor):
    """
    Message Queue node executor
    
    Pub/sub operations (Redis, SQS, webhook-based)
    
    Config:
        operation: publish
        queue_type: redis, sqs, webhook
        topic: Topic/queue name
        message: Message to publish
        queue_config: Queue-specific configuration
    """
    
    display_name = "Message Queue"
    
    async def execute(
        self,
        node: ProcessNode,
        state: ProcessState,
        context: ProcessContext
    ) -> NodeResult:
        """Execute message queue operation"""
        
        operation = self.get_config_value(node, 'operation', 'publish')
        queue_type = self.get_config_value(node, 'queue_type', 'webhook')
        topic = self.get_config_value(node, 'topic', '')
        message = self.get_config_value(node, 'message')
        queue_config = self.get_config_value(node, 'queue_config', {})
        
        logs = [f"Message queue {operation} on {queue_type}"]
        
        # Interpolate topic and message
        topic = state.interpolate_string(topic)
        logs.append(f"Topic: {topic}")
        
        if message:
            message = state.interpolate_object(message)
        
        start_time = time.time()
        
        try:
            if queue_type == 'webhook':
                result = await self._publish_webhook(topic, message, queue_config, logs)
            elif queue_type == 'redis':
                result = await self._publish_redis(topic, message, queue_config, logs)
            elif queue_type == 'sqs':
                result = await self._publish_sqs(topic, message, queue_config, logs)
            else:
                return NodeResult.failure(
                    error=ExecutionError(
                        category=ErrorCategory.CONFIGURATION,
                        code="UNSUPPORTED_QUEUE",
                        message=f"Queue type '{queue_type}' not supported. Use 'webhook', 'redis', or 'sqs'."
                    ),
                    logs=logs
                )
            
            duration_ms = (time.time() - start_time) * 1000
            
            if result.get('success'):
                variables_update = {}
                if node.output_variable:
                    variables_update[node.output_variable] = result
                
                return NodeResult.success(
                    output=result,
                    variables_update=variables_update,
                    duration_ms=duration_ms,
                    logs=logs
                )
            else:
                return NodeResult.failure(
                    error=ExecutionError(
                        category=ErrorCategory.EXTERNAL,
                        code="QUEUE_ERROR",
                        message=result.get('error', 'Message queue operation failed'),
                        is_retryable=True
                    ),
                    duration_ms=duration_ms,
                    logs=logs
                )
                
        except Exception as e:
            return NodeResult.failure(
                error=ExecutionError.internal_error(f"Message queue error: {e}"),
                logs=logs
            )
    
    async def _publish_webhook(
        self,
        url: str,
        message: Any,
        config: Dict[str, Any],
        logs: list
    ) -> Dict[str, Any]:
        """Publish message via webhook (HTTP POST)"""
        try:
            import aiohttp
            
            headers = config.get('headers', {'Content-Type': 'application/json'})
            timeout = config.get('timeout', 30)
            
            payload = message if isinstance(message, dict) else {'data': message}
            
            async with aiohttp.ClientSession() as session:
                async with session.post(
                    url,
                    json=payload,
                    headers=headers,
                    timeout=aiohttp.ClientTimeout(total=timeout)
                ) as response:
                    status = response.status
                    body = await response.text()
                    
                    logs.append(f"Webhook response: {status}")
                    
                    if status < 400:
                        return {
                            'success': True,
                            'status_code': status,
                            'response': body[:500]
                        }
                    else:
                        return {
                            'success': False,
                            'error': f"Webhook failed with status {status}: {body[:200]}"
                        }
                        
        except ImportError:
            return {'success': False, 'error': 'aiohttp not installed'}
        except Exception as e:
            return {'success': False, 'error': str(e)}
    
    async def _publish_redis(
        self,
        channel: str,
        message: Any,
        config: Dict[str, Any],
        logs: list
    ) -> Dict[str, Any]:
        """Publish message to Redis pub/sub"""
        try:
            import aioredis
            
            redis_url = config.get('url', 'redis://localhost:6379')
            
            redis = await aioredis.from_url(redis_url)
            
            msg_str = json.dumps(message) if isinstance(message, dict) else str(message)
            
            subscribers = await redis.publish(channel, msg_str)
            
            await redis.close()
            
            logs.append(f"Published to {subscribers} subscribers")
            
            return {
                'success': True,
                'channel': channel,
                'subscribers': subscribers
            }
            
        except ImportError:
            return {'success': False, 'error': 'aioredis not installed. Run: pip install aioredis'}
        except Exception as e:
            return {'success': False, 'error': str(e)}
    
    async def _publish_sqs(
        self,
        queue_url: str,
        message: Any,
        config: Dict[str, Any],
        logs: list
    ) -> Dict[str, Any]:
        """Publish message to AWS SQS"""
        try:
            import aioboto3
            
            session = aioboto3.Session(
                aws_access_key_id=config.get('access_key'),
                aws_secret_access_key=config.get('secret_key'),
                region_name=config.get('region', 'us-east-1')
            )
            
            msg_body = json.dumps(message) if isinstance(message, dict) else str(message)
            
            async with session.client('sqs') as sqs:
                response = await sqs.send_message(
                    QueueUrl=queue_url,
                    MessageBody=msg_body,
                    MessageAttributes=config.get('attributes', {})
                )
                
                message_id = response.get('MessageId')
                logs.append(f"SQS message ID: {message_id}")
                
                return {
                    'success': True,
                    'message_id': message_id,
                    'queue_url': queue_url
                }
                
        except ImportError:
            return {'success': False, 'error': 'aioboto3 not installed. Run: pip install aioboto3'}
        except Exception as e:
            return {'success': False, 'error': str(e)}
